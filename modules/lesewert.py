# -*- coding: utf-8 -*-
#%% Import Libraries
# IMPORT LIBRARIES

import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import matplotlib
import seaborn as sns
import datetime as dt

# Deutsch als mathematischen Standard einsetzen
import locale
locale.setlocale(locale.LC_ALL, "deu_deu")


# import dt. Mathezeichen (Komma statt punkt bei Dezimalzahlen)
matplotlib.rcParams['axes.formatter.use_locale'] = True


#import der Schrift für Matplotlib
from matplotlib import rcParams
import sys
import os
import matplotlib.font_manager as fm
#Schrift laden
fpath1 = 'C:\\Windows\\Fonts\\Campton-Light.otf'
campton_light = fm.FontProperties(fname=fpath1)


import timeit
import time
import requests
#für die Bildbearbeitung:
from PIL import Image
from PIL import ImageDraw, ImageFont
from io import StringIO
import textwrap

# alles von pptx importieren
from pptx import Presentation 
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt, Mm, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR
import time 

# Glätten der Line-Plots
from scipy.interpolate import interp1d
from scipy.signal import savgol_filter
from scipy.interpolate import spline
from scipy.interpolate import interp1d

# für die Lineare Regression bei der Entwicklung der Lesewerte
from sklearn.linear_model import Ridge

#%% LISTEN UND ONJEKTE SETZEN

#%% globale Variablen
# var für lokales Ressort
_lokales_ = "Lokales"


#%% glad

ausgaben = ["BIB", "FHA", "RV"]
werte = ["Lesewert", "Blickwert", "Durchlesewert"]
werte_dict = {"Lesewert":"LW", "Blickwert": "BW", "Durchlesewert":"DW"}
wertedict = {"Lesewert":"LW", "Blickwert": "BW", "Durchlesewert":"DW"}


platzierung = ["AA", "ZA", "SK", "ÜA", "TS1", "TS2", "TS3", "TS4", "TS5"]
platzierung_dict = {"AA": "Aufmacher", "ZA": "Zweitaufmacher", "SK": "Seitenkeller", "ÜA": "Überaufmacher", "TS1":"Einspalter", 
                    "TS2":"Zweispalter", "TS3":"Dreispalter", "TS4":"Vierspalter", "TS5": "Fünfspalter", "TS6": "Sechsspalter"}
darstellung_dict = {"NA": "Nachricht", "AK": "Ankündigung", "SA": "Struktur. Artikel", "ZU": "Zitattext/Umfrage",\
                        "SE":"Serie", "GA": "Gastartikel", "FF": "Freies Format", 'LB':'Leserbrief','ES': 'Essay', \
                        'RK': 'Rezension/Kritik', 'KM': 'Komm./Kolumne', 'BB': 'Bildbericht', 'BN': 'Bildnachricht', \
                        'GT': 'Geschichtstext', 'KI': 'Karikatur/Illustration/Grafik', 'ST': 'Serviceelement/Tabelle', \
                        'HG': 'Hintergrundkasten', 'IV': 'Interview', "RP": "Reportage/Porträt", "BF":"Bericht/Feature", "BI": "Bildnachricht" }
darstellungsform = ["BF", "RP", "IV", "KM", "BF", "LB", "SA", "NA", "RK", "ZU", "BN"]

rubrik_dict = {}
lokale_liste=[]

seitentitel_dict = {}

#seitentitel_dict = {'Veranstaltungen':'Veranstaltungen', 'Kirchen':"Kirchen", 
#                    'Lokales':"Lokales", 'Biberach':"Biberach", 'Rund um Biberach':'Rund um Biberach', 
#                    'Von der Schussen zur Rot': "Schussen zur Rot", 'Von der Rottum zur Iller':"Rottum zur Iller",
#                    "Kultur": "Kultur", "Friedrichshafen":"Friedrichshafen", "Region": "Region", "Biberach / Service":"Biberach/Service", 
#                    "Rund um Biberach / Schussen - Rot":"Biberach/Sch.-Rot",'Schussen-Rot / Rottum-Iller':"Sch.-Rot/Rottum-Ill.", 
#                    'Wir in Kreis und Region':'Wir in Kreis + Region', 'Friedrichshafen / Service':'Friedrichsh./Service', 
#                    'Markdorf und der westliche Bodenseekreis':"Markdorf/westl. BSK", 'Friedrichshafen / Service':'Friedrich./Service', 
#                    'LokalesImmenstaad / Oberteuringen':'Immens./Obert.', 'Wir am See':"Wir am See", 
#                    'Ravensburg / Weingarten':'Ravensburg/Weing.', 'Ravensburg':"Ravensburg", "Weingarten":"Weingarten", 
#                    'Oberschwaben & Allgäu': 'Oberschw. & Allgäu', 'Langenargen / Eriskirch / Kressbronn':'Lang./Eris./Kressbr.', 
#                    'Immenstaad / Oberteuringen / Meckenbeuren': 'Immenst./Obert./Meck.', 'Die Seite mit der Maus':"Seite mit der Maus", 
#                    "Extraseiten":"Extraseiten", 'Rottum-Iller / Kreis & Region': r'Rottum-Iller/K+R',
#                    'Langenargen / Eriskirch':"Langenargen/Eris.","Oberteuringen /Immenstaad/Meckenbeuren":"Ober./Imm./Mecken.", 
#                     'Friedrichshafen / Tettnang':"Fried./Tettnang","Garten": "Garten", "Kreismusikfest":"Kreismusikfest",   
#                     'Aktion Engagieren & kassieren':'Aktion Engagieren & kassieren', "Wissen": "Wissen", 
#                     'Biberach / Rund um Biberach':'Rund um Biberach', "Panorama":"Panorama",
#                     'Immenstaad / Meckenbeuren':'Imm. / Meckenb.', "Ernährung":"Ernährung", 
#                     'Meine Heimat. Mein Verein. Extraseiten':"Heimat/Vereine", "Bauen & Wohnen": "Bauen & Wohnen", 
#                     'Langenargen/Kressbronn':'Langenargen/Kress.', "Gemeinden":"Gemeinden",
#                     "Umland":"Umland", "Lokalsport":"Lokalsport", 'Kultur Lokal/Service':'Kultur Lokal/Service', 
#                     'rund um Biberach':'Rund um Biberach', 'Kultur Lokal':'Kultur Lokal', 'Immenstaad':'Immenstaad', 
#                     'Reise & Erholung':'Reise & Erholung', "Auto & Verkehr":"Auto & Verkehr", 'Oberteuringen / Immenstaad':'Oberteuringen / Immenstaad', 
#                     "Vermischtes":"Vermischtes", "Medien":"Medien", "Multimedia":"Multimedia", "Sternenhimmel":"Sternenhimmel",  
#                     "Kino":"Kino", "Fernsehen & Freizeit":"Fernsehen/Freizeit", "Geld & Service":"Geld & Service",
#                     "Fernsehen": "Fernsehen", "Leitartikel":"Leitartikel", "Familie":"Familie", "Von der Schussen zur Umlach":"Schussen z. Umlach", 
#                     'Langenargen / Kressbronn':'Langenargen / Kressbronn', 
#                     "Markdorf / Oberteuringen / Immenstaad / Salem":"Mark./Obert./Imm./Sal.", 
#                     "Oberteuringen / Meckenbeuren":"Oberteuringen / Meckenbeuren", "Langenargen":"Langenargen", 
#                     "Regionalsport / Lokalsport":"Reg./Lokalsport", "Regionalsport":"Regionalsport", "Szene":"Szene", 
#                     "Mode":"Mode", "Literatur": "Literatur", "Tiere":"Tiere", "Fernsehen":"Fernsehen", 
#                     "Lokale Eins BIB": "Lokale Eins BIB","Lokale Eins FHA": "Lokale Eins FHA", "Lokale Eins RV": "Lokale Eins RV", 
#                     "Wochenende": "Wochenende", 'Menschen':"Menschen", 'Lebensart':"Lebensart", 'Unterhaltung':"Unterhaltung",
#                     "Szene am Wochenende": "Szene am WE", 'Meine Seite':"Meine Seite" }
                    
orte_region=[]
orte_deutschland =[ 'Aach', 'Aachen', 'Aalen', 'Abenberg', 'Abensberg', 'Achern', 'Achim', 'Adelsheim', 'Adenau',
                     'Adorf/Vogtl.', 'Ahaus', 'Ahlen', 'Ahrensburg', 'Aichach', 'Aichtal', 'Aken (Elbe)', 'Albstadt',
                     'Alfeld (Leine)', 'Allendorf (Lumda)', 'Allstedt', 'Alpirsbach', 'Alsdorf', 'Alsfeld',
                     'Alsleben (Saale)', 'Altdorf bei Nürnberg', 'Altena', 'Altenberg', 'Altenburg', 
                     'Altenkirchen (Westerwald)', 'Altensteig', 'Altentreptow', 'Altlandsberg', 'Altötting', 'Alzenau', 
                     'Alzey', 'Amberg', 'Amöneburg', 'Amorbach', 'Andernach', 'Angermünde', 'Anklam',
                     'Annaberg-Buchholz', 'Annaburg', 'Annweiler am Trifels', 'Ansbach', 'Apolda', 'Arendsee (Altmark)',
                     'Arneburg', 'Arnis', 'Arnsberg', 'Arnstadt', 'Arnstein', 'Arnstein', 'Artern/Unstrut', 'Arzberg', 
                     'Aschaffenburg', 'Aschersleben', 'Asperg', 'Aßlar', 'Attendorn', 'Aub', 'Aue', 'Auerbach in der Oberpfalz',
                     'Auerbach/Vogtl.', 'Augsburg', 'Augustusburg', 'Aulendorf', 'Auma-Weidatal', 'Aurich', 'Babenhausen',
                     'Bacharach', 'Backnang', 'Bad Aibling', 'Bad Arolsen', 'Bad Belzig', 'Bad Bentheim', 'Bad Bergzabern', 
                     'Bad Berka', 'Bad Berleburg', 'Bad Berneck im Fichtelgebirge', 'Bad Bevensen', 'Bad Bibra', 
                     'Bad Blankenburg', 'Bad Bramstedt', 'Bad Breisig', 'Bad Brückenau', 'Bad Buchau', 'Bad Camberg', 
                     'Bad Colberg-Heldburg', 'Bad Doberan', 'Bad Driburg', 'Bad Düben', 'Bad Dürkheim', 'Bad Dürrenberg',
                     'Bad Dürrheim', 'Bad Elster', 'Bad Ems', 'Baden-Baden', 'Bad Fallingbostel', 'Bad Frankenhausen/Kyffhäuser',
                     'Bad Freienwalde (Oder)', 'Bad Friedrichshall', 'Bad Gandersheim', 'Bad Gottleuba-Berggießhübel',
                     'Bad Griesbach im Rottal', 'Bad Harzburg', 'Bad Herrenalb', 'Bad Hersfeld', 'Bad Homburg vor der Höhe', 
                     'Bad Honnef', 'Bad Hönningen', 'Bad Iburg', 'Bad Karlshafen', 'Bad Kissingen', 'Bad König', 
                     'Bad Königshofen im Grabfeld', 'Bad Köstritz', 'Bad Kötzting', 'Bad Kreuznach', 'Bad Krozingen', 
                     'Bad Laasphe', 'Bad Langensalza', 'Bad Lauchstädt', 'Bad Lausick', 'Bad Lauterberg im Harz', 
                     'Bad Liebenstein', 'Bad Liebenwerda', 'Bad Liebenzell', 'Bad Lippspringe', 'Bad Lobenstein', 
                     'Bad Marienberg (Westerwald)', 'Bad Mergentheim', 'Bad Münder am Deister', 'Bad Münstereifel', 
                     'Bad Muskau', 'Bad Nauheim', 'Bad Nenndorf', 'Bad Neuenahr-Ahrweiler', 'Bad Neustadt an der Saale', 
                     'Bad Oeynhausen', 'Bad Oldesloe', 'Bad Orb', 'Bad Pyrmont', 'Bad Rappenau', 'Bad Reichenhall', 
                     'Bad Rodach', 'Bad Sachsa', 'Bad Säckingen', 'Bad Salzdetfurth', 'Bad Salzuflen', 'Bad Salzungen',
                     'Bad Saulgau', 'Bad Schandau', 'Bad Schmiedeberg', 'Bad Schussenried', 'Bad Schwalbach', 'Bad Schwartau',
                     'Bad Segeberg', 'Bad Sobernheim', 'Bad Soden am Taunus', 'Bad Soden-Salmünster', 'Bad Sooden-Allendorf', 
                     'Bad Staffelstein', 'Bad Sulza', 'Bad Sülze', 'Bad Teinach-Zavelstein', 'Bad Tennstedt', 'Bad Tölz', 
                     'Bad Urach', 'Bad Vilbel', 'Bad Waldsee', 'Bad Wildbad', 'Bad Wildungen', 'Bad Wilsnack', 'Bad Wimpfen', 
                     'Bad Windsheim', 'Bad Wörishofen', 'Bad Wünnenberg', 'Bad Wurzach', 'Baesweiler', 'Baiersdorf', 'Balingen', 
                     'Ballenstedt', 'Balve', 'Bamberg', 'Barby', 'Bargteheide', 'Barmstedt', 'Bärnau', 'Barntrup', 'Barsinghausen',
                     'Barth', 'Baruth/Mark', 'Bassum', 'Battenberg (Eder)', 'Baumholder', 'Baunach', 'Baunatal', 'Bautzen', 
                     'Bayreuth', 'Bebra', 'Beckum', 'Bedburg', 'Beelitz', 'Beerfelden', 'Beeskow', 'Beilngries', 'Beilstein', 
                     'Belgern-Schildau', 'Bendorf', 'Bensheim', 'Berching', 'Berga/Elster', 'Bergen', 'Bergen auf Rügen', 
                     'Bergheim', 'Bergisch Gladbach', 'Bergkamen', 'Bergneustadt', 'Berka/Werra', 'Berlin', 'Bernau bei Berlin', 
                     'Bernburg (Saale)', 'Bernkastel-Kues', 'Bernsdorf', 'Bernstadt a. d. Eigen', 'Bersenbrück', 'Besigheim', 
                     'Betzdorf', 'Betzenstein', 'Beverungen', 'Bexbach', 'Biberach an der Riß', 'Biedenkopf', 'Bielefeld',
                     'Biesenthal', 'Bietigheim-Bissingen', 'Billerbeck', 'Bingen am Rhein', 'Birkenfeld', 
                     'Bischofsheim an der Rhön', 'Bischofswerda', 'Bismark (Altmark)', 'Bitburg', 'Bitterfeld-Wolfen', 
                     'Blankenburg (Harz)', 'Blankenhain', 'Blaubeuren', 'Blaustein', 'Bleckede', 'Bleicherode', 'Blieskastel', 
                     'Blomberg', 'Blumberg', 'Bobingen', 'Böblingen', 'Bocholt', 'Bochum', 'Bockenem', 'Bodenwerder', 'Bogen',
                     'Böhlen', 'Boizenburg/Elbe', 'Bonn', 'Bonndorf im Schwarzwald', 'Bönnigheim', 'Bopfingen', 'Boppard',
                     'Borgentreich', 'Borgholzhausen', 'Borken', 'Borken', 'Borkum', 'Borna', 'Bornheim', 'Bottrop', 'Boxberg',
                     'Brackenheim', 'Brake (Unterweser)', 'Brakel', 'Bramsche', 'Brandenburg an der Havel', 'Brand-Erbisdorf',
                     'Brandis', 'Braubach', 'Braunfels', 'Braunlage', 'Bräunlingen', 'Braunsbedra', 'Braunschweig', 'Breckerfeld',
                     'Bredstedt', 'Breisach am Rhein', 'Bremen', 'Bremerhaven', 'Bremervörde', 'Bretten', 'Breuberg', 'Brilon', 
                     'Brotterode-Trusetal', 'Bruchköbel', 'Bruchsal', 'Brück', 'Brüel', 'Brühl', 'Brunsbüttel', 'Brüssow', 
                     'Buchen (Odenwald)', 'Buchholz in der Nordheide', 'Buchloe', 'Bückeburg', 'Buckow (Märkische Schweiz)', 
                     'Büdelsdorf', 'Büdingen', 'Bühl', 'Bünde', 'Büren', 'Burg', 'Burgau', 'Burgbernheim', 'Burgdorf', 'Bürgel', 
                     'Burghausen', 'Burgkunstadt', 'Burglengenfeld', 'Burgstädt', 'Burg Stargard', 'Burgwedel', 'Burladingen', 
                     'Burscheid', 'Bürstadt', 'Buttelstedt', 'Buttstädt', 'Butzbach', 'Bützow', 'Buxtehude', 'Calau',
                     'Calbe (Saale)', 'Calw', 'Castrop-Rauxel', 'Celle', 'Cham', 'Chemnitz', 'Clausthal-Zellerfeld', 'Clingen', 
                     'Cloppenburg', 'Coburg', 'Cochem', 'Coesfeld', 'Colditz', 'Coswig', 'Coswig (Anhalt)', 'Cottbus', 'Crailsheim',
                     'Creglingen', 'Creußen', 'Creuzburg', 'Crimmitschau', 'Crivitz', 'Cuxhaven', 'Dachau', 'Dahlen', 'Dahme/Mark',
                     'Dahn', 'Damme', 'Dannenberg (Elbe)', 'Dargun', 'Darmstadt', 'Dassel', 'Dassow', 'Datteln', 'Daun', 
                     'Deggendorf', 'Deidesheim', 'Delbrück', 'Delitzsch', 'Delmenhorst', 'Demmin', 'Dessau-Roßlau', 'Detmold', 
                     'Dettelbach', 'Dieburg', 'Diemelstadt', 'Diepholz', 'Dierdorf', 'Dietenheim', 'Dietfurt an der Altmühl', 
                     'Dietzenbach', 'Diez', 'Dillenburg', 'Dillingen an der Donau', 'Dillingen/Saar', 'Dingelstädt', 'Dingolfing',
                     'Dinkelsbühl', 'Dinklage', 'Dinslaken', 'Dippoldiswalde', 'Dissen am Teutoburger Wald', 'Ditzingen', 'Döbeln',
                     'Doberlug-Kirchhain', 'Döbern', 'Dohna', 'Dömitz', 'Dommitzsch', 'Donaueschingen', 'Donauwörth', 'Donzdorf',
                     'Dorfen', 'Dormagen', 'Dornburg-Camburg', 'Dornhan', 'Dornstetten', 'Dorsten', 'Dortmund', 'Dransfeld',
                     'Drebkau', 'Dreieich', 'Drensteinfurt', 'Dresden', 'Drolshagen', 'Duderstadt', 'Duisburg', 'Dülmen',
                     'Düren', 'Düsseldorf', 'Ebeleben', 'Eberbach', 'Ebermannstadt', 'Ebern', 'Ebersbach an der Fils', 
                     'Ebersbach-Neugersdorf', 'Ebersberg', 'Eberswalde', 'Eckartsberga', 'Eckernförde', 'Edenkoben', 'Egeln',
                     'Eggenfelden', 'Eggesin', 'Ehingen (Donau)', 'Ehrenfriedersdorf', 'Eibelstadt', 'Eibenstock', 'Eichstätt', 
                     'Eilenburg', 'Einbeck', 'Eisenach', 'Eisenberg', 'Eisenberg (Pfalz)', 'Eisenhüttenstadt', 'Eisfeld', 'Eisleben',
                     'Eislingen/Fils', 'Ellingen', 'Ellrich', 'Ellwangen (Jagst)', 'Elmshorn', 'Elsdorf', 'Elsfleth', 'Elsterberg', 
                     'Elsterwerda', 'Elstra', 'Elterlein', 'Eltmann', 'Eltville am Rhein', 'Elzach', 'Elze', 'Emden', 'Emmelshausen',
                     'Emmendingen', 'Emmerich am Rhein', 'Emsdetten', 'Endingen am Kaiserstuhl', 'Engen', 'Enger', 'Ennepetal', 
                     'Ennigerloh', 'Eppelheim', 'Eppingen', 'Eppstein', 'Erbach', 'Erbach (Odenwald)', 'Erbendorf', 'Erding', 
                     'Erftstadt', 'Erfurt', 'Erkelenz', 'Erkner', 'Erkrath', 'Erlangen', 'Erlenbach am Main', 'Erlensee', 'Erwitte',
                     'Eschborn', 'Eschenbach in der Oberpfalz', 'Eschershausen', 'Eschwege', 'Eschweiler', 'Esens', 'Espelkamp',
                     'Essen', 'Esslingen am Neckar', 'Ettenheim', 'Ettlingen', 'Euskirchen', 'Eutin', 'Falkenberg/Elster', 
                     'Falkensee', 'Falkenstein/Harz', 'Falkenstein/Vogtl.', 'Fehmarn', 'Fellbach', 'Felsberg', 'Feuchtwangen', 
                     'Filderstadt', 'Finsterwalde', 'Fladungen', 'Flensburg', 'Flöha', 'Flörsheim am Main', 'Florstadt', 'Forchheim',
                     'Forchtenberg', 'Forst (Lausitz)', 'Frankenau', 'Frankenberg (Eder)', 'Frankenberg/Sa.', 'Frankenthal (Pfalz)', 
                     'Frankfurt am Main', 'Frankfurt (Oder)', 'Franzburg', 'Frauenstein', 'Frechen', 'Freiberg am Neckar', 'Freiberg',
                     'Freiburg im Breisgau', 'Freilassing', 'Freinsheim', 'Freising', 'Freital', 'Freren', 'Freudenberg', 
                     'Freudenberg', 'Freudenstadt', 'Freyburg (Unstrut)', 'Freystadt', 'Freyung', 'Fridingen an der Donau', 
                     'Friedberg', 'Friedberg', 'Friedland', 'Friedland', 'Friedrichroda', 'Friedrichsdorf', 'Friedrichshafen',
                     'Friedrichstadt', 'Friedrichsthal', 'Friesack', 'Friesoythe', 'Fritzlar', 'Frohburg', 'Fröndenberg/Ruhr', 
                     'Fulda', 'Fürstenau', 'Fürstenberg/Havel', 'Fürstenfeldbruck', 'Fürstenwalde/Spree', 'Fürth', 'Furth im Wald', 
                     'Furtwangen im Schwarzwald', 'Füssen', 'Gadebusch', 'Gaggenau', 'Gaildorf', 'Gammertingen', 'Garbsen', 
                     'Garching bei München', 'Gardelegen', 'Garding', 'Gartz (Oder)', 'Garz/Rügen', 'Gau-Algesheim', 'Gebesee', 
                   'Gedern', 'Geesthacht', 'Geestland', 'Gefell', 'Gefrees', 'Gehrden', 'Gehren', 'Geilenkirchen', 'Geisa', 
                   'Geiselhöring', 'Geisenfeld', 'Geisenheim', 'Geisingen', 'Geislingen', 'Geislingen an der Steige', 'Geithain', 
                   'Geldern', 'Gelnhausen', 'Gelsenkirchen', 'Gemünden am Main', 'Gemünden (Wohra)', 'Gengenbach', 'Genthin', 
                   'Georgsmarienhütte', 'Gera', 'Gerabronn', 'Gerbstedt', 'Geretsried', 'Geringswalde', 'Gerlingen', 'Germering',
                   'Germersheim', 'Gernsbach', 'Gernsheim', 'Gerolstein', 'Gerolzhofen', 'Gersfeld (Rhön)', 'Gersthofen', 'Gescher',
                   'Geseke', 'Gevelsberg', 'Geyer', 'Giengen an der Brenz', 'Gießen', 'Gifhorn', 'Ginsheim-Gustavsburg', 'Gladbeck',
                   'Gladenbach', 'Glashütte', 'Glauchau', 'Glinde', 'Glücksburg (Ostsee)', 'Glückstadt', 'Gnoien', 'Goch', 'Goldberg',
                   'Goldkronach', 'Golßen', 'Gommern', 'Göppingen', 'Görlitz', 'Goslar', 'Gößnitz', 'Gotha', 'Göttingen', 'Grabow',
                   'Grafenau', 'Gräfenberg', 'Gräfenhainichen', 'Gräfenthal', 'Grafenwöhr', 'Grafing bei München', 'Gransee', 'Grebenau',
                   'Grebenstein', 'Greding', 'Greifswald', 'Greiz', 'Greußen', 'Greven', 'Grevenbroich', 'Grevesmühlen', 'Griesheim', 
                   'Grimma', 'Grimmen', 'Gröditz', 'Groitzsch', 'Gronau (Leine)', 'Gronau (Westf.)', 'Gröningen', 'Großalmerode', 
                   'Groß-Bieberau', 'Großbottwar', 'Großbreitenbach', 'Großenehrich', 'Großenhain', 'Groß-Gerau', 'Großräschen', 
                   'Großröhrsdorf', 'Großschirma', 'Groß-Umstadt', 'Grünberg', 'Grünhain-Beierfeld', 'Grünsfeld', 'Grünstadt',
                   'Guben', 'Gudensberg', 'Güglingen', 'Gummersbach', 'Gundelfingen an der Donau', 'Gundelsheim', 'Günzburg', 
                   'Gunzenhausen', 'Güsten', 'Güstrow', 'Gütersloh', 'Gützkow', 'Haan', 'Hachenburg', 'Hadamar', 'Hagen', 'Hagenbach',
                   'Hagenow', 'Haiger', 'Haigerloch', 'Hainichen', 'Haiterbach', 'Halberstadt', 'Haldensleben', 'Halle (Saale)', 
                   'Halle (Westf.)', 'Hallenberg', 'Hallstadt', 'Haltern am See', 'Halver', 'Hamburg', 'Hameln', 'Hamm', 'Hammelburg',
                   'Hamminkeln', 'Hanau', 'Hannover', 'Hann. Münden', 'Harburg (Schwaben)', 'Hardegsen', 'Haren (Ems)', 
                   'Harsewinkel', 'Hartenstein', 'Hartha', 'Harzgerode', 'Haselünne', 'Haslach im Kinzigtal', 'Haßfurt', 
                   'Hattersheim am Main', 'Hattingen', 'Hatzfeld (Eder)', 'Hausach', 'Hauzenberg', 'Havelberg', 'Havelsee',
                   'Hayingen', 'Hechingen', 'Hecklingen', 'Heide', 'Heideck', 'Heidelberg', 'Heidenau', 'Heidenheim an der Brenz', 
                   'Heilbad Heiligenstadt', 'Heilbronn', 'Heiligenhafen', 'Heiligenhaus', 'Heilsbronn', 'Heimbach', 'Heimsheim',
                   'Heinsberg', 'Heitersheim', 'Heldrungen', 'Helmbrechts', 'Helmstedt', 'Hemau', 'Hemer', 'Hemmingen', 'Hemmoor',
                   'Hemsbach', 'Hennef (Sieg)', 'Hennigsdorf', 'Heppenheim (Bergstraße)', 'Herbolzheim', 'Herborn', 'Herbrechtingen', 
                   'Herbstein', 'Herdecke', 'Herdorf', 'Herford', 'Heringen/Helme', 'Heringen (Werra)', 'Hermeskeil', 
                   'Hermsdorf', 'Herne', 'Herrenberg', 'Herrieden', 'Herrnhut', 'Hersbruck', 'Herten', 'Herzberg am Harz',
                   'Herzberg (Elster)', 'Herzogenaurach', 'Herzogenrath', 'Hessisch Lichtenau', 'Hessisch Oldendorf', 
                   'Hettingen', 'Hettstedt', 'Heubach', 'Heusenstamm', 'Hilchenbach', 'Hildburghausen', 'Hilden', 
                   'Hildesheim', 'Hillesheim', 'Hilpoltstein', 'Hirschau', 'Hirschberg', 'Hirschhorn (Neckar)', 'Hitzacker (Elbe)', 
                   'Hochheim am Main', 'Höchstadt an der Aisch', 'Höchstädt an der Donau', 'Hockenheim', 'Hof', 'Hofgeismar', 
                   'Hofheim am Taunus', 'Hofheim in Unterfranken', 'Hohenberg an der Eger', 'Hohenleuben', 'Hohenmölsen', 
                   'Hohen Neuendorf', 'Hohenstein-Ernstthal', 'Hohnstein', 'Höhr-Grenzhausen', 'Hollfeld', 'Holzgerlingen', 
                   'Holzminden', 'Homberg (Efze)', 'Homberg (Ohm)', 'Homburg', 'Horb am Neckar', 'Hornbach', 'Horn-Bad Meinberg', 
                   'Hornberg', 'Hörstel', 'Horstmar', 'Höxter', 'Hoya', 'Hoyerswerda', 'Hückelhoven', 'Hückeswagen', 'Hüfingen',
                   'Hünfeld', 'Hungen', 'Hürth', 'Husum', 'Ibbenbüren', 'Ichenhausen', 'Idar-Oberstein', 'Idstein', 'Illertissen',
                   'Ilmenau', 'Ilsenburg (Harz)', 'Ilshofen', 'Immenhausen', 'Immenstadt im Allgäu', 'Ingelfingen', 'Ingelheim am Rhein',
                   'Ingolstadt', 'Iphofen', 'Iserlohn', 'Isny im Allgäu', 'Isselburg', 'Itzehoe', 'Jarmen', 'Jena', 'Jerichow', 
                   'Jessen (Elster)', 'Jever', 'Joachimsthal', 'Johanngeorgenstadt', 'Jöhstadt', 'Jülich', 'Jüterbog', 'Kaarst',
                   'Kahla', 'Kaisersesch', 'Kaiserslautern', 'Kalbe (Milde)', 'Kalkar', 'Kaltenkirchen', 'Kaltennordheim', 
                   'Kamen', 'Kamenz', 'Kamp-Lintfort', 'Kandel', 'Kandern', 'Kappeln', 'Karben', 'Karlsruhe', 'Karlstadt', 
                   'Kassel', 'Kastellaun', 'Katzenelnbogen', 'Kaub', 'Kaufbeuren', 'Kehl', 'Kelbra (Kyffhäuser)', 
                   'Kelheim', 'Kelkheim (Taunus)', 'Kellinghusen', 'Kelsterbach', 'Kemberg', 'Kemnath', 'Kempen', 
                   'Kempten (Allgäu)', 'Kenzingen', 'Kerpen', 'Ketzin/Havel', 'Kevelaer', 'Kiel', 'Kierspe', 'Kindelbrück',
                   'Kirchberg', 'Kirchberg an der Jagst', 'Kirchberg (Hunsrück)', 'Kirchen (Sieg)', 'Kirchenlamitz', 
                   'Kirchhain', 'Kirchheimbolanden', 'Kirchheim unter Teck', 'Kirn', 'Kirtorf', 'Kitzingen', 'Kitzscher',
                   'Kleve', 'Klingenberg am Main', 'Klingenthal', 'Klötze', 'Klütz', 'Knittlingen', 'Koblenz', 'Kohren-Sahlis',
                   'Kolbermoor', 'Kölleda', 'Köln', 'Königsberg in Bayern', 'Königsbrück', 'Königsbrunn', 'Königsee-Rottenbach',
                   'Königslutter am Elm', 'Königstein im Taunus', 'Königstein (Sächsische Schweiz)', 'Königswinter', 
                   'Königs Wusterhausen', 'Könnern', 'Konstanz', 'Konz', 'Korbach', 'Korntal-Münchingen', 'Kornwestheim',
                   'Korschenbroich', 'Köthen (Anhalt)', 'Kraichtal', 'Krakow am See', 'Kranichfeld', 'Krautheim', 
                   'Krefeld', 'Kremmen', 'Krempe', 'Kreuztal', 'Kronach', 'Kronberg im Taunus', 'Kröpelin', 'Kroppenstedt',
                   'Krumbach (Schwaben)', 'Kühlungsborn', 'Kulmbach', 'Külsheim', 'Künzelsau', 'Kupferberg', 'Kuppenheim',
                   'Kusel', 'Kyllburg', 'Kyritz', 'Laage', 'Laatzen', 'Ladenburg', 'Lage', 'Lahnstein', 'Lahr/Schwarzwald',
                   'Laichingen', 'Lambrecht (Pfalz)', 'Lampertheim', 'Landau an der Isar', 'Landau in der Pfalz', 
                   'Landsberg am Lech', 'Landsberg', 'Landshut', 'Landstuhl', 'Langelsheim', 'Langen', 'Langenau', 
                   'Langenburg', 'Langenfeld (Rheinland)', 'Langenhagen', 'Langenselbold', 'Langenzenn', 'Langewiesen', 'Lassan', 
                   'Laubach', 'Laucha an der Unstrut', 'Lauchhammer', 'Lauchheim', 'Lauda-Königshofen', 'Lauenburg/Elbe', 
                   'Lauf an der Pegnitz', 'Laufen', 'Laufenburg (Baden)', 'Lauffen am Neckar', 'Lauingen (Donau)', 'Laupheim', 
                   'Lauscha', 'Lauta', 'Lauter-Bernsbach', 'Lauterbach', 'Lauterecken', 'Lauterstein', 'Lebach', 'Lebus', 
                   'Leer (Ostfriesland)', 'Lehesten', 'Lehrte', 'Leichlingen (Rheinland)', 'Leimen', 'Leinefelde-Worbis', 
                   'Leinfelden-Echterdingen', 'Leipheim', 'Leipzig', 'Leisnig', 'Lemgo', 'Lengenfeld', 'Lengerich', 'Lennestadt',
                   'Lenzen', 'Leonberg', 'Leun', 'Leuna', 'Leutenberg', 'Leutershausen', 'Leutkirch im Allgäu', 'Leverkusen', 
                   'Lich', 'Lichtenau', 'Lichtenau', 'Lichtenberg', 'Lichtenfels', 'Lichtenfels', 'Lichtenstein/Sa.', 'Liebenau', 
                   'Liebenwalde', 'Lieberose', 'Liebstadt', 'Limbach-Oberfrohna', 'Limburg an der Lahn', 'Lindau (Bodensee)', 
                   'Linden', 'Lindenberg im Allgäu', 'Lindenfels', 'Lindow (Mark)', 'Lingen (Ems)', 'Linnich', 'Linz am Rhein',
                   'Lippstadt', 'Löbau', 'Löffingen', 'Lohmar', 'Lohne (Oldenburg)', 'Löhne', 'Lohr am Main', 'Loitz', 'Lollar', 
                   'Lommatzsch', 'Löningen', 'Lorch', 'Lorch', 'Lörrach', 'Lorsch', 'Lößnitz', 'Löwenstein', 'Lotte', 'Lübbecke', 
                   'Lübben (Spreewald)', 'Lübbenau/Spreewald', 'Lübeck', 'Lübtheen', 'Lübz', 'Lüchow (Wendland)', 'Lucka', 
                   'Luckau', 'Luckenwalde', 'Lüdenscheid', 'Lüdinghausen', 'Ludwigsburg', 'Ludwigsfelde', 'Ludwigshafen am Rhein', 
                   'Ludwigslust', 'Ludwigsstadt', 'Lugau', 'Lügde', 'Lüneburg', 'Lünen', 'Lunzenau', 'Lütjenburg', 'Lützen', 
                   'Lychen', 'Magdala', 'Magdeburg', 'Mahlberg', 'Mainbernheim', 'Mainburg', 'Maintal', 'Mainz', 'Malchin', 
                   'Malchow', 'Mannheim', 'Manderscheid', 'Mansfeld', 'Marbach am Neckar', 'Marburg', 'Marienberg', 'Marienmünster',
                   'Markdorf', 'Markgröningen', 'Märkisch Buchholz', 'Markkleeberg', 'Markneukirchen', 'Markranstädt', 
                   'Marktbreit', 'Marktheidenfeld', 'Marktleuthen', 'Marktoberdorf', 'Marktredwitz', 'Marktsteft', 'Marl', 
                   'Marlow', 'Marne', 'Marsberg', 'Maulbronn', 'Maxhütte-Haidhof', 'Mayen', 'Mechernich', 'Meckenheim', 'Medebach', 
                   'Meerane', 'Meerbusch', 'Meersburg', 'Meinerzhagen', 'Meiningen', 'Meisenheim', 'Meißen', 'Meldorf', 'Melle',
                   'Mellrichstadt', 'Melsungen', 'Memmingen', 'Menden (Sauerland)', 'Mendig', 'Mengen', 'Meppen', 'Merkendorf',
                   'Merseburg', 'Merzig', 'Meschede', 'Meßkirch', 'Meßstetten', 'Mettmann', 'Metzingen', 'Meuselwitz', 'Meyenburg', 
                   'Michelstadt', 'Miesbach', 'Miltenberg', 'Mindelheim', 'Minden', 'Mirow', 'Mittenwalde', 'Mitterteich', 'Mittweida',
                   'Möckern', 'Möckmühl', 'Moers', 'Mölln', 'Mönchengladbach', 'Monheim am Rhein', 'Monheim', 'Monschau', 'Montabaur', 
                   'Moosburg an der Isar', 'Mörfelden-Walldorf', 'Moringen', 'Mosbach', 'Mössingen', 'Mücheln (Geiseltal)', 'Mügeln',
                   'Mühlacker', 'Mühlberg/Elbe', 'Mühldorf am Inn', 'Mühlhausen/Thüringen', 'Mühlheim am Main', 'Mühlheim an der Donau',
                   'Mülheim an der Ruhr', 'Mülheim-Kärlich', 'Müllheim', 'Müllrose', 'Münchberg', 'Müncheberg', 'München', 
                   'Münchenbernsdorf', 'Munderkingen', 'Münnerstadt', 'Münsingen', 'Munster', 'Münster', 'Münstermaifeld', 
                   'Münzenberg', 'Murrhardt', 'Nabburg', 'Nagold', 'Naila', 'Nassau', 'Nastätten', 'Nauen', 'Naumburg', 
                   'Naumburg (Saale)', 'Naunhof', 'Nebra (Unstrut)', 'Neckarbischofsheim', 'Neckargemünd', 'Neckarsteinach', 
                   'Neckarsulm', 'Neresheim', 'Netphen', 'Nettetal', 'Netzschkau', 'Neu-Anspach', 'Neubrandenburg', 'Neubukow', 
                   'Neubulach', 'Neuburg an der Donau', 'Neudenau', 'Neuenbürg', 'Neuenburg am Rhein', 'Neuenhaus', 'Neuenrade',
                   'Neuenstadt am Kocher', 'Neuenstein', 'Neuerburg', 'Neuffen', 'Neuhaus am Rennweg', 'Neu-Isenburg', 'Neukalen',
                   'Neukirchen', 'Neukirchen-Vluyn', 'Neukloster', 'Neumark', 'Neumarkt in der Oberpfalz', 'Neumarkt-Sankt Veit', 
                   'Neumünster', 'Neunburg vorm Wald', 'Neunkirchen', 'Neuötting', 'Neuruppin', 'Neusalza-Spremberg', 'Neusäß', 
                   'Neuss', 'Neustadt an der Aisch', 'Neustadt an der Donau', 'Neustadt an der Waldnaab', 'Neustadt am Kulm', 
                   'Neustadt am Rübenberge', 'Neustadt an der Orla', 'Neustadt an der Weinstraße', 'Neustadt bei Coburg', 
                   'Neustadt (Dosse)', 'Neustadt-Glewe', 'Neustadt', 'Neustadt in Holstein', 'Neustadt in Sachsen', 'Neustrelitz',
                   'Neutraubling', 'Neu-Ulm', 'Neuwied', 'Nidda', 'Niddatal', 'Nidderau', 'Nideggen', 'Niebüll', 'Niedenstein', 
                   'Niederkassel', 'Niedernhall', 'Nieder-Olm', 'Niederstetten', 'Niederstotzingen', 'Nieheim', 'Niemegk', 
                   'Nienburg (Saale)', 'Nienburg/Weser', 'Nierstein', 'Niesky', 'Nittenau', 'Norden', 'Nordenham', 'Norderney',
                   'Norderstedt', 'Nordhausen', 'Nordhorn', 'Nördlingen', 'Northeim', 'Nortorf', 'Nossen', 'Nürnberg', 
                   'Nürtingen', 'Oberasbach', 'Oberharz am Brocken', 'Oberhausen', 'Oberhof', 'Oberkirch', 'Oberkochen', 
                   'Oberlungwitz', 'Obermoschel', 'Obernburg am Main', 'Oberndorf am Neckar', 'Obernkirchen', 'Ober-Ramstadt',
                   'Oberriexingen', 'Obertshausen', 'Oberursel (Taunus)', 'Oberviechtach', 'Oberweißbach/Thür. Wald', 'Oberwesel',
                   'Oberwiesenthal', 'Ochsenfurt', 'Ochsenhausen', 'Ochtrup', 'Oderberg', 'Oebisfelde-Weferlingen', 'Oederan',
                   'Oelde', 'Oelsnitz/Erzgeb.', 'Oelsnitz/Vogtl.', 'Oer-Erkenschwick', 'Oerlinghausen', 'Oestrich-Winkel', 
                   'Oettingen in Bayern', 'Offenbach am Main', 'Offenburg', 'Ohrdruf', 'Öhringen', 'Olbernhau', 'Olching', 
                   'Oldenburg (Oldb)', 'Oldenburg in Holstein', 'Olfen', 'Olpe', 'Olsberg', 'Oppenau', 'Oppenheim', 
                   'Oranienbaum-Wörlitz', 'Oranienburg', 'Orlamünde', 'Ornbau', 'Ortenberg', 'Ortrand', 'Oschatz', 
                   'Oschersleben (Bode)', 'Osnabrück', 'Osterburg (Altmark)', 'Osterburken', 'Osterfeld', 'Osterhofen', 
                   'Osterholz-Scharmbeck', 'Osterode am Harz', 'Osterwieck', 'Ostfildern', 'Ostheim vor der Rhön', 'Osthofen', 
                   'Östringen', 'Ostritz', 'Otterberg', 'Otterndorf', 'Ottweiler', 'Overath', 'Owen', 'Paderborn', 'Papenburg',
                   'Pappenheim', 'Parchim', 'Parsberg', 'Pasewalk', 'Passau', 'Pattensen', 'Pausa-Mühltroff', 'Pegau', 'Pegnitz', 
                   'Peine', 'Peitz', 'Penig', 'Penkun', 'Penzberg', 'Penzlin', 'Perleberg', 'Petershagen', 'Pfaffenhofen an der Ilm',
                   'Pfarrkirchen', 'Pforzheim', 'Pfreimd', 'Pfullendorf', 'Pfullingen', 'Pfungstadt', 'Philippsburg', 'Pinneberg',
                   'Pirmasens', 'Pirna', 'Plattling', 'Plau am See', 'Plaue', 'Plauen', 'Plettenberg', 'Pleystein', 'Plochingen',
                   'Plön', 'Pockau-Lengefeld', 'Pocking', 'Pohlheim', 'Polch', 'Porta Westfalica', 'Pößneck', 'Potsdam',
                   'Pottenstein', 'Preetz', 'Premnitz', 'Prenzlau', 'Pressath', 'Preußisch Oldendorf', 'Prichsenstadt', 
                   'Pritzwalk', 'Prüm', 'Puchheim', 'Pulheim', 'Pulsnitz', 'Putbus', 'Putlitz', 'Püttlingen', 'Quakenbrück', 
                   'Quedlinburg', 'Querfurt', 'Quickborn', 'Rabenau', 'Radeberg', 'Radebeul', 'Radeburg', 'Radevormwald', 
                   'Radolfzell am Bodensee', 'Raguhn-Jeßnitz', 'Rahden', 'Rain', 'Ramstein-Miesenbach', 'Ranis', 
                   'Ransbach-Baumbach', 'Rastatt', 'Rastenberg', 'Rathenow', 'Ratingen', 'Ratzeburg', 'Rauenberg', 
                   'Raunheim', 'Rauschenberg', 'Ravensburg', 'Ravenstein', 'Recklinghausen', 'Rees', 'Regen', 
                   'Regensburg', 'Regis-Breitingen', 'Rehau', 'Rehburg-Loccum', 'Rehna', 'Reichelsheim (Wetterau)', 
                   'Reichenbach im Vogtland', 'Reichenbach/O.L.', 'Reinbek', 'Reinfeld (Holstein)', 'Reinheim', 'Remagen',
                   'Remda-Teichel', 'Remscheid', 'Remseck am Neckar', 'Renchen', 'Rendsburg', 'Rennerod', 'Renningen',
                   'Rerik', 'Rethem (Aller)', 'Reutlingen', 'Rheda-Wiedenbrück', 'Rhede', 'Rheinau', 'Rheinbach', 
                   'Rheinberg', 'Rheinböllen', 'Rheine', 'Rheinfelden (Baden)', 'Rheinsberg', 'Rheinstetten', 'Rhens',
                   'Rhinow', 'Ribnitz-Damgarten', 'Richtenberg', 'Riedenburg', 'Riedlingen', 'Riedstadt', 'Rieneck', 
                   'Riesa', 'Rietberg', 'Rinteln', 'Röbel/Müritz', 'Rochlitz', 'Rockenhausen', 'Rodalben', 'Rodenberg',
                   'Rödental', 'Rödermark', 'Rodewisch', 'Rodgau', 'Roding', 'Römhild', 'Romrod', 'Ronneburg', 
                   'Ronnenberg', 'Rosbach vor der Höhe', 'Rosenfeld', 'Rosenheim', 'Rosenthal', 'Rösrath', 'Roßleben',
                   'Roßwein', 'Rostock', 'Rotenburg an der Fulda', 'Rotenburg (Wümme)', 'Roth', 'Rötha',
                   'Röthenbach an der Pegnitz', 'Rothenburg/O.L.', 'Rothenburg ob der Tauber', 'Rothenfels',
                   'Rottenburg am Neckar', 'Rottenburg a.d.Laaber', 'Röttingen', 'Rottweil', 'Rötz', 'Rüdesheim am Rhein',
                   'Rudolstadt', 'Ruhla', 'Ruhland', 'Runkel', 'Rüsselsheim am Main', 'Rutesheim', 'Rüthen', 
                   'Saalburg-Ebersdorf', 'Saalfeld/Saale', 'Saarbrücken', 'Saarburg', 'Saarlouis', 'Sachsenhagen',
                   'Sachsenheim', 'Salzgitter', 'Salzkotten', 'Salzwedel', 'Sandau (Elbe)', 'Sandersdorf-Brehna', 
                   'Sangerhausen', 'Sankt Augustin', 'Sankt Goar', 'Sankt Goarshausen', 'Sarstedt', 'Sassenberg',
                   'Sassnitz', 'Sayda', 'Schalkau', 'Schauenstein', 'Scheer', 'Scheibenberg', 'Scheinfeld', 
                   'Schelklingen', 'Schenefeld', 'Scheßlitz', 'Schieder-Schwalenberg', 'Schifferstadt', 'Schillingsfürst',
                   'Schiltach', 'Schirgiswalde-Kirschau', 'Schkeuditz', 'Schkölen', 'Schleiden', 'Schleiz', 
                   'Schleswig', 'Schlettau', 'Schleusingen', 'Schlieben', 'Schlitz', 'Schloß Holte-Stukenbrock', 
                   'Schlotheim', 'Schlüchtern', 'Schlüsselfeld', 'Schmalkalden', 'Schmallenberg', 'Schmölln', 
                   'Schnackenburg', 'Schnaittenbach', 'Schneeberg', 'Schneverdingen', 'Schömberg', 'Schönau', 
                   'Schönau im Schwarzwald', 'Schönberg', 'Schönebeck (Elbe)', 'Schöneck/Vogtl.', 'Schönewalde', 
                   'Schongau', 'Schöningen', 'Schönsee', 'Schönwald', 'Schopfheim', 'Schöppenstedt', 'Schorndorf', 
                   'Schortens', 'Schotten', 'Schramberg', 'Schraplau', 'Schriesheim', 'Schrobenhausen', 'Schrozberg', 
                   'Schüttorf', 'Schwaan', 'Schwabach', 'Schwäbisch Gmünd', 'Schwäbisch Hall', 'Schwabmünchen', 
                   'Schwaigern', 'Schwalbach am Taunus', 'Schwalmstadt', 'Schwandorf', 'Schwanebeck', 'Schwarzenbach am Wald',
                   'Schwarzenbach an der Saale', 'Schwarzenbek', 'Schwarzenberg/Erzgeb.', 'Schwarzenborn', 
                   'Schwarzheide', 'Schwedt/Oder', 'Schweich', 'Schweinfurt', 'Schwelm', 'Schwentinental', 'Schwerin', 
                   'Schwerte', 'Schwetzingen', 'Sebnitz', 'Seehausen (Altmark)', 'Seeland', 'Seelow', 'Seelze', 'Seesen',
                   'Sehnde', 'Seifhennersdorf', 'Selb', 'Selbitz', 'Seligenstadt', 'Selm', 'Selters (Westerwald)', 'Senden',
                   'Sendenhorst', 'Senftenberg', 'Seßlach', 'Siegburg', 'Siegen', 'Sigmaringen', 'Simbach am Inn',
                   'Simmern/Hunsrück', 'Sindelfingen', 'Singen (Hohentwiel)', 'Sinsheim', 'Sinzig', 'Soest', 'Solingen', 
                   'Solms', 'Soltau', 'Sömmerda', 'Sondershausen', 'Sonneberg', 'Sonnewalde', 'Sonthofen', 'Sontra', 
                   'Spaichingen', 'Spalt', 'Spangenberg', 'Speicher (Eifel)', 'Spenge', 'Speyer', 'Spremberg', 'Springe',
                   'Sprockhövel', 'Stade', 'Stadtallendorf', 'Stadtbergen', 'Stadthagen', 'Stadtilm', 'Stadtlengsfeld',
                   'Stadtlohn', 'Stadtoldendorf', 'Stadtprozelten', 'Stadtroda', 'Stadtsteinach', 'Stadt Wehlen', 
                   'Starnberg', 'Staßfurt', 'Staufen im Breisgau', 'Staufenberg', 'Stavenhagen', 'St. Blasien', 'Stein', 
                   'Steinach', 'Steinau an der Straße', 'Steinbach-Hallenberg', 'Steinbach (Taunus)', 'Steinfurt', 
                   'Steinheim an der Murr', 'Steinheim', 'Stendal', 'Sternberg', 'St. Ingbert', 'St. Georgen im Schwarzwald',
                   'Stockach', 'Stolberg (Rheinland)', 'Stollberg/Erzgeb.', 'Stolpen', 'Storkow (Mark)', 'Stößen',
                   'Straelen', 'Stralsund', 'Strasburg (Uckermark)', 'Straubing', 'Strausberg', 'Strehla', 'Stromberg',
                   'Stühlingen', 'Stutensee', 'Stuttgart', 'St. Wendel', 'Suhl', 'Sulingen', 'Sulz am Neckar', 'Sulzbach/Saar',
                   'Sulzbach-Rosenberg', 'Sulzburg', 'Sundern (Sauerland)', 'Südliches Anhalt', 'Süßen', 'Syke', 
                   'Tambach-Dietharz', 'Tangerhütte', 'Tangermünde', 'Tann (Rhön)', 'Tanna', 'Tauberbischofsheim', 'Taucha',
                   'Taunusstein', 'Tecklenburg', 'Tegernsee', 'Telgte', 'Teltow', 'Templin', 'Tengen', 'Tessin', 
                   'Teterow', 'Tettnang', 'Teublitz', 'Teuchern', 'Teupitz', 'Teuschnitz', 'Thale', 'Thalheim/Erzgeb.', 
                   'Thannhausen', 'Tharandt', 'Themar', 'Thum', 'Tirschenreuth', 'Titisee-Neustadt', 'Tittmoning', 
                   'Todtnau', 'Töging am Inn', 'Tönisvorst', 'Tönning', 'Torgau', 'Torgelow', 'Tornesch', 'Traben-Trarbach',
                   'Traunreut', 'Traunstein', 'Trebbin', 'Trebsen/Mulde', 'Treffurt', 'Trendelburg', 'Treuchtlingen', 
                   'Treuen', 'Treuenbrietzen', 'Triberg im Schwarzwald', 'Tribsees', 'Trier', 'Triptis', 'Trochtelfingen',
                   'Troisdorf', 'Trossingen', 'Trostberg', 'Tübingen', 'Tuttlingen', 'Twistringen', 'Übach-Palenberg',
                   'Überlingen', 'Uebigau-Wahrenbrück', 'Ueckermünde', 'Uelzen', 'Uetersen', 'Uffenheim', 'Uhingen', 
                   'Ulm', 'Ulmen', 'Ulrichstein', 'Ummerstadt', 'Unkel', 'Unna', 'Unterschleißheim', 'Usedom', 
                   'Usingen', 'Uslar', 'Vacha', 'Vaihingen an der Enz', 'Vallendar', 'Varel', 'Vechta', 'Velbert', 
                   'Velburg', 'Velden', 'Vellberg', 'Velen', 'Vellmar', 'Velten', 'Verden (Aller)', 'Veringenstadt',
                   'Verl', 'Versmold', 'Vetschau/Spreewald', 'Viechtach', 'Viernheim', 'Viersen', 'Villingen-Schwenningen',
                   'Vilsbiburg', 'Vilseck', 'Vilshofen an der Donau', 'Visselhövede', 'Vlotho', 
                   'Voerde (Niederrhein)', 'Vogtsburg im Kaiserstuhl', 'Vohburg an der Donau', 'Vohenstrauß',
                   'Vöhrenbach', 'Vöhringen', 'Volkach', 'Völklingen', 'Volkmarsen', 'Vreden',
                   'Wachenheim an der Weinstraße', 'Wächtersbach', 'Wadern', 'Waghäusel', 'Wahlstedt', 'Waiblingen',
                   'Waibstadt', 'Waischenfeld', 'Waldbröl', 'Waldeck', 'Waldenbuch', 'Waldenburg', 'Waldenburg', 
                   'Waldershof', 'Waldheim', 'Waldkappel', 'Waldkirch', 'Waldkirchen', 'Waldkraiburg', 'Waldmünchen', 
                   'Waldsassen', 'Waldshut-Tiengen', 'Walldorf', 'Walldürn', 'Wallenfels', 'Walsrode', 
                   'Waltershausen', 'Waltrop', 'Wanfried', 'Wangen im Allgäu', 'Wanzleben-Börde', 'Warburg',
                   'Waren (Müritz)', 'Warendorf', 'Warin', 'Warstein', 'Wassenberg', 'Wasserburg am Inn', 
                   'Wassertrüdingen', 'Wasungen', 'Wedel', 'Weener', 'Wegberg', 'Wegeleben', 'Wehr', 'Weida', 
                   'Weiden in der Oberpfalz', 'Weikersheim', 'Weil am Rhein', 'Weilburg', 'Weil der Stadt', 
                   'Weilheim an der Teck', 'Weilheim in Oberbayern', 'Weimar', 'Weingarten', 'Weinheim', 'Weinsberg', 
                   'Weinstadt', 'Weismain', 'Weißenberg', 'Weißenburg in Bayern', 'Weißenfels', 'Weißenhorn', 
                   'Weißensee', 'Weißenstadt', 'Weißenthurm', 'Weißwasser/O.L.', 'Weiterstadt', 'Welzheim', 'Welzow',
                   'Wemding', 'Wendlingen am Neckar', 'Werben (Elbe)', 'Werdau', 'Werder (Havel)', 'Werdohl', 
                   'Werl', 'Werlte', 'Wermelskirchen', 'Wernau (Neckar)', 'Werne', 'Werneuchen', 'Wernigerode',
                   'Wertheim', 'Werther (Westf.)', 'Wertingen', 'Wesel', 'Wesenberg', 'Wesselburen', 'Wesseling', 
                   'Westerburg', 'Westerstede', 'Wetter (Ruhr)', 'Wetter', 'Wettin-Löbejün', 'Wetzlar', 'Widdern', 
                   'Wiehe', 'Wiehl', 'Wiesbaden', 'Wiesmoor', 'Wiesensteig', 'Wiesloch', 'Wildau', 'Wildberg',
                   'Wildenfels', 'Wildeshausen', 'Wilhelmshaven', 'Wilkau-Haßlau', 'Willebadessen', 'Willich', 
                   'Wilsdruff', 'Wilster', 'Wilthen', 'Windischeschenbach', 'Windsbach', 'Winnenden', 
                   'Winsen (Luhe)', 'Winterberg', 'Wipperfürth', 'Wirges', 'Wismar', 'Wissen', 'Witten', 
                   'Wittenberg', 'Wittenberge', 'Wittenburg', 'Wittichenau', 'Wittlich', 'Wittingen', 
                   'Wittmund', 'Wittstock/Dosse', 'Witzenhausen', 'Woldegk', 'Wolfach', 'Wolfenbüttel', 
                   'Wolfhagen', 'Wolframs-Eschenbach', 'Wolfratshausen', 'Wolfsburg', 'Wolfstein', 'Wolgast', 
                   'Wolkenstein', 'Wolmirstedt', 'Worms', 'Wörrstadt', 'Wörth am Rhein', 'Wörth an der Donau', 
                   'Wörth am Main', 'Wriezen', 'Wülfrath', 'Wunsiedel', 'Wunstorf', 'Wuppertal', 'Würselen', 
                   'Wurzbach', 'Würzburg', 'Wurzen', 'Wustrow (Wendland)', 'Wyk auf Föhr', 'Xanten', 
                   'Zahna-Elster', 'Zarrentin am Schaalsee', 'Zehdenick', 'Zeil am Main', 'Zeitz', 
                   'Zell am Harmersbach', 'Zell im Wiesental', 'Zell (Mosel)', 'Zella-Mehlis', 'Zerbst/Anhalt', 
                   'Zeulenroda-Triebes', 'Zeven', 'Ziegenrück', 'Zierenberg', 'Ziesar', 'Zirndorf', 'Zittau', 
                   'Zörbig', 'Zossen', 'Zschopau', 'Zülpich', 'Zweibrücken', 'Zwenkau', 'Zwickau', 'Zwiesel', 
                   'Zwingenberg', 'Zwönitz', 
                   "Schleswig-Holstein", "Mecklenburg-Vorpommern", "Hamburg", "Bremen", "Niedersachsen", "NRW", 
                   "Nordrhein-Westfalen", "Hessen", "Rheinland-Pfalz", "Saarland", "Baden-Würrtemberg", "Bayern",
                   "Sachsen", "Brandenburg", "Mitteldeutschland", "Neuenburg", "Freiburg", "Deutschland"]



kolumnen_liste = []
kolumnen_dict = {'Auf gut Schwäbisch':'Auf gut Schwäbisch',
 'Leitartikel':"Leitartikel",
 'Unser Mann am Ball':'Unser Mann am Ball',
 'Knitz':'Knitz',
 '333 notiert':'333 notiert',
 'Joe Bauer in der Stadt':'Joe Bauer',
 'Der feine Unterschied':"Der feine Unterschied",
 'Der Wochenkehrer':'Wochenkehrer',
 'Bundesliga-Kolumne':'Bundesliga-Kolumne',
 'Unten rechts':'Unten rechts',
 'Oskar Beck':'Oskar Beck',
 'Meine Buchtipps':'Buchtipps',
 'Gerhard Raff':'Gerhard Raff',
 'Die "Tatort"-Fledderei':'Tatort-Fledderei',
 'Lesestoff':'Lesestoff',
 'Mein Wochenende':'Mein Wochenende', 
 '"Das Kapital" in Kürze':'Kapital in Kürze'}



ressort_list=[]
#ressort_list = ['Titel', 'Wir im Süden', 'Seite Drei', 'Nachrichten & Hintergrund',
#       'Meinung & Dialog', 'Wirtschaft', 'Journal', 'Kultur',
#       'Ratgeber', 'Panorama', 'Lokales', 'Lokalsport',
#       'Sport', 'Wochenende']               

#mantel_ressorts = ['Titel', 'Wir im Süden', 'Seite Drei', 'Nachr. & Hint.',
       #'Mein. & Dialog', 'Wirtschaft', 'Journal', 'Kultur', 'Fernsehen',
       #'Reise', 'Ratgeber', 'Panorama',
       #'Vermischtes', 'Sport', 'Wochenende']     
mantel_ressorts = ['Titelseite', 'Politik', 'Seite 3',
       'Zwischen Weser und Rhein', 'Wirtschaft', 'Aus aller Welt', 'Sport']
lokale_ressorts = [_lokales_]     
ressort_dict={'Titel':'Titel', 'Wir im Süden':'Wir im Süden', 'Seite Drei':'Seite Drei', 'Nachrichten & Hintergrund':'Nachr. & Hint.',
       'Meinung & Dialog':'Mein. & Dialog' , 'Wirtschaft':'Wirtschaft', 'Journal':'Journal', 'Kultur':'Kultur',
       'Ratgeber':'Ratgeber', 'Panorama':'Panorama', 'Lokales':'Lokales', 'Lokalsport':'Lokalsport',
       'Sport':'Sport', 'Wochenende':'Wochenende', 'Fernsehen': "Fernsehen", 
       'Titelseite':"Titelseite", 'Tagesthema':"Tagesthema", 'Die Dritte Seite':"Die Dritte Seite", 'Politik':"Politik",
       'Aus aller Welt':"Aus aller Welt",'Wirtschaft':"Wirtschaft", 'Kultur': "Kultur", 'Sport':"Sport", 'Entdecken':"Entdecken",
       'Multimed. Reportage': "Multim. Reportage",'Leserforum':"Leserforum", 'Wochenende':"Wochenende", "Die Seite Drei":"Die Seite Drei", 
       "Panorama": "Panorama", "Wissenswert":"Wissenswert", "Unsere Leser und wir":"Unsere Leser + wir", 
       'Lokales Stuttgart STZ':"Lokales Stutt. STZ", 'Lokales/Region STN':"Lokales/Region STN",  
       'Region/Baden-Württ. STZ':"Region/BaWü STZ",  'Stuttgart 5. Buch':"Stuttgart 5. Buch",
                                    'Filderstadt 5. Buch':"Filderstadt 5. Buch",
                                    'Lokales Ludwigsburg':'Lokales Ludwigsburg', 'Lokalsport':"Lokalsport", 
                                     "Landesnachrichten":"Landesnachr."}                    

ausgaben_liste=[]     
ausgaben_dict ={}


seitentitel_list = []

# erst Dax, dann MDAX, dann DOW, dann International
liste_unternehmen = ["Daimler", "Porsche", "Bosch", "Volkswagen", "Deutsche Bank", 
                     "LBBW", "EnBW", "Audi", "BMW", "Facebook",  "BASF", "BMW", 
                     "SAP", "Deutsche Bahn", "Heidelberger Druckmaschinen", "Nissan", 
                     "Adidas", "Allianz", "Bayer", "Beiersdorf", "Continental", 
                     "Covestro", "Deutsche Bank", "Deutsche Börse", "Post", "Telekom", 
                     "E.ON", "Fresenius", "HeidelbergCement", "Henkel", "Infineon", 
                     "Linde", "Lufthansa", "Merck", "Münchener Rück", "RWE", "Siemens", 
                     "Thyssen", "Vonovia", "Wirecard", 
                     "Deutsche Wohnen","Drillisch", "Aareal Bank", "N26", "Airbus", "Aurubis", 
                     "Aroundtown", "alstria", "Axel Springer", "Bechtle", "Brenntag", "Carl Zeiss", 
                     "Commerzbank", "Delivery Hero", "Euroshop", "Dialog Semiconductor", 
                     "Dürr", "Evonik", "Evotec", "Fielmann", "Fraport", "freenet", "Fuchs Petrolub", 
                     "GEA", "Gerresheimer", "Grand City Properties", "Hannover Rück", "Hella", "Hochtief", 
                     "Hugo Boss", "Innogy", "K+S", "Kion", "Knorr-Bremse", "Lanxess", "LEG Immobilien", 
                     "Metro", "MorphoSys", "MTU", "Nemetschek", "Norma Group", "Osram", "pbb", 
                     "ProSiebenSat", "Puma", "Qiagen", "Rheinmetall", "Rocket Internet", "RTL", 
                     "Sartorius", "Scout24", 
                     "3M", "American Express", "Apple", "Boeing", "Caterpillar", "Harley Davidson", 
                     "Chevron", "Cisco", "Coca Cola", "DowDuPont", "Exxon", "Goldman Sachs", "Home Depot", 
                     "IBM", "Intel", "Johnson Johnson", "JPMorgan", "McDonalds", "Merck", "Microsoft", 
                     "Nike", "Pfizer", "Procter Gamble", "Travelers", "United Technologies", "UnitedHealth", 
                     "Verizon", "Visa", "Walgreens", "Walmart", "Disney", 
                     "Sony", "Toyota", "Skoda", "Honda", "Renault", "Peugeot", "Fiat", "General Motors", "Ford", 
                     "Chrysler", "Samsung", "Huawei", "Alibaba"
                     ]

liste_sportarten = ["Frauentennis", "Wasserball", 'American Football','Football', 
                    "Judo",'Leichtathletik',"Ski alpin", "Gewichtheben", \
                   'Reitsport', 'Formel E', 'Fußball','Marathon', 'Badminton', 
                   'Golfsport','Snooker', 'Eishockey', \
                   'Langlauf', 'Radfahren', 'Handball', 'Skilanglauf', 'Tennis',
                   'Biathlon', 'Rudern','Rallye',\
                   'Laufen', 'Skispringen','Motorsport', 'Tischtennis','Radsport',
                   'Eiskunstlauf', 'Springsport', 'Ringen',\
                   'Basketball', 'Golf', 'Dressurreiten', 'Frauenfußball',
                   'Fünfkampf', 'Frauenvolleyball', 'Gewichtheben',\
                   'Rugby', 'Frauenhandball', 'Boxen', 'Volleyball', 'Fechten',
                   'Turnen','Beachvolleyball', "Kanu", 'Beachvolleball',
                   'Bouldern', 'Dressur', 'Formel 1', 'Kraftdreikampf', 
                   'Paralympics', 'Zehnkampf', 'Zeitfahren', 'Rudern']   

liste_vereine = ['TSG 1899 Hoffenheim', 'VfB Stuttgart II', 'FSV Mainz II','Fenerbahçe SK', 'Austria Wien', \
                     'Besiktas Istanbul', 'Viktoria Köln', 'RB Leipzig', 'Real Madrid', 'Olympique Marseille',\
                     'Manchester City', 'SC Freiburg II', 'VfB Stuttgart', 'Stuttgarter Kickers', 'Sampdoria Genua',\
                     'AS Monaco','FC Augsburg', 'Chemnitzer FC', 'Holstein Kiel', 'AFC Bournemouth','1. FC Heidenheim', \
                     'Paris SG', 'SV Elversberg', 'TSV 1860 München','Fortuna Düsseldorf', 'Hertha BSC Berlin',\
                     'VfL Wolfsburg', '1. FC Nürnberg', 'SV Werder Bremen', 'SC Freiburg', 'FC Ingolstadt', 'SV Sandhausen',\
                     'SC Paderborn', 'Manchester United', 'Werder Bremen II','FC Sevilla', 'FC Bayern München','FC Barcelona',\
                     'FC Schalke 04', 'SpVgg Greuther Fürth', 'FC Liverpool', 'Offenbacher Kickers', '1. FC Saarbrücken',\
                     'FC Rot-Weiß Erfurt', 'Juventus Turin', 'KSV Hessen Kassel e.V.', 'Atlético Madrid', \
                     'VfL Bochum', 'Eintracht Braunschweig', '1. FC Magdeburg', 'Eintracht Frankfurt', 'Dinamo Zagreb' \
                     '1. FC Kaiserslautern', 'FC Arsenal', '1. FSV Mainz 05', 'BVB', 'Wormatia Worms', \
                     '1. FC Köln', 'Mönchengladbach', 'Arminia Bielefeld', 'FC Chelsea', 'Hannover 96', 'Hannover', 'HSV',\
                     'VfL Wolfsburg', 'SG Sonnenhof Großaspach', 'SV Röchling Völklingen', 'Hamburger SV',\
                     'Bayer 04 Leverkusen','FSV Frankfurt', "1. FC Kaiserslautern", "FCK", "Viktoria Rebensburg", "FC Bayern München", \
                     "SC Neustadt", "KSV Durlach", "TSG Kaiserslautern", "SV Waldhof Mannheim", "FC Barcelona", "Manchester United",\
                     "Victoria Carl", "FC Speyer", "Wormatia Worms", "Hannover 96", 'Dynamo Dresden' 
              ]
#nur für die nw:  , 
liste_vereine = ['SC Paderborn', 'FC Schalke 04', 'Hannover 96', "Hamburger SV", "BVB", "Bayern München", 
                 "Arminia Bielefeld", "FC Liverpool", "Eintracht Frankfurt"]
zeitung_attribute_dict={'Ihre Zeitung ist unabhängig':"unabhängig", 
               'Ihre Zeitung ist informativ': "informativ",
               'Ihre Zeitung ist niveauvoll':"niveauvoll", 
               'Ihre Zeitung ist interessant':"interessant",
               'Ihre Zeitung ist übersichtlich':"übersichtlich",
               'Ihre Zeitung ist verständlich geschrieben':"verständlich",
               'Ihre Zeitung ist zeitgemäß':"zeitgemäß",
               'Ihre Zeitung glaubwürdig':"glaubwürdig",
               'Ihre Zeitung ist glaubwürdig':"glaubwürdig", 
               'Ihre Zeitung ist nah am Leser':"nah am Leser", 
               'Ihre Zeitung ist unterhaltsam':"unterhaltsam", 
               'Ihre Zeitung ist objektiv':"objektiv"}

zeitung_attribute_dict2 = {'Wie zufrieden sind Sie mit der Zeitung insgesamt?':'Zeitung insgesamt',
       'Wie zufrieden sind Sie mit dem redaktionellen Inhalt Ihrer Zeitung?':'redakt. Inhalt',
       'Wie zufrieden sind Sie mit der lokalen Berichterstattung Ihrer Zeitung?':'lokale Berichte',
       'Wie zufrieden sind Sie mit der Kundenbetreuung (telefonischer Kundenservice, Online-Servicebereich) Ihrer Zeitung?':'Kundenbetreuung',
       'Wie zufrieden sind Sie mit der Zustellung Ihrer Zeitung?':'Zustellung',
       'Wie zufrieden sind Sie mit der Kundenkarte Ihrer Zeitung?':'Kundenkarte'}

zeitung_attribute_dict3 = {
        'Wie häufig nutzen Sie die gedruckte Tageszeitung?':'gedruckte \nZeitung',
       'Wie häufig nutzen Sie das Onlineportal Ihrer Zeitung?':'Online-\nPortal',
       'Wie häufig nutzen Sie die Premium-App/das ePaper Ihrer Zeitung?':'Premium-App/\nE-Paper',
       'Wie häufig nutzen Sie die News-App Ihrer Zeitung?':'News-App'      
        }

#UMFRAGE-LISTEN
# Liste um die Umfrageergebnisse "Welche Themen interessieren Sie..." zu ordnen
# und auf verschiedene Folien zu verteilen
umfrage_sport = ['Berichte über Fußball', 'Berichte über Basketball', 'Berichte über Biathlon',
       'Berichte über Handball','Berichte über Schwimmen oder Wasserspringen',
       'Berichte über Volleyball', 'Berichte über American Football',
       'Berichte über Bergsport', 'Berichte über Eishockey',
       'Berichte über Leichtathletik', 'Berichte über Motorsport',
       'Berichte über Pferdesport', 'Berichte über Radsport',
       'Berichte über Schach', 'Berichte über Tennis', 'Berichte über Boxen' ]

umfrage_sport_dict = {'Berichte über Fußball':"Fußball", 
                      'Berichte über Basketball':"Basketball", 
                      'Berichte über Biathlon':"Biathlon",
                      'Berichte über Handball':"Handball",
                      'Berichte über Schwimmen oder Wasserspringen':"Schwimmen",
                      'Berichte über Volleyball': 'Volleyball',
                      'Berichte über Bergsport':"Bergsport", 
                      'Berichte über Eishockey':"Eishockey",
                      'Berichte über Leichtathletik':"Leichtathletik",
                      'Berichte über Motorsport':"Motorsport",
                      'Berichte über Pferdesport':"Pferdesport", 
                      'Berichte über Radsport':"Radsport",
                      'Berichte über Schach':"Schach",
                      'Berichte über Tennis':"Tennis",
                      'Berichte über Boxen':"Boxen", 
                      "Berichte über American Football": "Am. Football"}

#umfrage_gesell_dict = {"Berichte über Autos":"Autos", 
#                       "Berichte über Banken, Versicherungen und Geld":"Geld", 
#                       "Berichte über Bauprojekte in Ihrem Ort": "Bauprojekte", 
#                       "Berichte über Erziehungsfragen": "Erziehungsfragen", 
#                       "Berichte über Hochschulen und Forschungsinstitute": "Forschung", 
#                       "Berichte über Kirchen": "Kirchen", 
#                       "Berichte über klassische Musik": "Klassische Musik",
#                       "Informationen über Rock, Pop oder Jazz":"Rock/Pop/Jazz", 
#                       "Berichte über Landespolitik": "Landespolitik",
#                       "Berichte über lokale Schulen, Kindertageseinrichtungen":"Schulen/Kitas", 
#                       "Berichte über Messen": "Messen"}

umfrage_gesell_dict = {"Berichte über Autos":"Autos", 
                       "Berichte über Banken, Versicherungen und Geld":"Geld", 
                       "Berichte über Bauprojekte in Ihrem Ort": "Bauprojekte", 
                       "Berichte über Erziehungsfragen": "Erziehungsfragen", 
                       "Berichte über Hochschulen und Forschungsinstitute": "Forschung", 
                       "Berichte über Kirchen": "Kirchen", 
                       "Berichte über klassische Musik": "Klassische Musik",
                       "Informationen über Rock, Pop oder Jazz":"Rock/Pop/Jazz", 
                       "Berichte über Landespolitik": "Landespolitik",
                       "Berichte über lokale Schulen, Kindertageseinrichtungen":"Schulen/Kitas", 
                       "Berichte über Messen": "Messen", 
                         "Berichte über Prominente und Klatsch": "Promis/Klatsch", 
                       "Berichte über Raumfahrt und Astronauten": "Raumfahrt", 
                       "Berichte über Umwelt und Klima": "Umwelt/Klima", 
                       "Berichte über Unglück und Verbrechen": "Blaulicht", 
                       "Berichte über Wissenschaft und Forschung": "Wissenschaft"}

umfrage_ressorts_dict ={"Kulturteil in der Zeitung":"Kulturteil",
                        "Berichte über Innenpolitik": "Innenpolitik",
                        "Lokalteil in der Zeitung": "Lokalteil", 
                        "Politikteil in der Zeitung": "Politik",
                        "Berichte über Innenpolitik": "Innenpolitik",
                        "Sportteil in der Zeitung": "Sportteil", 
                        "Wirtschaftsteil in der Zeitung": "Wirtschaft",
                        "Berichte über Unternehmen, die nicht aus Ihrer Region sind": "nichtlokale Firmen", 
                        "Wochenendbeilage": "Wochenende" 
                        }

umfrage_themen_dict = {"Berichte über Mode":"Mode", 
                       "Berichte über Partnerschaft und Beziehung": "Partnerschaft", 
                       "Berichte über Prominente und Klatsch": "Promis/Klatsch", 
                       "Berichte über Raumfahrt und Astronauten": "Raumfahrt", 
                       "Berichte über Umwelt und Klima": "Umwelt/Klima", 
                       "Berichte über Unglück und Verbrechen": "Blaulicht", 
                       "Berichte über Wissenschaft und Forschung": "Wissenschaft", 
                       "Berichte über Renten": "Rente"}
                       
umfrage_themen_dict2 = {"Buchrezensionen":"Buchrezensionen", 
                        "Fernsehkritik und -tipps":"Fernsehen", 
                        "Humor":"Humor", 
                        "Informationen über Busse und Bahnen in Ihrem Ort": "ÖPNV vor Ort", 
                        "Informationen über einzelne Stadtteile": "Stadtteile", 
                        "Informationen über Essen und Trinken": "Essen + Trinken", 
                        "Informationen über Wellness und Kosmetik": "Wellness", 
                        "Informationen über Computer, Handys, Fernsehen etc.":"Computer/Technik"}

umfrage_themen_dict3 = {"Kino-Kritik": "Kino-Kritik", 
                        "Kreuzworträtsel/Sudoku":"Rätsel/Sudoku", 
                        "Polizei- und Gerichtsberichte aus Ihrem Ort": "Polizei/Gericht", 
                        "Ratgeber und Verbraucherthemen": "Ratgeber", 
                        "Reportagen":"Reportagen",
                        "Theaterrezensionen":"Theater", 
                        "Veranstaltungshinweise":"Veranstaltungen", 
                        "Berichte über Reisen und Tourismus": "Reisen/Tourismus", 
                        "Berichte über Autos und die Automobilbranche": "Autos", 
                        "Fernsehkritik und -tipps":"Fernsehen", 
                        "Humor":"Humor", 
                        "Informationen über Busse und Bahnen in Ihrem Ort": "ÖPNV vor Ort"}   



#umfrage_themen_dict3 = {"Kino-Kritik": "Kino-Kritik", 
#                        "Kreuzworträtsel/Sudoku":"Rätsel/Sudoku", 
#                        "Polizei- und Gerichtsberichte aus Ihrem Ort": "Polizei/Gericht", 
#                        "Ratgeber und Verbraucherthemen": "Ratgeber", 
#                        "Reportagen":"Reportagen",
#                        "Theaterrezensionen":"Theater", 
#                        "Veranstaltungshinweise":"Veranstaltungen", 
#                        "Berichte über Reisen und Tourismus": "Reisen/Tourismus", 
#                        "Berichte über Autos und die Automobilbranche": "Autos"}   

umfrage_themen_dict4 = {"Berichte über Außenpolitik":"Außenpolitik", 
                        "Berichte über Innenpolitik": "Innenpolitik",
                        "Berichte über Kommunalpolitik": "Kommunalpolitik", 
                        "Berichte über Stuttgart 21": "Stuttgart 21", 
                        "Berichte über Luftverschmutzung in Ihrem Ort": "Luftverschmutz.", 
                        "Berichte über Start-up-Unternehmen": "Start-ups", 
  
                        }

umfrage_themen_dict5 = {"Berichte über Freizeitsport": "Freizeitsport", 
                        "Berichte über Lokalsport": "Lokalsport", 
                        "Berichte über Gartenthemen":"Gartenthemen", 
                        "Berichte über Wohnen und Lifestyle": "Wohnen/Lifestyle", 
                        
                        "Berichte über Gesundheit":"Gesundheit", 
                        "Berichte über Ernährung": "Ernährung", 
                        "Berichte über Recht und juristische Themen":"Jura/Recht", 
                        "Berichte über Haustiere": "Haustiere", 
                        "Berichte über Energiepreise": "Energiepreise", 
                        "Berichte über gesellschaftliche Trends":"gesell. Trends"}

umfrage_verzicht = {"Telefontarife":"Telefontarife", "Börsenteil":"Börsenteil",
                    "Fernsehprogramm":"TV-Programm", "Wetter":"Wetter", 
                    "Veranstaltungskalender":"Terminkalender"}

#über diese Liste wird dann iteriert
umfrage_gesamt_liste = [umfrage_sport_dict,umfrage_gesell_dict, umfrage_ressorts_dict, 
                        umfrage_themen_dict,umfrage_themen_dict2, umfrage_themen_dict3, 
                        umfrage_themen_dict4, umfrage_themen_dict5, umfrage_verzicht]

# Zum Testen sollten bei einem neuen Kunden alle Dicts einmal mit der Update-Methode 
# zusammengezogen werden und mit den Inhalten der Kunden-Tabelle abgeglichen werden. 

# Die Fragen für die Pie-Charts
umfrage_piecharts = ["Seit wann haben Sie Ihre Zeitung abonniert?",
                     "Seit wann wohnen Sie am heutigen Wohnort?", 
                     "Zu welcher Tageszeit lesen Sie die Zeitung?",
                     "Wie viele Personen im Haushalt lesen Ihre Zeitung?",
                     
                     "Wie oft nehmen Sie durchschnittlich eine Ausgabe pro Tag in die Hand?", 
                     "Wie lange lesen Sie Ihre Zeitung durchschnittlich pro Tag?", 
                     "Wie viele Seiten der Zeitung lesen Sie durchschnittlich?", 
                     "Wo lesen Sie in der Regel Ihre Zeitung?", 
                     "An wie viele Personen außerhalb Ihres Haushalts geben Sie Ihre Zeitung weiter?", 
                     ]

umfrage_mini_bars_TV = ["Wie lange sehen Sie durchschnittlich an einem Tag fern?", 
                     "Wie oft sehen Sie Politik- und Wirtschaftsmagazine im TV?", 
                     "Wie oft sehen Sie Unterhaltungs-Shows im TV?", 
                     "Wie oft sehen Sie Spielfilme im TV?", 
                     "Welche Nachrichtensendungen sehen Sie sich hauptsächlich im TV an?", 
                     "Wie oft sehen Sie Nachrichten im TV?"]

umfrage_mini_bars_netz = ["Verfügen Sie über einen Internet-Zugang zu Hause?", 
                          "Verfügen Sie über ein internetfähiges Smartphone?", 
                          "Verfügen Sie über einen Internet-Zugang am Arbeitsplatz / Uni?", 
                          "Verfügen Sie über einen Tablet-PC?", 
                          "Wie lange surfen Sie täglich im Internet?",
                          "Wie oft nutzen Sie das Internet?" 
                          ]

umfrage_demografie_mini_bars = ["Sie sind ...", "Wie alt sind Sie?", "Wie groß ist Ihr Haushalt?"]

umfrage_demografie_bars = ["Welchen letzten Schulabschluss haben Sie?", "Sie sind zur Zeit ...", 
                           "Welcher Berufsgruppe ordnen Sie sich zu?"] 
 
# umfrage_pie wird hier mit Sortierung verknüpft                          
sortierung_umfrage = {}
        
        
# Liste mit Kolumnen
kolumnen_list = ['Guten Morgen', 'Leserbriefe', 'Kommentar','Radio Bielefeld', "6 Richtige", "Zum Sonntag", "Mein Tipp", #BI
                 "Campus Aktuell", "Radio Hochstift", "Auf ein Wort, Herr Pfarrer!", "Schon gehört?", # PB inklusive der Rubriken BI (z.B Guten Morgen)
                 "Schön gemacht!", "Notiert", "Tipp des Tages", # GT
                 "Marktplatz", "Kommentar", "Wort der Besinnung", "Persönlich", "Achtung Blitzer", 
                 "Persönlich", "Kurz notiert", 'Kommentar', 'Einwurf',
       'Briefe an die Redaktion', "Von Mittwoch zu Mittwoch"] #LÜB
 
#kolumnen_dict = {'Guten Morgen':"Guten Morgen", 'Leserbriefe':"Leserbriefe", 'Kommentar':"Kommentar",'Radio Bielefeld':"Radio Bielefeld",
#                 "6 Richtige":"6 Richtige", "Zum Sonntag":"Zum Sonntag", "Mein Tipp":"Mein Tipp", #BI
#                 "Campus Aktuell":"Campus Aktuell", "Radio Hochstift":"Radio Hochstift", "Auf ein Wort, Herr Pfarrer!":"Auf ein Wort!",
#                 "Schon gehört?":"Schon gehört?", # PB inklusive der Rubriken BI (z.B Guten Morgen)
#                 "Schön gemacht!":"Schön gemacht!", "Notiert":"Notiert", "Tipp des Tages":"Tipp des Tages", # GT
#                 "Marktplatz":"Marktplatz", "Kommentar":"Kommentar", "Wort der Besinnung":"Wort der Besinn.", "Persönlich":"Persönlich", 
#                 "Achtung Blitzer":"Achtung Blitzer"}

 
# ID-Nr für Downloads der Screenshots
id_nr = "1014" # DRP

ressortliste_stz = ['Titelseite', 'Tagesthema', 'Die Dritte Seite', 'Die Seite Drei', 'Politik',  'Landesnachrichten', "Die Brücke", 
       'Aus aller Welt',  'Panorama', 'Wirtschaft', 'Entdecken', 'Wissenswert', 'Leserforum', 'Unsere Leser und wir', 
                'Lokales Stuttgart STZ', 'Lokales/Region STN',   
                'Lokales Ludwigsburg', 'Region/Baden-Württ. STZ', 'Multimed. Reportage', 'Kultur', 'Sport',  'Stuttgart 5. Buch',
                'Filderstadt 5. Buch', 'Lokalsport','Wochenende']


ressortliste_stn = ['Titelseite', 'Tagesthema', 'Die Dritte Seite', 'Die Seite Drei', 'Politik', 'Landesnachrichten', "Die Brücke", 
       'Aus aller Welt',  'Panorama', 'Wirtschaft', 'Entdecken', 'Wissenswert', 'Leserforum', 'Unsere Leser und wir', 
                'Lokales Stuttgart STZ', 'Lokales/Region STN',   
                'Lokales Ludwigsburg', 'Region/Baden-Württ. STZ', 'Multimed. Reportage', 'Kultur', 'Sport',  'Stuttgart 5. Buch',
                'Lokalsport','Filderstadt 5. Buch', 'Wochenende']

ressortliste_spezial_res_stz = ['Titelseite', 'Tagesthema', 'Die Dritte Seite', 'Politik',
       "Die Brücke", 'Aus aller Welt', 'Wirtschaft', 'Entdecken','Leserforum', 
       'Lokale 1 Stutt. STZ', 'Lokales Stuttgart STZ','Lokale Eins Ludwigsburg', 'Lokales Ludwigsburg',
       'Region/Baden-Württ. STZ', 'Multimed. Reportage','Kultur','Sport', 'Lokale 1 Fild. STZ',
       'Filderstadt 5. Buch', 'Stuttgart 5. Buch', 'Lokalsport', 'Wochenende',
       'Die Seite Drei', 'Landesnachrichten', 'Panorama',
       'Lokale 1 Stutt. STN', 'Lokales/Region STN', 'Wissenswert',
       'Lokale 1 Fild. STN', 
        'Unsere Leser und wir']

ressortliste_spezial_res_stn = ['Titelseite', 'Tagesthema', 'Die Dritte Seite','Die Seite Drei',  'Politik',
       'Aus aller Welt', 'Panorama','Wirtschaft', 'Kultur', 'Entdecken', 'Unsere Leser und wir',
       'Lokale 1 Stutt. STZ', 'Lokales Stuttgart STZ',
       'Lokale 1 Stutt. STN','Lokales/Region STN', 
       'Sport',  'Wissenswert','Lokale 1 Fild. STN', 'Stuttgart 5. Buch',
       'Region/Baden-Württ. STZ',  'Lokale 1 Fild. STZ', 'Filderstadt 5. Buch',
        'Lokalsport', 'Wochenende']

textlänge = ["bis 200", "201-400", "401-600", "601-800", "über 800"]
      
#%% Function setlist() - Listen/Dicts verändern

# Die Funktion ersetzt/erweitert im Bedarfsfall bestehende Listen. 


def setlist(platzierung={}, darstellung={}, darstellungsf=[], 
            seitentitel={}, ressort=[], mantel = [], ressortdict={}, 
            sportarten=[], ausgaben=[], ausgabendict={}, id_nummer="1009", seiten_lokal= [], 
            lokal=[], kolumnen=[], zeitung_attribute ={}, 
            umfrage_pie={}, umfrage_tv=[], umfrage_demografie=[], umfrage_netz=[], 
            sort_umfrage = {}, region=[], deutschland=[], international=[], unternehmen= [], 
            rubriken={}, rub_dict={}, seiten_list=[]):
    
    # checken, ob die Parameter länger als 0 sind = neue Parameter
    # in dem Fall müssen die globalen Variablen verändert werden
    if len(platzierung)>0:
        print("Neue Platzierung_dict")
        global platzierung_dict
        platzierung_dict = platzierung
    
    if len(darstellung)>0:
        print("Neue Darstellungs_dict")
        global darstellung_dict
        darstellung_dict = platzierung_dict

    if len(darstellungsf)>0:
        print("Neue Liste Darstellungsform")
        global darstellungsform 
        darstellungsform = darstellungsf
        
    if len(seitentitel)>0:
        print("Neue Seitentitel gesetzt")
        global seitentitel_dict
        seitentitel_dict = seitentitel
    
    if len(ressort)>0:
        print("Neue Ressortliste gesetzt")
        global ressort_list
        ressort_list=ressort
    
    if len(mantel)>0: 
        print("Neue Mantelressorts gesetzt")
        global mantel_ressorts
        mantel_ressorts = mantel
    
    if len(ressortdict)>0:
        print("Neues Ressortdict gesetzt")
        global ressort_dict
        ressort_dict = ressortdict
    
    if len(sportarten)>0: 
        print("Neue Sportarten hinzugefügt")
        global liste_sportarten
        for elem in sportarten:
            liste_sportarten.append(elem)
    
    if len(ausgaben)>0: 
        print("Neue Ausgaben-Liste hinzugefügt")
        global ausgaben_liste
        
        for elem in ausgaben: 
            ausgaben_liste.append(elem)
    
    if len(ausgabendict)>0:
        global ausgaben_dict
        ausgaben_dict = ausgabendict
        
    if id_nummer !="0": 
        print ("Neue ID für Downloads gesetzt")
        global id_nr
        id_nr = id_nummer
        
    if len(seiten_lokal) >0:
        print("Neue Liste für Lokalseiten gesetzt.")
        global seitentitel_lokal
        seitentitel_lokal = seiten_lokal
    
#    if len(zeitung_attribute) >0: 
#        print("Neue Liste für Zeitungsattribute in Umfrage-Sheets gesetzt.")
#        global zeitung_attribute_dict
#        zeitung_attribute_dict = zeitung_attribute
        
    if len(lokal) >0:
         print("Neue Liste mit lokalen Ressorts gesetzt")
         global lokale_liste
         lokale_liste = lokal
         
    if len(kolumnen) >0:
        print("Neue Liste mit Kolumnen gesetzt")
        global kolumnen_liste
        kolumnen_liste = kolumnen
        
     
    if len(zeitung_attribute)>0:
        print("Neues Dict Zeitungsattribute für Umfrage gesetzt.")
        global zeitung_attribute_dict
        zeitung_attribute_dict = zeitung_attribute
        
    if len(umfrage_pie) > 0:
        print("Neues Dict Umfrage Pie gesetzt.")
        global umfrage_piecharts
        umfrage_piecharts = umfrage_pie
        
        
        
    if len(umfrage_tv) > 0: 
        print("Neue Liste Umfrage TV angelegt.")
        global umfrage_mini_bars_TV
        umfrage_mini_bars_TV = umfrage_tv
        
    if len(umfrage_demografie) > 0:
        print("Neue Liste Umfrage Demografie gesetzt")
        global umfrage_demografie_bars
        umfrage_demografie_bars = umfrage_demografie
        
    if len(umfrage_netz) > 0:
        print("Neue Liste Umfrage Internet gesetzt")
        global umfrage_mini_bars_netz
        umfrage_mini_bars_netz = umfrage_netz
        
        
    if len(sort_umfrage) > 0:  
        print("Neues Dict für Sortierung")
        global sortierung_umfrage
        sortierung_umfrage = sort_umfrage
    ###
    if len(region)>0:
        print("Neue Orte (regional) gesetzt")
        global orte_region
        orte_region = region
        
    if len(deutschland)>0:
        print("Neue Orte (Deutschland) gesetzt")
        global orte_deutschland
        orte_deutschland = deutschland
    
    if len(international)>0:
        print("Neue Orte (international) gesetzt")
        global orte_international
        orte_international = international
    
    if len(kolumnen)>0: 
        print("Neue Liste für Kolumnen gesetzt")
        global kolumnen_list
        kolumnen_list = kolumnen
        
    if len(unternehmen)>0: 
        print("Neue Liste für Unternehmen angelegt")
        global liste_unternehmen
        liste_unternehmen = unternehmen
    
    if len(rubriken)>0: 
        print("Neues Rubrikendict gesetzt")
        global rubrik_dict
        rubrik_dict = rubriken
    
    if len(seiten_list)>0:#
        global seitentitel_list
        seitentitel_list = seiten_list
        
    if len(rub_dict)>0: 
        global kolumnen_dict
        kolumnen_dict = rub_dict
#TODO irgendwie müssen alle einzeln 

    
    
    
#%%
        
        
        
#%% ALLGEMEINES

#%% Entwicklung WErte im Messverlauf        
#Grafik Entwicklung - alle zwei Tage ein Balken + Gleitender Mittelwert
        # Unter target nimmt die Funktion die Argumente Lesewert, Blickwert, Durchlesewert entgegen. 
        # Unter title_text kann eine Überschrift für die Grafik übergeben werden
        #Funktion nimmt Dataframe wie es kommt, wird nicht mehr sortiert z.B. nach Ressort.
        
        
        
def grafik_entwicklung(prs, df, target="Lesewert", mean_line=0, legend="large", 
                       grid=True, title_text = False, steps=1, ma=False, limit_y=False, steps_x_label = 1):
   
    if target == "Lesewert": 
        group_param = "Artikel-Lesewert (Erscheinung) in %"
    elif target == "Blickwert":  
        group_param = "Artikel-Blickwert (Erscheinung) in %"
    elif target == "Durchlesewert": 
        group_param = "Artikel-Durchlesewerte (Erscheinung) in %"
    df_ma = df.groupby("Erscheinungsdatum", as_index=False).mean()
    df_ = df.groupby("Erscheinungsdatum", as_index=False).mean()
    if steps !=1: 
        df_ = df_.iloc[0::steps, :]
    # Grafik zeichnen
     #Schriftfarbe und Farbe der Ticks festlegen
    set_font_color ="#8c8f91" 
    
    # Werte für die Achsen werden festgelegt
    if ma==False:
        x = df_["Erscheinungsdatum"].apply(lambda x: x.strftime("%d.%m.%Y"))
        xn = np.arange(len(x))
        labels = x
        y = df_[group_param]
    if ma==True:
        # hier wird der Original-DF auf die Länge des moving-average-df eingekürzt, 
        # meist um vier Zeilen für moving avg von 5
        df_x = df_.iloc[4:]
       
        x = df_x["Erscheinungsdatum"].apply(lambda x: x.strftime("%d.%m.%Y"))
        xn = np.arange(len(x))
        labels = x
        y = df_x[group_param]
        
    
    
    # Seaborn-Style und Größe des Plots festlegen
    sns.set_style("white")
    fig, ax1 = plt.subplots(figsize=(20,8))
    
    #setzt die linke Y-Achse in Campton light
    # rechte Y-Achse können wir erst zum Code-Ende ansteuern
    plt.yticks(fontproperties=campton_light)

     # Achsen, Ticks und alles andere festlegen
    
    if grid==True:
        ax1.grid(color= set_font_color, linestyle='-', linewidth=1, axis="y")
    
     # Barcharts einzeichnen
    bars = ax1.bar(xn,y, color="#f77801", width=0.3, label="Lesewert")
    
    # zeichnet die Regressionslinie               
#    lr = Ridge()
#    lr.fit(xn, y)
#    plt.plot(xn, lr.coef_*x+lr.intercept_, color="orange")        
    if ma == False: 
        sns.regplot(x=xn, y=y, ax=ax1, color="grey", ci=None)
    if ma == True:
        try: 
            ma_name_05 = target + "05"
            ma_name_10 = target + "10"
            ma_name_15 = target + "10"
            # columns für moving average festlegen
            df_ma[ma_name_05] = df_[group_param].rolling(window=5).mean()
            df_ma[ma_name_10] = df_[group_param].rolling(window=10).mean()
            df_ma[ma_name_15] = df_[group_param].rolling(window=15).mean()
            df_ma = df_ma[~df_ma[ma_name_05].isnull()]
    #         #Glättungsversuch durch svgol-Filter
    #        y_10 = savgol_filter(df_ma[ma_name_10], 11, 5)
            y_5 = savgol_filter(df_ma[ma_name_05], 11, 3)
    #        y_15 = savgol_filter(df_ma[ma_name_15], 11, 3)
    #        y_10_pure = df_ma[ma_name_10]
            #x = df_ma["Erscheinungsdatum"].apply(lambda x: x.strftime("%d.%m.%Y"))
            y = df_ma[ma_name_05]
            #xn = np.arange(len(x))
            #labels = x
            
            #plt.plot(xn,  y_5, color="grey", linewidth=4)
            plt.plot(y_5, color="grey", linewidth=4)
        except:
            print("Fehler bei Entwicklungslinie")
        #Glättung mit interpolate.interp1d
#        x_new = np.linspace(xn.min(), xn.max(), 3000)
#        print(x_new.shape)
#        f = interp1d(xn, y, kind="quadratic")
#        y_smooth = f(x_new)
##        
#        plt.plot(x_new, y_smooth, color="gray", linewidth=2)
       #
        
        # Glättung mit neuen X-und Y-Werten
        
        
        
        
        
        
#        #Glättung mit interpolate.interp1d
#        x_new = np.linspace(xn.min(), xn.max(), 10)
#        
#        f = interp1d(xn, y_ma, kind="quadratic")
#        y_smooth = f(x_new)
#        
#        plt.plot(x_new, y_smooth, color="green")
#        
        # Glättungsversuch mit spline
       
#        x_smooth = np.linspace(xn.min(), xn.max(), 10)
#        
#        y_smooth = spline(xn,y,x_smooth,  order=2)
#        plt.plot(x_smooth, y_smooth, color="green")
##       
#        
#        x_smooth2 = np.linspace(xn.min(), xn.max(), 10)
#        
#        y_smooth2 = spline(xn,y,x_smooth,  order=3)
#        plt.plot(x_smooth2, y_smooth2, color="red")
        
#       
        
        
        
        #sns.regplot(x=xn, y=df_[ma_name_05], ax = ax1, color="red", ci=None)
        #sns.regplot(x=xn, y=y, ax=ax1, color="grey", ci=None)       
        #plt.plot(df_ma[ma_name_10], color="green")
       # plt.plot(y_10, color="black")
        #plt.plot(y_5, linewidth=5, color='grey')
        #plt.plot(y_10, color="black")
      #  plt.plot(y_15, color="blue")
       # plt.plot(df_ma[ma_name_05], color="red", label ="Gleitender Mittelwert")
       # plt.plot(df_ma[ma_name_15], color="blue") 
       # plt.plot(y_15, color="green")    
        
    ax1.set_ylabel('Ø ' +target +' in Prozent', color= set_font_color, \
                   fontproperties=campton_light, fontsize=50)
    #  Anzahl der Ticks 
    #if steps_x_label == 1:
    ax1.xaxis.set(ticks=range(0, len(xn)))
     
    ax1.set_xticklabels(labels = labels, rotation=90, ha="center",  weight=800,\
                        color= set_font_color, fontproperties=campton_light, \
                        fontsize=20) # Labels werden ausgerichtet
    if steps_x_label != 1:
        for label in ax1.xaxis.get_ticklabels()[::2]:
            label.set_visible(False)
    ax1.patch.set_facecolor('white') # Hintergrundfarbe auf weiß, dann... 
    ax1.patch.set_alpha(0.0) # Hintergrund ausblenden, damit zweite Grafik 
                                                #   (der Plot) sichtbar wird
    ax1.set_zorder(2) # erste Grafik wird vor die zweite Grafik geschoben
    
    ax1.yaxis.label.set_size(22)
    
    #begrenzt die Y-Achse
    if limit_y == True: 
        bottom = y.min() - y.min()*0.3
        top = y.max()
        if y.min() > 5:
            ax1.set_ylim(bottom, top)
    
    
    # Abstände Bars zur Achse (standardmäßig bei 0.5)
    plt.margins(x=0.03) # ziehen Bars näher an die Achse
     #obere Linie ausblenden
    ax1.spines["top"].set_visible(False)
    #ax1.spines["left"].set_color("gray")
    ax1.spines["top"].set_visible(False)
    ax1.spines["bottom"].set_visible(False)
    ax1.spines["left"].set_visible(False)
    ax1.spines["right"].set_visible(False)
    
    
    
    # jetzt werden die Y-Ticks links in Campton Light gefasst
    plt.yticks(fontproperties=campton_light)
    
    
    # Bei Bedarf Linie mit dem Durchschnitt einziehen
    if mean_line !=0:
        
        labeltext = "Ø LW Seitentitel: {:1.1f}".format(float(mean_line)).replace(".", ",")
        linie = ax1.axhline(y=mean_line, xmin=0.01, xmax=0.99, color=set_font_color, label=labeltext)
    
    
     # Zahlen an den Y-Achsen verändern, Größe und Farbe
    ax1.tick_params(axis='y', labelsize=25, colors= set_font_color)
    
    legend_height = 1.24
    if legend=="normal":
            legend_height = 1.24
    elif legend == "strange": 
            legend_height =1.44
    else:
        legend_height = 1.24
             
   
    
    if mean_line != 0:
        leg = plt.legend(bbox_to_anchor=(1, legend_height), handles=[bars, linie], markerscale=140)
    else:
        leg = plt.legend(bbox_to_anchor=(1, legend_height), handles=[bars], markerscale=140)

    for text in leg.get_texts(): 
        plt.setp(text, color= set_font_color, size=21)
    
   
             
    #plt.tight_layout()
    
    # Canvaseinstellung / Position des Plots
    # Function nutzt Voreinstellung aus Parametern
    pos_left = 0.1 # 0.2
    pos_right=0.9 #0.8
    pos_top=0.90
    pos_bottom = 0.3
    
    
    plt.subplots_adjust(left=pos_left, right=pos_right, top=pos_top, 
                        bottom=pos_bottom)
    
    filename = "grafik_entwicklung_lesewert.png"
    plt.savefig(filename)
    
    # Prüfung ob Titel als Paramenter mitgegeben wurde, ansonsten automatisch eintragen
    if title_text == False: 
        title_text = "Entwicklung " + target 
    else: 
        title_text = title_text
    plt.close()
    # Plot wird auf PPTX-Sheet gezogen
    
    picture_sheet(prs, filename, title_text=title_text)            
    
    return prs

      
        
        
#%% Initial-Funktion für Fragebogen 
        
        
        
def fragebogen(prs, df, liste = ausgaben_liste):
    
    
    # Identifizierung der Fragetexte ("Ihre Zeitung ist...")
    # geschieht über das dict zeitung_attribute_dict
    # die Anlage einer zusätzlichen Liste ist nicht notwendig, 
    # wir suchen über das dict
    
    # check ob attribute_df
    
    #Falls nötig: GESAMT
    analyse_mantel_abschluss(prs, df, liste_ressorts=mantel_ressorts)
    
    for elem in liste:
        
        elem_df = df[df["ZTG"]==elem]
        attribute_df = elem_df[elem_df["Fragetext"].isin(zeitung_attribute_dict)]
        #Berechnung der Werte, Erstellung des Diagramms, Schreiben auf PPTX
        #zeitungsattribute_berechnung(prs, attribute_df)
        #zeitung_themen(prs, elem_df)
        #deckblatt(prs, "Umfrage" + elem)
        #umfrage_pie(prs, elem_df)
        
        #mini_bars(prs, elem_df, title_text = "Zeitungsnutzung", liste=umfrage_mini_bars_TV)
        #mini_bars(prs, elem_df, title_text = "Internetnutzung", liste=umfrage_mini_bars_netz)
        
        
        
    return prs
        



#%% Berechnung Attribute_Werte der Zeitung  
# Die Funktion erstellt mehrfarbige, gestapelte  Bar-Charts für die Darstellung
# der Attribut auf den Satz "Ihre Zeitung ist..." 
# Als Parameter benötigt die Funktion das pptx-Objekt und die
    # Umfragedaten
# Die Datei benötigt ein vorsortiertes Dataframe, bedeutet: sortiert nicht mehr selbst, 
    # zum Beispiel nach Ausgaben etc. 

def zeitungsattribute_berechnung(prs, df, title="", frage_dict = zeitung_attribute_dict):
    print(df.shape[0])
    print(df.Fragetext)
    df = df[df["Fragetext"].isin(frage_dict)]
    print(df.shape)
    
    df = df[["Welle", "Fragetext", "Antworttext"]]
    df_group = df.groupby(["Fragetext", "Antworttext"], as_index=False).count() 
    # DAten von long auf wide umbauen, um Grafik leichter zeichnen zu können
    df_group = df_group.pivot(index= "Fragetext", columns="Antworttext", 
                              values= "Welle")
    
    df_group.reset_index(inplace=True)
    
    df_group.fillna(0, inplace=True)
   # if frage_dict==zeitung_attribute_dict:
    df_pivot = df_group.reindex(columns=["Fragetext", "trifft nicht zu", "trifft eher nicht zu", "trifft eher zu", "trifft voll zu"])
    #if frage_dict==zeitung_attribute_dict2:
    #    df_pivot = df_group.reindex(columns=["Fragetext", "überhaupt nicht zufrieden", "weniger zufrieden", "zufrieden", "voll zufrieden"])
    #if frage_dict==zeitung_attribute_dict3:
    #    df_pivot = df_group.reindex(columns=["Fragetext", "nie", "monatlich", "wöchentlich", "täglich", "mehrmals täglich"])
    #else:
     #   print(frage_dict)
    # START PLOTTING
    # Variablen festlegen
    x = df_pivot["Fragetext"]
    df_pivot["shortnames"] = df_pivot["Fragetext"].apply(lambda x: frage_dict[x])
    labels = df_pivot["shortnames"]
    xn = range(len(x))
    width = 0.3
    #Schriftfarbe und Farbe der Ticks festlegen, fig/ax aufrufen
    set_font_color ="#8c8f91" 
    sns.set_style("white")
    fig, ax1 = plt.subplots(figsize=(18,8))
    #setzt Schriftart und Schriftgrößé für linke Y-Achse
    plt.yticks(fontproperties=campton_light)
    ax1.tick_params(axis='y', labelsize=25, colors= set_font_color)
    rotation=45
    ha="right"
    fontsize_ticks = 30
    p1 = ax1.bar(xn, df_pivot["trifft nicht zu"], width, 
                     label=("trifft nicht zu"), color="#fdbe85")
    p2 = ax1.bar(xn, df_pivot["trifft eher nicht zu"], width, \
                 bottom=df_pivot["trifft nicht zu"], label=("trifft eher nicht zu"), color="#fd8d3c")
    p3 = ax1.bar(xn, df_pivot["trifft eher zu"],width,\
                 bottom=df_pivot["trifft nicht zu"]+
                 df_pivot["trifft eher nicht zu"], label=("trifft eher zu"), 
                 color="#e6550d")
    p4 = ax1.bar(xn, df_pivot["trifft voll zu"], width, \
                 bottom=df_pivot["trifft nicht zu"]+df_pivot["trifft eher nicht zu"]
                 +df_pivot["trifft eher zu"],color="#a63603", 
                 label=("trifft voll zu"))
    if frage_dict == zeitung_attribute_dict:
        rotation=45
        ha="right"
        fontsize_ticks = 30
        p1 = ax1.bar(xn, df_pivot["trifft nicht zu"], width, 
                     label=("trifft nicht zu"), color="#fdbe85")
        p2 = ax1.bar(xn, df_pivot["trifft eher nicht zu"], width, \
                 bottom=df_pivot["trifft nicht zu"], label=("trifft eher nicht zu"), color="#fd8d3c")
        p3 = ax1.bar(xn, df_pivot["trifft eher zu"],width,\
                 bottom=df_pivot["trifft nicht zu"]+
                 df_pivot["trifft eher nicht zu"], label=("trifft eher zu"), 
                 color="#e6550d")
        p4 = ax1.bar(xn, df_pivot["trifft voll zu"], width, \
                 bottom=df_pivot["trifft nicht zu"]+df_pivot["trifft eher nicht zu"]
                 +df_pivot["trifft eher zu"],color="#a63603", 
                 label=("trifft voll zu"))
    
    if frage_dict == zeitung_attribute_dict2:
        rotation=45
        ha="right"
        fontsize_ticks = 30
        p1 = ax1.bar(xn, df_pivot["überhaupt nicht zufrieden"], width, 
                     label=("überhaupt nicht zufrieden"), color="#fdbe85")
        p2 = ax1.bar(xn, df_pivot["weniger zufrieden"], width, \
                 bottom=df_pivot["überhaupt nicht zufrieden"], label=("weniger zufrieden"), color="#fd8d3c")
        p3 = ax1.bar(xn, df_pivot["zufrieden"],width,\
                 bottom=df_pivot["überhaupt nicht zufrieden"]+
                 df_pivot["weniger zufrieden"], label=("zufrieden"), 
                 color="#e6550d")
        p4 = ax1.bar(xn, df_pivot["voll zufrieden"], width, \
                 bottom=df_pivot["überhaupt nicht zufrieden"]+df_pivot["weniger zufrieden"]
                 +df_pivot["zufrieden"],color="#a63603", 
                 label=("voll zufrieden"))
    
    if frage_dict == zeitung_attribute_dict3:
        rotation = 45
        ha="right"
        fontsize_ticks = 30
        p1 = ax1.bar(xn, df_pivot["nie"], width, 
                     label=("nie"), color="#feedde")
        p2 = ax1.bar(xn, df_pivot["monatlich"], width, \
                 bottom=df_pivot["nie"], label=("monatlich"), color="#fdbe85")
        p3 = ax1.bar(xn, df_pivot["wöchentlich"],width,\
                 bottom=df_pivot["nie"]+
                 df_pivot["monatlich"], label=("wöchentlich"), 
                 color="#fd8d3c")
        p4 = ax1.bar(xn, df_pivot["täglich"], width, \
                 bottom=df_pivot["nie"]+df_pivot["monatlich"]
                 +df_pivot["wöchentlich"],color="#e6550d", 
                 label=("täglich"))
        p5 = ax1.bar(xn, df_pivot["mehrmals täglich"], width, \
                 bottom=df_pivot["nie"]+df_pivot["monatlich"]
                 +df_pivot["wöchentlich"]+df_pivot["täglich"],color="#a63603", 
                 label=("mehrmals\ntäglich"))
    
    #Labels und Ticks festlegen
    ax1.set_ylabel('', color= set_font_color, \
                   fontproperties=campton_light, fontsize=30)
    ax1.xaxis.set(ticks=range(0, len(xn))) #  Anzahl der Ticks 
    ax1.set_xlabel("")
    ax1.set_xticklabels(labels = labels, rotation=rotation, ha=ha,  weight=800,\
                        color= set_font_color, fontproperties=campton_light, \
                        fontsize=fontsize_ticks) # Labels werden ausgerichtet
    # Legende
    
    handles,labels = ax1.get_legend_handles_labels()
    
    if frage_dict == zeitung_attribute_dict3:
        handles = [handles[4],handles[3], handles[2], handles[1], handles[0]]
        labels = [labels[4], labels[3], labels[2], labels[1], labels[0]]
    else:
        handles = [handles[3], handles[2], handles[1], handles[0]]
        labels = [labels[3], labels[2], labels[1], labels[0]]
        
        
        
    leg = ax1.legend(handles, labels, bbox_to_anchor=(1.05, 0.90), loc=2, borderaxespad=0., markerscale=140) # 
    for text in leg.get_texts(): 
        plt.setp(text, color= set_font_color, size=30)
        
    #Lage der Bars / Sichtbarkeit der Spines  
    plt.margins(x=0.03) # ziehen Bars näher an die Achse
    #obere Linie ausblenden
    ax1.spines["top"].set_visible(False)
    #ax1.spines["left"].set_color("gray")
    ax1.spines["top"].set_visible(False)
    ax1.spines["bottom"].set_visible(False)
    ax1.spines["left"].set_visible(False)
    ax1.spines["right"].set_visible(False)
    
    #Grid-Linien
    ax1.grid(color= "#e0e0e0", linestyle='-', linewidth=1, axis="y")
          
    # Titel hinzufügen
    figure_title= "Ihre Zeitung ist..."
    fontsize=50
    
    if frage_dict == zeitung_attribute_dict:
        figure_title= "Ihre Zeitung ist..."
        fontsize=50
    if frage_dict == zeitung_attribute_dict2:
        figure_title= "Wie zufrieden sind Sie mit diesen Leistungen?"
        fontsize=40
    if frage_dict == zeitung_attribute_dict3:
        figure_title= "Wie häufig nutzen Sie folgende Produkte?"
        fontsize=37
    ax1.text(0.5, 1, figure_title, horizontalalignment='center',
             verticalalignment='baseline', 
             transform = ax1.transAxes,fontproperties=campton_light,
             color = set_font_color, fontsize=fontsize, 
             bbox=dict(facecolor='none', edgecolor='none', pad=30) )  
            #[ 'center' | 'top' | 'bottom' | 'baseline' ]
    p_left = 0.15
    p_right = 0.92
    p_top = 0.85
    p_bottom = 0.3 
#if label_position == "large": 
#       p_left = 0.12
#       p_right = 0.92
#       p_top = 0.85
#       p_bottom = 0.4
#   
#   if label_position == "xlarge": 
#       p_left = 0.14
#       p_right = 0.92
#       p_top = 0.85
#       p_bottom = 0.45
    
    if frage_dict == zeitung_attribute_dict: 
        p_left = 0.15
        p_right = 0.92
        p_top = 0.85
        p_bottom = 0.3
    
    if frage_dict == zeitung_attribute_dict2: 
        p_left = 0.35
        p_right = 0.92
        p_top = 0.70
        p_bottom = 0.3
    
    if frage_dict == zeitung_attribute_dict3: 
        p_left = 0.25
        p_right = 0.90
        p_top = 0.63
        p_bottom = 0.3
    
    
    plt.subplots_adjust(left=p_left, right=p_right, top=p_top, 
                        bottom=p_bottom)   
    plt.margins(0.05,0.19)
    #plt.tight_layout()
    #plt.tight_layout(pad=4, w_pad=0, h_pad= 1.0)
    # Canvaseinstellung / Position des Plots
    #
   
    
    
    filename = "grafik_ztg_attribute.png"
    plt.savefig(filename, bbox_inches="tight")
   
    plt.close()
    # Plot wird auf PPTX-Sheet gezogen
    # Vergabe des Titeltextes
    if len(title)>0:
        title_text = "Zeitungsnutzung " + ausgaben_dict[title]
    else:
        title_text = "Zeitungsnutzung"
    picture_sheet(prs, filename, title_text=title_text)         
    #savefig("filename.pdf", bbox_inches = 'tight',
    #pad_inches = 0)
    
    #gca().xaxis.set_major_locator(NullLocator())
    #gca().yaxis.set_major_locator(NullLocator())
    return prs     
        
        


#%% Welche Themen - Balkengrafiken quer für Fragebogen
    # Die Funktion benötigt das PPTX-Objekt und das Dataframe der Umfrage. 
    # Erstellt werden Querbalkengrafiken für die Frage: "Welche der folgenden Themen 
    # interessieren Sie besonders?"
    # Die Antworten,  die jeweils auf eine Folie sollen, sind unter Listen und Objekten
    # jeweils in Listen angegeben. Die müssen bei Bedarf geändert werden. 
    
def zeitung_themen_horizontal(prs, df, title="", fragetext = "Welche der folgenden Themen interessieren Sie besonders in der Zeitung?"):
    
    
    df_ = df[df["Fragetext"]==fragetext]
    
    for elem in umfrage_gesamt_liste: 
        
        df_loop = df_[df_["Antworttext"].isin(elem)]
        
        df_group = df_loop.groupby(["Antworttext"], as_index=False).count().sort_values(by="Welle")
        
        
        # Grafik zeichnen
     
        df_group["Antwort_short"] = df_group["Antworttext"].apply(lambda x: elem[x])
        answer = df_group["Antwort_short"]
        y_pos = range(len(answer))
        nennung = df_group["Welle"]
        bar_color = "#f77801"
       
        width = 0.3
        set_font_color = "#8c8f91"
    
    
        #Schriftfarbe und Farbe der Ticks festlegen
        
        sns.set_style("white")
        fig, ax1 = plt.subplots(figsize=(19,7))
        #setzt Schriftart und Schriftgrößé für linke Y-Achse
        plt.yticks(fontproperties=campton_light)
        plt.xticks(fontproperties=campton_light)
        ax1.tick_params(axis='x', labelsize=25, colors= set_font_color)
        ax1.tick_params(axis='y', labelsize=25, colors= set_font_color)
        ax1.set_yticks(y_pos)
        ax1.set_yticklabels(answer)
    
        #für horizontale Balken
        ax1.barh(y_pos, nennung, color=bar_color, align="center")
        
        
        plt.margins(x=0.05) # ziehen Bars näher an die Achse
             #obere Linie ausblenden
        ax1.spines["top"].set_visible(False)
            #ax1.spines["left"].set_color("gray")
        ax1.spines["top"].set_visible(False)
        ax1.spines["bottom"].set_visible(False)
        ax1.spines["left"].set_visible(False)
        ax1.spines["right"].set_visible(False)
    
    
    
        #Grid-Linien
        ax1.grid(color= "#e0e0e0", linestyle='-', linewidth=1, axis="x")
    
        # Titel hinzufügen
        figure_title= fragetext
        plt.text(0, 1.0, figure_title,
                 #horizontalalignment='center',
    
                transform = ax1.transAxes, 
                fontproperties=campton_light, 
                color = set_font_color, 
                fontsize=30)
        
        
        #plt.subplots_adjust(left=0.5, right=0.65, top=0.85, 
                            #bottom=0.4)   
        plt.margins(0.05,0.19)
        plt.tight_layout()
        
        
        filename = "grafik_ztg_themen .png"
        plt.savefig(filename, bbox_inches="tight")
        plt.close()
        # Plot wird auf PPTX-Sheet gezogen
        # Vergabe des Titeltextes
        if len(title)>0:
            title_text = "Zeitungsnutzung " + ausgaben_dict[title]
        else: 
            title_text = "Zeitungsnutzung"
        picture_sheet(prs, filename, title_text=title_text)  
    return prs



def zeitung_themen(prs, df, title="", fragetext = "Welche der folgenden Themen interessieren Sie besonders in der Zeitung?", special=""):
    
    df_ = df[df["Fragetext"]==fragetext]
    
    for elem in umfrage_gesamt_liste: 
        
        if elem =="umfrage_verzicht":
            fragetext = "Auf was könnten Sie in Ihrer Zeitung verzichten?"
            df_ = df[df["Fragetext"]==fragetext]
            print("Umfrage_verzicht erkannt")
        
        df_loop = df_[df_["Antworttext"].isin(elem)]
        df_group = df_loop.groupby(["Antworttext"], as_index=False).count().sort_values(by="Welle", ascending=False)
        
        
        # Grafik zeichnen
     
        df_group["Antwort_short"] = df_group["Antworttext"].apply(lambda x: elem[x])
        x = df_group["Antwort_short"]
        labels = x
        xn = range(len(x))
        y = df_group["Welle"]
        bar_color = "#f77801"
       
        width = 0.8
        set_font_color = "#8c8f91"
        ticksize = 30
    
        #Schriftfarbe und Farbe der Ticks festlegen
        
        sns.set_style("white")
        fig, ax1 = plt.subplots(figsize=(19,7))
        
        
        #setzt Schriftart und Schriftgrößé für linke Y-Achse
        bars = ax1.bar(xn,y, color="#f77801", width=0.8, label="Anzahl Antworten")
        rotation = 45
        
        ax1.xaxis.set(ticks=range(0, len(xn))) #  Anzahl der Ticks 
        
        
        ax1.set_xticklabels(labels = labels, rotation=rotation, ha="right",  weight=800,\
                        color= set_font_color, fontproperties=campton_light, \
                        fontsize=ticksize) # Labels werden ausgerichtet
        
        plt.yticks(fontproperties=campton_light)
        ax1.tick_params(axis='y', labelsize=0, colors= set_font_color)
         #Grid wird eingebaut
        ax1.grid(color= "#dddddd", linestyle='-', linewidth=1, axis="y")
        
        
        for p in ax1.patches:
        # Problem: Ist ein Balken nur vier groß, ist der Abstand zu gering
        # TODO... NOCH EIN BISSCHEN EXPERIMENTIEREN
            height = p.get_height()
        
            txt = '{:1.0f}'.format(height).replace(".", ",")
            ax1.text(p.get_x()+p.get_width()/2., height + 5, txt, ha="center",\
                     fontproperties=campton_light, color= set_font_color, rotation=0\
                     ,fontsize= 25, weight = 1000)
        
        
        #ax1.yaxis.label.set_size(22)
        plt.margins(x=0.03) # ziehen Bars näher an die Achse
        
        ax1.spines["top"].set_visible(False)
        #ax1.spines["left"].set_color("gray")
        ax1.spines["top"].set_visible(False)
        ax1.spines["bottom"].set_visible(False)
        ax1.spines["left"].set_visible(False)
        ax1.spines["right"].set_visible(False)
        
        
         # Überschrift
        # ist sie länger als 40, wird sie mit textwrap auf zwei Zeilen verteilt
        
        
        #ax1.axes.get_yaxis().set_visible(False)
        
        # Titel hinzufügen
        figure_title= fragetext
        ax1.text(0.5, 1.2, figure_title, horizontalalignment='center',
             verticalalignment='baseline', 
             transform = ax1.transAxes,fontproperties=campton_light,
             color = set_font_color, fontsize=32, 
             bbox=dict(facecolor='none', edgecolor='none', pad=30) )  
            #[ 'center' | 'top' | 'bottom' | 'baseline' ]
        
        
       
        
        
    
        pos_left = 0.1 # 0.2
        pos_right=0.9 #0.8
        pos_top=0.8
        pos_bottom = 0.42
        
        plt.subplots_adjust(left=pos_left, right=pos_right, top=pos_top, 
                       bottom=pos_bottom)
        
        filename = "grafik_ztg_themen.png"
        
       
        plt.savefig(filename)
        picture_sheet(prs, filename, title_text="Zeitungsnutzung")
        plt.close()
        
        
    
    
    
    return prs





#%% Pie-Charts (klein) Zeitungsnutzung
    # Diese Funktion schreibt jeweils zwei Pie-Charts auf ein PPTX-Dokument. 
    # Die Reihenfolge der Pie-Charts ist unter Listen und Objekte in der Liste 
    # umfrage_piecharts festgelegt. 
def umfrage_pie(prs, df, title=""):
    
    counter=0
    for elem in umfrage_piecharts: #umfrage_piecharts enthält die Fragen für die Pies.
        counter += 1
        
        df_ = df[df["Fragetext"]==elem]
        
        df_group = df_.groupby(["Antworttext"], as_index=False).count().sort_values(by="Welle", ascending= False)
        print("df_group vor Sortierung")
        print(df_group)
        
        # Sortierung, falls Werte laut sortierungsliste sortiert werden sollen 
        if elem in sortierung_umfrage:  
        # Labels werden jetzt sortiert
        # Sortierung findet sich in Liste sortierung_umfrage
            df_group["Antworttext"] = pd.Categorical(df_group["Antworttext"], sortierung_umfrage[elem])
            df_group = df_group.sort_values(by="Antworttext")
            print("group nach sortierung: ")
            print(df_group)
            #labels = sorted(labels, key=sorter.index)
        
        
        
        sizes = df_group["Welle"]
        set_font_color ="#8c8f91"
         
        # legt Position der Zahlen innerhalb der Pie-Chart fest
        def autodistance(autopct):
            
            if autopct >=5:
            
                return 1.15 # 1.2 bei kleinen Zahlen
            else:
                return 0.8
        
        # Erstellt die tatsächlichen Nutzerzahlen anhand der automatisch generierten Prozentzahlen
        def make_label(sizes):
            total = sizes.sum()
            def my_label(pct):
                return round(total*(pct/100)).astype(int)
            mylabel = my_label
            
            #if mylabel >=6:
               # pct_distance = 0.8
            #if mylabel <6:
               # pct_distance = 1.2
            return mylabel
        
        # Falls nötig: Explode wird für jedes Element errechnet
        explode = ()
        for i in range(len(sizes)):
            explode = explode +(0.03,)
        
        colors =["#ff9900", "#ffc570", "#fce3bf", "#fcf6d1"]
        fig, ax1 = plt.subplots(figsize=(8,8))
        handles, labels = ax1.get_legend_handles_labels()
        
        labels = df_group["Antworttext"].values
      
       
        
        
        ax1.pie(sizes, radius=1, frame=False, shadow=False, autopct=make_label(sizes), startangle=90, colors=colors,\
                textprops={'fontsize': 20, "color": set_font_color},pctdistance =0.8, 
                wedgeprops={"linewidth":2, "edgecolor":"white"})
        ax1.axis("equal")
        
            
        
        
       
        
       
        
        
       
    
        
        
        leg = ax1.legend(handles, labels=labels, #bbox_to_anchor=(1, 0.2), 
                         markerscale=14, mode="expand", 
                         borderaxespad=0., loc=8) #1.Wert x, 2. Wert y
        for text in leg.get_texts(): 
            plt.setp(text, color= set_font_color, size=20)
        
        # Überschrift
        # ist sie länger als 40, wird sie mit textwrap auf zwei Zeilen verteilt
        figure_title= elem
        posx=0.5
        posy=0.82
        if len(figure_title) >40:
            figure_title = "\n".join(textwrap.wrap(elem, width=35))
            posy=0.8
          
        
        plt.text(posx, posy, figure_title,
             horizontalalignment='center',
             
             transform = ax1.transAxes, 
            fontproperties=campton_light, 
            color = set_font_color, 
            fontsize=28)
        
        
        filename = "grafik_pie.png"
        filename2 = "grafik_pie2.png"
        
        #plt.tight_layout()
        #plt.savefig(filename)
    # Canvaseinstellung / Position des Plots
    # Function nutzt Voreinstellung aus Parametern
    
        
        pos_left = 0.25 # 0.2
        pos_right=0.75 #0.8
        pos_top=0.98
        pos_bottom = 0.1
        plt.subplots_adjust(left=pos_left, right=pos_right, top=pos_top, 
                       bottom=pos_bottom)
        
        
        # plt.savefig(filename, bbox_inches="tight")
        if len(title)>0:
            title_text = "Zeitungsnutzung " +ausgaben_dict[title]
        else: 
            title_text = "Zeitungsnutzung"
        
        
        if counter%2!=0:
            plt.savefig(filename)
            #double_picture_sheet(prs, filename, filename2, title_text=title_text)
        elif counter%2==0:
            plt.savefig(filename2)
            double_picture_sheet(prs, filename, filename2, title_text=title_text)
        plt.close()
        
        #plt.savefig(filename, bbox_inches="tight")
        
        #picture_sheet(prs, filename, title_text=title_text) 
        #double_picture_sheet(prs, filename, title_text=title_text)
        
    return prs





#%% Fragebogen Mini-Barchart
    # Die Fuktion legt die Fragen fest, die per Mini-Barchart auf PPTX-Docs gezogen
    # werden (je zwei auf eine Seite). 
    # Die Funktion benötigt eine Liste,  in der die Fragen angegeben sind, die 
    # genutzt werden sollen. Voreingestellt ist umfrage_mini_bars_TV, es gibt auch die 
    # liste umfrage_mini_bars_netz. Die Listen finden sich unter Listen und Objekte. 
    
    # Die Überschrift kann mit title_text individuell eingestellt werden. 
   
        
def mini_bars(prs, df, title_text = "Fernsehnutzung", liste=umfrage_mini_bars_TV, title=""):
    counter = 0
    
    
    
    
    if len(liste)==0:
        print("Achtung, Funktion mini_bars ohne Liste mit Umfrage-Themen")
    for elem in liste:
        counter +=1
        df_ = df[df["Fragetext"]==elem]
        df_group = df_.groupby(["Antworttext"], as_index=False).count().sort_values(by="Welle", ascending= False)
        
        # Größe der X-Ticks, wird bei längeren Texten Internetnutzung geändert    
        ticksize=30  
        len_group = df_group.Antworttext.map(len).max()
        if len_group >=20: 
            ticksize = 23
        
        
      
        
        # Sorteriung vornehmen nach vorgegebenen REihenfolgen in sortierung_umfrage
        if elem in sortierung_umfrage:
            sorter = sortierung_umfrage[elem]
            
            df_group.Antworttext = df_group.Antworttext.astype("category")
            df_group.Antworttext.cat.set_categories(sorter, inplace=True)
            df_group.sort_values(["Antworttext"], inplace=True)
        else: 
            print(elem + " ... is not in dict")
            
            
        set_font_color ="#8c8f91"
        
        x = df_group["Antworttext"]
        
            
        
        
        def check_labels(elem):
            if len(elem)>19:
                return elem[:16] + "..."
            else:
                return elem
        
        def check_labels_pers(elem):
            if len(elem)>30:
                return elem[:26] + "..."
            else:
                return elem
        
        def check_labels_internet(elem):
            if len(elem)> 8:
                
                return_label = "\n".join(textwrap.wrap(elem, width=8))
                global ticksize
                ticksize = 5
                return return_label
            else:
                return elem
          
        #ticksize = 30
        figsize = (11,8.8)
        #figsize = (3.94, 3.15) # umgerechnet (10,8) in cm
        if title_text == "Fernsehnutzung":  
            labels = x.apply(check_labels)
            rotation=45
            ha="right"
            # Canvaseinstellung / Position des Plots
            # Function nutzt Voreinstellung aus Parametern
            p_left = 0.25#  vorher: 0.3
            p_right = 0.75# vorher: 0.8
            p_top = 0.7
            p_bottom = 0.4
           
            
            #p_left = 0.08
            #p_right = 0.92
            #p_top = 0.85
            #p_bottom = 0.33
            
            
        elif title_text == "Internetnutzung": 
            labels = x.apply(check_labels_internet)
            
            
            rotation= 0
            
            ha="center"
            # Canvaseinstellung / Position des Plots
            # Function nutzt Voreinstellung aus Parametern
            p_left = 0.1
            p_right = 0.9
            p_top = 0.65
            p_bottom = 0.3
            
            if len(labels)>9:
                ticksize=5
            if len(labels)<=9:
                ticksize=30
        #TODO ticksize muss sich verkleinern, wenn der Text der X-Ticks länger wird    
        
        
        
        
        elif title_text == "Persönliches":
            labels = x.apply(check_labels_pers)
            rotation=45
            ha="right"
            # Canvaseinstellung / Position des Plots
            # Function nutzt Voreinstellung aus Parametern
            p_left = 0.25#  vorher: 0.3
            p_right = 0.75# vorher: 0.8
            p_top = 0.7
            p_bottom = 0.4
            
            
            
            
        xn = range(len(x))
        y = df_group["Welle"]
        
         # Seaborn-Style und Größe des Plots festlegen
        sns.set_style("white")
        fig, ax1 = plt.subplots(figsize=figsize)
    
        #setzt die linke Y-Achse in Campton light
        # rechte Y-Achse können wir erst zum Code-Ende ansteuern
        #plt.yticks(fontproperties=campton_light)
        
        #Grid einbauen
        
        
        bars = ax1.bar(xn,y, color="#f77801", width=0.8, label="Anzahl Antworten")
        
        
        ax1.xaxis.set(ticks=range(0, len(xn))) #  Anzahl der Ticks 
        
        
        ax1.set_xticklabels(labels = labels, rotation=rotation, ha=ha,  weight=800,\
                        color= set_font_color, fontproperties=campton_light, \
                        fontsize=ticksize) # Labels werden ausgerichtet
        
        for p in ax1.patches:
        # Problem: Ist ein Balken nur vier groß, ist der Abstand zu gering
        # TODO... NOCH EIN BISSCHEN EXPERIMENTIEREN
            height = p.get_height()
            var_height = 5
            if height <=8:
                var_height = 2
            txt = '{:1.0f}'.format(height).replace(".", ",")
            ax1.text(p.get_x()+p.get_width()/2., height + var_height, txt, ha="center",\
                     fontproperties=campton_light, color= set_font_color, rotation=0\
                     ,fontsize= 30, weight = 1000)
        
        #ax1.yaxis.label.set_size(22)
        plt.margins(x=0.03) # ziehen Bars näher an die Achse
        
        ax1.spines["top"].set_visible(False)
        #ax1.spines["left"].set_color("gray")
        ax1.spines["top"].set_visible(False)
        ax1.spines["bottom"].set_visible(False)
        ax1.spines["left"].set_visible(False)
        ax1.spines["right"].set_visible(False)
        #ax1.grid(color= set_font_color, linestyle='-', linewidth=1, axis="x")
        
         # Überschrift
        # ist sie länger als 40, wird sie mit textwrap auf zwei Zeilen verteilt
        figure_title= elem
        posx=0.5
        posy=1.5
        if len(figure_title) >40:
            figure_title = "\n".join(textwrap.wrap(elem, width=35))
            
       
        
        plt.text(posx, posy, figure_title,
             horizontalalignment='center',
             
             transform = ax1.transAxes, 
            fontproperties=campton_light, 
            color = set_font_color, 
            fontsize=38)
        
        #ax1.axes.get_yaxis().set_visible(False)
        
        #plt.tight_layout()
        #grid einbauen
        #Grid einbauen
        ax1.tick_params(axis='y', labelsize=0, colors= set_font_color)
        ax1.grid(color= "#dddddd", linestyle='-', linewidth=1, axis="y")
       
        
        plt.subplots_adjust(left=p_left, right=p_right, top=p_top, 
                        bottom=p_bottom)
        
        filename1 = "grafik_lesewert_mini_bars.png"
        filename2 = "grafik_lesewert_mini_bars2.png"
        
        if len(title)>0:
            title_txt = title_text + " " + ausgaben_dict[title]
        
        else: 
            title_txt = title_text
        
        
        if counter%2 !=0:
            
            plt.savefig(filename1)
        elif counter%2 ==0:
            
            plt.savefig(filename2)
            double_picture_sheet(prs, filename1, filename2, title_text=title_txt)
        plt.close()
        
    return prs


#%% Grafik Umfrage Verdienst
def umfrage_geld(df_, fragetext = 'Wie hoch ist das monatliche Nettoeinkommen, das alle zusammen in Ihrem Haushalt haben - nach Abzug von Steuer und Sozialversicherung?'):
    geld_list = ["unter 1.500 €", "zwischen 1.500 und 2.500 €", "über 2.500 €"]
    
    df_geld = df_[df_["Fragetext"]==fragetext]
    df_geld = df_geld.groupby("Antworttext", as_index=False).count()
    
    
    
    x = df_geld       

#%%    

#%% ALLGEMEINE ERKENNTNISSE
    

#%% Kennzahlen der Messung

def kennzahlen(df, df_scans):
    print()
    print("---------------- Kennzahlen der Messung --------------------------")
    scans = df_scans.shape
    artikel = df.shape
    lw_hoch = df[df["Artikel-Lesewert (Erscheinung) in %"]>25].shape
    lw_niedrig = df[df["Artikel-Lesewert (Erscheinung) in %"]<5].shape
    
    print("Anzahl Scans: {}".format(scans))
    print("Anzahl ausgewertete Artikel: {}".format(artikel))
    print("Anzahl Artikel über LW von 25: {}".format(lw_hoch))
    print("Anzahl Artikel unter LW von 5: {}".format(lw_niedrig))
    
    
    
#%% Lesetage 
    # Funktion erstellt Barchart über die Lesetage und die durchschnittliche Zahl von Lesern 
    # an den jeweiligen Tagen

# year = Voreingestellt 2018, damit wertet die Funktion nur treffer aus dem Jahr 2018 aus. Bei False werden alle Treffer ausgewertet
    # TODO Funktion erweitern mit von-bis-Angaben

def lesetage(prs, df_scans, title_text="Bitte Titel eingeben", year=False, 
             multi_line = False, multi_line_ZTG_list = []):
    # Filtern nach Scans = Wahr
    
    df_time = df_scans.copy()
   # df_time = df_scans[df_scans["Treffer"]==True].copy()
    
    df_time["time"] = df_time["Erfassungsdatum"] + "-" + df_time["Erfassungsuhrzeit"]
    
    #Erstelle Datetime-Format
    df_time["time"] = pd.to_datetime(df_time["time"], format="%Y-%m-%d-%H:%M:%S")
    
#    # Filtere Jahre, nur 2018 erlaubt
#    if year != False: 
#        mask = df_time["time"].dt.year == year
#        df_time = df_time[mask]      
#        
    # Erstellen einer Spalte mit den Lesestunden    
    df_time["hour"] = df_time["time"].dt.hour
    df_time["weekday"] = df_time["time"].dt.weekday # 0 = Montag, 6 = Sonntag
    
   
    # creating new Dataframe
    df_new = pd.DataFrame(columns=["Erfassungsdatum", "weekday", "Leserzahl"])
    # alle untersuchten ETs identifizieren
    list_days = df_time["Erfassungsdatum"].unique()
    
    for elem in list_days: 
        df_ = df_time[df_time["Erfassungsdatum"]==elem]
        #jeder Teilnehmer wird pro Tag nur einmal gezählt
        leserzahl = df_.WellenteilnahmenId.nunique()
        wochentag = df_.iloc[0].weekday
        # neue Daten werden mit pd.Series an das neue Dataframe angehängt
        df_new = df_new.append(pd.Series([elem, wochentag, leserzahl], index=df_new.columns), ignore_index=True)
    
    #Wichtig! Ansonsten hat der folgende Groupby keine Zahlen zum rechnen
    df_new["Leserzahl"] = pd.to_numeric(df_new["Leserzahl"]).astype(float)
    
    weekday_group = df_new.groupby("weekday", as_index=False).mean()
    week = ["Montag", "Dienstag", "Mittwoch", "Donnerstag", "Freitag", "Samstag", "Sonntag"]
    weekday_group["weekday"] = weekday_group["weekday"].apply(lambda x: week[x])
    
   
    
    x = weekday_group["weekday"]
    labels = x
    xn = range(len(x))
    
    y = weekday_group["Leserzahl"]
    
    
    
    
    plot_axis(prs, x=x, labels = labels, xn = xn, y = y, yn = 0,\
                 pos_left = 0.08, pos_right=0.92, pos_top=0.85, \
                 pos_bottom = 0.33, article = "total", grid=False,\
                 title_text=title_text, axis=1, mean_line = 0,\
                 legend="normal", special = "Lesetage")
    
    return prs

#%% Lesezeiten    
''' Diese Helper-Function erstellt Analysen über den Lesezeitraum (Lesezeitpuinkt am Tag. Sie benötigt 
das Dataframe ScanAuswertungen aus der Lesewert-Datenbank. Ansonsten kümmert sie 
sich um sich selbst. 

'''

def lesezeit(prs, df_scans, title_text="Um welche Uhrzeit wird gelesen?"):
    
    df_time = df_scans[df_scans["Treffer"]==True]
    
    
    # Zusammenführen der Datumsangaben
    
    df_time["time"] = df_time["Erfassungsdatum"] + "-" + df_time["Erfassungsuhrzeit"]
    
    #Erstelle Datetime-Format
    df_time["time"] = pd.to_datetime(df_time["time"], format="%Y-%m-%d-%H:%M:%S")
    
    # Filtere Jahre, nur 2018 erlaubt
    #mask = df_time["time"].dt.year == 2018
    #df_time = df_time[mask]

    # Erstellen einer Spalte mit den Lesestunden
    df_time["hour"] = df_time["time"].dt.hour
    
    
    # neues Dataframe mit 24 Columns, um für jeden einzelnen Erscheinungstag 
    # zu zählen
    df_new = pd.DataFrame(columns = np.arange(24))
    list_days = df_time["Erfassungsdatum"].unique()
    
    for elem in list_days:
        df_ = df_time[df_time["Erfassungsdatum"]==elem]
        daily_list = []
        for i in range (24):
            reader = df_[df_["hour"]==i].WellenteilnahmenId.nunique()
            daily_list.append(reader)
        df_new = df_new.append(pd.Series(daily_list, index=df_new.columns), 
                               ignore_index= True)
    
    # neues Dataframe um die Reader der einzelnen Tage pro Stunde zusammenzurechnen
    col_names = ["hour", "reader"] 
    df_sum = pd.DataFrame(columns=col_names)
    
    for i in range(24):
        hour = i
        reader = df_new[i].sum()
        #df_sum = df_sum.append(pd.Series([hour,reader], index=df_sum.columns), ignore_index= True)
        
        df_sum = df_sum.append(pd.Series([hour,reader], index=df_sum.columns), ignore_index= True)
    
    # converting df_sum-columns into int
    df_sum["reader"] = pd.to_numeric(df_sum["reader"]).astype(float)
    
    
    
   # Anzahl Erscheinungstage: 
    ET_tage = df_time["Erfassungsdatum"].nunique()
    
    
    x= df_sum["hour"]
    x_ = np.array(x)
    labels = x
    xn = range(len(x))
    y = df_sum["reader"] /ET_tage
    y_ = np.array(y)
    yn = df_sum["reader"] /ET_tage
    
    x_new = np.linspace(x_.min(), x_.max(),500)
    f = interp1d(x_, y_)
    y_smooth=f(x_new)
    
    
    
    
    # Seaborn-Style und Größe des Plots festlegen
    sns.set_style("white")
    fig, ax1 = plt.subplots(figsize=(20,8))
    
    #setzt die linke Y-Achse in Campton light
    # rechte Y-Achse können wir erst zum Code-Ende ansteuern
    plt.yticks(fontproperties=campton_light)
    
    #ax1.plot(x_new,y_smooth, color="#f77801", label="Lesewert")
    ax1.set_ylabel('Ø Leserzahl', color="#aeb0b2", fontproperties=campton_light, fontsize=50)
    ax1.set_xlabel('Uhrzeit', color="#aeb0b2", fontproperties=campton_light, fontsize=20)
    ax1.xaxis.set(ticks=range(0, len(xn))) # wir müssen die Anzahl der Ticks händisch festlegen
    ax1.set_xticklabels(labels = x, rotation=45, ha="right",  weight=800, color="#aeb0b2", \
                        fontproperties=campton_light, fontsize=30) # Labels werden ausgerichtet
                        #fontname=campton_light) # Labels werden ausgerichtet  
    fill = ax1.fill_between(xn, yn, alpha=1, color="#f77801",label="Ø Artikel/Tag") # Raum unter Yn-Linie wird gefüllt, Farbe wird transparent
    
    ax1.yaxis.label.set_size(22)
    
    ax1.spines["top"].set_visible(False)
    ax1.spines["left"].set_color("gray")
    ax1.spines["bottom"].set_color("gray")
    ax1.spines["right"].set_color("gray")
    ax1.spines["top"].set_visible(False)
    ax1.spines["bottom"].set_visible(False)
    ax1.spines["left"].set_visible(False)
    ax1.spines["right"].set_visible(False)
    
    # Zahlen an den Y-Achsen verändern, Größe und Farbe
    ax1.tick_params(axis='y', labelsize=25, colors="#aeb0b2", )
    
    
    plt.tight_layout()
    plt.subplots_adjust(left=0.08, right=0.92, top=0.85, bottom=0.3)
    filename = "grafik_zeiten3.jpg"
    plt.savefig(filename)
    
    picture_sheet(prs, filename, title_text= title_text)
    
    return prs

#%% LW in Zahlen - Textmarken
    # TODO  hier noch Funktion schreibne, die Zahl der Scans, 
    # Anzahl der Artikel, Artikel mit LW über 30%, ARtikel mit LW unter 5%, davon 
    # Zahl Ankündigungen und Veranstaltungsthemen



#%% Ort finden 
#Diese Funktion gleicht Listen mit Orten mit den Handlungsorten der Datensätze ab
#seite3["Ort"] = ""
#seite3["Handlungsorte"] == seite3["Handlungsorte"].fillna("x")
#def check_city(df):
#    
#    for i in df.index:
#        sentence = df.get_value(i, "Handlungsorte")
#       
#        if any(ext in sentence for ext in lw.orte_region):
#            df.set_value(i, "Ort", "Region")
#        elif any(ext in sentence for ext in lw.orte_deutschland):
#            df.set_value(i, "Ort", "National")
#        elif sentence == "x": 
#            df.set_value(i, "Ort", "keine Angaben")
#        else: 
#            df.set_value(i, "Ort", "International")
#check_city(data_analyse)

#%% LW-Marken - Textanzeige der Daten
    
''' Diese Funktion wirft lediglich die Werte als Text heraus. Funktioniert
nur für die Schwäbische, muss für jede Zeitung neu gemacht werden. 
'''

def marken_analyse(df, df_doublesplitid=False, df_nichtkum=False, kunde=""):
    # Werte Gesamt
    print("")
    print("__________________________________________________")
    print("--------------------------------------------------")
    print("WERTE GESAMT .... die hier nehmen!!!..ö STUTT WERTE GESAMT STZ und STN und Lokalteile")
    gesamt_lw = df['Artikel-Lesewert (Erscheinung) in %'].mean()
    gesamt_bw = df['Artikel-Blickwert (Erscheinung) in %'].mean()
    gesamt_dw = df['Artikel-Durchlesewerte (Erscheinung) in %'].mean()
    print ("LW gesamt: {:1.1f}%".format(gesamt_lw).replace(".", ","))
    print ("BW gesamt: {:1.1f}%".format(gesamt_bw).replace(".", ","))
    print ("DW gesamt: {:1.1f}%".format(gesamt_dw).replace(".", ","))
    
    print("WERTE GESAMT .... mit multiplen Split_IDs")
    gesamt_lw = df_doublesplitid['Artikel-Lesewert (Erscheinung) in %'].mean()
    gesamt_bw = df_doublesplitid['Artikel-Blickwert (Erscheinung) in %'].mean()
    gesamt_dw = df_doublesplitid['Artikel-Durchlesewerte (Erscheinung) in %'].mean()
    print ("LW gesamt: {:1.1f}%".format(gesamt_lw).replace(".", ","))
    print ("BW gesamt: {:1.1f}%".format(gesamt_bw).replace(".", ","))
    print ("DW gesamt: {:1.1f}%".format(gesamt_dw).replace(".", ","))
    
    
    
    print("Werte GESAMT NICHTKUMULIERT - für STUTT AUSGABEN, z.B. STN Filderstadt")
    gesamt_lw = df_nichtkum['Artikel-Lesewert (Erscheinung) in %'].mean()
    gesamt_bw = df_nichtkum['Artikel-Blickwert (Erscheinung) in %'].mean()
    gesamt_dw = df_nichtkum['Artikel-Durchlesewerte (Erscheinung) in %'].mean()
    print ("LW gesamt: {:1.1f}%".format(gesamt_lw).replace(".", ","))
    print ("BW gesamt: {:1.1f}%".format(gesamt_bw).replace(".", ","))
    print ("DW gesamt: {:1.1f}%".format(gesamt_dw).replace(".", ","))
    
    
    
    
    if kunde != "MHS":
    
        # Werte Lokales
        #df_l = df[(df["Ressortbeschreibung"]=="Lokales") | (df["Ressortbeschreibung"]=="Lokalsport")]
        print("----------------------end--------------------------")
        print("")
        print("")
        print("WERTE LOKALTEILE")
        # check ob df_doublesplitid vorhanden ist
        if isinstance(df_doublesplitid, pd.DataFrame):
            print("Datensatz mit doppelten SplitId gefunden....... die hier nehmen!!!")
            df_lokal = df_doublesplitid.copy()
            
        else: 
            print("Kein Datensatz mit doppelten SplitId gefunden")
            print("Analyse mit kumulierten Daten.")
            df_lokal = df.copy()
        print(df_lokal.shape)
        df_l = df_lokal[df_lokal["Ressortbeschreibung"]== "Lokales"]
        
        ausgaben = ausgaben_liste
           
        for ausg in ausgaben: 
            df_ = df_l[df_l["ZTG"]==ausg]
            lw = df_['Artikel-Lesewert (Erscheinung) in %'].mean()
            bw = df_['Artikel-Blickwert (Erscheinung) in %'].mean()
            dw = df_['Artikel-Durchlesewerte (Erscheinung) in %'].mean()
            print ("LW Lokalteil " + ausg +": {:1.1f}%".format(lw).replace(".", ","))
            print ("BW Lokalteil " + ausg +": {:1.1f}%".format(bw).replace(".", ","))
            print ("DW Lokalteil " + ausg +": {:1.1f}%".format(dw).replace(".", ","))
        
        
        print("Lokale Werte mit nchtkumulierten Daten")   
        print(df.shape)
        df_lok = df_nichtkum.copy()
        df_lok = df_lok[df_lok["Ressortbeschreibung"]=="Lokales"]
        df_lok2 = df[df["Ressortbeschreibung"]=="Lokales"]
        
        for ausg in ausgaben: 
            
            df__ = df_lok[df_lok["ZTG"]==ausg]
            lw = df__['Artikel-Lesewert (Erscheinung) in %'].mean()
            bw = df__['Artikel-Blickwert (Erscheinung) in %'].mean()
            dw2 = df__['Artikel-Durchlesewerte (Erscheinung) in %'].mean()
            print ("LW Lokalteil " + ausg +": {:1.1f}%".format(lw).replace(".", ","))
            print ("BW Lokalteil " + ausg +": {:1.1f}%".format(bw).replace(".", ","))
            print ("DW Lokalteil " + ausg +": {:1.1f}%".format(dw2).replace(".", ","))
        print("-------------------------end-------------------")
        print("")
        print("")
        print("Letzer Lok mit kumulierten Daten .... ODER die hier nehmen, gleiches ERgebnis")
        for ausg in ausgaben: 
            
            df___ = df_lok2[df_lok2["ZTG"]==ausg]
            lw = df___['Artikel-Lesewert (Erscheinung) in %'].mean()
            bw = df___['Artikel-Blickwert (Erscheinung) in %'].mean()
            dw2 = df___['Artikel-Durchlesewerte (Erscheinung) in %'].mean()
            print ("LW Lokalteil " + ausg +": {:1.1f}%".format(lw).replace(".", ","))
            print ("BW Lokalteil " + ausg +": {:1.1f}%".format(bw).replace(".", ","))
            print ("DW Lokalteil " + ausg +": {:1.1f}%".format(dw2).replace(".", ","))
        
        # Mantel
         
         
        print("WERTE MANTEL")
        mantel = df[df["Ressortbeschreibung"].isin(mantel_ressorts)]
        
        mantel_lw = mantel['Artikel-Lesewert (Erscheinung) in %'].mean()
        mantel_bw = mantel['Artikel-Blickwert (Erscheinung) in %'].mean()
        mantel_dw = mantel['Artikel-Durchlesewerte (Erscheinung) in %'].mean()
        print ("LW Mantel: {:1.1f}%".format(mantel_lw).replace(".", ","))
        print ("BW Mantel: {:1.1f}%".format(mantel_bw).replace(".", ","))
        print ("DW Mantel: {:1.1f}%".format(mantel_dw).replace(".", ","))
        print("-------------------------end-------------------")
        print("")
        print("")
        
        
        print("WERTE EINZELAUSGABEN (NICHTKUM)")
        
        # Check ob nichtkumulierte Werte übergeben worden sind. 
        # Ansonsten nehmen wir die kumulierten
        # Check über Type und isinstance()
        
        if isinstance(df_nichtkum, pd.DataFrame):
            print("Eigener Datensatz für nichtkumulierte Werte gefunden...DIE HIER NEHMEN...")
            df_total = df_nichtkum.copy()
            
        else: 
            print("Keine nichtkumulierten Daten vorhanden.")
            print("Analyse mit kumulierten Daten.")
            df_total = df.copy()
            
            
            
            
        for ausg in ausgaben: 
            df_ = df_total[df_total["ZTG"]==ausg]
            lw = df_['Artikel-Lesewert (Erscheinung) in %'].mean()
            bw = df_['Artikel-Blickwert (Erscheinung) in %'].mean()
            dw = df_['Artikel-Durchlesewerte (Erscheinung) in %'].mean()
            
            print("LW Ausgabe " +ausg + "gesamt:  {:1.1f}%".format(lw).replace(".", ","))
            print("BW Ausgabe " +ausg + "gesamt:  {:1.1f}%".format(bw).replace(".", ","))
            print("DW Ausgabe " +ausg + "gesamt:  {:1.1f}%".format(dw).replace(".", ","))
         
            print("-------------------------WErte für Ausgaben mit kumulierten Werten (Mantel) und doppelten Splits------------")
            
            df_2 = df_doublesplitid[df_doublesplitid["ZTG"]==ausg]
            lw = df_2['Artikel-Lesewert (Erscheinung) in %'].mean()
            bw = df_2['Artikel-Blickwert (Erscheinung) in %'].mean()
            dw = df_2['Artikel-Durchlesewerte (Erscheinung) in %'].mean()
            print("LW Ausgabe " +ausg + "gesamt:  {:1.1f}%".format(lw).replace(".", ","))
            print("BW Ausgabe " +ausg + "gesamt:  {:1.1f}%".format(bw).replace(".", ","))
            print("DW Ausgabe " +ausg + "gesamt:  {:1.1f}%".format(dw).replace(".", ","))
            print()
            print()
           
        print("-------------------------end-------------------")
        print("")
        print("")


#%% Analyse Ressrots/Seitentitel nach Lesewert
    
'''
Diese Helperfunction wirft eine Grafik zum Lesewert und der Erscheinungshäufigkeit
einzelner Ressorts oder Seitentitel heraus. 


Target="ressort" ist die Voreinstellung. Sollen Seitentitel analysiert werden, 
target="seitentitel"

Die Funktion benötigt für Target="ressort" nur den normalen, bereinigte Datensatz. 
Für Target = "seitentitel" benötigen wir den Datensatz auf Ressortebene. 

special: Hier fasse ich verschiedene Anwendungsmöglichkeiten zusammen. 
1 = Grafik mit einer Y-Achse links, LW für einzelne Ressorts, Linie mit allg.DS
2 = Grafik mit zwei Y-Achsen, Artikelhäufigkeit pro Ressort (absolute Zahl)



minimize, Standard 5, zeigt an ab welcher Zahl von Artikeln ein Ressort/Seitentitel
in die Auswertung aufgenommen wird.

label_position gibt an, wie viel Platz die Labels unten und links bekommen. 
Voreingestellt auf normal. Auch möglich:  large und xlarge für besonders lange 
Labelnamen. strange sind spezialformate

Über die ressort_liste kann der Funktion noch explizit mitgeteilt werden, 
welche Ressorts/Seitentitel in die Analyse einfließen sollen

Unter Special können Spezialfälle für einzelne Zeitungstitel angelegt werden. 

order = "ok" meint, dass wir an der Sortierung der Ressorts nach Seitenzahl /Median
nicht mehr händisch eingreifen müssen. order = "new" heißt, dass wir manuel eingreifen

mean_line = der Wert für eine eventuelle Durchschnittslinie muss der Funktion am 
Beginn mitgegeben werden, da sich der Durchschnitt nicht zwangsläufig aus dem 
Untersuchungs-DF ergibt. 

Anzahl Lokales = Gibt die Anzahl der Lokalteile an, durch die bei der Gesamt-Betrachtung
die Artikelzahl Lokales (bzw. der Ressorts, die unter lokale_liste stehen, s.u.) geteilt werden soll. 
In der Regel könnte man hier die Länge der Ausgaben_liste angeben... aber lieber händisch
'''

# TODO Aus irgendeinem Grund lässt sich die Liste lokale_ressorts nicht ohne weiteres als parameter benutzen. 
# wenn ich in der Funktion loklist = lokale_ressorts setze, bleibt loklist leer... obwohl lokale_ressorts gefüllt sind. ????
#%%Kolumen-Finder
# Diese Funktion durchsucht einen Dataframe nach Kolumnen und spielt diese in einer Lesewertgrafik aus. 

# anzahl = hier kann die Mindestanzahl an Veröffentlichungen angegeben werden. 
# ausnahme = eine List mit Ausnahmen, z.B. einzelnen Kolumnen wir "Kommentar", die man nicht drin haben möchte
# limit = number of kolumns

def kolumnen_finder(prs, df, anzahl=5, ausnahmem=[], axis = 2, label_position = "large", title_text = "Kolumnen", grid=True, limit=10):
    
    # neues Dataframe wird erstellt, hier werden die Kolumnen später aufgelistet
    columns = ["Name", "LW", "BW", "DW", "Anzahl"]
    df_kolumne = pd.DataFrame(columns=columns)
    
    liste = df.Beschreibung.unique()
    for elem in liste: 
        de = df[df["Beschreibung"]==elem]
        shape = de.shape[0]
        lw = de["Artikel-Lesewert (Erscheinung) in %"].mean()
        bw = de["Artikel-Blickwert (Erscheinung) in %"].mean()
        dw = de["Artikel-Durchlesewerte (Erscheinung) in %"].mean()
    
    
        df_kolumne = df_kolumne.append({"Name":elem, "LW":lw, "BW":bw, "DW":dw, "Anzahl":shape}, ignore_index=True)
   
    df_kolumne.sort_values(by="LW", ascending=False, inplace=True)
    df_kolumne = df_kolumne[df_kolumne["Anzahl"]>=anzahl]
    df_kolumne = df_kolumne.head(limit)
    #  Grafik-Funktion wird aktiviert
         
   
    x  = df_kolumne["Name"]
    

    labels = x
    xn = range(len(x))
    y = df_kolumne["LW"]
    et = df.Erscheinungsdatum.nunique()
    yn = df_kolumne["Anzahl"]
    #yn_total = df_kolumne["Anzahl"]
    #yn_mean = df_kolumne["Anzahl"] / et
   # Titel der Grafik einstellen
   
   
   
       
    p_left = 0.08
    p_right = 0.92
    p_top = 0.85
    p_bottom = 0.33
   
    if label_position == "large": 
       p_left = 0.12
       p_right = 0.92
       p_top = 0.85
       p_bottom = 0.4
   
    if label_position == "xlarge": 
       p_left = 0.14
       p_right = 0.92
       p_top = 0.85
       p_bottom = 0.50

    if label_position == "xxlarge": 
       p_left = 0.14
       p_right = 0.92
       p_top = 0.85
       p_bottom = 0.55
       
    if label_position == "strange":
       p_left = 0.14
       p_right = 0.80
       p_top = 0.75
       p_bottom = 0.45
    
   
   
   
   
   # Grafik erstellen
    plot_axis(prs, x=x, labels = labels, xn = xn, y = y, \
                  pos_left = p_left, pos_right=p_right, pos_top=p_top, \
                  pos_bottom = p_bottom, grid=grid,\
                  title_text=title_text, axis=1)
    
   #Grid-Linien für die zweite Grafik ausschalten
   
   # STEP 2 - Berechnung durchschnittlicher Lesewert plus Häufigkeit je Ressort
   # Grafik erstellen
    if axis==2:
       
        plot_axis(prs, x=x, labels = labels, xn = xn, y = y, yn=yn, \
                 pos_left = p_left, pos_right=p_right, pos_top=p_top, \
                 pos_bottom = p_bottom, grid=grid,\
                 title_text=title_text, axis=2, mean_line = 0)
   
   
    
#%%
    

#%%
    

def grafik_lesewert(prs, df, target="ressort", minimize=5, label_position="normal",
                    ressort_liste=[], special=False, title_text="", order="ok", legend="normal", sort="Seitennummer", 
                    article="total", mean_line = 0, grid=False, anzahl_lokales=1, mean_line_title = "", axis=2, 
                    limit_y = False, extension_df = 0, font_size=30): 
    
  
   print("Hier ist grafik Lesewerte 2553... die Anzahl der Datensätze beträgt..." + str(df.shape[0]))
   
   group_param=""
   
   #Festlegung welche Werte hier gesucht werden
   
   
   
   if target=="ressort":
       group_param = "Ressortbeschreibung"
       mean_line_title = "Ø LW Gesamt"
   
   elif target == "Textlänge": 
       group_param = "Textlänge"
       mean_line_title = "Ø LW Gesamt"
   elif target == "darstellungsformen": 
       group_param = "Darstellungsformen"
       mean_line_title = "Ø LW Gesamt"
   elif target == "platzierungen":
       group_param = "Platzierungen"
       mean_line_title = "Ø LW Gesamt"
   elif target == "seitentitel":
       group_param = "Seitentitel"
       mean_line_title = "Ø LW Gesamt"
   elif target == "rubriken": 
       group_param="Beschreibung"
       mean_line_title = "Ø LW Gesamt"
   elif target == "ort": 
       group_param = "Ort"
       mean_line_title = "Ø LW Gesamt"
       sort = ["Region", "National", "International"]
   elif target == "special": 
       group_param = "Spezial"
       mean_line_title = "Ø LW Gesamt"
   elif target == "special_res":
       
       group_param = "Spezial_Res"
       mean_line_title = "Ø LW Gesamt"
   else: 
       group_param = target
       mean_line_title = "Ø LW Gesamt"
     
   # Abfrage, ob eine Auswahl-Liste mitgegeben wurde. Falls ja: DF bereinigen.      
   if len(ressort_liste)>0:
       
       df = df[df[group_param].isin(ressort_liste)]    
  
   
   df_group = df.groupby(group_param, as_index=False)
   df_group = df_group.agg({"SplitId":"size", "Seitennummer":"median", "Artikel-Lesewert (Erscheinung) in %":"mean"})
  
   df_group = df_group[df_group["SplitId"]>=minimize]
  
   #WErte für Lok1 (in einer Series) werden manuell übergeben
   if isinstance(extension_df, pd.Series):
       print("Series als Extension erkannt")
       df_group = df_group.append(extension_df, ignore_index=True)
   if isinstance(extension_df, pd.DataFrame):
       print("Dataframe als Extension erkannt")
       df_group = df_group.append(extension_df, ignore_index=True)
   print("--------------------------") 
   print(" und das hier ist das erweiterte DF")
   print(df_group.head(20))
   
   
   
   
   
   if sort == "Seitennummer": 
       #niedrigste Seitennummer nach vorne
       df_group = df_group.sort_values(by=sort)
   if sort == "Lesewert":
       # hnöchster Lesewert nach vorne
       df_group = df_group.sort_values(by="Artikel-Lesewert (Erscheinung) in %", ascending=False)
   
   if sort == ["Region", "National", "International"]:
       df_group.Ort = df_group.Ort.astype("category")
       df_group.Ort.cat.set_categories(sort, inplace=True)
       df_group.sort_values(["Ort"], inplace = True)
   
   
   if sort == "Kategorie":
       
       df_group[group_param] = df_group[group_param].astype("category")
       df_group[group_param].cat.set_categories(ressort_liste, inplace=True)
       df_group.sort_values([group_param], inplace = True)
   
  
  
       
   # Reihenfolge ändern, damit Lok 1 vorne steht:
   if order == "new": 
       list_lok = df.Seitentitel.unique()
       if "Biberach" in list_lok:
           
           df_group = df_group.reindex([4,2,1,5,8,7,9, 6, 3])
           
       elif "Friedrichshafen" in list_lok:
           print("------------- TEST FHA -------------------------")
           
           
           
           df_group = df_group.reindex([4,1,3,6,5,9,8,2])
           
       
  
   df_group.dropna(inplace=True)
  
   # Lokale Werte durch Anzahl der ausgewerteten Lokalteile teilen... 
   if anzahl_lokales!=1:
      
      
       for i in range(len(lokale_ressorts)):
           
           #mask = df_group[group_param] == lokale_ressorts[i]
          
           #mask = lokale_ressorts[i]
           # für die NW
           mask = df_group[group_param] == "weitere Lokalseiten"
           a = df_group.loc[mask, "SplitId"]
           b = df_group.loc[mask, "SplitId"] / anzahl_lokales
           df_group.loc[mask, "SplitId"] = df_group.loc[mask, "SplitId"] / anzahl_lokales
           c = df_group.loc[mask, "SplitId"]
          
   df_group.reset_index(drop=True, inplace=True)
  
    # STEP 1 - Berechnung des durchschnittlichen Lesewerts aus den Einzelwerten
    # der Ressorts(Seitentitel)
   
   lw_durchschnitt = df_group["Artikel-Lesewert (Erscheinung) in %"].mean()
   
  
   # SPECIAL für Schwäbisch
   # Werte für Titelseite und Lokales werden gedrittelt
   # Grund: Da bei Titel und Lokalem jeweils alle drei 
   
   
   # TODO Viel zu gefärhlich... man muss die Index-Nummer der jeweiligen 
   # Werte angeben... totale Fehlerquelle
   
   if special==True: 
       #df_group[df_group["Ressortbeschreibung"]=="Lokales"]["]
       mask = df_group["Ressortbeschreibung"]=="Lokales"
       val1 = df_group[mask].SplitId
       val_lok = val1 / 3
       
       df_group.set_value(10, "SplitId", val_lok)
       
       mask2 = df_group["Ressortbeschreibung"]=="Titel"
       val2 = df_group[mask2].SplitId
       val_titel = val2 / 3
       df_group.set_value(0, "SplitId", val_titel)
       
       mask3 = df_group["Ressortbeschreibung"]=="Lokalsport"
       val3 = df_group[mask3].SplitId
       val_loksport = val3/3
       df_group.set_value(11, "SplitId", val_loksport)
       
       
      
     
    
   df_group_copy = df_group.copy()
   ersch = df.Erscheinungsdatum.nunique()
   df_group_copy["SplitId"] = df_group_copy["SplitId"].apply(lambda x: round(x/ersch, 2))
  
    # Werte für Achsen festlegen
   #Achsen-Werte festlegen
   x = df_group[group_param]
   
   if target == "ressort":
       
         # schiebt die beiden ersten Lokalseiten im Mantel ans Ende der lokalen 
       # SEiten, vor allem hinter die Lok 1
   
    
        df_group["shortnames"] = df_group[group_param].apply(lambda x: ressort_dict[x])
       
        labels = df_group["shortnames"]
        
   elif target ==  "special_res":
       
    
        df_group["shortnames"] = df_group[group_param].apply(lambda x: ressort_dict[x])
       
        labels = df_group["shortnames"]
        
    
   elif target == "rubriken":
       
       try: 
           
           df_group["shortnames"] = df_group[group_param].apply(lambda x: kolumnen_dict[x])
           
       except KeyError:
           df_group["shortnames"] = df_group[group_param]
       labels = df_group["shortnames"]
   
    
   # Seitentitel sind oft zu lang, daher ersetzen wir die Labels mit Kurzversionen    
   elif (target == "seitentitel") | (target == "special") :
     
       def checkseitentitel(e): 
           try:
               return seitentitel_dict[e]
           except: 
               return e
       df_group["shortnames"] = df_group[group_param].apply(checkseitentitel) 
       labels = df_group["shortnames"]
       
   
   elif target == "darstellungsformen": 
       df_group["shortnames"] = df_group[group_param].apply(lambda x: darstellung_dict[x])
       labels = df_group["shortnames"]
   
   elif target == "platzierungen": 
       df_group["shortnames"] = df_group[group_param].apply(lambda x: platzierung_dict[x])
       labels = df_group["shortnames"]
       
   elif target == "ort": 
       labels = df_group["Ort"]

   else: 
       labels = df_group[group_param]
    
   xn = range(len(x))
   y = df_group["Artikel-Lesewert (Erscheinung) in %"]
   et = df.Erscheinungsdatum.nunique()
   yn_total = df_group["SplitId"]
   yn_mean = df_group["SplitId"] / et
   
   # Titel der Grafik einstellen
   
   
   if target=="ressort": 
       title = "Lesewert nach Ressorts"
       

       
   elif target == "seitentitel":
       title = "Lesewert nach Seitentiteln"
   
   if len(title_text)>1:
       title = title_text
    
   # Postion und Lage der Labels auf der X-Achse bestimmen    
   
       
   p_left = 0.08
   p_right = 0.92
   p_top = 0.85
   p_bottom = 0.33
   
   if label_position == "large": 
       p_left = 0.12
       p_right = 0.92
       p_top = 0.85
       p_bottom = 0.4
   
   if label_position == "xlarge": 
       p_left = 0.14
       p_right = 0.92
       p_top = 0.85
       p_bottom = 0.50

   if label_position == "xxlarge": 
       p_left = 0.14
       p_right = 0.92
       p_top = 0.85
       p_bottom = 0.55
       
   if label_position == "strange":
       p_left = 0.14
       p_right = 0.80
       p_top = 0.75
       p_bottom = 0.45
    
   
   # Berechnung Durchschnitt für eventuelle Durchschnittslinie
   if mean_line !=0:
       mean_line = mean_line
   
   if article=="total": 
       yn = yn_total
   else:
       yn = yn_mean
   
   # Grafik erstellen
   plot_axis(prs, x=x, labels = labels, xn = xn, y = y, \
                  pos_left = p_left, pos_right=p_right, pos_top=p_top, \
                  pos_bottom = p_bottom, article = article, grid=grid,\
                  title_text=title, axis=1, mean_line=mean_line, 
                  mean_line_title = mean_line_title, legend=legend, limit_y=limit_y, 
                  font_size = font_size)
    
   #Grid-Linien für die zweite Grafik ausschalten
   
   # STEP 2 - Berechnung durchschnittlicher Lesewert plus Häufigkeit je Ressort
   # Grafik erstellen
   if axis==2:
       
       plot_axis(prs, x=x, labels = labels, xn = xn, y = y, yn=yn, \
                 pos_left = p_left, pos_right=p_right, pos_top=p_top, \
                 pos_bottom = p_bottom, article = article, grid=grid,\
                 title_text=title, axis=2, mean_line = 0, legend=legend, limit_y=limit_y, 
                 font_size = font_size)
   
   
   
   
   return prs 
 
    
#%% Ressorts nach Geschlecht
   # TODO: Daten fehlen noch

#%% Wochentage nach Lesewert und durchschnittlicher Artikelanzahl pro Tag
#Func Plot von Lesewert nach Erscheiungstag
'''
Lesewert nach Erscheinungstag

Diese Funktion  übernimmt einen vorgefilterten Datensatz (zum Beispiel alle 
Artikel eines Ressorts etc.) und zeigt dann alle Erscheinungstage mit 
durchschnittlichem Lesewert und durchschnittliche Anzahl der erschienen Artikel 
an diesem Tag. 

Optional kann eine neue Überschrift vergeben werden. 



'''

def lesewert_erscheinung(df,df_multi_split, prs, title_text="Lesewert nach Wochentagen", grid=False):
    
    df_ = df.copy()
    df_multi = df_multi_split.copy()
    
    erscheinungstage = df_["Erscheinungsdatum"].nunique()
    
    df_["Tag"] = df_["Erscheinungsdatum"].apply(lambda x: x.strftime('%A'))
    df_multi["Tag"] = df_multi["Erscheinungsdatum"].apply(lambda x: x.strftime('%A'))
    
    df_ = df_[["SplitId", "Ressortbeschreibung", "Erscheinungsdatum", "weekday", "Tag", \
               "Artikel-Lesewert (Erscheinung) in %", "Seitennummer", "Darstellungsformen"]]
    
    df_multi = df_multi[["SplitId", "Ressortbeschreibung", "Erscheinungsdatum", "weekday", "Tag", \
               "Artikel-Lesewert (Erscheinung) in %", "Seitennummer", "Darstellungsformen"]]
    
    # Werte für kumulierten Datensatz - hier filtern wir den durchschnittlichen Lesewert heraus
    df_number = df_.groupby(["Tag", "Erscheinungsdatum"], as_index=False).count()
    df_2 = df_number.groupby("Tag", as_index = False).count()
    df_2 = df_2.rename(columns={"Erscheinungsdatum": "Tageszahl"})
    df_2 = df_2[["Tag", "Tageszahl"]]
    # Groupby um Anzahl der Artikel und durchschnittliche Lesewerte zu erhalten
    df_group = df_.groupby("Tag").agg({"SplitId":"size", "Artikel-Lesewert (Erscheinung) in %":"mean", \
                                           "Erscheinungsdatum":"count"})
        
    df_group = df_group.reindex(["Montag", "Dienstag", "Mittwoch", "Donnerstag", "Freitag", "Samstag"])
    
    # Columns ordnen
    df_group = df_group.rename(columns={"SplitId":"Artikelanzahl", "Artikel-Lesewert (Erscheinung) in %":"LW"})
    df_group['Tag'] = df_group.index
    
    #Verschmelzung der beiden Tabellen
    df_group = pd.merge(df_group, df_2, on="Tag", how="left")
    df_group.dropna(inplace=True)
    
    
    # Werte für nichtkumulierten Datensatz - hier filtern wir den durchschnittlichen Lesewert heraus
    df_number_multi = df_multi.groupby(["Tag", "Erscheinungsdatum"], as_index=False).count()
    df_2_multi = df_number_multi.groupby("Tag", as_index = False).count()
    df_2_multi = df_2_multi.rename(columns={"Erscheinungsdatum": "Tageszahl"})
    df_2_multi = df_2_multi[["Tag", "Tageszahl"]]
    # Groupby um Anzahl der Artikel und durchschnittliche Lesewerte zu erhalten
    df_group_multi = df_multi.groupby("Tag").agg({"SplitId":"size", "Artikel-Lesewert (Erscheinung) in %":"mean", \
                                           "Erscheinungsdatum":"count"})
        
    df_group_multi = df_group_multi.reindex(["Montag", "Dienstag", "Mittwoch", "Donnerstag", "Freitag", "Samstag"])
    
    # Columns ordnen
    df_group_multi = df_group_multi.rename(columns={"SplitId":"Artikelanzahl", "Artikel-Lesewert (Erscheinung) in %":"LW"})
    df_group_multi['Tag'] = df_group_multi.index
    
    #Verschmelzung der beiden Tabellen
    df_group_multi = pd.merge(df_group_multi, df_2_multi, on="Tag", how="left")
    df_group_multi.dropna(inplace=True)
    
    
    
    
    
    
    #Achsen-Werte festlegen
    x = df_group["Tag"]
    labels = x
    xn = range(len(x))
    y = df_group["LW"]
    #Ausnahme für Gütersloher Kultur und Mitteldeutsche
    yn = df_group_multi["Artikelanzahl"]  / df_group["Tageszahl"]
    if (df.iloc[0].Seitentitel=="Gütersloher Kultur") or (df.iloc[0].Ressortbeschreibung=="Sport") : 
        yn = df_group["Artikelanzahl"] / df_group["Tageszahl"]
    
    
    
    plot_axis(prs, x=x, labels = labels, xn = xn, y = y, yn = yn,\
                  pos_left = 0.08, pos_right=0.92, pos_top=0.85, \
                  pos_bottom = 0.33, article = "mean", grid=grid,\
                  title_text = title_text, axis=2)


    return prs






   
#%% Lesewert/Artikelzahl nach Darstellungsformen

''' 
Die Funktion 
darstellungsformen()
zeichnet eine Grafik mit zwei y-Achsen zu den Darstellungsformen und der Anzahl 
der Artikel in diesen Formen. 

Der Parameter minimum sagt aus wie häufig eine Darstellungsform vorkommen muss, um 
gemessen zu werden. 
By default steht der Wert auf 5. 

Ansonsten muss der User nur das entsprechende Dataframe eingeben, das 
ausgewertet werden soll. 

'''

def darstellungsformen(prs, df, minimum = 5, geschlecht=False, title_text="Bitte eingeben", special="", grid=False):
    # liste heißt Darstellu´ngsform
    mask = df["Darstellungsformen"].notnull()
    df = df[mask]
    
    #Diese Funktion spaltet die Darstellungsformen-Zellen auf
    # außerdem werden hier die beiden Darstellungsformen  BB und BN zu BI zusammengefasst
    def splititup(elem):
        splitter = [x.strip() for x in elem.split(',')]
        
        for x in splitter: 
            if x in darstellungsform: 
                if x == "BB" or x == "BN":
                    return "BI"
                else: 
                    return x
            else:
                pass

    #neue Spalte DF wird angelegt, hier werden erwünschte Darstellungsformen verzeichnet
    df["DF"] = df["Darstellungsformen"].apply(splititup) 
    
    
    if geschlecht == False: 
        # Tabelle wird verkleinert, Zeilen ohne Eintrag in DF werden gelöscht
        df_ = df[["SplitId", "Erscheinungsdatum", "Darstellungsformen", "DF", "Artikel-Lesewert (Erscheinung) in %"]]
        
    if geschlecht == True:  
        df_ = df[["SplitId", "Erscheinungsdatum", "Darstellungsformen", "DF", "LW_w", "LW_m", "Artikel-Lesewert (Erscheinung) in %"]]
    
    mask_df = df_["DF"].notnull()
    
    df_ = df_[mask_df]
    
    if special == "kommis_nw":
        
        def cat_kommis(df): 
            for i in df.index:
                date = df.get_value(i, "Erscheinungsdatum")
                
                if date < pd.to_datetime("2018-05-24"):
                    df.set_value(i, "DF", "Kommentar bis\nzum 23. Mai")
                else: 
                    df.set_value(i, "DF", "Kommentar ab\ndem 24. Mai")
        cat_kommis(df_)

    #Anzahl Erscheinungstage wird ermittelt
    erscheinungstage = df_["Erscheinungsdatum"].nunique()
    
    
    if geschlecht == False :
        
        # Tabelle mit Darstellungsform, Anzahl Artikel und Durchschnittslesewert wird erstellt
        df_DF = df_.groupby(["DF"], as_index=False).agg({"SplitId":"size", "Artikel-Lesewert (Erscheinung) in %":"mean"})\
        .sort_values(by="Artikel-Lesewert (Erscheinung) in %", ascending=False)
    
    if geschlecht == True:
        df_DF = df_.groupby(["DF"], as_index=False).agg({"SplitId":"size", "LW_w":"mean", "LW_m":"mean",
                           "Artikel-Lesewert (Erscheinung) in %":"mean"})\
        .sort_values(by="Artikel-Lesewert (Erscheinung) in %", ascending=False)
        
    
    df_DF = df_DF[df_DF["SplitId"]>=minimum]
    
    shape_tester = df_DF.shape[0]
    # GRAFIK ANFERTIGEN
    print(df_DF.head())
    
    # Titel
    if title_text == "Bitte eingeben":
        title = "Darstellungsformen und Lesewerte"
    else: 
        title = title_text
    
    
    if shape_tester >0: 
        if geschlecht == False: 
            
        
        
            # Achsen festlegen
            x= df_DF["DF"]
            labels = df_DF["DF"].apply(lambda x: darstellung_dict[x])
            if special == "kommis_nw": 
                labels = df_DF["DF"]
                
            else: 
                labels = df_DF["DF"].apply(lambda x: darstellung_dict[x])
            print(labels)
            xn = np.arange(len(x)) #brauchen wir, weil Matplot keine Strings zur X-Achse verarbeiten kann, nur Zahlen
            y = df_DF["Artikel-Lesewert (Erscheinung) in %"]
            yn = df_DF["SplitId"]
        
            #Zeichenfunktion wird aufgerufen
            if special == "kommis_nw":
                plot_axis(prs, x=x, labels = labels, xn = xn, y = y, \
                      pos_left = 0.12, pos_right=0.92, pos_top=0.85, \
                      pos_bottom = 0.4, article = "total", grid=grid,\
                      title_text = title, axis=1, umbruch_x=False)
                
            else: 
                plot_axis(prs, x=x, labels = labels, xn = xn, y = y, yn = yn,\
                      pos_left = 0.12, pos_right=0.92, pos_top=0.85, \
                      pos_bottom = 0.4, article = "total", grid=grid,\
                      title_text = title, umbruch_x=False)
            
        if geschlecht == True: 
            
            #Achsen festlegen
            # Werte festlegen    
            x = df_DF["DF"]
            labels = df_DF["DF"].apply(lambda x: darstellung_dict[x])
            xn = np.arange(len(x))
            y_w = df_DF["LW_w"]
            y_m = df_DF["LW_m"]
            
            multiple_bars(prs, x=x, labels=labels, xn=xn,y_w=y_w, y_m = y_m, 
                          title_text=title_text, grid=grid)
        
    
    
    return prs

#%% alle Kategorien - Bar-Charts
    
# dise Funktion sorgt für die Barcharts einzelner Kategorien- 
def auswahl_kategorien(prs, df, minimum = 5, title_text="Bitte eingeben", special=""):
    print(df.iloc[0]["Ressortbeschreibung"])
    
    liste = rubrik_dict[df.iloc[0]["Ressortbeschreibung"]]
    print(len(liste))
    print(liste)
    print("------")
    df_ = df[df["Beschreibung"].isin(liste)]
    df_group = df_.groupby(["Beschreibung"], as_index=False).agg({"SplitId":"size", "Artikel-Lesewert (Erscheinung) in %":"mean"})\
        .sort_values(by="Artikel-Lesewert (Erscheinung) in %", ascending=False)
    
    
    df_group = df_group[df_group["SplitId"]>=minimum]
    
    shape_tester = df_group.shape[0]
    # GRAFIK ANFERTIGEN
    
    
    # Titel
    if title_text == "Bitte eingeben":
        title = "Rubriken " + df_.iloc[0].Ressortbeschreibung
    else: 
        title = title_text
    
    
    if shape_tester >0: 

        
            # Achsen festlegen
            x= df_group["Beschreibung"]
            if special == "kommis_nw": 
                labels = df_DF["DF"]
                
            else: 
                labels = df_group["Beschreibung"]
            
            xn = np.arange(len(x)) #brauchen wir, weil Matplot keine Strings zur X-Achse verarbeiten kann, nur Zahlen
            y = df_group["Artikel-Lesewert (Erscheinung) in %"]
            yn = df_group["SplitId"]
        
            #Zeichenfunktion wird aufgerufen
            if special == "kommis_nw":
                plot_axis(prs, x=x, labels = labels, xn = xn, y = y, \
                      pos_left = 0.12, pos_right=0.92, pos_top=0.85, \
                      pos_bottom = 0.4, article = "total", grid=True,\
                      title_text = title, axis=1)
                
            else: 
                plot_axis(prs, x=x, labels = labels, xn = xn, y = y, yn = yn,\
                      pos_left = 0.12, pos_right=0.92, pos_top=0.85, \
                      pos_bottom = 0.4, article = "total", grid=True,\
                      title_text = title, axis=1)
            
        
    
    return prs

#%% Lesewert/Artikelzahl nach Platzierungen

def platzierungen(prs, df, minimum = 5, geschlecht=False, title_text="Bitte eingeben", special=""):
    # alle Null-Werte ausblenden
    print("Hier ist Platzierungen, Z 3164, Dataframe hat {} Zeilen".format(df.shape[0]))
    
    mask = df["Platzierungen"].notnull()
    df = df[mask]
    
    if geschlecht == False: 
        # Tabelle wird verkleinert, Zeilen ohne Eintrag in DF werden gelöscht
        df_ = df[["SplitId", "Erscheinungsdatum", "Platzierungen", "Artikel-Lesewert (Erscheinung) in %"]]
        
    if geschlecht == True:  
        df_ = df[["SplitId", "Erscheinungsdatum", "Platzierungen", "LW_w", "LW_m", "Artikel-Lesewert (Erscheinung) in %"]]

    
    if special == "kommis_nw": 
        def cat_kommis(df): 
            for i in df.index:
                date = df.get_value(i, "Erscheinungsdatum")
                if date > "2018-05-23":
                    df.set_value(i, "Platzierungen", "Kommentar \nvor dem 24. Mai")
                else: 
                    df.set_value(i, "Platzierungen", "Kommentar \ab dem 24. Mai")
    
    
    #Anzahl Erscheinungstage wird ermittelt
    #erscheinungstage = df_["Erscheinungsdatum"].nunique()
    
    
    if geschlecht == False :
        
        # Tabelle mit Darstellungsform, Anzahl Artikel und Durchschnittslesewert wird erstellt
        df_pl = df_.groupby(["Platzierungen"], as_index=False).agg({"SplitId":"size", "Artikel-Lesewert (Erscheinung) in %":"mean"})\
        .sort_values(by="Artikel-Lesewert (Erscheinung) in %", ascending=False)
    
    if geschlecht == True:
        df_pl = df_.groupby(["Platzierungen"], as_index=False).agg({"SplitId":"size", "LW_w":"mean", "LW_m":"mean",
                           "Artikel-Lesewert (Erscheinung) in %":"mean"})\
        .sort_values(by="Artikel-Lesewert (Erscheinung) in %", ascending=False)
        
    
    df_pl = df_pl[df_pl["SplitId"]>=minimum]
    print("Hier ist Func Platzierungen, Zeiel 3202. Yn darf nicht 0 sein, ist hier {}".format(df_pl["SplitId"]))
    # GRAFIK ANFERTIGEN
    if geschlecht == False: 
        
    
    
        # Achsen festlegen
        x= df_pl["Platzierungen"]
        labels = df_pl["Platzierungen"].apply(lambda x: platzierung_dict[x])
        xn = np.arange(len(x)) #brauchen wir, weil Matplot keine Strings zur X-Achse verarbeiten kann, nur Zahlen
        y = df_pl["Artikel-Lesewert (Erscheinung) in %"]
        yn = df_pl["SplitId"]
    
        #Zeichenfunktion wird aufgerufen
        plot_axis(prs, x=x, labels = labels, xn = xn, y = y, yn = yn,\
                  pos_left = 0.12, pos_right=0.92, pos_top=0.85, \
                  pos_bottom = 0.4, article = "total", grid=False,\
                  title_text = "Platzierungen und Lesewerte")
        
    if geschlecht == True: 
        
        #Achsen festlegen
        # Werte festlegen    
        x = df_pl["Platzierungen"]
        labels = df_pl["Platzierungen"].apply(lambda x: platzierung_dict[x])
        xn = np.arange(len(x))
        y_w = df_pl["LW_w"]
        y_m = df_pl["LW_m"]
        
        multiple_bars(prs, x=x, labels=labels, xn=xn,y_w=y_w, y_m = y_m, 
                      title_text=title_text)
    
#%% Lesewert und aRtikelzahl allgemein
        
def lesewert_artikelzahl(prs, df, minimum = 5, geschlecht=False, title_text="Bitte eingeben", 
                         mode="Handlungsort", liste=[]):
    
    
    
    
    
    # alle Null-Werte ausblenden
    mask = df[mode].notnull()
    df = df[mask]
    
    if geschlecht == False: 
        # Tabelle wird verkleinert, Zeilen ohne Eintrag in DF werden gelöscht
        df_ = df[["SplitId", "Erscheinungsdatum", mode, "Artikel-Lesewert (Erscheinung) in %"]]
        
    if geschlecht == True:  
        df_ = df[["SplitId", "Erscheinungsdatum", mode, "LW_w", "LW_m", "Artikel-Lesewert (Erscheinung) in %"]]

    
    #Anzahl Erscheinungstage wird ermittelt
    erscheinungstage = df_["Erscheinungsdatum"].nunique()
    
    
    if geschlecht == False :
        
        # Tabelle mit Darstellungsform, Anzahl Artikel und Durchschnittslesewert wird erstellt
        df_pl = df_.groupby([mode], as_index=False).agg({"SplitId":"size", "Artikel-Lesewert (Erscheinung) in %":"mean"})\
        .sort_values(by="Artikel-Lesewert (Erscheinung) in %", ascending=False)
    
    if geschlecht == True:
        df_pl = df_.groupby([mode], as_index=False).agg({"SplitId":"size", "LW_w":"mean", "LW_m":"mean",
                           "Artikel-Lesewert (Erscheinung) in %":"mean"})\
        .sort_values(by="Artikel-Lesewert (Erscheinung) in %", ascending=False)
        
    
    df_pl = df_pl[df_pl["SplitId"]>=minimum]
    
    # GRAFIK ANFERTIGEN
    if geschlecht == False: 
        
    
    
        # Achsen festlegen
        x= df_pl[mode]
        labels = df_pl[mode].apply(lambda x: platzierung_dict[x])
        xn = np.arange(len(x)) #brauchen wir, weil Matplot keine Strings zur X-Achse verarbeiten kann, nur Zahlen
        y = df_pl["Artikel-Lesewert (Erscheinung) in %"]
        yn = df_pl["SplitId"]
    
        #Zeichenfunktion wird aufgerufen
        plot_axis(prs, x=x, labels = labels, xn = xn, y = y, yn = yn,\
                  pos_left = 0.12, pos_right=0.92, pos_top=0.85, \
                  pos_bottom = 0.4, article = "total", grid=False,\
                  title_text = "Platzierungen und Lesewerte")
        
    if geschlecht == True: 
        
        #Achsen festlegen
        # Werte festlegen    
        x = df_pl["Platzierungen"]
        labels = df_pl["Platzierungen"].apply(lambda x: platzierung_dict[x])
        xn = np.arange(len(x))
        y_w = df_pl["LW_w"]
        y_m = df_pl["LW_m"]
        
        multiple_bars(prs, x=x, labels=labels, xn=xn,y_w=y_w, y_m = y_m, 
                      title_text=title_text)        
        


#%% Vergleich Geschlechter

#target gibt an, welche Werte verglichen werden sollen, also Darstellungsform oder 
# Ressort     



#%% Artikellänge und Lesewert/Artikelanzahl
# TODO: Noch nicht geschrieben
def multiple_bars_geschlecht(prs, df_geschlecht, 
                             target="Ressorts", ressort_liste = [], 
                             grid = False, legend="normal", 
                             title_text= "Ressorts nach Geschlecht", 
                             sort="Seitennummer", 
                             width = "normal"):
    
    # Wahl nach Untersuchungsspalte
    if target == "Ressorts":
        group_param = "Ressortbeschreibung"
        sort = sort
    if target == "Spezial_Res":  
        group_param = "Spezial_Res"
    if len(ressort_liste)>0: 
        df_geschlecht = df_geschlecht[df_geschlecht["Ressortbeschreibung"].isin(ressort_liste)]
    
    
    
    
    df_group = df_geschlecht.groupby(group_param, as_index=False)
    df_group = df_group.agg({"LW_m":"mean", "LW_w": "mean", "Seitennummer":"median"})
    print(df_group)
    if sort=="Kategorie":
       
       df_group[group_param] = df_group[group_param].astype("category")
       df_group[group_param].cat.set_categories(ressort_liste, inplace=True)
       df_group.sort_values([group_param], inplace = True)
        
    else: 
        df_group.sort_values(by=sort, inplace=True)
    
    # Werte festlegen    
    x = df_group[group_param]
   
    df_group["short"] = df_group[group_param].apply(lambda x: ressort_dict[x])
    
    labels = df_group["short"]
    
    xn = np.arange(len(x))
    y_w = df_group["LW_w"]
    y_m = df_group["LW_m"]
    
    
    
    # PLOTTING
    
    #Schriftfarbe und Farbe der Ticks festlegen
    set_font_color ="#8c8f91" 
       
    #Style setzen, plots anlegen
    sns.set_style("white")
    fig, ax1 = plt.subplots(figsize=(20,8))
    bar1 = ax1.bar(xn-0.40, y_w, label = "Frauen", width=0.4, align="center", color="#f77801")
    bar2 = ax1.bar(xn, y_m, width=0.4, label = "Männer", align="center", color="#f7bb83")
    
    
    
    # setzt Grid wenn gewünscht
    if grid==True:
        ax1.grid(color= set_font_color, linestyle='-', linewidth=1, axis="y")
    
    
    
    # setzt die Y-Achse 
    plt.yticks(fontproperties=campton_light)
    y_label_text = 'Ø Lesewert in Prozent'
    ax1.set_ylabel(y_label_text, color= set_font_color, fontproperties=campton_light, fontsize=50)
    ax1.yaxis.label.set_size(22)
    ax1.tick_params(axis='y', labelsize=25, colors= set_font_color)
    
    
    
    # setzt die X-Achse
    
    ax1.xaxis.set(ticks=range(0, len(xn))) #  Anzahl der Ticks 
    ax1.set_xticklabels(labels = labels, rotation=45, ha="right",  weight=800,\
                        color= set_font_color, fontproperties=campton_light, \
                        fontsize=30) # Labels werden ausgerichtet
    
    
    #Legende einbauen
    legend_height = 1.24
    if legend=="normal":
            legend_height = 1.24
    elif legend == "strange": 
            legend_height =1.44
    else:
        legend_height = 1.24
        
    leg = plt.legend(bbox_to_anchor=(1, legend_height), handles=[bar1, bar2], markerscale=140)
    for text in leg.get_texts(): 
        plt.setp(text, color= set_font_color, size=25)
    
    # DESIGN
    # Abstände Bars zur Achse (standardmäßig bei 0.5)
    plt.margins(x=0.03) # ziehen Bars näher an die Achse
     #obere Linie ausblenden
    ax1.spines["top"].set_visible(False)
    #ax1.spines["left"].set_color("gray")
    ax1.spines["top"].set_visible(False)
    ax1.spines["bottom"].set_visible(False)
    ax1.spines["left"].set_visible(False)
    ax1.spines["right"].set_visible(False)
    
    pos_left = 0.08
    pos_right=0.92
   
    pos_top=0.85
    pos_bottom = 0.4
    
    if legend=="large": 
        pos_left = 0.08
        pos_right=0.92
        pos_top=0.85
        pos_bottom = 0.45
    
    if legend=="xlarge": 
        pos_left = 0.08
        pos_right=0.92
        pos_top=0.85
        pos_bottom = 0.48
    
    if width=="special": 
        pos_right=0.88
    
    plt.subplots_adjust(left=pos_left, right=pos_right, top=pos_top, 
                        bottom=pos_bottom)
    
    filename = "grafik_geschlecht_vergleich.jpg" 
   
    
     
                  
    plt.savefig(filename)
                    
    picture_sheet(prs, filename, title_text=title_text)     
   
    
    plt.close()  
    return prs
#%%
    

#%% Vergleich Lokales STZ & STN
 
def lokales_vergleich_tabelle(prs, df):    

#neue Tabelle anlegen
    slide_layout = prs.slide_layouts[2]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    
    #Tabelle anlegen
    rows = 4 
    cols = 3
    left = Inches(1.05)
    top = Inches(2)
    width = Inches (8.0)
    height = Inches(0.6)
    font_size = 12
        
    total_rows = 0
        
    # Titelzeile des Sheets festlegen
    title_placeholder = slide.shapes.title            
    title_placeholder.text = "Vergleich Lokales"
    title_placeholder.text_frame.paragraphs[0].font.bold = True
        
        #Layout der Tabelle festlegen und neues Dokument für Tabelle festlegen
        #nur dann Tabelle zeichnen, wenn mindestens ein Wert in der Tabelle abgefragt wird...    
    table = shapes.add_table(rows, cols, left, top, width, height).table
    table.columns[0].width = Inches(2)
    table.columns[1].width = Inches(2.7)
    table.columns[2].width = Inches(3.3)
        
        
    table.cell(0, 0).text = "Werte"
    table.cell(0, 1).text = "Lokales STZ" 
    table.cell(0, 2).text = "Lokales STN" 
        
    #table.last_row=True
    table.cell(1,0).text = "LW"
    table.cell(2,0).text = "BW"
    table.cell(3,0).text = "DW"
    
    # Daten Lokales STN & STZ
    df_lokal = df[df["Ressortbeschreibung"].isin(lokale_ressorts)] 
    df_lokal_stz = df_lokal[df_lokal["Zeitung"]== "STZ"]
    df_lokal_stn = df_lokal[df_lokal["Zeitung"]== "STN"]
    

    
    for i in range(1):
            
                
               
            table.cell(i+1,1).text = str(round(df_lokal_stz["Artikel-Lesewert (Erscheinung) in %"].mean(), 1)).replace(".", ",")
            table.cell(i+1,2).text = str(round(df_lokal_stn\
                                                ["Artikel-Lesewert (Erscheinung) in %"].mean(), 1)).replace(".", ",")
            table.cell(i+2,1).text = str(round(df_lokal_stz\
                                                ["Artikel-Blickwert (Erscheinung) in %"].mean(), 1)).replace(".", ",")
            table.cell(i+2,2).text = str(round(df_lokal_stn\
                                                ["Artikel-Blickwert (Erscheinung) in %"].mean(), 1)).replace(".", ",")
            table.cell(i+3,1).text = str(round(df_lokal_stz\
                                                ["Artikel-Durchlesewerte (Erscheinung) in %"].mean(), 1)).replace(".", ",")
            table.cell(i+3,2).text = str(round(df_lokal_stn\
                                                ["Artikel-Durchlesewerte (Erscheinung) in %"].mean(), 1)).replace(".", ",")
           
   # Größe + Schriftart der oberen Zeile festlegen
    # muss hier geschehen, da ich sie mit i nicht ansteuern kann
    for col in range(cols):
        text_frame = table.cell(0, col).text_frame
        p = text_frame.paragraphs[0]
        run = p.add_run()
        font = run.font
        p.font.name ="Campton-Light"
        p.font.size = Pt(11)
       
        p.alignment = PP_ALIGN.CENTER
    # table.last_row = True


    # data rows, Tabelle füllen
    font_size = 12
    for i in range(0, rows):
        for j in range(0,cols):
            #table.cell(i+1, j).text_frame.paragraphs[0].font.size=Pt(font_size)
            text_frame = table.cell(i, j).text_frame
            p = text_frame.paragraphs[0]
            run = p.add_run()
            font = run.font
            p.font.name ="Campton-Light"
            p.font.size = Pt(11)
       
            p.alignment = PP_ALIGN.CENTER
        
        for row in range(rows): 
            text_frame = table.cell(row,1).text_frame
            p = text_frame.paragraphs[0]
            #p.alignment = PP_ALIGN.LEFT
    return prs

#%% Top 5 Kolumnen
# target = zeigt an, ob nur in den lokalen Ressorts, den Mantelressort oder allen Ressorts gesucht werden soll (lokal, mantel, gesamt) 
    # minimum gibt Mindestanzahl von Veröffentlichungen an 

def top_kolumnen(prs, df, liste=kolumnen_liste, target="gesamt", minimum=4, title_text = "Bitte Titel eingeben"):
    if target == "lokal":
        df_ = df[df["Ressortbeschreibung"].isin(lokale_ressorts)]
    if target == "mantel":
        df_ = df[df["Ressortbeschreibung"].isin(mantel_ressorts)]
    else: 
        df_ = df.copy()
    
    # nur festgelegte Kolumnen untersuchen
    df_ = df_[df_["Beschreibung"].isin(liste)]
    
    # nach Kolumnen gruppieren
    df_group = df_.groupby("Beschreibung")
    df_group = df_group.agg({"SplitId":"size", "Artikel-Lesewert (Erscheinung) in %":"mean"})
    
    # Einzeichnen
    plot_axis(prs, x=0, labels = 0, xn = 0, y = 0, yn = 0,
                  pos_left = 0.08, pos_right=0.92, pos_top=0.85, 
                  pos_bottom = 0.33, article = "total", grid=False,
                  title_text=title_text, axis=1, 
                  mean_line = 0, legend="normal", special=False)
    
    return prs
    
#%% Kommentare Übersicht LW
def kommentare_tabelle(prs, df, version="STZ ", title_text = "default"):
    
    # Berechnung der einzelnen Kommentare
    
    #STZ:
    if version=="STZ":     
        df_z = df[df["Zeitung"]=="STZ"]
        
        # Kommentare Titelseite
        df_1 = df_z[(df_z["Ressortbeschreibung"]=="Titelseite") & 
                    (df_z["Darstellungsformen"]=="KM") &
                    ((df_z["Platzierungen"]=="TS1") |
                            (df_z["Platzierungen"]=="ZA"))]
        
        # Kommentar "Unten rechts" auf Seite 3
        df_2 = df_z[(df_z["Seitennummer"]==3) & (df_z["Beschreibung"]=="Unten rechts")]
        
        # weitere Kommentare auf Seite 3
        df_3 = df_z[(df_z["Platzierungen"]=="TS1") & (df_z["Beschreibung"]!="Unten rechts")]
        
        # Kommentare Wirtschaft 1
        df_4 = df_z[(df_z["Darstellungsformen"]=="KM")]
        
        # Wirtschaft Börsenwoche
        # TODO: Hier sind nicht alle richtig codiert. Börsenwoche erscheint immer montags, 
        # 16.4. ist z.B. nicht codiert. 
        
        df_5 = df_z[(df_z["Beschreibung"]=="Börsenwoche")]
        
        # Kommentare Lokal
        # Ein Stück Stadt
        df_6 = df_z[df_z["Beschreibung"]=="Ein Stück Stadt"]
        # Gerhard Raff
        df_7 = df_z[df_z["Beschreibung"]=="Gerhard Raff"]
        
        rows = 8
        
    if version == "STN": 
        df_z = df[df["Zeitung"]=="STN"]
        
        # Kommentare Titelseite
        df_1 = df_z[(df_z["Ressortbeschreibung"]=="Titelseite") & 
                    (df_z["Darstellungsformen"]=="KM")]
        
        # Leitartikel Seite 2
        df_2 = df_z[((df_z["Seitentitel"]=="MEINUNG ^. NACHRICHTEN") | 
                (df_z["Seitentitel"]=="MEINUNG  NACHRICHTEN"))] 
        df_2 = df_1[df_1["Beschreibung"]=="Leitartikel"]
        
        rows = 3
        
    
    
    #TABELLE
    
    # neues Sheet aufrufen
    slide_layout = prs.slide_layouts[2] #ist im Master Überschrift inkl aller Symbole, ansonsten leer
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    
    
    # Form der Tabelle Festlegen
    cols = 4
    
    left = Inches(1.05)
    top = Inches(1.55)
    width = Inches (8.3)
    height = Inches(0.6)
    
    # jetzt legen wir die Tabelle an, ein table-Objekt
    table = shapes.add_table(rows, cols, left, top, width, height).table
    
    
    # Breite der Spalten festlegen (für modus Gesamt und Ressort)
    
    table.columns[0].width = Inches(4.7)
    table.columns[1].width = Inches(1.2)
    table.columns[2].width = Inches(1.2)
    table.columns[3].width = Inches(1.2)
    
    
    #  Index benennen
    
    table.cell(0, 0).text = "Kommentar/Rubrik"
    table.cell(0, 1).text = "LW"
    table.cell(0, 2).text = "BW"
    table.cell(0, 3).text = "DW"
    
    # Überschrift festlegen
    
    title_placeholder = slide.shapes.title
    
    if title_text == "default":
        title = "Kommentare " + version
    
    title_placeholder.text = title
    title_placeholder.text_frame.paragraphs[0].font.bold = True
    title_placeholder.text_frame.paragraphs[0].font.name = "Campton-Bold"
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(32)
    
    
    # Größe + Schriftart der oberen Zeile festlegen
    # muss hier geschehen, da ich sie mit i nicht ansteuern kann
    for col in range(cols):
        text_frame = table.cell(0, col).text_frame
        p = text_frame.paragraphs[0]
        run = p.add_run()
        font = run.font
        p.font.name ="Campton-Light"
        p.font.size = Pt(11)
       
        p.alignment = PP_ALIGN.CENTER
        # table.last_row = True
    
    
    # data rows, Tabelle füllen
    if version == "STZ":
        
    
        table.cell(1, 0).text = "Kommentar Titelseite"
        table.cell(1, 1).text = str(round(df_1["Artikel-Lesewert (Erscheinung) in %"].mean(), 1)).replace(".", ",")
        table.cell(1, 2).text = str(round(df_1["Artikel-Blickwert (Erscheinung) in %"].mean(), 1)).replace(".", ",")
        table.cell(1, 3).text = str(round(df_1["Artikel-Durchlesewerte (Erscheinung) in %"].mean(), 1)).replace(".", ",")
        
        table.cell(2, 0).text = '"Unten rechts" - Seite 3'
        table.cell(2, 1).text = str(round(df_2["Artikel-Lesewert (Erscheinung) in %"].mean(), 1)).replace(".", ",")
        table.cell(2, 2).text = str(round(df_2["Artikel-Blickwert (Erscheinung) in %"].mean(), 1)).replace(".", ",")
        table.cell(3, 3).text = str(round(df_2["Artikel-Durchlesewerte (Erscheinung) in %"].mean(), 1)).replace(".", ",")
            
        table.cell(3, 0).text = 'Weitere Kommentare - Seite 3'
        table.cell(3, 1).text = str(round(df_3["Artikel-Lesewert (Erscheinung) in %"].mean(), 1)).replace(".", ",")
        table.cell(3, 2).text = str(round(df_3["Artikel-Blickwert (Erscheinung) in %"].mean(), 1)).replace(".", ",")
        table.cell(3, 3).text = str(round(df_3["Artikel-Durchlesewerte (Erscheinung) in %"].mean(), 1)).replace(".", ",")
        
        table.cell(4, 0).text = 'Kommentare Wirtschaft'
        table.cell(4, 1).text = str(round(df_4["Artikel-Lesewert (Erscheinung) in %"].mean(), 1)).replace(".", ",")
        table.cell(4, 2).text = str(round(df_4["Artikel-Blickwert (Erscheinung) in %"].mean(), 1)).replace(".", ",")
        table.cell(4, 3).text = str(round(df_4["Artikel-Durchlesewerte (Erscheinung) in %"].mean(), 1)).replace(".", ",")
        
        table.cell(5, 0).text = 'Börsenwoche Wirtschaft'
        table.cell(5, 1).text = str(round(df_5["Artikel-Lesewert (Erscheinung) in %"].mean(), 1)).replace(".", ",")
        table.cell(5, 2).text = str(round(df_5["Artikel-Blickwert (Erscheinung) in %"].mean(), 1)).replace(".", ",")
        table.cell(5, 3).text = str(round(df_5["Artikel-Durchlesewerte (Erscheinung) in %"].mean(), 1)).replace(".", ",")
        
        table.cell(6, 0).text = 'Ein Stück Stadt - Lokal'
        table.cell(6, 1).text = str(round(df_6["Artikel-Lesewert (Erscheinung) in %"].mean(), 1)).replace(".", ",")
        table.cell(6, 2).text = str(round(df_6["Artikel-Blickwert (Erscheinung) in %"].mean(), 1)).replace(".", ",")
        table.cell(6, 3).text = str(round(df_6["Artikel-Durchlesewerte (Erscheinung) in %"].mean(), 1)).replace(".", ",")
        
        table.cell(7, 0).text = 'Gerhard Raff - Lokal'
        table.cell(7, 1).text = str(round(df_5["Artikel-Lesewert (Erscheinung) in %"].mean(), 1)).replace(".", ",")
        table.cell(7, 2).text = str(round(df_5["Artikel-Blickwert (Erscheinung) in %"].mean(), 1)).replace(".", ",")
        table.cell(7, 3).text = str(round(df_5["Artikel-Durchlesewerte (Erscheinung) in %"].mean(), 1)).replace(".", ",")
     
    if version == "STN":
        
        table.cell(1, 0).text = "Kommentar Titelseite"
        table.cell(1, 1).text = str(round(df_1["Artikel-Lesewert (Erscheinung) in %"].mean(), 1)).replace(".", ",")
        table.cell(1, 2).text = str(round(df_1["Artikel-Blickwert (Erscheinung) in %"].mean(), 1)).replace(".", ",")
        table.cell(1, 3).text = str(round(df_1["Artikel-Durchlesewerte (Erscheinung) in %"].mean(), 1)).replace(".", ",")
        
        table.cell(2, 0).text = "Leitartikel Seite 2"
        table.cell(2, 1).text = str(round(df_1["Artikel-Lesewert (Erscheinung) in %"].mean(), 1)).replace(".", ",")
        table.cell(2, 2).text = str(round(df_1["Artikel-Blickwert (Erscheinung) in %"].mean(), 1)).replace(".", ",")
        table.cell(2, 3).text = str(round(df_1["Artikel-Durchlesewerte (Erscheinung) in %"].mean(), 1)).replace(".", ",")
    

    
    for i in range(0, rows-1):
        for j in range(0,cols):
            #table.cell(i+1, j).text_frame.paragraphs[0].font.size=Pt(font_size)
            text_frame = table.cell(i+1, j).text_frame
            p = text_frame.paragraphs[0]
            run = p.add_run()
            font = run.font
            p.font.name ="Campton-Light"
            p.font.size = Pt(11)
       
            p.alignment = PP_ALIGN.CENTER
        for row in range(rows): 
            text_frame = table.cell(row,1).text_frame
            p = text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
    
    return prs 
        
    
    

#%% ANALYSE RESSORTS
   
#%%  Initial Ressort-Analyse
# TODO: def analyse_ressorts(prs, df):


    
#%%  Deckblätter Function Zwischenbericht
    ''' Diese Funktion erstellt die Deckblätter inklusive der Screenshots
    für die Auswertungen (einzelne Ressorts, Lokalteile etc.).  
    
    '''
    
def deckblatt_macher(prs, df, ressort,  platzierung, darstellungsform, 
                     seitentitel, ZTG= "null"):

    
    
   
    
    #neuen Slide aufrufen
    
    slide_layout = prs.slide_layouts[2]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    
   
    ressort_neu = ressort.replace("/", "_").replace(" ", "_").replace("&", "und")
    
    if ZTG !="null": 
        page_path = "./seiten_schwaebische/Schw_" + ZTG + "_" + ressort_neu + ".jpg"
        
    elif ZTG=="null": 
        page_path = "./seiten_schwaebische/Schw_" + ressort_neu + ".jpg"
       
        
    page_left = Inches(7)
    page_top = Inches(0.9)
    page_width = Inches(2.7)
    try:
        page = slide.shapes.add_picture(page_path, page_left, page_top, width=page_width)
    except FileNotFoundError:
        print(ressort_neu + " ... JPEG Titelseite nicht gefunden")
    
    #Form der Tabelle festlegen
        
    rows=1+len(platzierung)+len(darstellungsform)+len(seitentitel)+1
    if len(seitentitel)>0:
        rows +=1
   

    
    cols = 4
    left = Inches(1.05)
    top = Inches(2)
    width = Inches (5)
    height = Inches(0.6)
    
    total_rows = 0

    
    # Titelzeile des Sheets festlegen
    title_placeholder = slide.shapes.title
    if ZTG != "null":
        title_placeholder.text = ressort + " - " + df.iloc[0]["ZTG"]
    elif ZTG == "null":
        title_placeholder.text = ressort 
    
    title_placeholder.text_frame.paragraphs[0].font.bold = True
    title_placeholder.text_frame.paragraphs[0].font.name = "Campton-Bold"
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(32)
    
    #Layout der Tabelle festlegen und neues Dokument für Tabelle festlegen
    #nur dann Tabelle zeichnen, wenn mindestens ein Wert in der Tabelle abgefragt wird... 
    if rows>1:
        
        table = shapes.add_table(rows, cols, left, top, width, height).table
        table.columns[0].width = Inches(2.6)
        table.columns[1].width = Inches(0.8)
        table.columns[2].width = Inches(0.8)
        table.columns[3].width = Inches(0.8)
        table.cell(0, 0).text = ""
        table.cell(0, 1).text = "LW"
        table.cell(0, 2).text = "BW"    
        table.cell(0, 3).text = "DW"    
        table.last_row=True 
    
        #Datensatz für die jeweilige Tabelle anfertigen
        #Wir haben jetzt das df einer Zeitung mit den Artikeln eines Ressorts: 
        data = df.copy()
    
    
        # for-Schleife zündet nur, wenn len(seitentitel) != 0 
        for h in range(len(seitentitel)):
            #table.cell(h+1, 0).text = seitentitel_dict[seitentitel[h]]
            table.cell(h+1, 0).text = seitentitel[h]
            table.cell(h+1, 1).text = str(round(data[data["Seitentitel"]==seitentitel[h]]\
                                           ["Artikel-Lesewert (Erscheinung) in %"].mean(),1)).replace(".", ",")
            table.cell(h+1, 2).text = str(round(data[data["Seitentitel"]==seitentitel[h]]\
                                           ["Artikel-Blickwert (Erscheinung) in %"].mean(),1)).replace(".", ",")
            table.cell(h+1, 3).text = str(round(data[data["Seitentitel"]==seitentitel[h]]\
                                           ["Artikel-Durchlesewerte (Erscheinung) in %"].mean(),1)).replace(".", ",")
        #if h==(len(seitentitel)-1):
            #table.cell(h)
            
            total_rows = len(seitentitel)+2
            
        for i in range(len(platzierung)):
            x=0
            y = len(platzierung)+1
            if len(seitentitel) >0: 
                y = y+len(seitentitel)+1
            if len(seitentitel) >= 1:
                x = len(seitentitel)+1
            table.cell(i+1+x,0).text = platzierung_dict[platzierung[i].strip()]
            table.cell(i+1+x,1).text = str(round(data[data["Platzierungen"]==platzierung[i].strip()]\
                                     ["Artikel-Lesewert (Erscheinung) in %"].mean(), 1)).replace(".", ",")
            table.cell(i+1+x,2).text = str(round(data[data["Platzierungen"]==platzierung[i].strip()]\
                                     ["Artikel-Blickwert (Erscheinung) in %"].mean(), 1)).replace(".", ",")
            table.cell(i+1+x,3).text = str(round(data[data["Platzierungen"]==platzierung[i].strip()]\
                                     ["Artikel-Durchlesewerte (Erscheinung) in %"].mean(), 1)).replace(".", ",")
            # TODO: Richtige Zeilennummer ausrechnen, falls keine Darstellungsform abgefragt wird
            total_rows = len(platzierung)+1
            if len(seitentitel) >= 1:
                total_rows = len(platzierung)+1+len(seitentitel)+1
            
        for k in range(len(darstellungsform)):
            y = len(platzierung)+1
            if len(seitentitel) >0: 
                y = y+len(seitentitel)+1
            table.cell(k+y,0).text  = darstellung_dict[darstellungsform[k]]
            table.cell(k+y,1).text = str(round(data[data["Darstellungsformen"]\
                                                .str.contains(darstellungsform[k], na=False)]\
                                           ["Artikel-Lesewert (Erscheinung) in %"].mean(),1)).replace(".", ",")
            table.cell(k+y,2).text = str(round(data[data["Darstellungsformen"]\
                                                .str.contains(darstellungsform[k], na=False)]\
                                           ["Artikel-Blickwert (Erscheinung) in %"].mean(),1)).replace(".", ",")
            table.cell(k+y,3).text = str(round(data[data["Darstellungsformen"]\
                                                .str.contains(darstellungsform[k], na=False)]\
                                           ["Artikel-Durchlesewerte (Erscheinung) in %"].mean(),1)).replace(".", ",")
            
            total_rows = k+y+1
        
        table.cell(total_rows, 0).text = "Gesamt"
        table.cell(total_rows, 1).text = str(round(data["Artikel-Lesewert (Erscheinung) in %"].mean(), 1)).replace(".", ",")
        table.cell(total_rows, 2).text = str(round(data['Artikel-Blickwert (Erscheinung) in %'].mean(), 1)).replace(".", ",")
        table.cell(total_rows, 3).text = str(round(data["Artikel-Durchlesewerte (Erscheinung) in %"].mean(), 1)).replace(".", ",")
    
    
    # Schrift in Campton umwandeln etc.  
    
    # Schrift in erster Zeile (Index) umstellen
    for col in range(cols):
        text_frame = table.cell(0, col).text_frame
        p = text_frame.paragraphs[0]
        run = p.add_run()
        font = run.font
        p.font.name ="Campton-Light"
        p.font.size = Pt(11)
        p.alignment = PP_ALIGN.CENTER
    
    for i in range(rows):
        for j in range (cols):
            text_frame = table.cell(i, j).text_frame
            p = text_frame.paragraphs[0]
            run = p.add_run()
            font = run.font
            p.font.name ="Campton-Light"
            p.font.size = Pt(11)
   
            p.alignment = PP_ALIGN.CENTER
        for row in range(rows): 
            text_frame = table.cell(row,0).text_frame
            p = text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
                
        
        #  Schrift verkleinern, wenn mehr als sieben Zeilen vorhanden sind
        #if rows > 7: 
        #    font_size = 10 
        #else:
        #    font_size = 12
        #for row in range(rows):
        #    for col in range(cols):
        #        table.cell(row, col).text_frame.paragraphs[0].font.size=Pt(font_size)
                
        
        
        
    return prs



#%% Deckblatt-Macher für Abschlussbericht
def deckblatt_abschluss(prs, df_, df_geschlecht=[], title_text="", lokales=False, kunde="nw", geschlecht=True):
    lw_mean = df_["Artikel-Lesewert (Erscheinung) in %"].mean()
    dw_mean = df_['Artikel-Durchlesewerte (Erscheinung) in %'].mean()
    bw_mean = df_["Artikel-Blickwert (Erscheinung) in %"].mean()
#    m_mean = df_geschlecht["LW_m"].mean()
#    w_mean = df_geschlecht["LW_w"].mean()
    # Anlegen des Powerpoint-Dokuments
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # if Schleife für händisch eingegebenen Titel
    if len(title_text)>0: 
        titletext = title_text
    else:
        titletext = df_.iloc[0]["Ressortbeschreibung"]

    
    
    title_placeholder = slide.shapes.title
    title_placeholder.text= titletext
    placeholder_object = slide.placeholders[10] #unterer Text, größtenteils vom Coach auszufüllen
    placeholder_body = slide.placeholders[12] # oberer Text, orange Box mit LW
    placeholder_picture = slide.placeholders[13]
    # richtige Bilddatei wählen (Lokal oder Mantel)
    #ressort = df_.iloc[0]["Ressortbeschreibung"].replace(" ", "_").replace("/", "_").replace(",", "_")
    ressort = "Lokalsport"
    #Bild laden und einsetzen
    
    if kunde=="nw": 
        if lokales == True: 
            ressort = df_.iloc[0]["Ressortbeschreibung"]+"_"+df_.iloc[0]["ZTG"]
        else: 
            ressort = df_.iloc[0]["Ressortbeschreibung"].replace(" ", "_").replace("/", "_").replace(",", "_")
        picture_url = "./seiten_nw/nw_"+ressort+".jpg"
    
    if kunde=="SWZ":
        ressort = df_.iloc[0]["Ressortbeschreibung"].replace(" ", "_").replace("/", "_").replace(",", "_").replace("&", "und")
        if lokales == True: 
            ressort = df_.iloc[0]["ZTG"] + "_" + df_.iloc[0]["Ressortbeschreibung"]
        picture_url = "./seiten_schwaebische/Schw_"+ressort+".jpg"
    
    if kunde == "MHS":
        
        if df_.iloc[0]["Ressortbeschreibung"] in lokale_liste:
            picture_url = "./seiten/lok_"+ressort+".jpg"
        elif df_.iloc[0]["Ressortbeschreibung"]=="Wochenende": 
            picture_url = "./seiten/Wochenende.jpg"
        elif df_.iloc[0]["Ressortbeschreibung"]=="Die Brücke": 
            picture_url = "./seiten/Die_Brücke.jpg"
        else:
            picture_url = "./seiten/"+df_.iloc[0]["Zeitung"]+"_"+ressort+".jpg"
    if kunde=="MZ": 
        ressort = df_.iloc[0]["Ressortbeschreibung"].replace(" ", "_").replace("/", "_").replace(",", "_").replace("&", "und")
        
        if df_.iloc[0]["Ressortbeschreibung"] == "Lokales":
            picture_url = "./seiten_mz/mz_" + df_.iloc[0].ZTG + "_Lokales.jpg"
        elif df_.iloc[0].Ressortbeschreibung == "Regionalsport":
            picture_url = "./seiten_mz/mz_" + df_.iloc[0].ZTG + "_Regionalsport.jpg"
        else:
            picture_url = "./seiten_mz/mz_" + ressort + ".jpg"
    
    if kunde == "NOZ":
       ressort = df_.iloc[0]["Ressortbeschreibung"].replace(" ", "_").replace("/", "_").replace(",", "_").replace("&", "und") 
       picture_url = "./seiten_noz/noz_"+df_.iloc[0]["ZTG"]+"_"+ressort+".jpg"
    
    if kunde == "DRP":
        ressort = df_.iloc[0]["Ressortbeschreibung"].replace(" ", "_").replace("/", "_").replace(",", "_").replace("&", "und")
        if (ressort=="Lokales")|(ressort=="Lokalsport"): 
            picture_url = "./seiten_drp/drp_"+df_.iloc[0]["ZTG"]+"_"+ressort+".jpg"     
        else: 
            picture_url = "./seiten_drp/drp_"+ressort+".jpg"    
    picture = placeholder_picture.insert_picture(picture_url)
    
    if kunde == "nw": 
        placeholder_body.text = "LW {:1.1f} - BW {:1.1f} - DW {:1.1f}".format(lw_mean, bw_mean, dw_mean).replace(".", ",") 
        placeholder_object.text = "Steigerung zu 2018 (+)"
    else:
        placeholder_body.text = "Lesewert Ø " + '{:1.1f}'.format(lw_mean).replace(".", ",") 
        placeholder_object.text = "Blickwert: " + '{:1.1f}'.format(bw_mean).replace(".", ",")+ "\n" + \
        "Durchlesewert: " + '{:1.1f}'.format(dw_mean).replace(".", ",")
        
#    + \
#    "\n" + "\n" + "Lesewert Frauen: " + "{:1.1f}".format(w_mean).replace(".", ",") + \
#    "\n" + "Lesewert Männer: " + "{:1.1f}".format(m_mean).replace(".", ",")

    return prs


#%% Tabelle Darstellungsformen / Platzierungen + Seitenbild
    
#Unter research_object verbirgt sich der Untersuchungsgegenstand
# (Darstellungsformen oder Platzierungen)
    
# Unter sort kann eingestellt werden, ob die Ergebnisse nach Seitennummer oder 
    # nach Lesewert sortiert werden sollen. 

def tabelle_ressortauswertung(prs, df, research_object="Darstellungsformen", 
                              sort="Lesewert", title_text="",minimum=5, 
                              lokales=False, kunde="nw"):
    
    
    df_ = df.copy()
    df_ = df_[~df_[research_object].isnull()]
    
        
    if research_object == "Darstellungsformen":
        # alle Zeilen mit nur einem Wert
        df_oneval = df_[df_[research_object].str.len()<=2]
        
        # alle Zeilen mit zwei Werten, hier erster Wert ausgewähl
        df_firstval = df_[df_[research_object].map(len)>2]
        df_firstval["Darstellungsformen"] = df_firstval["Darstellungsformen"].apply(lambda x: x[:2])
        
        # alle Zeilen mit zwei Werten, hier zweiter Wert ausgewählt
        df_secondval = df_[df_[research_object].map(len)>2]
        df_secondval["Darstellungsformen"] = df_secondval["Darstellungsformen"].apply(lambda x: x[-2:])
        
        df_ = df_oneval.append([df_firstval, df_secondval], ignore_index=True)
        element_list = darstellungsform
        element_dict = darstellung_dict
    elif research_object == "Platzierungen": 
        element_list = platzierung
        element_dict = platzierung_dict
    
    df_ = df_[df_[research_object].isin(element_list)]
    
    df_ = df_.groupby(research_object, as_index=False)
    df_ = df_.agg({"Seitennummer":"median", "SplitId":"size", 
                   "Artikel-Lesewert (Erscheinung) in %":"mean", 
                   "Artikel-Blickwert (Erscheinung) in %":"mean", 
                   "Artikel-Durchlesewerte (Erscheinung) in %":"mean"})
    df_ = df_[df_["SplitId"]>=minimum]
    
    
    if sort=="Seitennummer":
        df_ = df_.sort_values(by="Seitennummer")
    elif sort == "Lesewert": 
        df_ = df_.sort_values(by="Artikel-Lesewert (Erscheinung) in %", ascending = False)
    
    # Tabelle einzeichnen
    
    # neues Sheet aufrufen
    slide_layout = prs.slide_layouts[14] #ist im Master Überschrift inkl aller Symbole, ansonsten leer
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    
    #Bild einsetzen
    placeholder_picture = slide.placeholders[13]
    
    ressort = df.iloc[0]["Ressortbeschreibung"].replace(" ", "_").replace("/", "_").replace(",", "_")
    #Bild laden und einsetzen
    
    if kunde=="nw": 
        if lokales == True: 
            ressort = df.iloc[0]["Ressortbeschreibung"]+"_"+df.iloc[0]["ZTG"]
        else: 
            ressort = df.iloc[0]["Ressortbeschreibung"].replace(" ", "_").replace("/", "_").replace(",", "_")
        picture_url = "./seiten_nw/nw_"+ressort+".jpg"
    if kunde == "MHS":
        
        if df.iloc[0]["Ressortbeschreibung"] in lokale_liste:
            picture_url = "./seiten/lok_"+ressort+".jpg"
        elif df.iloc[0]["Ressortbeschreibung"]=="Wochenende": 
            picture_url = "./seiten/Wochenende.jpg"
        elif df.iloc[0]["Ressortbeschreibung"]=="Die Brücke": 
            picture_url = "./seiten/Die_Brücke.jpg"
        else:
            picture_url = "./seiten/"+df.iloc[0]["Zeitung"]+"_"+ressort+".jpg"
    
    if kunde == "SWZ":
         ressort = df.iloc[0]["Ressortbeschreibung"].replace(" ", "_").replace("/", "_").replace(",", "_").replace("&", "und")
         if lokales == True: 
            ressort = df.iloc[0]["ZTG"] + "_" + df.iloc[0]["Ressortbeschreibung"]
         picture_url = "./seiten_schwaebische/Schw_"+ressort+".jpg"
    
    if kunde == "MZ":
        ressort = df.iloc[0]["Ressortbeschreibung"].replace(" ", "_").replace("/", "_").replace(",", "_").replace("&", "und")
        if lokales == True: 
            ressort = df.iloc[0]["ZTG"] + "_" + df.iloc[0]["Ressortbeschreibung"]
        picture_url = "./seiten_mz/mz_"+ressort+".jpg"
    
    if kunde == "NOZ":
        ressort = df.iloc[0]["Ressortbeschreibung"].replace(" ", "_").replace("/", "_").replace(",", "_").replace("&", "und")
#        if lokales == True: 
#            ressort = df.iloc[0]["ZTG"] + "_" + df.iloc[0]["Ressortbeschreibung"]
        active_ztg = df.iloc[0]["ZTG"]
        picture_url = "./seiten_noz/noz_"+active_ztg + "_" +ressort+".jpg"
        
    if kunde == "DRP": 
        print("Zeile 4154 " + str(df_.columns))
        # ACHTUNG: HIER MUSS DF statt DF_ benutzt werden!!!! 
        ressort = df.iloc[0]["Ressortbeschreibung"].replace(" ", "_").replace("/", "_").replace(",", "_").replace("&", "und")
        if (ressort=="Lokales")|(ressort=="Lokalsport"): 
            picture_url = "./seiten_drp/drp_"+df.iloc[0]["ZTG"]+"_"+ressort+".jpg"     
        else: 
            picture_url = "./seiten_drp/drp_"+ressort+".jpg" 
    
    
    picture = placeholder_picture.insert_picture(picture_url)
    
   
    # Form der Tabelle Festlegen
    rows = df_.shape[0]+1
    # Nummer der Spalten abhängig von Gesamt- oder Ressort-Auswertung
    cols = 5
    
    left = Inches(1.05)
    top = Inches(1.6)
    width = Inches (5.8)
    height = Inches(0.6)
    
    # jetzt legen wir die Tabelle an, ein table-Objekt
    table = shapes.add_table(rows, cols, left, top, width, height).table
    
    
    # Breite der Spalten festlegen (für modus Gesamt und Ressort)
    
    table.columns[0].width = Inches(2.7)
    table.columns[1].width = Inches(0.7)
    table.columns[2].width = Inches(0.7)
    table.columns[3].width = Inches(0.7)
    table.columns[4].width = Inches(1)
    
    
    table.cell(0, 0).text = research_object[:-2]
    table.cell(0, 1).text = "LW"
    table.cell(0, 2).text = "BW"
    table.cell(0, 3).text = "DW"
    table.cell(0, 4).text = "Anzahl"
    
    
    # Überschrift festlegen
    title_placeholder = slide.shapes.title
    if len(title_text)>1:
        title = title_text
    else:  
        title = research_object
    title_placeholder.text = title
    title_placeholder.text_frame.paragraphs[0].font.bold = True
    title_placeholder.text_frame.paragraphs[0].font.name = "Campton-Bold"
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(32)
        
 
    
    # Größe + Schriftart der oberen Zeile festlegen
    # muss hier geschehen, da ich sie mit i nicht ansteuern kann
    for col in range(cols):
        text_frame = table.cell(0, col).text_frame
        p = text_frame.paragraphs[0]
        run = p.add_run()
        font = run.font
        p.font.name ="Campton-Light"
        p.font.size = Pt(11)
       
        p.alignment = PP_ALIGN.CENTER
        # table.last_row = True
    
    
    # data rows, Tabelle füllen
    for i in range(0, df_.shape[0]):
        
       
        table.cell(i+1, 0).text = element_dict[df_.iloc[i][research_object]]
        table.cell(i+1, 1).text = str(round(df_.iloc[i]["Artikel-Lesewert (Erscheinung) in %"], 1)).replace(".", ",")
        table.cell(i+1, 2).text = str(round(df_.iloc[i]["Artikel-Blickwert (Erscheinung) in %"], 1)).replace(".", ",")
        table.cell(i+1, 3).text = str(round(df_.iloc[i]["Artikel-Durchlesewerte (Erscheinung) in %"], 1)).replace(".", ",")
        table.cell(i+1, 4).text = str(df_.iloc[i]["SplitId"])
    
        for j in range(0,cols):
            #table.cell(i+1, j).text_frame.paragraphs[0].font.size=Pt(font_size)
            text_frame = table.cell(i+1, j).text_frame
            p = text_frame.paragraphs[0]
            run = p.add_run()
            font = run.font
            p.font.name ="Campton-Light"
            p.font.size = Pt(11)
       
            p.alignment = PP_ALIGN.CENTER
        for row in range(rows): 
            text_frame = table.cell(row,0).text_frame
            p = text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
   
    return prs
    
    


#%% Top-10 Function
    
''' 
Die Top-10 Helper-Function erstellt eine Tabelle mit den zehn besten Artikeln

Bei df_berechnet = False errechnet sich die Funktion die zehn besten Artikel 
anhand des Lesewerts selbst. Ansonsten = True bedeutet, dass das Dataframe
bereits in der richtigen Reihenfolge und mit der richtigen Anzahl an Datensätzen 
eingespielt wird. 

h = "default" bedeutet, dass das Programm sich die Überschrift selbst
aus dem Datensatz erstellt. In Einzelfällen (wie die 10 besten Artikel Gesamt),
kann sie auch händisch eingegeben werden.

Mode = "ressort" bedeutet, dass nur Artikel aus einem Ressort verwertet werden, 
eine zusätzlich Anzeige des Ressorts in der Tabelle ist also unnötig.  
Mode = "Gesamt" bedeutet, dass in der Tabelle eine Spalte mit der Bezeichnung des
Ressorts eingebaut wird. "Gesamt" zeigt an, dass die gesamte Ausgabe untersucht wird
(bedeutet: der gesamte DF wird eingegeben), "Mantel" bedeutet dasselbe für den Mantel. 
Bei den letzten beiden wird auch die ÜS entsprechend verändert. 

Screenshots = True bedeutet, dass die besten Top5-Artikel als Screenshots 
eingebaut werden. Die Zahl kann unter number_screenshots verändert werden. 

Bei Zeitung = True wird der Zeitungstitel mit in der ÜS ausgespielt. 

Bei Geschlecht kann männlich/weiblich angegeben werden. 
Alter = Jeweils die Altersklassen angeben

 
top_10(prs, df, df_berechnet = False, screenshots=True, number_screenshots = 5, 
           mode="gesamt", headline="Top 5 Gesamt", zeitung=False)
lokales_ressort - sollte einer der Top-10-Artikel aus dem Ressort Lokales kommen, 
dann erhält er hinter dem Ressort noch die ZTG-Kennung
'''

def top_10(prs, df, df_berechnet = False, screenshots=True, number_screenshots = 10, 
           mode="ressort", headline="default", zeitung=True, geschlecht="", alter="", 
           lokales_ressort="Lokales", kunden_id = 0):
    
    # Definition der Spalten, die später ausgewertet werden sollen
#    LW = "Artikel-Lesewert (Erscheinung) in %"
#    BW = "Artikel-Blickwert (Erscheinung) in %"
#    DW = "Artikel-Durchlesewerte (Erscheinung) in %"
    
    # Hier müssen die Variablen für Ressortbeschreibung etc. festgelegt werden, 
    # je nachdem ob 
    ressort = "Ressort"
    artikel = "Artikel"
    seite = "Seite"
    datum = "Datum"
    ausgabename = "Ausgabename"
    
    
    LW = "LW"
    BW="BW"
    DW= "DW"
        
    operation = "Gesamt"
    
    if len(geschlecht)>0: 
        if geschlecht == "männlich":
           LW = "LW_m"
           BW = "BW_m"
           DW = "DW_m"
         
        elif geschlecht == "weiblich":
           LW = "LW_w"
           BW = "BW_w"
           DW = "DW_w" 
        operation = geschlecht
        
    if len(alter)>0:
        if alter == "50 bis 59 Jahre": 
            LW = "LW_50bis59"
            BW = "BW_50bis59"
            DW = "DW_50bis59"
        
        if alter == "60 bis 69 Jahre": 
            LW = "LW_60bis69"
            BW = "BW_60bis69"
            DW = "DW_60bis69"
        
        if alter == "40 bis 59 Jahre":
            LW = "LW_40bis59"
            BW = "BW_40bis59"
            DW = "DW_40bis59"
            
        if alter == "60 bis 79 Jahre":
            LW = "LW_60bis79"
            BW = "BW_60bis79"
            DW = "DW_60bis79"  
        
        operation= alter
        
   
    # falls Top Ten noch nicht berechnet, hier berechnen
    sorter = LW
    if df_berechnet == False:
        df = df.sort_values(by=sorter, ascending=False).head(10)
        
    # neues Sheet aufrufen
    slide_layout = prs.slide_layouts[2] #ist im Master Überschrift inkl aller Symbole, ansonsten leer
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    
    
    # Form der Tabelle Festlegen
    rows = df.shape[0]+1
    # Nummer der Spalten abhängig von Gesamt- oder Ressort-Auswertung
    if mode == "ressort":
        cols = 7
    if (mode == "Gesamt") | (mode=="Mantel"): 
        cols = 7
    
    left = Inches(1.05)
    top = Inches(1.55)
    width = Inches (8.3)
    height = Inches(0.6)
    
    # jetzt legen wir die Tabelle an, ein table-Objekt
    table = shapes.add_table(rows, cols, left, top, width, height).table
    
    
    # Breite der Spalten festlegen (für modus Gesamt und Ressort)
    if mode == "ressort": 
        table.columns[0].width = Inches(0.4)
        table.columns[1].width = Inches(4.0)
        table.columns[2].width = Inches(1.1)
        table.columns[3].width = Inches(0.7)
        table.columns[4].width = Inches(0.7)
        table.columns[5].width = Inches(0.7)
        table.columns[6].width = Inches(0.7)
    
    if (mode == "Gesamt") | (mode=="Mantel"): 
        table.columns[0].width = Inches(0.4)
        table.columns[1].width = Inches(3.6)
        table.columns[2].width = Inches(1.1)
        table.columns[3].width = Inches(1.5)
        table.columns[4].width = Inches(0.566)
        table.columns[5].width = Inches(0.566)
        table.columns[6].width = Inches(0.566)
    #  Index benennen
    
    table.cell(0, 0).text = ""
    table.cell(0, 1).text = "Artikel"
    table.cell(0, 2).text = "Datum"
    # Vierte Spalte wird bei Gesamt zur Ressortspalte
    if mode == "ressort": 
        table.cell(0,3).text = "Seite"
    if (mode == "Gesamt") | (mode=="Mantel"): 
        table.cell(0,3).text = "Ressort"
        
    table.cell(0, 4).text = "LW"
    table.cell(0, 5).text = "BW"
    table.cell(0, 6).text = "DW"
    
    # Überschrift festlegen
    
    title_placeholder = slide.shapes.title
    title = ""
    if headline == "default":
        if zeitung == True:
            
            title = "Top 10 - " + df.iloc[0].loc[ressort] + " " + df.iloc[0].loc[ausgabename]
        elif zeitung == False: 
            title = "Top 10 - " + df.iloc[0].loc[ressort]
    
    else:
        title = headline
    
    if len(geschlecht)>0: 
        title = "Top 10 - " + df.iloc[0].loc[ressort] + " " + geschlecht
        if kunden_id == 1011: 
            title = headline
            
            
    if len(alter)>0:
        title = "Top 10 - " + df.iloc[0].loc[ressort] + " " + alter
        if kunden_id == 1011: 
            title = headline
    
    title_placeholder.text = title
    title_placeholder.text_frame.paragraphs[0].font.bold = True
    title_placeholder.text_frame.paragraphs[0].font.name = "Campton-Bold"
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(32)
    
    
        
    
    
    # Größe + Schriftart der oberen Zeile festlegen
    # muss hier geschehen, da ich sie mit i nicht ansteuern kann
    for col in range(cols):
        text_frame = table.cell(0, col).text_frame
        p = text_frame.paragraphs[0]
        run = p.add_run()
        font = run.font
        p.font.name ="Campton-Light"
        p.font.size = Pt(11)
       
        p.alignment = PP_ALIGN.CENTER
        # table.last_row = True
    
    
    # data rows, Tabelle füllen
    for i in range(0, df.shape[0]):
        
        table.cell(i+1, 0).text = str(i+1)
    
        # Checken wie lang der Überschrift-Text ist, bei Bedarf nach dem letzten vollst Wort abschneiden und ... setzen
        # erst checken wir, ob ÜS ein String ist
        val = df.iloc[i]["Artikel"]
        if isinstance(val,str):
            if len(df.iloc[i]["Artikel"]) > 45:
                kurze_üs = textwrap.shorten(df.iloc[i]["Artikel"], width=46, placeholder="...")
                table.cell(i+1, 1).text = kurze_üs
            elif len(df.iloc[i]["Artikel"]) <=1:
                table.cell(i+1,1).text = "KEINE ÜS! Ausg: " + df.iloc[i][ausgabename] +", Seite: " + str(df.iloc[i][seite])
            else:
                table.cell(i+1, 1).text = df.iloc[i]["Artikel"][:47]
         # wenn überschrift kein String: 
        elif isinstance(val, float):  
             table.cell(i+1,1).text = "KEINE ÜS! Ausg: " + df.iloc[i][ausgabename] +", Seite: " + str(df.iloc[i][seite])
        elif isinstance(val, int):  
             table.cell(i+1,1).text = "KEINE ÜS! Ausg: " + df.iloc[i][ausgabename] +", Seite: " + str(df.iloc[i][seite])
#        SPEZOIALITÄT NW... hier war wohl eine Nummer als ÜS
#        val = df.iloc[i]["Ueberschrifttext"]
#        if isinstance(val,str):
#            table.cell(i+1, 1).text = df.iloc[i]["Ueberschrifttext"][:47]
#        elif isinstance(val, float):  
#            table.cell(i+1, 1).text = "Komisch"
#        
        # jetzt aus dem Timestamp eine lesbare Datumsangabe machen
        print(df.iloc[i][datum])
        if (kunden_id == 1011) | (kunden_id == 1016): 
            datum_bearbeitet = df.iloc[i][datum].strftime("%d.%m.%Y")
        
            table.cell(i+1, 2).text = datum_bearbeitet
        
        # Beschrift bei Gesamt= REssort, bei Ressortauswertung = Seitennummer
        if mode=="ressort": 
            table.cell(i+1, 3).text = str(df.iloc[i][seite])
        if (mode == "Gesamt") | (mode=="Mantel"):
            if df.iloc[i][ressort] == lokales_ressort:
                table.cell(i+1, 3).text = str(df.iloc[i][ressort]+" "+df.iloc[i][ausgabename])
                if kunden_id == 1016:
                    table.cell(i+1, 3).text = str(df.iloc[i][ressort])
                    
            else:
                table.cell(i+1, 3).text = df.iloc[i][ressort]
                            
            
        
     
        
        
        table.cell(i+1, 4).text = str(round(df.iloc[i][LW], 1)).replace(".", ",")
        table.cell(i+1, 5).text = str(round(df.iloc[i][BW], 1)).replace(".", ",")
        table.cell(i+1, 6).text = str(round(df.iloc[i][DW], 1)).replace(".", ",")
    
        for j in range(0,cols):
            #table.cell(i+1, j).text_frame.paragraphs[0].font.size=Pt(font_size)
            text_frame = table.cell(i+1, j).text_frame
            p = text_frame.paragraphs[0]
            run = p.add_run()
            font = run.font
            p.font.name ="Campton-Light"
            p.font.size = Pt(11)
       
            p.alignment = PP_ALIGN.CENTER
        for row in range(rows): 
            text_frame = table.cell(row,1).text_frame
            p = text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
    
    if screenshots==True: 
        screenshot(prs, df, zeitung=zeitung, mode=mode, kunden_id=kunden_id)
      
    return prs 
    
    
    
    
#%% Screenshot-Function    
    
'''
screenshot() fertigt die Screenshots inklusive der LW-Marken an. 
Sie kann händisch aufgerufen oder direkt über die function Top-10 aktiviert 
werden. 

Rangliste = False bedeutet, dass die oberste Ziffer der Marke offen bleibt. 
Rangliste = True bedeutet, dass die obere Ziffer aus df["Rangliste"] bedient wird, 
also die Platzierung des Artikels unter allen Artikeln des Tages. 
Bei vielen Ausgaben ist das aber nur in den Lokalteilen mit den nichtkumulierten 
Artikeln möglich... 

Number of Screenshots (number) ist auf 5 voreingestellt, kann aber verändert werden. 

ausgabe... hier sind Spezialfälle der jeweiligen Ausgabe hinterlegt. Kann 
man händisch ändern oder der Funktion als Parameter mitgeben. 

'''


def screenshot(prs, df, number=5, rangliste=False, ausgabe = "", 
               zeitung=True, mode="ressort", outside=False, kunden_id = 0):
#    LW = "Artikel-Lesewert (Erscheinung) in %"
#    BW = "Artikel-Blickwert (Erscheinung) in %"
#    DW = "Artikel-Durchlesewerte (Erscheinung) in %"
#    seite = "Seitennummer"
#    if df["LW"]:
    LW = "LW"
    BW = "BW"
    DW = "DW"
    seite = "Seite"
    ressort = "Ressort"
    id_nr = kunden_id
    ausgabename = "Ausgabename" # sonst auch gerne ZTG
   # df["Ressortbeschreibung"] == df["Ressort"]
        
    df = df.sort_values(by=LW, ascending=False).head(number)
    for i in range(df.shape[0]):
        #id_nr = "1011" #für die Schwäbische Zeitung
        
        
        elem = df.iloc[i]
       
        # normale URL zum Laden der Screenshots
        url = "https://lesewert.azureedge.net/layer/"+str(id_nr)+"/"+ str(elem.AusgabeId)+ "/"+ str(elem.ErscheinungsId) + "/"+\
        str(elem.Seite) + "/" +  str(elem.ArtikelId) +  ".jpg"
        print(url)
        # Ausweich-URL zum Laden der Screenshots
        url2= "https://mein.lesewert.de/Offline/ContentProxy?url=" + str(id_nr) + "/"+ str(elem.AusgabeId) + "/"+ \
        str(elem.ErscheinungsId) + "/"+ str(elem[seite]) + "/" +  str(elem.ArtikelId) +  ".jpg"
        
        response = requests.get(url, stream=True)
        response.raw.decode_content = True
    
        # Try-Block prüft, ob URL antwortet
        # Wenn nicht wird alternative URL geladen
        try:
            im = Image.open(response.raw)
        except OSError:
            print("ACHTUNG FEHLER")
            print("Konnte folgende URL nicht laden:")
            print(url)
                       
            response2 = requests.get(url2, stream=True)
            response2.raw.decode_content = True
            im = Image.open(response2.raw)
    
        #Breite des Bildes verringern, wenn der Text höher als breit ist, z.B. ganze Reportage-Seite
        # nur dann ist garantiert, dass die ÜS im Bild ist. 
        if im.size[0] < im.size[1]:
            
            corr = im.size[1]/im.size[0]
            if ausgabe == "schwaebische":
                height = int(400*corr)
                im = im.resize((400, height), Image.ANTIALIAS)
            else:
                height = int(550*corr)
                im = im.resize((550, height), Image.ANTIALIAS)
        else: 
            corr = im.size[1]/im.size[0]
            if ausgabe == "schwaebische": 
                height = int(600*corr)
                im = im.resize((600, height), Image.ANTIALIAS)
        #TODO:_ Achtung, an den Fotos rumgefummelt
        
        # Bildbearbeitung mit PILLOW
        # Beschriften der LW-Marke
        marke = Image.open("lw_marke_png.png")
        draw = ImageDraw.Draw(marke)
        
        font = ImageFont.truetype('Campton-Light.otf', size=60) 
        font_bold = ImageFont.truetype('Campton-Bold.otf', size=60) 
        (x1,y1) = (340,200)
        (x2,y2) = (340, 300)
        (x3,y3) = (340, 400)
        (x4,y4) = (340, 490)
        #message1 = str(elem["Platzierung"])
        print(elem) 
        message1 = str(elem["Platzierung"]) + str(".")
        message2 = str(round(elem[LW],1)).replace(".", ",") + "%"
        message3 = str(round(elem[BW],1)).replace(".", ",") + "%"
        message4 = str(round(elem[DW],1)).replace(".", ",") + "%"
        color="rgb(255, 255, 255)"
        draw.text((x1,y1), message1, fill=color, font=font)
        draw.text((x2,y2), message2, fill=color, font=font_bold)
        draw.text((x3,y3), message3, fill=color, font=font)
        draw.text((x4,y4), message4, fill=color, font=font)
   
        # Größe der Marke anpassen
        marke = marke.resize((944, 708), Image.ANTIALIAS)
            
        
        #Zwischenspeichern der Marke
        final1 = Image.new("RGBA", (1000, 377), (255,255,255,255))
        # TODO Feineinstellung, Ausgang (0x,0y)
        final1.paste(im, (0,0))
        if im.size[0]<im.size[1]:        # ist der Artikel lang, wird er kleiner, Marke rutscht nach rechts oben  
            final1.paste(marke, (650,5), marke)
        else: # ansonsten bleibt die Marke unten rechts
            final1.paste(marke, (530,35), marke)

        final1.save("final1.png")  
    
    
        # neuen Slide mit Screenshot des Artikels anlegen
        # prs wurde mit Funktionsaufruf übergeben
        slide_layout = prs.slide_layouts[11]  # Layout mit Bild
        slide = prs.slides.add_slide(slide_layout)
        shapes = slide.shapes
    
        
        title_placeholder = slide.shapes.title
         
        
        if zeitung == True: 
            title = 'Top ' + str(i+1) + " - " + elem[ressort] + " " + elem.ausgabename
        elif zeitung == False:
            title = 'Top ' + str(i+1) + " - " + elem[ressort]
        
        if (mode=="Mantel") | (mode=="Gesamt"):
            # Ressortbeschreibung für Überschrift im Modus Gesamt und Mantel
            elem_ressort = df.iloc[i][ressort]
            if elem_ressort == _lokales_:
                title = "Top" + str(i+1) + " - " + mode + " (" + elem_ressort + " " + df.iloc[i][ausgabename] + ")"
            else:  
                title = "Top" + str(i+1) + " - " + mode + " (" + elem_ressort + ")"
        
        title_placeholder.text = title
        title_placeholder.text_frame.paragraphs[0].font.bold = True
        title_placeholder.text_frame.paragraphs[0].font.name = "Campton-Bold"
        title_placeholder.text_frame.paragraphs[0].font.size = Pt(32)
        
        placeholder = slide.placeholders[10]
       
        filename = "final1.png"
        if outside: 
            filename = "/screenshots/"+ elem[ressort] +"_" + i +".png"
        #title_placeholder.text_frame.paragraphs[0].font = font_bold
        picture = placeholder.insert_picture(filename)
    
    
    
    return prs







#%% Screenshot function für Power BI
    

def screenshot_powerbi(df, number=5, rangliste=False, ausgabe = "", kunden_id =  0, 
               zeitung=True, mode="ressort", outside=False, markenposition="rechts"):
    LW = "Artikel-Lesewert (Erscheinung) in %"
    BW = "Artikel-Blickwert (Erscheinung) in %"
    DW = "Artikel-Durchlesewerte (Erscheinung) in %"
    seite = "Seitennummer"
    LW = "LW"
    BW = "BW"
    DW = "DW"
    id_nr = kunden_id
    
    df["Seitennummer"] = df["Seite"]
    seite = "Seite"
    ressort = "Ressort"
    
        
    df = df.sort_values(by=LW, ascending=False).head(number)
    for i in range(df.shape[0]):
        #id_nr = "1011" #für die Schwäbische Zeitung
        
        
        elem = df.iloc[i]
       
        # normale URL zum Laden der Screenshots
        url = "https://lesewert.azureedge.net/layer/"+str(id_nr)+"/"+ str(elem.AusgabeId)+ "/"+ str(elem.ErscheinungsId) + "/"+\
        str(elem.Seitennummer) + "/" +  str(elem.ArtikelId) +  ".jpg"
        # Ausweich-URL zum Laden der Screenshots
        url2= "https://mein.lesewert.de/Offline/ContentProxy?url=" + str(id_nr) + "/"+ str(elem.AusgabeId) + "/"+ \
        str(elem.ErscheinungsId) + "/"+ str(elem[seite]) + "/" +  str(elem.ArtikelId) +  ".jpg"
        
        response = requests.get(url, stream=True)
        response.raw.decode_content = True
    
        # Try-Block prüft, ob URL antwortet
        # Wenn nicht wird alternative URL geladen
        try:
            im = Image.open(response.raw)
        except OSError:
            print("ACHTUNG FEHLER")
            print("Konnte folgende URL nicht laden:")
            print(url)
                       
            response2 = requests.get(url2, stream=True)
            response2.raw.decode_content = True
            im = Image.open(response2.raw)
    
        #Breite des Bildes verringern, wenn der Text höher als breit ist, z.B. ganze Reportage-Seite
        # nur dann ist garantiert, dass die ÜS im Bild ist. 
        if im.size[0] < im.size[1]:
            corr = im.size[1]/im.size[0]
            if ausgabe == "schwaebische":
                height = int(400*corr)
                im = im.resize((400, height), Image.ANTIALIAS)
            else:
                height = int(550*corr)
                im = im.resize((550, height), Image.ANTIALIAS)
        else: 
            corr = im.size[1]/im.size[0]
            if ausgabe == "schwaebische": 
                height = int(600*corr)
                im = im.resize((600, height), Image.ANTIALIAS)
        #TODO:_ Achtung, an den Fotos rumgefummelt
        
        # Bildbearbeitung mit PILLOW
        # Beschriften der LW-Marke
        marke = Image.open("lw_marke_png.png")
        draw = ImageDraw.Draw(marke)
        
        font = ImageFont.truetype('Campton-Light.otf', size=60) 
        font_bold = ImageFont.truetype('Campton-Bold.otf', size=60) 
        (x1,y1) = (340,200)
        (x2,y2) = (340, 300)
        (x3,y3) = (340, 400)
        (x4,y4) = (340, 490)
        #message1 = str(elem["Platzierung"])
        print(elem) 
        message1 = str(elem["Platzierung"]) + str(".")
        message2 = str(round(elem[LW],1)).replace(".", ",") + "%"
        message3 = str(round(elem[BW],1)).replace(".", ",") + "%"
        message4 = str(round(elem[DW],1)).replace(".", ",") + "%"
        color="rgb(255, 255, 255)"
        draw.text((x1,y1), message1, fill=color, font=font)
        draw.text((x2,y2), message2, fill=color, font=font_bold)
        draw.text((x3,y3), message3, fill=color, font=font)
        draw.text((x4,y4), message4, fill=color, font=font)
   
        # Größe der Marke anpassen
        if markenposition == "mittig":
            marke = marke.resize((629, 472), Image.ANTIALIAS)
        else: 
            marke = marke.resize((944, 708), Image.ANTIALIAS)
            
            
        
            
        
        #Zwischenspeichern der Marke
        final1 = Image.new("RGBA", (1000, 377), (255,255,255,255))
        # TODO Feineinstellung, Ausgang (0x,0y)
        final1.paste(im, (0,0))
        
        if markenposition == "mittig":
            print("Mittig")
            final1.paste(marke, (20, 150), marke)
        else:
            if im.size[0]<im.size[1]:        # ist der Artikel lang, wird er kleiner, Marke rutscht nach rechts oben  
                final1.paste(marke, (650,5), marke)
            else: # ansonsten bleibt die Marke unten rechts
                final1.paste(marke, (530,35), marke)
        
        

        filename = "C:/Users/Felix/Desktop/LW/screenshots/"+ elem[ressort] +"_" +  str(i) +".png"

        final1.save(filename)  
    
    
       
       
      




#%% Function Ankündigungen

''' 
Diese Funktion stellt Ankündigungen anderen Nachrichten gegenüber. Sie benötigt
ein vorsortiertes DF, den Rest erledigt sie selbst.
TODO: Hier kann man eine grundlegende Tabellen-Funktion überlegen.  
''' 
def lesewert_ankündigungen(prs, df, title_text =""):
    data = df.copy()
    
    #neue Tabelle anlegen
    slide_layout = prs.slide_layouts[2]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    
    #Tabelle anlegen
    rows = 4 
    cols = 3
    left = Inches(1.05)
    top = Inches(2)
    width = Inches (8.0)
    height = Inches(0.6)
    font_size = 12
        
    total_rows = 0
        
    # Titelzeile des Sheets festlegen
    title_placeholder = slide.shapes.title            
    if len(title_text)>0:  
        title_placeholder.text = title_text
    else:  
        title_placeholder.text = "Ankündigungen - "  + data.iloc[0].ZTG
    title_placeholder.text_frame.paragraphs[0].font.bold = True
        
        #Layout der Tabelle festlegen und neues Dokument für Tabelle festlegen
        #nur dann Tabelle zeichnen, wenn mindestens ein Wert in der Tabelle abgefragt wird...    
    table = shapes.add_table(rows, cols, left, top, width, height).table
    table.columns[0].width = Inches(2)
    table.columns[1].width = Inches(2.7)
    table.columns[2].width = Inches(3.3)
        
        
    table.cell(0, 0).text = ""
    table.cell(0, 1).text = "Ankündigungen" 
    table.cell(0, 2).text = "alle anderen Texte" 
        
    #table.last_row=True
    table.cell(1,0).text = "LW"
    table.cell(2,0).text = "BW"
    table.cell(3,0).text = "DW"
    
    # TODO Liste muss für schwäbische gekennzeichnet werden
    liste_ak = ["AK", 'AK, BF','AK, NA', 'AK, HG', 'AK, BN', 'BF, AK', 'BB, AK', 'AK, RK', 'NA, AK', 'AK, BB', 'AK, IV', 'AK, RP', 'AK, FF' ]
    
    df_ak = data[data["Darstellungsformen"].isin(liste_ak)]
    df_notak = data[~data["Darstellungsformen"].isin(liste_ak)]
       

    
    for i in range(1):
            
                
               
            table.cell(i+1,1).text = str(round(df_ak\
                                                ["Artikel-Lesewert (Erscheinung) in %"].mean(), 1)).replace(".", ",")
            table.cell(i+1,2).text = str(round(df_notak\
                                                ["Artikel-Lesewert (Erscheinung) in %"].mean(), 1)).replace(".", ",")
            table.cell(i+2,1).text = str(round(df_ak\
                                                ["Artikel-Blickwert (Erscheinung) in %"].mean(), 1)).replace(".", ",")
            table.cell(i+2,2).text = str(round(df_notak\
                                                ["Artikel-Blickwert (Erscheinung) in %"].mean(), 1)).replace(".", ",")
            table.cell(i+3,1).text = str(round(df_ak\
                                                ["Artikel-Durchlesewerte (Erscheinung) in %"].mean(), 1)).replace(".", ",")
            table.cell(i+3,2).text = str(round(df_notak\
                                                ["Artikel-Durchlesewerte (Erscheinung) in %"].mean(), 1)).replace(".", ",")
           
   # Größe + Schriftart der oberen Zeile festlegen
    # muss hier geschehen, da ich sie mit i nicht ansteuern kann
    for col in range(cols):
        text_frame = table.cell(0, col).text_frame
        p = text_frame.paragraphs[0]
        run = p.add_run()
        font = run.font
        p.font.name ="Campton-Light"
        p.font.size = Pt(11)
       
        p.alignment = PP_ALIGN.CENTER
    # table.last_row = True


    # data rows, Tabelle füllen
    font_size = 12
    for i in range(0, rows):
        for j in range(0,cols):
            #table.cell(i+1, j).text_frame.paragraphs[0].font.size=Pt(font_size)
            text_frame = table.cell(i, j).text_frame
            p = text_frame.paragraphs[0]
            run = p.add_run()
            font = run.font
            p.font.name ="Campton-Light"
            p.font.size = Pt(11)
       
            p.alignment = PP_ALIGN.CENTER
        
        #for row in range(rows): 
            #text_frame = table.cell(row,1).text_frame
            #p = text_frame.paragraphs[0]
            #p.alignment = PP_ALIGN.LEFT
                
       
    return prs  
    


#%% Übersicht 3 - Lokales Schwäbische
'''
Helper-Funktion für die Schwäbische Zeitung - Ressort Lokales mit Analysen 
auf Seitentitel-Ebene im Lokalen. 

Die Funktion benötigt ein gereinigtes Dataframe, den Rest erledigt sie selbst. 


''' 
  
    

def analyse_lokales(prs, df, df_raw="noch einzusetzen"):
    
    # Liste der drei Ausgaben
    
    ausgabe_liste = ausgaben_liste
    
    st_liste = seitentitel_lokal
    
    
    df_ = df[df["Ressortbeschreibung"]=="Lokales"]
    
    for ausgabe in ausgabe_liste: 
        
        # für jeden Durchlauf ein Dataframe erstellen
        df_ausgabe = df_[df_["ZTG"]==ausgabe]
        
        
         # alle Grundwerte im jedem Schleifendurchlauf auf Null setzen
        platzierung = ["AA", "SK"]
        darstellungsform = ["NA"]
        seitentitel=[]
        
        
        
        
        deckblatt_macher(prs, df_ausgabe, "Lokales", platzierung, 
                         darstellungsform, seitentitel, ZTG= ausgabe)
        
        
        top_10(prs, df_ausgabe, df_berechnet = False, screenshots=True, number_screenshots = 5, 
          mode="ressort", headline="Top 10 Lokales " + df_ausgabe["Ausgabenteil"].unique()[0])
        
        
        # Möglichkeit für einzelne keys oder Ausgaben bestimmte Werte einzustellen
        if ausgabe == "FHA": 
            key = "strange"
        else:
            key = "xlarge"
        
        grafik_lesewert(prs, df_ausgabe, target="seitentitel", minimize=5, label_position=key,
                        ressort_liste = st_liste, order="new", legend="xlarge")
        
        
        lesewert_erscheinung(df_ausgabe, prs, title_text="LW nach Wochentagen - " + df_ausgabe.iloc[0].Ausgabenteil)
        
        
        lesewert_ankündigungen(prs, df_ausgabe)
        
    # Jetzt die Übersicht Lokalsport
    
    df_loksport = df[df["Ressortbeschreibung"]=="Lokalsport"]
    
    for ausgabe in ausgabe_liste: 
        
        df_ausgabe_sport = df_loksport[df_loksport["ZTG"]==ausgabe]
        
        platzierung = ["AA", "SK"]
        darstellungsform = ["NA"]
        seitentitel=[]
        
        deckblatt_macher(prs, df_ausgabe_sport, "Lokalsport", platzierung, 
                         darstellungsform, seitentitel, ZTG= ausgabe)
        
        top_10(prs, df_ausgabe_sport, df_berechnet = False, screenshots=True, number_screenshots = 5, 
            mode="ressort", headline="Top 10 Lokalsport " + df_ausgabe_sport["Ausgabenteil"].unique()[0])
        
        #grafik_lesewert(prs, df_ausgabe_sport, target="seitentitel", minimize=5, label_position="xlarge", order="new")
    
        #lesewert_erscheinung(df_ausgabe_sport, prs, title_text="LW nach Wochentagen - " + df_ausgabe.iloc[0].Zeitung)
        
        
    
    return prs


  
    
#%% 1.Lokalseite finden - Schwäbische
    
''' Diese Funktion sucht bei der Schwäbischen Zeitung die erste Lokalseite
in allen drei Ausgaben heraus. 
Läuft normalerweise direkt im Analysecode im Notebook. 

''' 
# Function um Lokale 1 zu identifizieren



def lokale_eins(df):
    ausgaben = ["BIB", "FHA", "RV"]
    list_datum = data2["Erscheinungsdatum"].unique()
    for ausgabe in ausgaben: 
        
        df_ausgabe = df[df["ZTG"]==ausgabe]
       
        for date in list_datum: 
            
            df_ = df_ausgabe[df_ausgabe["Erscheinungsdatum"]==date]
            df_ = df_[df_["Ressortbeschreibung"]== "Lokales"]
            df_ = df_[(df_["Seitentitel"]!="Veranstaltungen") & (df_["Seitentitel"]!="Kirchen")\
                     & (df_["Seitentitel"]!= "Wir in Kreis und Region") & (df_["Seitentitel"]!="Kultur Lokal")\
                     & (df_["Seitentitel"]!= "Region") & (df_["Seitentitel"]!="Umland")\
                     & (df_["Seitentitel"]!="Oberschwaben & Allgäu")]
            
            df_ = df_.sort_values(by="Seitennummer")
            first_page = df_.iloc[0].Seitennummer
            df_ = df_[df_["Seitennummer"]==first_page]
            for i in df_.index: 
                neuer_st = "Lokale Eins " + ausgabe
                df.set_value(i, "Seitentitel", neuer_st)
        
                
            
            
#lokale_eins(data2)
        



#%% Function Tabelle Infokästen

''' 
Diese Funktion schaut sich die Darstellungsform Infokästen (HG)

'''
def tabelle_infokästen(prs, df):
    df_ = df[df["Darstellungsformen"]=="HG"]
    df_not = df[df["Darstellungsformen"]!="HG"]
    
    # Erstellung Tabelle
     #neue Tabelle anlegen
    slide_layout = prs.slide_layouts[2]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    
    #Tabelle anlegen
    rows = 4 
    cols = 3
    left = Inches(1.05)
    top = Inches(2)
    width = Inches (8.0)
    height = Inches(0.6)
    font_size = 12
        
    total_rows = 0
        
    # Titelzeile des Sheets festlegen
    title_placeholder = slide.shapes.title            
    title_placeholder.text = "Lesewert Infokästen"
    title_placeholder.text_frame.paragraphs[0].font.bold = True
        
        #Layout der Tabelle festlegen und neues Dokument für Tabelle festlegen
        #nur dann Tabelle zeichnen, wenn mindestens ein Wert in der Tabelle abgefragt wird...    
    table = shapes.add_table(rows, cols, left, top, width, height).table
    table.columns[0].width = Inches(2)
    table.columns[1].width = Inches(2.7)
    table.columns[2].width = Inches(3.3)
        
        
    table.cell(0, 0).text = ""
    table.cell(0, 1).text = "Infokästen" 
    table.cell(0, 2).text = "alle anderen Artikel" 
        
    #table.last_row=True
    table.cell(1,0).text = "LW"
    table.cell(2,0).text = "BW"
    table.cell(3,0).text = "DW"
    
        
    for i in range(1):
            
                
               
            table.cell(i+1,1).text = str(round(df_["Artikel-Lesewert (Erscheinung) in %"].mean(), 1)).replace(".", ",")
            table.cell(i+1,2).text = str(round(df_not\
                                                ["Artikel-Lesewert (Erscheinung) in %"].mean(), 1)).replace(".", ",")
            table.cell(i+2,1).text = str(round(df_\
                                                ["Artikel-Blickwert (Erscheinung) in %"].mean(), 1)).replace(".", ",")
            table.cell(i+2,2).text = str(round(df_not\
                                                ["Artikel-Blickwert (Erscheinung) in %"].mean(), 1)).replace(".", ",")
            table.cell(i+3,1).text = str(round(df_\
                                                ["Artikel-Durchlesewerte (Erscheinung) in %"].mean(), 1)).replace(".", ",")
            table.cell(i+3,2).text = str(round(df_not\
                                                ["Artikel-Durchlesewerte (Erscheinung) in %"].mean(), 1)).replace(".", ",")
           
   # Größe + Schriftart der oberen Zeile festlegen
    # muss hier geschehen, da ich sie mit i nicht ansteuern kann
    for col in range(cols):
        text_frame = table.cell(0, col).text_frame
        p = text_frame.paragraphs[0]
        run = p.add_run()
        font = run.font
        p.font.name ="Campton-Light"
        p.font.size = Pt(11)
       
        p.alignment = PP_ALIGN.CENTER
    # table.last_row = True


    # data rows, Tabelle füllen
    font_size = 12
    for i in range(0, rows):
        for j in range(0,cols):
            #table.cell(i+1, j).text_frame.paragraphs[0].font.size=Pt(font_size)
            text_frame = table.cell(i, j).text_frame
            p = text_frame.paragraphs[0]
            run = p.add_run()
            font = run.font
            p.font.name ="Campton-Light"
            p.font.size = Pt(11)
       
            p.alignment = PP_ALIGN.CENTER
        
        for row in range(rows): 
            text_frame = table.cell(row,1).text_frame
            p = text_frame.paragraphs[0]
            #p.alignment = PP_ALIGN.LEFT
                
    return prs


    


        
#%% Func gelesene vs veröffentlichte Artikel
# TODO muss noch für die Library angepasst werden, ist jetzt erstmal nur rüberkopiert. 
    

def gelesene_artikel(prs, df_scan, df_nichtkum):
    df_s = df_scan.copy()
    df_nk = df_nichtkum.copy()
    
    #ID in String verwandeln
    df_s["TeilartikelVeroeffentlichungenId"] = df_s["TeilartikelVeroeffentlichungenId"].apply(str)
    
    # Umbennen der ID-Spalte
    df_s.rename(columns={"TeilartikelVeroeffentlichungenId":"ArtikelId"}, inplace=True)
    
    # merge mit dem Datensatz, in dem die doppelten Split-IDs noch vorhanden sind
    df_merge = pd.merge(df_s, df_nk, how="left", on="ArtikelId")
    df_m = df_merge.copy()
    
    df_m = df_m[["ArtikelId", "Erscheinungsdatum", "WellenteilnahmenId", "Ressortbeschreibung", "ZTG", "Ausgabename"]]
    
    # Alle Scans ohne entsprechendes Printgegenüber löschen
    # neue Spalte mit Wochentag anlegen (bei manchen Analysen gefordert)
    df_m = df_m[~df_m["Erscheinungsdatum"].isnull()]
    df_m["day"] = df_m["Erscheinungsdatum"].dt.weekday

    
    # Berechnung durschschnittliche Anzahl Artikel pro Tag
    # Anzahl Erscheinungstage:
    anzahl_et = df_m.Erscheinungsdatum.nunique() 

    # Anzahl Leser pro Tag und Ausgabe bzw. Gesamt
    df_group_leserpT_Gesamt = df_m.groupby(["Erscheinungsdatum", "WellenteilnahmenId"], as_index=False).count()
    df_group_leserpT_Gesamt = df_group_leserpT_Gesamt.groupby("Erscheinungsdatum", as_index=False).count()
    df_group_leserpT_Gesamt.rename(columns={"WellenteilnahmenId":"Leser_Tag_gesamt"}, inplace=True)
    df_group_leserpT_Gesamt = df_group_leserpT_Gesamt[["Erscheinungsdatum", "Leser_Tag_gesamt"]]

    df_group_leserpT_Ausgabe = df_m.groupby(["Erscheinungsdatum", "Ausgabename", "WellenteilnahmenId"], as_index=False).count()
    df_group_leserpT_Ausgabe = df_group_leserpT_Ausgabe.groupby(["Erscheinungsdatum", "Ausgabename"], as_index=False).count()
    df_group_leserpT_Ausgabe.rename(columns={"WellenteilnahmenId":"Leser_Tag_Ausgabe"}, inplace=True)
    df_group_leserpT_Ausgabe = df_group_leserpT_Ausgabe[["Erscheinungsdatum", "Ausgabename", "Leser_Tag_Ausgabe"]]

    # Merge Leser Ausgabe und Gesamt
    df_group_leser = pd.merge(df_group_leserpT_Ausgabe, df_group_leserpT_Gesamt, how="left", on="Erscheinungsdatum")


    # Anzahl gelesene Artikel pro Tag und Ausgabe bzw. Gesamt
    df_group_artikelpT_Gesamt = df_m.groupby("Erscheinungsdatum", as_index=False).count()
    df_group_artikelpT_Gesamt.rename(columns={"ArtikelId":"Scans_Gesamt"}, inplace=True)
    df_group_artikelpT_Gesamt = df_group_artikelpT_Gesamt[["Erscheinungsdatum", "Scans_Gesamt"]]


    df_group_artikelpT_Ausgabe = df_m.groupby(["Erscheinungsdatum", "Ausgabename"], as_index=False).count()
    df_group_artikelpT_Ausgabe.rename(columns={"ArtikelId":"Scans_Ausgabe"}, inplace=True)
    df_group_artikelpT_Ausgabe = df_group_artikelpT_Ausgabe[["Erscheinungsdatum", "Ausgabename", "Scans_Ausgabe"]]

    # Zusammenführen der Artikel-Count-Dateien
    df_group_scans = pd.merge(df_group_artikelpT_Ausgabe, df_group_artikelpT_Gesamt, how="left", on="Erscheinungsdatum")

    # Anzahl veröffentlichter Artikel 
    df_Artikel_Gesamt = data2.groupby(["Erscheinungsdatum"], as_index=False).count()
    df_Artikel_Gesamt.rename(columns={"SplitId":"Artikelzahl_Tag_Gesamt"}, inplace=True)
    df_Artikel_Gesamt= df_Artikel_Gesamt[["Erscheinungsdatum", "Artikelzahl_Tag_Gesamt"]]

    df_Artikel_Ausgabe = data2.groupby(["Erscheinungsdatum", "Ausgabename"], as_index=False).count()
    df_Artikel_Ausgabe.rename(columns={"SplitId":"Artikelzahl_Tag_Ausgabe"}, inplace=True)
    df_Artikel_Ausgabe = df_Artikel_Ausgabe[["Erscheinungsdatum", "Ausgabename", "Artikelzahl_Tag_Ausgabe"]]

    #Final merge
    df_group = pd.merge(df_group_artikel, df_group_leser, how="inner",  on=["Erscheinungsdatum", "Ausgabename"])
    df_group_final = pd.merge(df_group, df_Artikel_Ausgabe, on=["Erscheinungsdatum", "Ausgabename"], how="inner")
    
    # Spalten mit Analyse-Zahlen
    df_group_final["Artikel_Leser_Tag_Gesamt"] = df_group_final["Scans_Gesamt"] / df_group_final["Leser_Tag_gesamt"]
    df_group_final["Artikel_Leser_Tag_Ausgabe"] = df_group_final["Scans_Ausgabe"] / df_group_final["Leser_Tag_Ausgabe"] 
    
    # df fpür Grafik ausspielen
    df_grafik = df_group_final.groupby("Ausgabename", as_index=False).mean()
    
    
    
    # GRAFIK PLOTTEN
    
    # Achsen festlegen
    x= df_grafik["Ausgabename"]
    df_grafik["Label"] = df_grafik["Ausgabename"].apply(lambda x: dict_ausgaben[x])
    labels = df_grafik["Label"]
    xn = range(len(x))
    y = df_grafik["Artikel_Leser_Tag_Ausgabe"]
    yn = df_grafik["Artikelzahl_Tag_Ausgabe"]
    set_font_color ="#8c8f91" 
    
     # Seaborn-Style und Größe des Plots festlegen
    sns.set_style("white")
    fig, ax1 = plt.subplots(figsize=(20,8))
    
    #setzt die linke Y-Achse in Campton light
    # rechte Y-Achse können wir erst zum Code-Ende ansteuern
    plt.yticks(fontproperties=campton_light)

     # Achsen, Ticks und alles andere festlegen
    
   
    # Barcharts einzeichnen
    bars = ax1.bar(xn,y, color="#f77801", width=0.3, label="Ø gelesene Artikel/Tag")
            
    
                   
                   
        
    ax1.set_ylabel('Ø Artikel/Tag', color= set_font_color, \
                   fontproperties=campton_light, fontsize=50)
    ax1.xaxis.set(ticks=range(0, len(xn))) #  Anzahl der Ticks 
    ax1.set_xticklabels(labels = labels, rotation=45, ha="right",  weight=800,\
                        color= set_font_color, fontproperties=campton_light, \
                        fontsize=30) # Labels werden ausgerichtet
    
    ax1.patch.set_facecolor('white') # Hintergrundfarbe auf weiß, dann... 
    ax1.patch.set_alpha(0.0) # Hintergrund ausblenden, damit zweite Grafik 
                                                #   (der Plot) sichtbar wird
    ax1.set_zorder(2) # erste Grafik wird vor die zweite Grafik geschoben
    
    # Werte über die Balken setzen
    for p in ax1.patches:
        # Problem: Ist ein Balken nur vier groß, ist der Abstand zu gering
        # TODO... NOCH EIN BISSCHEN EXPERIMENTIEREN
        height = p.get_height()
        
        txt = '{:1.1f}'.format(height).replace(".", ",")
        ax1.text(p.get_x()+p.get_width()/2., height + 3, txt, ha="center",\
                     fontproperties=campton_light, color= set_font_color, rotation=0\
                     ,fontsize= 30, weight = 1000)
        
       
   
    ax1.plot(xn, yn, alpha=0.2) # Linie wird gefadet
    fill = ax1.fill_between(xn, yn, alpha=1, color="#ffcd92", label="Ø Artikel/Tag") #, label=ax2_y_label
                            # Raum unter Yn-Linie wird gefüllt, Farbe wird transparent
    ax1.set_ylim(0,yn.max()+1) # manuelles Setzen der Y-Achse rechts, 
                                #damit die auch bei 0 anfängt. 
   
    
    #ax2.set_ylabel(ax2_y_labeltext, color= set_font_color, \
                   #fontproperties=campton_light, labelpad=15) # labelpad =margin

    # Größe der Achsen-Beschriftung festlegen
        #ax2.yaxis.label.set_size(22)
        
    ax1.yaxis.label.set_size(22)
    
    
    
    # Abstände Bars zur Achse (standardmäßig bei 0.5)
    plt.margins(x=0.03) # ziehen Bars näher an die Achse
     #obere Linie ausblenden
    ax1.spines["top"].set_visible(False)
    #ax1.spines["left"].set_color("gray")
    ax1.spines["top"].set_visible(False)
    ax1.spines["bottom"].set_visible(False)
    ax1.spines["left"].set_visible(False)
    ax1.spines["right"].set_visible(False)
    
    
    
    
    
    
    # jetzt werden die Y-Ticks links in Campton Light gefasst
    plt.yticks(fontproperties=campton_light)
    
    
    #
     # Zahlen an den Y-Achsen verändern, Größe und Farbe
    ax1.tick_params(axis='y', labelsize=25, colors= set_font_color)
    
    legend_height = 1.24
    
    
    # Legende einbauen
    
    leg = plt.legend(bbox_to_anchor=(1, legend_height), handles=[bars, fill], markerscale=140)
        
        
        
        
  
    for text in leg.get_texts(): 
        plt.setp(text, color= set_font_color, size=21)
    
   
             
    plt.tight_layout()
    
    # Canvaseinstellung / Position des Plots
    # Function nutzt Voreinstellung aus Parametern
    pos_left = 0.08
    pos_right=0.92
    pos_top=0.85
    pos_bottom = 0.4
    
    plt.subplots_adjust(left=pos_left, right=pos_right, top=pos_top, 
                        bottom=pos_bottom)
    
    filename = "grafik_lesewert_plot_two_axis.png"
    plt.savefig(filename)
    
    plt.close()
    # Plot wird auf PPTX-Sheet gezogen
    title_text = "Artikel - Erschienen und gelesen"
    lw.picture_sheet(prs, filename, title_text=title_text)            
    
    return prs
    
        

#%% Kolumnen-Finder 




#%%
#%%#%%
#%% HELPER-FUCNTIONS

#%% Deckblatt-Generator
    # klappt noch nicht so richtig mit dem Design... 
    
def deckblatt(prs, title_text):
    slide_layout = prs.slide_layouts[15]  # Layout mit zwei Platzhaltern idx 10 + 12
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    
    
    title_placeholder = slide.shapes.title
    
    #title_placeholder.text = "Hallo"
    title_placeholder.text = title_text 
    title_placeholder.text_frame.paragraphs[0].font.bold = True
    title_placeholder.text_frame.paragraphs[0].font.name = "Campton-Bold"
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(32)
    return prs
    
#%% Func Bilder auf PPTX-Folien

'''
HELPER-FUNCTIONEN



POWERPOINT-ERSTELLUNG

picture_sheet(prs, filename)

Diese Function übernimmt das PRS-Objekt und eine erstellte Bilddatei (z.B. über
Matplotlib) und fertigt damit ein neues Sheet an. 


'''
def picture_sheet(prs, filename, title_text):
    
    slide_layout = prs.slide_layouts[11]  # Layout mit Bild
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    
    placeholder = slide.placeholders[10]
    title_placeholder = slide.shapes.title
    
    pic_path = filename
    pic_left = Inches(0.4)
    pic_top = Inches(2)
    #pic_height = Inches(8)
    pic_width = Inches(9)
    #pic = slide.shapes.add_picture(pic_path, pic_left, pic_top, width=pic_width)
    picture = placeholder.insert_picture(pic_path)
    
    #Titel festlegen
    title_placeholder = slide.shapes.title
    #title_placeholder.text = "Hallo"
    title_placeholder.text = title_text 
    title_placeholder.text_frame.paragraphs[0].font.bold = True
    title_placeholder.text_frame.paragraphs[0].font.name = "Campton-Bold"
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(32)
    return prs
        
     
#%% Zwei Bilder auf einer Folie
    
def double_picture_sheet(prs, filename1, filename2, title_text):
    
    slide_layout = prs.slide_layouts[12]  # Layout mit zwei Platzhaltern idx 10 + 12
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    
    placeholder1 = slide.placeholders[10]
    placeholder2 = slide.placeholders[12]
    title_placeholder = slide.shapes.title
    
    pic_path1 = filename1
    pic_path2 = filename2
    
    pic_left = Inches(0.4)
    pic_top = Inches(2)
    #pic_height = Inches(8)
    pic_width = Inches(9)
    #pic = slide.shapes.add_picture(pic_path, pic_left, pic_top, width=pic_width)
    picture1 = placeholder1.insert_picture(pic_path1)
    picture2 = placeholder2.insert_picture(pic_path2)
    
    #Titel festlegen
    title_placeholder = slide.shapes.title
    #title_placeholder.text = "Hallo"
    title_placeholder.text = title_text 
    title_placeholder.text_frame.paragraphs[0].font.bold = True
    title_placeholder.text_frame.paragraphs[0].font.name = "Campton-Bold"
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(32)
    return prs


#%% Func Zeichnen von LW-Bar-Charts
'''
MATPLOTLIB-ERSTELLUNG

plot_two_axis()

Diese Funktion übernimmt vorgefertigte Werte und erstellt daraus eine 
Grafik mit zwei verschiedenen Y-Achsen (links und rechts), einem Balkendiag
und einer Farbfläche im Hintergrund. 

Dafür muss angegeben werden, wie sich x, y, yn zusammensetzen soll. 
Voreingestellt sind hierzu die x, y, yn etc. Werte, die in der übergeordneten
Funktion errechnet werden müssen. Eigentlich sollten sie dann automatisch 
übernommen werden. TESTEN!!!! 

Außerdem muss es möglich sein, die Position der Grafik im Canvas zu verschieben.
Voreingestellt werden die Standardwerte übergeben, der Nutzer kann dann bei 
Bedarf andere Werte einstellen.  

Mit article="total" rechnet die Funktion die Anzahl aller Artikel im 
Erscheinungsraum zusammen. 
article="mean" zeigt die Zahl der Artikle pro Tag im Durchschnitt an. 

Mit grid=False werden die Gridlinien auf der Y-Achse ausgeschaltet. Können 
durch grid=True wieder eingeschaltet werden. 

Parameter axis zeigt an, ob mit nur einer oder Y-zwei Achsen gearbeitet 
werden soll

Wenn mean_line != 0 gesetzt wird, wird eine Linie eingezeichnet, die den 
Durchschnittswert anzeigt.

mit legend="strange" lassen sich besondere formate bearbeiten

'''


def plot_axis(prs, x=0, labels = 0, xn = 0, y = 0, yn = 0,
                  pos_left = 0.08, pos_right=0.92, pos_top=0.85, 
                  pos_bottom = 0.33, article = "total", grid=False,
                  title_text="Bitte Titel eingeben", axis=2, 
                  mean_line = 0, legend="normal", special=False, 
                  mean_line_title = "", umbruch_x=True, 
                  limit_y=False, font_size = 30):
    
    #Schriftfarbe und Farbe der Ticks festlegen
    set_font_color ="#8c8f91" 
    
    # Schriftgröße festlegen
    font_size = font_size
    print("Hier ist plot axis, zeile 5315. Festgelegte Schriftgröeß: {}".format(font_size))
    print("Hier die Werte: x:{}, xn:{}".format(x, xn))
    # Werte für die Achsen werden festgelegt
    x = x
    xn = xn
    y = y
    yn = yn
    
    def check_labels(elem):
        if len(elem)>=18:
            return_label = "\n".join(textwrap.wrap(elem, width=8))
            return return_label
        else: 
            return elem
    
            
#    if umbruch_x:
#        labels = x.apply(check_labels)
        
    
    
    # Seaborn-Style und Größe des Plots festlegen
    sns.set_style("white")
    fig, ax1 = plt.subplots(figsize=(20,8))
    
    #setzt die linke Y-Achse in Campton light
    # rechte Y-Achse können wir erst zum Code-Ende ansteuern
    plt.yticks(fontproperties=campton_light)

     # Achsen, Ticks und alles andere festlegen
    
    if axis==2: 
        ax2 = ax1.twinx() # Zwillings-Y-Achse anlegen, X-Achse wird geteilt

        # Wenn nötig - Gridlinien
        if grid==True: 
            ax1.grid(color= set_font_color, linestyle='-', linewidth=1, axis="y")
    else:
        if grid==True: 
            ax1.grid(color= set_font_color, linestyle='-', linewidth=1, axis="y")
    
     # Barcharts einzeichnen
    # Label für die X-Achse-Legendenbeschriftung festlegen
    legende_text = "Lesewert"
    if special=="Lesetage":
        legende_text ="Ø Leserzahl"
    bars = ax1.bar(xn,y, color="#f77801", width=0.3, label=legende_text)
            
    
                   
    # setzt Text der Y-Achse fest
               
    y_label_text = 'Ø Lesewert in Prozent'
    if special == "Lesetage": 
        y_label_text = "Ø Leserzahl pro Tag"
        
    ax1.set_ylabel(y_label_text, color= set_font_color, \
                   fontproperties=campton_light, fontsize=50)
    ax1.xaxis.set(ticks=range(0, len(xn))) #  Anzahl der Ticks 
    ax1.set_xticklabels(labels = labels, rotation=45, ha="right",  weight=800,\
                        color= set_font_color, fontproperties=campton_light, \
                        fontsize=30) # Labels werden ausgerichtet
    
    ax1.patch.set_facecolor('white') # Hintergrundfarbe auf weiß, dann... 
    ax1.patch.set_alpha(0.0) # Hintergrund ausblenden, damit zweite Grafik 
                                                #   (der Plot) sichtbar wird
    ax1.set_zorder(2) # erste Grafik wird vor die zweite Grafik geschoben
    
   
    
    for p in ax1.patches:
        # Problem: Ist ein Balken nur vier groß, ist der Abstand zu gering
        # TODO... NOCH EIN BISSCHEN EXPERIMENTIEREN
        height = p.get_height()
         # Werte über die Balken setzen
    # height_text gibt an, wie vile höher das ganze Gesetzt werden soll
        height_text = height +1
        if special == "Lesetage": 
            height_text = height + 10
        if height <8:
            height_text = height+1
        txt = '{:1.1f}'.format(height).replace(".", ",")
        ax1.text(p.get_x()+p.get_width()/2., height_text, txt, ha="center",\
                     fontproperties=campton_light, color= set_font_color, rotation=0\
                     ,fontsize= font_size, weight = 1000)
        
       
    if axis==2:
        
     # Korrekte Y2-Beschriftung, Label und Legende
     # Parameter article wird hier verarbeiteta
        if article=="total": 
             ax2_y_labeltext = "Artikel im Messzeitraum"
             ax2_y_label ="Artikelanzahl"
        elif article=="mean":
             ax2_y_labeltext ="Ø Anzahl Artikel pro Tag"
             ax2_y_label = "Ø Artikel/Tag"
    
        ax2.plot(xn, yn, alpha=0.2) # Linie wird gefadet
        fill = ax2.fill_between(xn, yn, alpha=1, color="#ffcd92", label=ax2_y_label)
                            # Raum unter Yn-Linie wird gefüllt, Farbe wird transparent
        ax2.set_ylim(0,yn.max()+1) # manuelles Setzen der Y-Achse rechts, 
                                #damit die auch bei 0 anfängt. 
   
    
        ax2.set_ylabel(ax2_y_labeltext, color= set_font_color, \
                   fontproperties=campton_light, labelpad=15) # labelpad =margin

    # Größe der Achsen-Beschriftung festlegen
        ax2.yaxis.label.set_size(22)
        
    ax1.yaxis.label.set_size(22)
    
    
    
    # Abstände Bars zur Achse (standardmäßig bei 0.5)
    plt.margins(x=0.03) # ziehen Bars näher an die Achse
     #obere Linie ausblenden
    ax1.spines["top"].set_visible(False)
    #ax1.spines["left"].set_color("gray")
    ax1.spines["top"].set_visible(False)
    ax1.spines["bottom"].set_visible(False)
    ax1.spines["left"].set_visible(False)
    ax1.spines["right"].set_visible(False)
    
    
    if axis==2: 
        ax2.spines["top"].set_visible(False)
        ax2.spines["top"].set_visible(False)
        ax2.spines["bottom"].set_visible(False)
        ax2.spines["left"].set_visible(False)    
        ax2.spines["right"].set_visible(False)
    
    
    
    # jetzt werden die Y-Ticks links in Campton Light gefasst
    plt.yticks(fontproperties=campton_light)
    
    
    # Bei Bedarf Linie mit dem Durchschnitt einziehen
    if mean_line !=0:
        print(mean_line)
        labeltext = mean_line_title +": {:1.1f}".format(float(mean_line)).replace(".", ",")
        print("Achtung, Labeltext")
        print(labeltext)
        linie = ax1.axhline(y=mean_line, xmin=0.01, xmax=0.99, color=set_font_color, label=labeltext)
    
    
     # Zahlen an den Y-Achsen verändern, Größe und Farbe
    ax1.tick_params(axis='y', labelsize=25, colors= set_font_color)
    
    legend_height = 1.24
    if legend=="normal":
            legend_height = 1.24
    elif legend == "strange": 
            legend_height =1.35
            pos_bottom = 0.39
            pos_top = 0.8
    else:
        legend_height = 1.24
    print(mean_line)
    if axis==2: 
        ax2.tick_params(axis='y', labelsize=25, colors= set_font_color)
    
    
    # Legende einbauen
        print(mean_line)
        print ("----------------------")
        if mean_line != 0:
            print("mean_line !=0, axis =2")
            print(mean_line)
            leg = plt.legend(bbox_to_anchor=(1, legend_height), handles=[bars, fill, linie], markerscale=140)
        else:
            leg = plt.legend(bbox_to_anchor=(1, legend_height), handles=[bars, fill], markerscale=140)
        
        
        
        
        
    elif axis==1:
        print(mean_line)
        if mean_line != 0:
            print("mean_line !=0, axis =1")
            print(mean_line)
            leg = plt.legend(bbox_to_anchor=(1, legend_height), handles=[bars, linie], markerscale=140)
            
        else:
            leg = plt.legend(bbox_to_anchor=(1, legend_height), handles=[bars], markerscale=140)
    
    for text in leg.get_texts(): 
        plt.setp(text, color= set_font_color, size=21)
    
   
             
    plt.tight_layout()
    
    # Canvaseinstellung / Position des Plots
    # Function nutzt Voreinstellung aus Parametern
    plt.subplots_adjust(left=pos_left, right=pos_right, top=pos_top, 
                        bottom=pos_bottom)
    
    filename = "grafik_lesewert_plot_two_axis.png"
    plt.savefig(filename)
    
    plt.close()
    # Plot wird auf PPTX-Sheet gezogen
    picture_sheet(prs, filename, title_text=title_text)            
    
    return prs


#%% Schlagworte nach Begriffen durchsuchen 
    
''' Value=<string> legt fest, ob die Funktion nach Sportarten, Vereinen
oder Unternehmen sucht


sort = "LW" oder "Artikelzahl" zeigt an, wie ddie Top 10 erstellt und wie sie georndet werden soll
min_artikel = Mindestanzahl an Artikel 
''' 
    
def schlagworte_finden(prs, df, value="sportart", sort= "LW", min_artikel = 5, title=""):
    
    
    #Erstellen ein neues Dict, um daraus später ein DF zu machen
    result_dict = {}
    #Labels festlegen,die später dem Dataframe übergeben werden
    if value == "sportart":
        labels = ["Sportart", "LW", "BW", "DW", "Artikelanzahl"]
        col = "Themen"
        liste = liste_sportarten
    if value == "vereine":
        labels = ["Verein", "LW", "BW", "DW", "Artikelanzahl"]
        col = "Akteure"
        liste = liste_vereine
    if value == "unternehmen": 
        labels = ["Unternehmen", "LW", "BW", "DW", "Artikelanzahl"]
        col = "Akteure"
        liste = liste_unternehmen
    # jetzt durchlaufen wir die Liste mit Analysewörtern
    for elem in liste: 
        # check, ob das Analysewort irgendwo im String vorhanden ist
        # nan werden ausgeblendet
        df_elem = df[df[col].str.contains(elem, na=False)]
        #Dict füllen, Analysewort als Key, die Werte als Values
        result_dict[elem] = round(df_elem["Artikel-Lesewert (Erscheinung) in %"].mean(),1), \
        round(df_elem["Artikel-Blickwert (Erscheinung) in %"].mean(),1),\
        round(df_elem["Artikel-Durchlesewerte (Erscheinung) in %"].mean(),1), \
        df_elem.shape[0] # letzter Wert shape[0] = Anzahl gefundener Artikel
    
    # pd-Dataframe erstellen, orient="Index" bedeutet, Keys werden Zeilen, Werte werden Spalten
    sw_df = pd.DataFrame.from_dict(result_dict, orient="index").reset_index()
    
    #Spalten umbenennen
    sw_df.columns=labels
    #sw_df.loc[sw_df.Sportart=="American Football", ["Sportart"]] = "Am. Football"
    # Werte sortieren, höchster Lesewert zuerst
    sw_df = sw_df[sw_df["Artikelanzahl"]>=min_artikel]
    sw_df = sw_df.sort_values(by=sort, ascending=False).head(10)
    
    x_col = sw_df.columns[0]
    x = sw_df[x_col]
    labels = x
    xn = range(len(x))
    y = sw_df["LW"]
    yn = sw_df["Artikelanzahl"]
    if len(title)==0:
        title_string = "Lesewert nach " + x_col
    if len(title)>=1: 
        title_string = title
    plot_axis(prs, x=x, labels = labels, xn = xn, y = y, yn = yn,\
                  pos_left = 0.17, pos_right=0.92, pos_top=0.85, \
                  pos_bottom = 0.48, article = "total", grid=False,\
                  title_text = title_string, axis=2, )
    

    
    
    
    return prs

#%% Muss-Kann-Soll
def mks(df_ereignis, df, liste=ausgaben_liste):
    df_ = df_ereignis[(df_ereignis["Key"]== "Priority.High") | 
                                  (df_ereignis["Key"] == "Priority.Low") | 
                                    (df_ereignis["Key"] == "Priority.Medium")]
  
    
    df_["TeilartikelVeroeffentlichungsId"] = df_["TeilartikelVeroeffentlichungsId"].astype(int)
    df["ArtikelId"] = df["ArtikelId"].astype(int)
    #ausgaben = ["BIB", "FHA", "RV"]
    
     
    #df = df[df["Ressortbeschreibung"]=="Lokales"]
    
   
   
    df_merge = pd.merge(df, df_, left_on="ArtikelId", right_on="TeilartikelVeroeffentlichungsId", how="left")
    
   
    
    #df_mks = df_merge[~df_merge["TeilartikelVeroeffentlichungsId"].isnull()]
    #df_mks = df_merge.groupby("Key").count()
   
    
#    sizes = df_mks["ArtikelId"]
#    set_font_color ="#8c8f91"
    
     # legt Position der Zahlen innerhalb der Pie-Chart fest
#    def autodistance(autopct): 
#        if autopct >=5:
#            return 1.15 # 1.2 bei kleinen Zahlen
#        else:
#            return 0.8
        
#    # Erstellt die tatsächlichen Nutzerzahlen anhand der automatisch generierten Prozentzahlen
#    def make_label(sizes):
#        total = sizes.shape[0]
#        def my_label(pct):
#            return round(total*(pct/100)).astype(int)
#        mylabel = my_label
#        
#        #if mylabel >=6:
#           # pct_distance = 0.8
#        #if mylabel <6:
#           # pct_distance = 1.2
#        return mylabel
#    
#        # Falls nötig: Explode wird für jedes Element errechnet
#        explode = ()
#        for i in range(len(sizes)):
#            explode = explode +(0.03,)
#        
#        colors =["#ff9900", "#ffc570", "#fce3bf", "#fcf6d1"]
#        fig, ax1 = plt.subplots(figsize=(8,8))
#        test = make_label(sizes)
#        
#        
#        ax1.pie(sizes, radius=1, frame=False, shadow=False, autopct=make_label(sizes), startangle=90, colors=colors,\
#                textprops={'fontsize': 20, "color": set_font_color},pctdistance =0.8, 
#                wedgeprops={"linewidth":2, "edgecolor":"white"})
#        ax1.axis("equal")
#        
#            
#        
#        handles, labels = ax1.get_legend_handles_labels()
#        labels = df_group["Antworttext"].values     
    df_muss = df_merge[df_merge["Key"]=="Priority.High"].count()
    
    df_kann = df_merge[df_merge["Key"]=="Priority.Medium"].count()
    df_soll = df_merge[df_merge["Key"]=="Priority.Low"].count()
    #total = df_muss["ArtikelId"]+ df_kann["ArtikelId"] + df_soll["ArtikelId"]
    total = df_muss["AusgabeId"] + df_kann["AusgabeId"] + df_soll["AusgabeId"]
    muss = df_muss["AusgabeId"] / total * 100
    kann = df_kann["AusgabeId"] / total * 100
    soll = df_soll["AusgabeId"] / total * 100
   
    muss_lw = df_merge[df_merge["Key"]=="Priority.High"]["Artikel-Lesewert (Erscheinung) in %"].mean()
    kann_lw = df_merge[df_merge["Key"]=="Priority.Medium"]["Artikel-Lesewert (Erscheinung) in %"].mean()
    soll_lw = df_merge[df_merge["Key"]=="Priority.Low"]["Artikel-Lesewert (Erscheinung) in %"].mean()
    aus = ""
    print ("-------------------------  " + aus + "  ----------------")
    print("muss für Lok " + aus +": {}  in Prozent: ".format(df_muss))
    print("in Prozent {:1.1f}".format(muss))
    print("kann für Lok " + aus +": {} in Prozent: ".format(df_kann))
    print("in Prozent {:1.1f}".format(kann))
    print("soll für Lok " + aus +": {} in Prozent: ".format(df_soll))
    print("in Prozent {:1.1f}".format(soll))
   
    print ("LW muss für " + aus + ": {:1.1f}".format(muss_lw))
    print ("LW kann für " + aus + ": {:1.1f}".format(kann_lw))
    print ("LW soll für " + aus + ": {:1.1f}".format(soll_lw))
    
#
#    
#    
#        
#        colors =["#ff9900", "#ffc570", "#fce3bf", "#fcf6d1"]
#        fig, ax1 = plt.subplots(figsize=(8,8))
#        test = make_label(sizes)
#        
#        
#        ax1.pie(sizes, radius=1, frame=False, shadow=False, autopct=make_label(sizes), startangle=90, colors=colors,\
#                textprops={'fontsize': 20, "color": set_font_color},pctdistance =0.8, 
#                wedgeprops={"linewidth":2, "edgecolor":"white"})
#        ax1.axis("equal")
#        
#            
#        
#        handles, labels = ax1.get_legend_handles_labels()
#        labels = df_group["Antworttext"].values
#        
#           
#        if elem in sortierung_umfrage:
#            sorter = sortierung_umfrage[elem]
#            labels = sorted(labels, key=sorter.index)
#        else: 
#            print(elem + " ... is not in dict")
#        
#        leg = ax1.legend(handles, labels=labels, #bbox_to_anchor=(1, 0.2), 
#                         markerscale=14, mode="expand", 
#                         borderaxespad=0., loc=8) #1.Wert x, 2. Wert y
#        for text in leg.get_texts(): 
#            plt.setp(text, color= set_font_color, size=20)
#        
#        # Überschrift
#        # ist sie länger als 40, wird sie mit textwrap auf zwei Zeilen verteilt
#        figure_title= elem
#        posx=0.5
#        posy=0.82
#        if len(figure_title) >40:
#            figure_title = "\n".join(textwrap.wrap(elem, width=35))
#            posy=0.8
#          
#        
#        plt.text(posx, posy, figure_title,
#             horizontalalignment='center',
#             
#             transform = ax1.transAxes, 
#            fontproperties=campton_light, 
#            color = set_font_color, 
#            fontsize=28)
#        
#        
#        filename = "grafik_mks.png"
#        filename2 = "grafik_pie2.png"
#       
#        # Labels werden jetzt sortiert
#        # Sortierung findet sich in Liste sortierung_umfrage
#        
#        
#        if elem in sortierung_umfrage:
#            sorter = sortierung_umfrage[elem]
#            labels = sorted(labels, key=sorter.index)
#        else: 
#            print(elem + " ... is not in dict")
#        
#        leg = ax1.legend(handles, labels=labels, #bbox_to_anchor=(1, 0.2), 
#                         markerscale=14, mode="expand", 
#                         borderaxespad=0., loc=8) #1.Wert x, 2. Wert y
#        for text in leg.get_texts(): 
#            plt.setp(text, color= set_font_color, size=20)
#        
#        # Überschrift
#        # ist sie länger als 40, wird sie mit textwrap auf zwei Zeilen verteilt
#        figure_title= elem
#        posx=0.5
#        posy=0.82
#        if len(figure_title) >40:
#            figure_title = "\n".join(textwrap.wrap(elem, width=35))
#            posy=0.8
#          
#        
#        plt.text(posx, posy, figure_title,
#             horizontalalignment='center',
#             
#             transform = ax1.transAxes, 
#            fontproperties=campton_light, 
#            color = set_font_color, 
#            fontsize=28)
#        
#        
#        filename = "grafik_pie.png"
#        filename2 = "grafik_pie2.png"
#        
#        #plt.tight_layout()
#        #plt.savefig(filename)
#    # Canvaseinstellung / Position des Plots
#    # Function nutzt Voreinstellung aus Parametern
#    
#        
#        pos_left = 0.25 # 0.2
#        pos_right=0.75 #0.8
#        pos_top=0.98
#        pos_bottom = 0.1
#        plt.subplots_adjust(left=pos_left, right=pos_right, top=pos_top, 
#                       bottom=pos_bottom)
#        
#        
#        # plt.savefig(filename, bbox_inches="tight")
#        if len(title)>0:
#            title_text = "Zeitungsnutzung " +ausgaben_dict[title]
#        else: 
#            title_text = "Zeitungsnutzung"
#        
#        
#        if counter%2!=0:
#            plt.savefig(filename)
#        elif counter%2==0:
#            plt.savefig(filename2)
#            double_picture_sheet(prs, filename, filename2, title_text=title_text)
#        plt.close()
#    
#        
        # Grafik erstellen
#   plot_axis(prs, x=x, labels = labels, xn = xn, y = y, \
#                  pos_left = p_left, pos_right=p_right, pos_top=p_top, \
#                  pos_bottom = p_bottom, article = article, grid=False,\
#                  title_text=title_text, axis=1, mean_line=mean_line, legend=legend)
#    
#   
#   # STEP 2 - Berechnung durchschnittlicher Lesewert plus Häufigkeit je Ressort
#   # Grafik erstellen
#   plot_axis(prs, x=x, labels = labels, xn = xn, y = y, yn=yn, \
#             pos_left = p_left, pos_right=p_right, pos_top=p_top, \
#             pos_bottom = p_bottom, article = article, grid=False,\
#             title_text=title_text, axis=2, mean_line = 0, legend=legend)
#    

        
#%% Multiple Bars, reine Zeichenfunktion
def multiple_bars(prs, x=0, labels=0, xn=0, y_w=0, y_m=0,  grid = False, legend="normal", 
                  title_text = "Bitte ÜS eingeben"):
    
    
    
    # PLOTTING
    
    #Schriftfarbe und Farbe der Ticks festlegen
    set_font_color ="#8c8f91" 
       
    #Style setzen, plots anlegen
    sns.set_style("white")
    fig, ax1 = plt.subplots(figsize=(20,8))
    bar1 = ax1.bar(xn-0.40, y_w, label = "Frauen", width=0.4, align="center", color="#f77801")
    bar2 = ax1.bar(xn, y_m, width=0.4, label = "Männer", align="center", color="#f7bb83")
    
    
    
    # setzt Grid wenn gewünscht
    if grid==True:
        ax1.grid(color= set_font_color, linestyle='-', linewidth=1, axis="y")
    
    
    
    # setzt die Y-Achse 
    plt.yticks(fontproperties=campton_light)
    y_label_text = 'Ø Lesewert in Prozent'
    ax1.set_ylabel(y_label_text, color= set_font_color, fontproperties=campton_light, fontsize=50)
    ax1.yaxis.label.set_size(22)
    ax1.tick_params(axis='y', labelsize=25, colors= set_font_color)
    
    
    
    # setzt die X-Achse
    
    ax1.xaxis.set(ticks=range(0, len(xn))) #  Anzahl der Ticks 
    ax1.set_xticklabels(labels = labels, rotation=45, ha="right",  weight=800,\
                        color= set_font_color, fontproperties=campton_light, \
                        fontsize=30) # Labels werden ausgerichtet
    
     # Achsen, Ticks und alles andere festlegen
    
    if grid==True:
        ax1.grid(color= set_font_color, linestyle='-', linewidth=1, axis="y")
    
    
    #Legende einbauen
    legend_height = 1.24
    if legend=="normal":
            legend_height = 1.24
    elif legend == "strange": 
            legend_height =1.44
    else:
        legend_height = 1.24
        
    leg = plt.legend(bbox_to_anchor=(1, legend_height), handles=[bar1, bar2], markerscale=140)
    for text in leg.get_texts(): 
        plt.setp(text, color= set_font_color, size=25)
    
    # Text über den Balken
    for p in ax1.patches:
        # Problem: Ist ein Balken nur vier groß, ist der Abstand zu gering
        # TODO... NOCH EIN BISSCHEN EXPERIMENTIEREN
        height = p.get_height()
         # Werte über die Balken setzen
    # height_text gibt an, wie vile höher das ganze Gesetzt werden soll
        height_text = height +1
#        if special == "Lesetage": 
#            height_text = height + 10
        if height <8:
            height_text = height+0.3
        txt = '{:1.1f}'.format(height).replace(".", ",")
        ax1.text(p.get_x()+p.get_width()/2., height_text, txt, ha="center",\
                     fontproperties=campton_light, color= set_font_color, rotation=0\
                     ,fontsize= 20, weight = 1000)
        
    
    
    
    # DESIGN
    # Abstände Bars zur Achse (standardmäßig bei 0.5)
    plt.margins(x=0.03) # ziehen Bars näher an die Achse
     #obere Linie ausblenden
    ax1.spines["top"].set_visible(False)
    #ax1.spines["left"].set_color("gray")
    ax1.spines["top"].set_visible(False)
    ax1.spines["bottom"].set_visible(False)
    ax1.spines["left"].set_visible(False)
    ax1.spines["right"].set_visible(False)
    
    pos_left = 0.08
    pos_right=0.92
    pos_top=0.85
    pos_bottom = 0.4
    
    if legend=="normal": 
        pos_left = 0.08
        pos_right=0.92
        pos_top=0.85
        pos_bottom = 0.4
        
    if legend=="large":
        pos_left = 0.08
        pos_right=0.92
        pos_top=0.85
        pos_bottom = 0.45
    
    
    plt.subplots_adjust(left=pos_left, right=pos_right, top=pos_top, 
                        bottom=pos_bottom)
    
    filename = "geschlecht_vergleich_bars.jpg" 
   
    
     
                  
    plt.savefig(filename)
                    
    picture_sheet(prs, filename, title_text=title_text)     
   
    
    plt.close()  
    return prs
           
    
        



#%% Check Listen Umfrage
# Diese Funktion gleicht die Listen mit den tatsächlichen Fragen und Antworten 
    # des jeweiligen Datensatzes ab. 
    
#def list_checker(df):
#    test = df[df["Fragetext"].isin()]
#%% ÄLTERE CODESCHNIPSEL


#%% Übersicht 2 - Mantelressorts             

 
#übersicht_ressort erstellt die Deckblätter für die einzelnen Ressort, 
#ruft dann Top 10 und fünf Screenshots auf

#Problem: Diese Funktion muss theoretisch für jede Ausgabe neu geschrieben
#werden..

#n df_raw sind noch alle Split-IDs enthalten. Wird benötigt, damit die 
# Titelseiten einzeln ausgewertet werden können. 
# ODER: df_raw = nichtkumulierte Werte




def analyse_mantel(prs, df, df_raw="noch einzusetzen", ressort_liste = [], ausgabe_liste= []): 
    
    # check ob Ressortliste händisch angelegt wurde, ansonsten 
    
   
    # Liste mit allen Ressorts des Dataframes wird angelegt
    ressort_liste = mantel_ressorts
    
    
    # Liste mit den drei Ausgaben
    ausgabe_liste = ausgaben_liste
    
    # Liste mit Ressorts wird durchlaufen
    # Anforderungen der Coaches werden individuell in der for-Schleife festgelegt
    
    for ressort in ressort_liste: 
        
        # alle Grundwerte im jedem Schleifendurchlauf auf Null setzen
        platzierung = ["AA", "SK"]
        darstellungsform = ["NA"]
        seitentitel = df[df["Ressortbeschreibung"]==ressort]
        seitentitel = seitentitel.Seitentitel.unique()
        ressort_gefunden = True
        
        #Festlegung, welche Werte analysiert werden sollen
        # Version Schwäbische Zeitung
#        if ressort == "Titel":
#            platzierung = ["AA"]
#            darstellungsform = ["NA"]
#            seitentitel = ["Unterm Strich", "Leitartikel"]
#        elif ressort == "Wir im Süden":
#            platzierung = ["AA", "SK"]
#            darstellungsform = ["NA"]
#        elif ressort == "Seite Drei":
#            platzierung = ["AA", "SK"]
#        elif ressort == "Nachrichten & Hintergrund":
#            platzierung = ["AA", "SK"]
#            darstellungsform = ["NA"]
#        elif ressort == "Meinung & Dialog": 
#            platzierung = ["AA"]
#            darstellungsform = ["NA"]
#        elif ressort == "Wirtschaft":
#            platzierung = ["AA", "SK"]
#            darstellungsform = ["NA"]
#        elif ressort == "Journal":
#            platzierung = ["AA", "SK"]
#            darstellungsform = ["NA"]
#        elif ressort == "Kultur":
#            platzierung = ["AA", "SK"]
#            darstellungsform = ["NA"]
#        elif ressort == "Sport":
#            platzierung = ["AA", "SK"]
#            darstellungsform = ["NA"]
#        elif ressort == "Ratgeber":
#            platzierung = ["AA", "SK"]
#            darstellungsform = ["NA"]
#        elif ressort == "Wochenende":
#            seitentitel = ['Wochenende', 'Menschen', 'Lebensart', 'Unterhaltung',
#                           'Szene am Wochenende', 'Meine Seite']
#        
#            
#        else:
#            print ("Ressort " + ressort + " nicht für Auswertung angefordert")
#            ressort_gefunden = False
    
    
       # weitere Spezialanforderungen an einzelne Ressorts
       # zuerst: check ob ein Ressort aus dem Katalog vorhanden ist
        
        if ressort_gefunden: 
            
            # Datensatz mit Daten nur aus dem betreffenden Ressort anlegen
            df_ = df[df["Ressortbeschreibung"]== ressort]
            #df_r = df_raw[df_raw["Ressortbeschreibung"]==ressort]
            
            # Version Schwäbisch Zeitung
            #Spezialfälle
            
#            if ressort == "Titel":
#                
#                for title in ausgabe_liste:
#                    
#                    #Datensatz für jeweiligen Zeitungstitel
#                    df_changed = df_r[df_r["ZTG"]==title]
#                    deckblatt_macher(prs, df_changed, ressort, platzierung, 
#                                     darstellungsform, seitentitel, ZTG=title)
#                    
#                    top_10(prs, df_changed, df_berechnet = False)
#                    
#            elif ressort == "Wochenende":
#                #ausblenden eines Teilstücks eines bereits gewerteten 
#                #Pro und Contra-Artikels
#                df_we = df_[df_["SplitId"]!="23496"]
#                # für func grafik_lesewert wird Liste mit allen Seitentiteln
#                    # erstellt
#                liste_we = df_we["Seitentitel"].unique()
#                
#                deckblatt_macher(prs, df_we, ressort, platzierung, darstellungsform, 
#                     seitentitel, ZTG="null")
#                top_10(prs, df_we, df_berechnet = False, zeitung=False)
#                print("Wochenende Graifk Lesewert wird aufgerufen")
#                print (liste_we)
#                grafik_lesewert(prs, df_we, target="seitentitel", minimize=0, ressort_liste = liste_we, 
#                                order="ok", legend="xlarge")
#            elif ressort == "Sport": 
#                deckblatt_macher(prs, df_, ressort, platzierung, darstellungsform, 
#                     seitentitel, ZTG="null")
#                top_10(prs, df_, df_berechnet = False, zeitung=False)
#                schlagworte_finden(prs, df_, "sportart", sort = "Artikelanzahl")
#            elif ressort== "Ratgeber":
#                #Artiekl ist mit zwei versch. IDs doppelt vorhanden
#                df_ra=df_[df["SplitId"]!="24505"]
#                deckblatt_macher(prs, df_ra, ressort, platzierung, darstellungsform, 
#                     seitentitel, ZTG="null")
#                top_10(prs, df_ra, df_berechnet = False, zeitung=False)
#            elif ressort == "Seite Drei":
#                # Artikel ist doppelt vorhanden
#                df_sd=df_[df["SplitId"]!="25837"]
#                deckblatt_macher(prs, df_sd, ressort, platzierung, darstellungsform, 
#                     seitentitel, ZTG="null")
#                top_10(prs, df_sd, df_berechnet = False, zeitung=False)
#            
#            
#            else:
#                # Achtung: Immer df_ übergeben, damit immer nur die Daten des
#                # aktuellen Ressorts verarbeitet werden
#                deckblatt_macher(prs, df_, ressort, platzierung, darstellungsform, 
#                     seitentitel, ZTG="null")
#                top_10(prs, df_, df_berechnet = False, zeitung=False)
            
            deckblatt_macher(prs, df_, ressort, platzierung, darstellungsform, 
                    seitentitel, ZTG="null")
            top_10(prs, df_, df_berechnet = False, zeitung=False)

            



    return prs
#%% Werte berechnen
# Funktion wirft für ein Dataframe die Lese-, Blick- Und Durchlesewerte aus. Vor allem interessant bei händischen Arbeiten für kleine Analysen. 

def werte_rechnen(df, text="", et=0):
    lesewert = "Artikel-Lesewert (Erscheinung) in %"
    
    blickwert = "Artikel-Blickwert (Erscheinung) in %"
   
    durchlesewert = "Artikel-Durchlesewerte (Erscheinung) in %"
    
    
    lw = df[lesewert].mean()
    bw = df[blickwert].mean()
    dw = df[durchlesewert].mean()
    shape = df.shape[0]
    if et==0: 
        erscht = df.Erscheinungsdatum.nunique()
    if et!=0: 
        erscht = et
    durchschnitt = shape/erscht
    zusatztext = text
    
    print("LESEWERT "+zusatztext+ "durchschnittlich: {0:.2f} Prozent".format(lw))
    print("BLICKWERT "+zusatztext+ "durchschnittlich: {0:.2f} Prozent".format(bw))
    print("DURCHLESEWERT "+zusatztext+ "durchschnittlich: {0:.2f} Prozent".format(dw))
    print("Größe des Datensatzes "+zusatztext+": {0:.2f}".format(shape))
    print("Durchschnitt pro Erscheinungstag "+zusatztext+": {0:.2f}".format(durchschnitt))
    
    
#%% Werte berechnen
# Funktion wirft für ein Dataframe die Lese-, Blick- Und Durchlesewerte aus. Vor allem interessant bei händischen Arbeiten für kleine Analysen. 

def werte_berechnen(df, text="", et=0):
    lesewert = "LW"
    
    blickwert = "BW"
   
    durchlesewert = "DW"
    
    
    lw = df[lesewert].mean()
    bw = df[blickwert].mean()
    dw = df[durchlesewert].mean()
    shape = df.shape[0]
    if et==0: 
        erscht = df.Datum.nunique()
    if et!=0: 
        erscht = et
    durchschnitt = shape/erscht
    zusatztext = text
    
    print("LESEWERT "+zusatztext+ "durchschnittlich: {0:.2f} Prozent".format(lw))
    print("BLICKWERT "+zusatztext+ "durchschnittlich: {0:.2f} Prozent".format(bw))
    print("DURCHLESEWERT "+zusatztext+ "durchschnittlich: {0:.2f} Prozent".format(dw))
    print("Größe des Datensatzes "+zusatztext+": {0:.2f}".format(shape))
    print("Durchschnitt pro Erscheinungstag "+zusatztext+": {0:.2f}".format(durchschnitt))
    

#%% Funktion zur Analyse der Mantel-Daten  
####    ### MANTEL ASUWERT
def mantel_auswertung(prs, df, df_geschlecht, kunde="", lokales=False, ztg = ""):
        print("l6188 - mantel_auswertung gestartet")
        titletext=df.iloc[0]["Ressortbeschreibung"]
        titletext_lang=df.iloc[0]["Ressortbeschreibung"]
        titletext_kurz=df.iloc[0]["Ressortbeschreibung"]
        if len(ztg)>1:
            titletext_lang=df.iloc[0]["Ressortbeschreibung"] + " " + ausgaben_dict[ztg]
            titletext_kurz = df.iloc[0]["Ressortbeschreibung"] + " " + ztg
        deckblatt_abschluss(prs, df, df_geschlecht, title_text=titletext_lang, lokales=lokales, kunde=kunde, geschlecht=True)
        print(titletext_kurz)
#        
#     # Ressortübersicht
        tabelle_ressortauswertung(prs, df, research_object="Darstellungsformen", sort="Lesewert", 
                                  title_text = "Darstellungsformen " + titletext, minimum = 2, kunde=kunde, lokales=lokales)
        darstellungsformen(prs, df, minimum = 2, geschlecht=False, title_text="Darstellungsformen " + titletext_kurz, grid=False)
        darstellungsformen(prs, df_geschlecht, minimum = 2, geschlecht=True, title_text="Darstellungsformen " + titletext_kurz)
        tabelle_ressortauswertung(prs, df, research_object="Platzierungen", sort="Lesewert", 
                                  title_text = "Platzierungen " + titletext_kurz, minimum = 5, kunde=kunde, lokales=lokales)
        platzierungen(prs, df, minimum = 5, geschlecht=False, title_text="Platzierungen " + titletext_kurz)
        platzierungen(prs, df_geschlecht, minimum = 5, geschlecht=True, title_text="Platzierungen " + titletext_kurz)
       
         #14 - Entwicklung LW, BW, DW  CHECK
        
        for el in werte:
            grafik_entwicklung(prs, df, target=el, mean_line=0, legend="large", 
                               grid=True, title_text = el + " " + titletext_lang, 
                               steps_x_label=1, limit_y = True)
        # Entwicklung AA und SK - also nur die großen Stücke
#        df_großetexte = df[(df["Platzierungen"]=="AA") | (df["Platzierungen"]=="SK")]
#        for el in werte: 
#            grafik_entwicklung(prs, df_großetexte, target=el, mean_line=0, legend="large", grid=True, title_text = el + " AA und SK", steps_x_label=2)
        if kunde == "NOZ": 
            altersangabe1 = "50 bis 59 Jahre"
            altersangabe2 = "60 bis 69 Jahre"
        elif kunde == "DRP": 
            altersangabe1 = "40 bis 59 Jahre"
            altersangabe2 = "60 bis 79 Jahre"
        # nehme hier bei der Rheinpfalz ein Lesewert-Erklärstück raus            
        if (kunde == "DRP") & (titletext == "Hintergrund"): 
            df = df[(df["ArtikelId"]!=57947) & (df["ArtikelId"]!=86426)&\
                    (df["ArtikelId"]!=72408) & (df["ArtikelId"]!=44006)]
            df_geschlecht = df_geschlecht[(df_geschlecht["ArtikelId"]!=57947) & (df_geschlecht["ArtikelId"]!=86426)&\
                                          (df_geschlecht["ArtikelId"]!=72408) & (df_geschlecht["ArtikelId"]!=44006)]
        top_10(prs, df, df_berechnet = False, screenshots=True, number_screenshots = 5, 
           mode="ressort", headline="Top 10 - "+titletext_lang, zeitung=False)
        top_10(prs, df_geschlecht, df_berechnet = False, screenshots=False, number_screenshots = 5, 
           mode="ressort", headline="Top 10 - "+titletext, zeitung=False, geschlecht="weiblich")
        top_10(prs, df_geschlecht, df_berechnet = False, screenshots=False, number_screenshots = 5, 
           mode="ressort", headline="Top 10 - "+titletext, zeitung=False, geschlecht="männlich")
        top_10(prs, df_geschlecht, df_berechnet = False, screenshots=False, number_screenshots = 5, 
           mode="ressort", headline="Top 10 - "+titletext, zeitung=False, alter=altersangabe1)        
        top_10(prs, df_geschlecht, df_berechnet = False, screenshots=False, number_screenshots = 5, 
           mode="ressort", headline="Top 10 - "+titletext, zeitung=False, alter=altersangabe2)
#     
#%%
        
#%% Funktion lokale Analyse

def lokale_auswertung(prs, df, df_geschlecht, liste_kolumnen=[], kunde="", ztg = ""):
    # Titel der jeweiligen Grafiken festlegen
    titletext=df.iloc[0]["Ressortbeschreibung"]
    titletext_lang=df.iloc[0]["Ressortbeschreibung"]
    titletext_kurz=df.iloc[0]["Ressortbeschreibung"]
    if len(ztg)>1:
        titletext_lang=df.iloc[0]["Ressortbeschreibung"] + " " + ausgaben_dict[ztg]
        titletext_kurz = df.iloc[0]["Ressortbeschreibung"] + " " + ztg
    
    # Deckblatt und Tabellen
    deckblatt_abschluss(prs, df, kunde=kunde, title_text = titletext_lang)
    tabelle_ressortauswertung(prs, df, research_object="Darstellungsformen", sort="Lesewert", 
                                  title_text = "Darstellungsformen " + ausgaben_dict[ztg], minimum = 5, kunde="NOZ", lokales=True)
    tabelle_ressortauswertung(prs, df, research_object="Platzierungen", sort="Lesewert", 
                                  title_text = "Platzierungen " + ausgaben_dict[ztg], minimum = 5, kunde="NOZ", lokales=True)
        
    platzierungen(prs, df, minimum = 5, geschlecht=False, title_text="Lesewert Platzierungen " + ztg, special="")
    platzierungen(prs, df_geschlecht, minimum = 5, geschlecht=True, title_text="Platzierungen nach Geschlecht", special="")
    darstellungsformen(prs, df, minimum = 5, geschlecht=False, title_text="Lesewert Darstellungsformen "+ztg, special="")
    darstellungsformen(prs, df_geschlecht, minimum = 5, geschlecht=True, title_text="Darstellungsformen nach Geschlecht", special="")
    # Entwicklungsgrafiken LW, BW udn DL für alle Texte + Aufmacher + Seitenkeller
    for el in werte:
        grafik_entwicklung(prs, df, target=el, mean_line=0, legend="large",
                               grid=True, title_text = el + " " + ausgaben_dict[ztg],
                               steps=1,  limit_y=True, ma = False)
        df_ = df[(df["Platzierungen"]=="AA")|(df["Platzierungen"]=="SK")]
        grafik_entwicklung(prs, df_, target=el, mean_line=0, legend="large",
                               grid=True, title_text = el + " Aufmacher und Seitenkeller", steps=2)
        
## SONDERWÜNSCHE LOKALES
    
        #Auswertung Seitentitel in Column Special
    grafik_lesewert(prs, df, target="special", minimize=5, label_position="large",
                    ressort_liste=[], special=False, order="ok", legend="normal", sort="Seitennummer", 
                    article="total", mean_line = 0, grid=True, anzahl_lokales=1, axis=1, \
                    title_text="Lesewert nach Seitentiteln " + ausgaben_dict[ztg])
        
       # Auswertung nach Textlänge
    grafik_lesewert(prs, df, target="Textlänge", minimize=1, label_position="large",
                    ressort_liste=textlänge, special=False, order="ok", legend="normal", sort="Kategorie", 
                    article="total", mean_line = 0, grid=False, anzahl_lokales=1, axis=2, 
                    title_text="Textlängen " + ausgaben_dict[ztg])
#        if elem == "HAL": 
#            res_li = ["]
    grafik_lesewert(prs, df, target="seitentitel", minimize=5, label_position="large",
                    ressort_liste=[], special=False, order="ok", legend="normal", sort="Seitennummer", 
                    article="total", mean_line = 0, grid=False, anzahl_lokales=1, axis=2, 
                    title_text="Seitentitel " + ztg)
     
    grafik_lesewert(prs, df, target="special", minimize=5, label_position="large",
                    ressort_liste=[], special=False, order="ok", legend="normal", sort="Seitennummer", 
                    article="total", mean_line = 0, grid=False, anzahl_lokales=1, axis=2, 
                    title_text="Seitentitel " + ztg)
        
    lesewert_erscheinung(df,df, prs, title_text="LW nach Wochentagen")
    lesewert_ankündigungen(prs, df, title_text="Ankündigungen - " + ausgaben_dict[ztg])
    top_10(prs, df, df_berechnet = False, screenshots=True, number_screenshots = 5, 
           mode="ressort", headline="Top 10 - Lokales "+ausgaben_dict[ztg], zeitung=True, geschlecht="", alter="", 
           lokales_ressort="Lokales")
        # Top 10 Geschlechter
    top_10(prs, df_geschlecht, df_berechnet = False, screenshots=False, number_screenshots = 5, 
           mode="ressort", headline="default", zeitung=True, geschlecht="weiblich")
    top_10(prs, df_geschlecht, df_berechnet = False, screenshots=False, number_screenshots = 5, 
           mode="ressort", headline="default", zeitung=True, geschlecht="männlich")
        
 #        Top 10 Altersgruppen
    top_10(prs, df_geschlecht, df_berechnet = False, screenshots=False, number_screenshots = 5, 
           mode="ressort", headline="default", zeitung=True, alter="40 bis 59 Jahre")
    top_10(prs, df_geschlecht, df_berechnet = False, screenshots=False, number_screenshots = 5, 
           mode="ressort", headline="default", zeitung=True, alter="60 bis 79 Jahre")
   
# 
#          Kommentare und Kolumnen
    liste_kommis = liste_kolumnen
    df_liste = df[df["Beschreibung"].isin(liste_kommis)]
         
    df_kommentare = df[df["Darstellungsformen"].str.contains("KM")]
         
    besch_lok = "Kommentare gesamt"
    split_lok = df_kommentare.shape[0]
    seite_lok = df_kommentare.Seitennummer.median()
    lw_lokal = df_kommentare["Artikel-Lesewert (Erscheinung) in %"].mean()
    col = ["Beschreibung", "SplitId", "Seitennummer", "Artikel-Lesewert (Erscheinung) in %"]
    extension_df = pd.DataFrame({"Beschreibung":[besch_lok], 
                                 "SplitId":[split_lok], 
                                 "Seitennummer":[seite_lok], 
                                 "Artikel-Lesewert (Erscheinung) in %":[lw_lokal]})
            
#          
    print(extension_df)
    grafik_lesewert(prs, df_liste, target="rubriken", minimize=5, label_position="xlarge",
                    ressort_liste=liste_kommis, special=False, title_text="Kolumnen/Rubriken/Meinung " + ztg, order="ok",
                    legend="normal", sort="Lesewert", article="mean", extension_df=extension_df, grid=False)
    
    
        
        

#%% Übersicht 1 - Func Analyse Gesamt   
    
'''
df_scan ist das Dataframe mit den Messdaten aus der ScanAuswertung. 

df_doublesplitid ist ein optionaler Datensatz, bei dem die multiplen Split_ids
nicht entfernt wurden. Manchmal notwendig. 

''' 

    
def analyse_gesamt(prs, df, df_scan=False, df_doublesplitid=False, 
                   df_nichtkum=False):
    
    
    # Demografie - kümmert sich DD drum
    
    # Lesewert in Zahlen - händisch anlegen
    
    # Chart mit Gesamtwert - händisch anlegen
    
    # Chart mit den drei Lokalteilen - händisch anlegen
    
    # Chart mit Mantelteil - händisch anlegen
    # bzw. für einzelne Titel anlegen
    marken_analyse(df, df_doublesplitid=df_doublesplitid, df_nichtkum=df_nichtkum)
    # Übersicht Sonderseiten
   
    # Anzeige Lesezeit
    # TODO ScanAuswertung table runterladen
    #lesezeit(prs, df_scans)
    
       
     
    
    #df_sonderseite = df[df["Seitentitel"].isin(sonderseiten_liste)]
    
    grafik_lesewert(prs, df, target="ressort", minimize=5, label_position="large", 
                    sort = 'Seitennummer', ressort_liste=mantel_ressorts)
    
    # Chart mit Artikelanzahl und LW nach Wochentage
    ##lesewert_erscheinung(df, prs, title_text="Lesewert nach Wochentagen")
    
    # Chart verwendete Darstellungsform und Anzahl Artikel 
    darstellungsformen(prs, df, minimum = 5)
    
    # Chart mit ressort und lesewert
    ##grafik_lesewert(prs, df, target="ressort", minimize=5, special=True, ressort_liste=ressort_list)
    
    # Chart mit Lesezeit
    ##lesezeit(prs, df_scan)
    
    # Vergleich Infokästen
    tabelle_infokästen(prs, df)
    
    
 #%% Gruppen berechnen
# diese helper-Funktion wirft die Lesewerte zum Beispiel für alle SEitentitel eines Dataframes heraus. 

def gruppe_rechnen(df, param):
    df_group = df.groupby(param).agg({"Artikel-Lesewert (Erscheinung) in %":"mean", "SplitId":"count"})
    print(df_group)


#%%
    
    
#%% Lesewert nach Ort
def lw_nach_ort(df, orte_region=orte_region):
    df["Handlungsorte"] = df["Handlungsorte"].fillna("x")
    df["Ort"] = ""
    for i in df.index: 
        sentence = df.get_value(i, "Handlungsorte")
        if any(ext in sentence for ext in orte_region):
            df.set_value(i, "Ort", "Region")
        elif any(ext in sentence for ext in orte_deutschland):
            df.set_value(i, "Ort", "Deutschland")
        elif sentence == "x": 
            df.set_value(i, "Ort", "keine Angaben")
        else: 
            df.set_value(i, "Ort", "International")
        
#    lw_regional = df[df["Ort"]=="Region"]["Artikel-Lesewert (Erscheinung) in %"].mean()
#    lw_deutschland = df[df["Ort"]=="Deutschland"]["Artikel-Lesewert (Erscheinung) in %"].mean()
#    lw_international = df[df["Ort"]=="International"]["Artikel-Lesewert (Erscheinung) in %"].mean()
#    
#    print("LW regional: {0:.2f}".format(lw_regional))
#    print("LW Deuschland: {0:.2f}".format(lw_deutschland))
#    print("LW International: {0:.2f}".format(lw_international))
    
    
    
    
#%% EINZELNE ZEITUNGEN - INITIAL-FUNKTIONEN
# TODO:  Funktion schreiben    
    
def neue_westfaelische(df, df_nk = 0, df_double_split = 0, df_scans=0, df_umfrage=0,
                       df_geschlecht = 0, liste = ausgaben_liste, df_ereignisse=0): 
    
    # Timer aktivieren
    t_0 = time.time()
    
    #prs-Objekt erstellen
    prs = Presentation("LW_neu_2.pptx")
    #*******************************************************************
    #ALLGEMEINES
    
#    # 1 - Entwicklungsfolie Lesewert/DW/BW im Messverlauf - CHECK
#    print("Starte Auswertung Entwicklung über Messzeitraum...")
#    c
#    for elem in werte:
#        grafik_entwicklung(prs, df, target=elem, mean_line=0, legend="large", grid=True, title_text = False, ma=False)
    #kultur = df[df["Ressortbeschreibung"]=="Kultur / Medien"]
    #themen_liste = ["Society/Promis/Klatsch & Tratsch","Kunst/Kunsthandwerke", "Musik", "Theater/Oper", "Kino/Film", "Fernsehen/TV"]
#    grafik_lesewert(prs, kultur,minimize=5, label_position="large",
#                    ressort_liste=themen_liste, special=False, title_text="Kulturthemen nach Lesewert", order="ok", legend="normal", sort="Lesewert", 
#                    article="total", mean_line = 0, grid=False, anzahl_lokales=1, mean_line_title = "", axis=2)
#    
#    
#    # 2 - Fragebogen - CHECK
#    print ("Initialisiere Auswertung der Fragebögen")
#    # Aufruf jeder einzelnen Ausgabe
#    for elem in ausgaben_liste:
#        print(ausgaben_liste)
#        
#        df_ = df_umfrage[df_umfrage["ZTG"]==elem]
#       
#        #Aufruf gestapelte Bars "Ihre Zeitung ist..."
#       # zeitungsattribute_berechnung(prs, df_, title = elem) # elem wird später über ausgaben_dict umgewandelt
#        zeitungsattribute_berechnung(prs, df_, title = elem) # elem wird später über ausgaben_dict umgewandelt
#        #zeitungsattribute_berechnung(prs, df_, title = elem, frage_dict = zeitung_attribute_dict2)
#        #zeitungsattribute_berechnung(prs, df_, title = elem, frage_dict = zeitung_attribute_dict3)
#        #Aufruf Bars "Welche Themen
#        zeitung_themen(prs, df_, title=elem)
#        
#        # Aufruf kleine Pie-Charts
#        umfrage_pie(prs, df_, title=elem)
#        
#        # Aufruf kleine Barcharts, TV- und Internetnutzung
#        
#        mini_bars(prs, df_, title_text = "Fernsehnutzung", liste=umfrage_mini_bars_TV, title=elem)
#        mini_bars(prs, df_, title_text = "Internetnutzung", liste=umfrage_mini_bars_netz, title=elem)
#        mini_bars(prs, df_, title_text = "Persönliches", liste=umfrage_demografie_bars, title=elem)
#        
#    # 3 - Kennzahlen der Messung - CHECK
#    
#    kennzahlen(df, df_scans)
#    
#        
#        
#    # 4 - TEXTMARKEN - Analyse LW, BW, DW Gesamt, Ausgaben, Lokales, Mantel - CHECK
#    #print("-------------------------------------------")
#   #print("WERTE FÜR LW-MARKEN:")
#    marken_analyse(df, df_doublesplitid=df_double_split, df_nichtkum=df_nk)
##    
##
##     5 -Screenshots Gesamt + Mantel Top 5 - CHECK
#    print("Top 10 Gesamt werden ausgespielt...")
#    top_10(prs, df, df_berechnet = False, screenshots=True, number_screenshots = 5, 
#           mode="Gesamt", headline="Top 10 Gesamt", zeitung=False, lokales_ressort="Lokales")
#    # Top 10 Mantelressorts CHECK
#    print("Top 10 Mantel werden ausgespielt...")
#    df_mantel = df[df["Ressortbeschreibung"].isin(mantel_ressorts)]
#    top_10(prs, df_mantel, df_berechnet = False, screenshots=True, number_screenshots = 5, 
#           mode="Mantel", headline="Top 10 Mantelressorts", zeitung=False)
    
 # 6 - Leseraktivität nach Tageszeit und Wochentag (aus Scandaten) - CHECK
   # TODO Werte für alle vier Wellen ausspielen, dafür muss ich Scan-Daten mit Asugabedaten verknüpfen... 
   #  ... muss also wissen, woher der einzelne Scan herkommt. 
#    print("Analyse Lesetage und Lesezeiten gestartet...." )
   # lesetage(prs, df_scans, title_text="Lesetage - Gesamt")
   # lesezeit(prs, df_scans, title_text="Lesetage - Gesamt")
#    for elem in ausgaben_liste:
#        df_ = df_scans[df_scans["ZTG"]==elem]
#        lesetage(prs, df_, title_text="Lesetage - " + ausgaben_dict[elem]) 
#        lesezeit(prs, df_, title_text="Lesezeiten - " + ausgaben_dict[elem])
   
    
##    
###
##    # 8 - Ressort-Übersicht Lesewerte plus Artikelzahl - CHECK
###    
    mean_line = df["Artikel-Lesewert (Erscheinung) in %"].mean()   
    grafik_lesewert(prs, df, target="ressort", minimize=5, label_position="xlarge",
                    ressort_liste=ressort_list, special=False, title_text="Ressorts Gesamt", order="ok",
                    legend="normal", sort="Seitennummer", article="mean", anzahl_lokales=4, 
                    mean_line=mean_line)
  
    ressort_list_erweitert = ['Titelseite',
 'Politik und Meinung',
 'Seite 3',
 'Zwischen Weser und Rhein',
 'Wirtschaft',
 'Aus aller Welt',
 'Sport',
 'Lokalsport',
 
 "Lok 1 BI", "Lok 1 GT", "Lok 1 HF", "Lok 1 PB",
 
 'weitere Lokalseiten', 
 'Kultur / Medien',
 'Rätsel / Roman',
 'Kinder / TV-Programm',
 'Das Magazin',
 'Das Magazin - wochentags', 
'Magazin' ]
    grafik_lesewert(prs, df, target="special_res", minimize=5, label_position="xlarge",
                    ressort_liste=ressort_list_erweitert, special=False, title_text="Ressorts Gesamt", order="ok",
                    legend="normal", sort="Kategorie", article="mean", anzahl_lokales=4, 
                    mean_line=mean_line)    
##    
#    # für die einzelnen Wellen: - CEHCK
#    for elem in ausgaben_liste: 
#        df_ = df_nk[df_nk["ZTG"]==elem]
#        mean_line = df_["Artikel-Lesewert (Erscheinung) in %"].mean()
#        grafik_lesewert(prs, df_, target="ressort", minimize=5, label_position="xlarge",
#                    ressort_liste=ressort_list, special=False, title_text="Ressorts - " + ausgaben_dict[elem] , order="ok",
#                    legend="normal", sort="Seitennummer", article="mean", mean_line=mean_line) 
#    
##    
##    # 9 - Ressorts nach Geschlecht und Tutek - CHECK
##    #gesamt
#    print("Ressorts nach Geschlechtern untersuchen...")
    #multiple_bars_geschlecht(prs, df_geschlecht, target="Ressorts", grid = True, legend="large", 
#                             ressort_liste = ressort_list, title_text= "Ressorts nach Geschlecht - Gesamt")
##    # für einzelne Wellen
#    for elem in ausgaben_liste: 
#        df_geschlecht_ = df_geschlecht[df_geschlecht["ZTG"]==elem]
#        print(elem)
#        print(df_geschlecht_.shape[0])
#        multiple_bars_geschlecht(prs, df_geschlecht_, target="Ressorts", 
#                                 grid = True, legend="large", 
#                                 title_text = "Ressorts nach Geschlecht - " + elem, 
#                                 ressort_liste = ressort_list)
#    #prs-Objekt abspeichern
##    
##    # 10 - Lesewert nach TAgen im Verhältnis zum Angebot - CHECK
    #lesewert_erscheinung(df, df, prs, title_text="LW nach Wochentagen - Gesamt", grid=False)
#    for elem in ausgaben_liste: 
#        df_ = df_nk[df_nk["ZTG"]==elem]
#        df_double_split_  = df_double_split[df_double_split["ZTG"]==elem]
#        lesewert_erscheinung(df_,df_double_split_, prs, title_text="LW nach Wochentagen - " + elem)
   
    # 11 - Lesewert nach darstellungsform im Verhältnis zur Artikelzahl + DF und Geschlecht- CHECK
    # TODO - nichtkumuliert mit Geschlechterdatei verbinden
#    df_bi = df[(df["Ressortbeschreibung"]=="Lokales") & (df["ZTG"]=="BI")]
#    df_bi_gesch = df_geschlecht[(df_geschlecht["Ressortbeschreibung"]=="Lokales") & (df_geschlecht["ZTG"]=="BI")]
  #  darstellungsformen(prs, df, minimum = 4, geschlecht=False, title_text="Darstellungsformen Gesamt")
    #darstellungsformen(prs, df_geschlecht, minimum = 4, geschlecht=True, title_text="Darstellungsformen nach Geschlecht", grid=True)
#    
#    for elem in ausgaben_liste: 
#        df_ = df_nk[df_nk["ZTG"]==elem]
#        df_geschlecht_ = df_geschlecht[df_geschlecht["ZTG"]==elem]
#        darstellungsformen(prs, df_geschlecht_, minimum = 5, geschlecht=True, title_text="Darstellungsformen " + ausgaben_dict[elem])
#    
    #12 - Lesewert Kolumnen
#    liste_kolumnen = ['Die schräge Meldung', 'Kommentar Seite 2', 'Meinungsbörse', 'Einwurf', 'Steinwurf', 'TV-Kritik','Lesezeichen']
#    grafik_lesewert(prs, df, target="rubriken", minimize=5, label_position="xlarge",
#                    ressort_liste=liste_kolumnen, special=False, title_text="Lesewerte Rubriken/Kolumnen ", order="ok", legend="normal", sort="Lesewert", 
#                    article="mean", mean_line=0, grid=True)
#    #ÜBERSICHT SEITEN MANTEL
#    
#    print("Beginne Übersicht Mantelseiten....")
##    #mantel_ressorts_plus_loksport = ["Titelseite", "Zwischen Weser und Rhein"]
#    mantel_ressorts_plus_loksport = ['Titelseite', 'Politik und Meinung', 'Seite 3',
#       'Zwischen Weser und Rhein', 'Wirtschaft', 'Aus aller Welt', 'Sport', 'Kinder / TV-Programm',
#       'Rätsel / Roman', 'Daily']
#    mantel_ressorts_plus_loksport = ["Kultur"]
#    for elem in mantel_ressorts_plus_loksport:
#        print("Analysiere " + elem +"...")
#        df_ = df[df["Ressortbeschreibung"]==elem]
#        df_geschlecht_ = df_geschlecht[df_geschlecht["Ressortbeschreibung"]==elem]
#        df_aufmacher = df_[df_["Darstellungsformen"]=="AA"] 
#        df_großetexte = df_[df_["Platzierungen"].isin(["AA", "SK"])]
#        #print (df_großetexte.shape)
#     # 12 - Deckblatt Ressort - CHECK    
#        
#        deckblatt_abschluss(prs, df_)
####        
###     # 13 - Ressortübersicht
#       
#        tabelle_ressortauswertung(prs, df_, research_object="Platzierungen", sort="Lesewert", 
#                                  title_text = "Platzierungen " + ressort_dict[elem], minimum = 5)
#        platzierungen(prs, df_, minimum = 5, geschlecht=False, title_text="Platzierungen " + ressort_dict[elem])
#        platzierungen(prs, df_geschlecht_, minimum = 5, geschlecht=True, title_text="Platzierungen " + ressort_dict[elem])
#        tabelle_ressortauswertung(prs, df_, research_object="Darstellungsformen", sort="Lesewert", 
#                                  title_text = "Darstellungsformen " + ressort_dict[elem], minimum = 5)
#        darstellungsformen(prs, df_, minimum = 5, geschlecht=False, title_text="Darstellungsformen " + ressort_dict[elem])
#        darstellungsformen(prs, df_geschlecht_, minimum = 5, geschlecht=True, title_text="Darstellungsformen " + ressort_dict[elem])
#       
#       
#         #15 Top 10 mit Screenshots gesamtes Ressort CHECK
#        print(elem)
#        top_10(prs,  df_, df_berechnet = False, screenshots=True, number_screenshots = 5, 
#           mode="ressort", headline="Top 10 - " + elem, zeitung=True)
#        
#        # Top 10 Geschlechter
#        top_10(prs, df_geschlecht_, df_berechnet = False, screenshots=False, number_screenshots = 5, 
#           mode="ressort", headline="default", zeitung=True, geschlecht="weiblich")
#        top_10(prs, df_geschlecht_, df_berechnet = False, screenshots=False, number_screenshots = 5, 
#           mode="ressort", headline="default", zeitung=True, geschlecht="männlich")
#        
#        # Top 10 Altersgruppen
#        top_10(prs, df_geschlecht_, df_berechnet = False, screenshots=False, number_screenshots = 5, 
#           mode="ressort", headline="default", zeitung=True, alter="40 bis 59 Jahre")
#        top_10(prs, df_geschlecht_, df_berechnet = False, screenshots=False, number_screenshots = 5, 
#           mode="ressort", headline="default", zeitung=True, alter="60 bis 79 Jahre")
#   
#          #14 - Entwicklung LW, BW, DW  CHECK
#        for el in werte:
#            grafik_entwicklung(prs, df_, target=el, mean_line=0, 
#                               legend="large", grid=True, 
#                               title_text = el + " " + elem, ma=False, steps=1)
#        # Entwicklung AA und SK - also nur die großen Stücke
#        for el in werte: 
#            grafik_entwicklung(prs, df_großetexte, target=el, mean_line=0, 
#                               legend="large", grid=True,
#                               title_text = wertedict[el] + " " + elem + " - AA und SK", 
#                               ma = False, steps=1)
#        
#        
#        # LW nach Handlungsorten
#        grafik_lesewert(prs, df_, target="ort", minimize=5, label_position="normal",
#                    ressort_liste=["Region", "National", "International"], special=False, title_text=elem + " - Artikel nach Orten", order="ok", legend="normal", sort="Seitennummer", 
#                    article="total", mean_line = 0, grid=True, anzahl_lokales=1, mean_line_title = "")
#        df__ = df_[df_["Platzierungen"]=="AA"]
#        grafik_lesewert(prs, df__, target="ort", minimize=1, label_position="normal",
#                    ressort_liste=["Region", "National", "International"], special=False, title_text=elem + " - Aufmacher nach Orten", order="ok", legend="normal", sort="Seitennummer", 
#                    article="total", mean_line = 0, grid=True, anzahl_lokales=1, mean_line_title = "")
#        print(rubrik_dict[elem])
#        if len(rubrik_dict[elem]) > 0:
#            rubrik_liste = rubrik_dict[elem]
#            print(rubrik_liste)
#            grafik_lesewert(prs, df_, target="rubriken", minimize=5, label_position="large",
#                    ressort_liste=rubrik_liste, special=False, title_text="Rubriken " + ressort_dict[elem], order="ok", legend="large", sort="Lesewert", 
#                    article="total", mean_line = 0, grid=False, anzahl_lokales=1, mean_line_title = "", axis=2)
#            auswahl_kategorien(prs, df_, minimum = 5, title_text="Rubriken " + elem, special="")
            
#            
##            
##        
#        if elem == "Politik und Meinung":
#            
#            df_kommi = df_[df_["Darstellungsformen"]=="KM"]
#            df_kommi2 = df_[(df_["Beschreibung"]=="Kommentare") | (df_["Beschreibung"]=="Kommentar")]
#            darstellungsformen(prs, df_kommi, geschlecht=False, 
#                               title_text= "Kommentare DF - vor/nach dem 23. Mai", 
#                               special="kommis_nw")
#            darstellungsformen(prs, df_kommi2, geschlecht=False, 
#                               title_text= "Kommentare Rubrik - vor/nach dem 23. Mai", 
#                               special="kommis_nw")
#        
#            
#        
#   
#     
        
        
        
       # SPEZIAL SPORT
        
        
         
#    df_sport = df[df["Ressortbeschreibung"]=="Sport"]
#    lesewert_erscheinung(df_sport, df_sport, prs, title_text="Lesewert nach Wochentagen Sport")
#    # Grafik Sportarten
#    schlagworte_finden(prs, df_sport, value="sportart", sort= "LW", min_artikel = 5, title="")
#    schlagworte_finden(prs, df_sport, value="vereine", sort= "LW", min_artikel = 1, title="")
#    grafik_lesewert(prs, df_sport, minimize=5, label_position="large",
#                        ressort_liste=[], special=False, title_text="Beliebteste Sportarten", order="ok", legend="normal", sort="Lesewert", 
#                        article="total", mean_line = 0, grid=False, anzahl_lokales=1, mean_line_title = "", axis=2) 
#        
#        
#       #   Sport Angebot und Lesewert nach Erscheinungstagen
#    df_sport = df[df["Ressortbeschreibung"]=="Sport"]
#    
#    # Grafik Sportarten
#    schlagworte_finden(prs, df_sport, value="sportart", sort= "LW", min_artikel = 5, title="")
#    schlagworte_finden(prs, df_sport, value="vereine", sort= "LW", min_artikel = 5, title="")
#   
#        
##    # ÜBERSICHT LOKALES
##    
    #Loksport
#    ausgaben_liste_ohne_pb  =   ["BI", "HF", "GT"]
#    for ausgabe in ausgaben_liste_ohne_pb:
#        df_loksport = df[df["Ressortbeschreibung"]== "Lokalsport"]
#        df_lokaus = df_loksport[df_loksport["ZTG"]==ausgabe]
#        
#        deckblatt_abschluss(prs, df_lokaus, title_text = "Lokalsport " + ausgabe, lokales=False)
#        tabelle_ressortauswertung(prs, df_lokaus, research_object="Darstellungsformen", sort="Lesewert", 
#                                  title_text = "Darstellungsformen " + ausgabe, minimum = 5)
#        darstellungsformen(prs, df_lokaus, minimum = 5, geschlecht=False, title_text="Darstellungsformen " +  ausgaben_dict[ausgabe])
#        tabelle_ressortauswertung(prs, df_lokaus, research_object="Platzierungen", sort="Lesewert", 
#                                  title_text = "Platzierungen " + ausgabe, minimum = 5)
#        platzierungen(prs, df_lokaus, minimum = 5, geschlecht=False, title_text="Platzierungen " +  ausgaben_dict[ausgabe])
#        lesewert_erscheinung(df_lokaus, df_lokaus, prs, title_text="Lesewert nach Wochentagen Lokalsport" + ausgabe)
#        schlagworte_finden(prs, df_lokaus, value="sportart", sort= "LW", min_artikel = 5, title="sportarten Lokalsport "+ausgabe)
#            

    # für Lokallkutlr
#    ausgaben_liste =["BI", "GT", "HF", "PB"]
#    
#    for ausgabe in ausgaben_liste:
#        df_ = df[(df["ZTG"]==ausgabe) & (df["Ressortbeschreibung"]=="Lokales")]
#        # für Lokale Kultur
#        df_ = df[(df["ZTG"]==ausgabe) & ((df["Ressortname"]=="Lokale Kultur") | (df["Seitentitel"]=="Lokale Kultur") | (df["Seitentitel"]=="Gütersloher Kultur"))]
#        
#        df_geschlecht_ = df_geschlecht[(df_geschlecht["ZTG"]==ausgabe) & (df_geschlecht["Ressortbeschreibung"]=="Lokales")]
##        #  für Lokale Kultur
##        df_geschlecht_ = df_geschlecht[( df_geschlecht["ZTG"]==ausgabe) & (( df_geschlecht["Ressortname"]=="Lokale Kultur") | ( df_geschlecht["Seitentitel"]=="Lokale Kultur") | ( df_geschlecht["Seitentitel"]=="Gütersloher Kultur"))]
#        df_großetexte = df_[df_["Platzierungen"].isin(["AA", "SK"])]
##        print(df_.shape)
##        print(ausgabe)
####    # 16 Deckblatt
#        deckblatt_abschluss(prs, df_, title_text = "Lokales " + ausgaben_dict[ausgabe], lokales=True)
####        
####    # 17 Entwicklung im Messverlauf
##        print(df_.shape)
#        df_1 = df_[(df_["Platzierungen"]=="AA") | (df_["Platzierungen"]=="SK")]
#        
#        for el in werte:
#            grafik_entwicklung(prs, df_, target=el, mean_line=0, legend="large",
#                               grid=True, title_text = el + " Lokales " + ausgaben_dict[ausgabe], ma=False, limit_y = False)
#            grafik_entwicklung(prs, df_großetexte, target=el, mean_line=0, legend="large",
#                               grid=True, title_text = el + " AA und SK " + ausgaben_dict[ausgabe], 
#                               ma=False, limit_y = False)
####
### 23 Darstellungsformen und Platzierungen
#        tabelle_ressortauswertung(prs, df_, research_object="Darstellungsformen", sort="Lesewert", 
#                                  title_text = "Darstellungsformen " + ausgaben_dict[ausgabe], minimum = 5, lokales=True)
#        darstellungsformen(prs, df_, minimum = 5, geschlecht=False, title_text="Darstellungsformen " +  ausgaben_dict[ausgabe])
#        #darstellungsformen(prs, df_geschlecht_, minimum = 5, geschlecht=True, title_text="Darstellungsformen " +  ausgaben_dict[ausgabe])
#        tabelle_ressortauswertung(prs, df_, research_object="Platzierungen", sort="Lesewert", 
#                                  title_text = "Platzierungen " +  ausgaben_dict[ausgabe], minimum = 5, lokales=True)
#        platzierungen(prs, df_, minimum = 5, geschlecht=False, title_text="Platzierungen " +  ausgaben_dict[ausgabe])
#        platzierungen(prs, df_geschlecht_, minimum = 5, geschlecht=True, title_text="Platzierungen " +  ausgaben_dict[ausgabe])
###        
##
##        
####    # 18 Lesewert nach Seitentiteln
#        df_st = df_[df_["Seitentitel"].isin(seitentitel_list)] 
#        print(df_st.Seitentitel.unique())
#        df_mag = df[df["Ressortbeschreibung"]=="Daily"]
#        grafik_lesewert(prs, df_mag, target="seitentitel", minimize=5, label_position="xlarge",
#                    special=False, title_text="LW nach Seitentiteln Magazin", order="ok", legend="normal", sort="Seitennummer", 
#                    article="total")
#        grafik_lesewert(prs, df_, target="seitentitel", minimize=5, label_position="xlarge",
#                    ressort_liste=seitentitel_list, special=False, title_text="LW mit df_ nach Seitentiteln " + ausgabe, order="ok", legend="normal", sort="Seitennummer", 
#                    article="total")
#        grafik_lesewert(prs, df_, target="special", minimize=5, label_position="xlarge",
#                    ressort_liste=[], special=False, title_text="LW Special nach Seitentiteln " + ausgabe, order="ok", legend="normal", sort="Seitennummer", 
#                    article="total")
###         
####    # 19 Lesewert nach Tagen und im Vergleich zum Angebot
       #lesewert_erscheinung(df_, df_, prs, title_text="Lesewert nach Wochentagen " + ausgabe)
####    
####    # 20 TODO Muss-kann-soll
####    
####    # 21 Meldung vs Ankündiger
       # lesewert_ankündigungen(prs, df_)
#####    
####    # 22 Top 10
#        top_10(prs, df_, df_berechnet = False, screenshots=True, number_screenshots = 5, 
#           mode="ressort", headline="Top 10 - " + ausgaben_dict[ausgabe], zeitung=True)
#        # nach Geschlechtern
#        top_10(prs, df_geschlecht_, df_berechnet = False, screenshots=False, number_screenshots = 5, 
#           mode="ressort", headline="Top 10 - " + ausgaben_dict[ausgabe] + " - Frauen", zeitung=True, geschlecht="weiblich")
#        top_10(prs, df_geschlecht_, df_berechnet = False, screenshots=False, number_screenshots = 5, 
#           mode="ressort", headline="Top 10 - " + ausgaben_dict[ausgabe] + " - Frauen", zeitung=True, geschlecht="männlich")
#        # nach Altersgruppen
#        top_10(prs, df_geschlecht_, df_berechnet = False, screenshots=False, number_screenshots = 5, 
#           mode="ressort", headline="default", zeitung=True, alter="40 bis 59 Jahre")
#        top_10(prs, df_geschlecht_, df_berechnet = False, screenshots=False, number_screenshots = 5, 
#           mode="ressort", headline="default", zeitung=True, alter="60 bis 79 Jahre") 
#         
###    
#    # 24 Übersicht Kolumnen + Rubriken 
#        if ausgabe == "BI": 
##           kolumnen_list = ['Guten Morgen BI', 'Das war die Woche', 'Kirchliche Nachrichten', 
##                            'Leserbriefe', 'Persönlich', 'Radio Bielefeld', 'Sperrungen', 
##                            'Zum Sonntag']
#            kolumnen_list = ['Flottmann', 'Kommentar', 'Meinung']
#        if ausgabe == "GT": 
#           kolumnen_list = ['Guten Morgen GT', 'Kennzeichen-GT', 'Kommentar', 'Leserbriefe', 'Testphase Mensch',
#                            'Tipp der Redaktion']
#        
#        if ausgabe == "HF": 
#            kolumnen_list = ['Das Mittwochsrätsel', 'Das sach auch man', 'Das Schaufenster', 'Iut De Noberskopp', 
#                             'Kommentar', 'Leserbriefe', 'Radio Herford', 'So gesehen']
#        
#        if ausgabe == "PB":
#            kolumnen_list = ['Auf ein Wort', 'Campus Aktuell', 'Guten Morgen PB', 'Kennzeichen PB',
#                             'Klartext von Gläubigen', 'Kommentar', 'Leserbriefe', 'Persönlich', 'Zwischenruf']
#       
#        grafik_lesewert(prs, df_, target="rubriken", minimize=1, label_position="xlarge",
#                            ressort_liste=kolumnen_list, special=False, title_text="Lesewerte Rubriken/Kolumnen ", order="ok", legend="normal", sort="Lesewert", 
#                        article="mean", mean_line=0, grid=True)
            
#        df_kultur=df[df["Ressortbeschreibung"]=="Kultur"]
#        kult_rub = ['Stern der Woche', 'Lesezeichen', 'TV-Kritik', 'Hennig liest',
#           'Ducksch spielt', 'Persönlich']
#        grafik_lesewert(prs, df_kultur, target="rubriken", minimize=1, label_position="xlarge",
#                        ressort_liste=kult_rub, special=False, title_text="Lesewerte Rubriken/Kolumnen Kultur", order="ok", legend="normal", sort="Lesewert", 
#                        article="mean", mean_line=0, grid=True)
##             
#    # 25 Lokale Übersichten 1. Lokseite etc. 
#        if ausgabe == "BI": 
#            lok_eins = "Lok 1 BI"
#            lok_zwei = "Stadtteile 1 BI"
#            üs_zwei = "Stadtteile 1"
#        if ausgabe == "GT": 
#            lok_eins = "Lok 1 GT"
#            lok_zwei = "Kreis 1 GT"
#            üs_zwei = "Kreisseite 1"
#        if ausgabe == "PB": 
#            lok_eins = "Lok 1 PB"
#            lok_zwei = "Kreis 1 PB"
#            üs_zwei = "Kreisseite 1"
#        if ausgabe == "LÜB": 
#            lok_eins = "Lok 1 LÜB"
#            
#            
#        df_lok1 = df_[df_["Spezial"]==lok_eins]
#        deckblatt_abschluss(prs, df_lok1, title_text = "Lokale Eins " + ausgaben_dict[ausgabe], lokales=True)
#        tabelle_ressortauswertung(prs, df_lok1, research_object="Darstellungsformen", sort="Lesewert", 
#                                  title_text = "Darstellungsformen " + ausgaben_dict[ausgabe], minimum = 5, lokales=True)
#        darstellungsformen(prs, df_lok1, minimum = 5, geschlecht=False, title_text="Darstellungsformen " +  ausgaben_dict[ausgabe])
#        
#        tabelle_ressortauswertung(prs, df_lok1, research_object="Platzierungen", sort="Lesewert", 
#                                  title_text = "Platzierungen " +  ausgaben_dict[ausgabe], minimum = 5, lokales=True)
#        platzierungen(prs, df_lok1, minimum = 5, geschlecht=False, title_text="Platzierungen " +  ausgaben_dict[ausgabe])
#        
#        
#        # jetzt, falls vorhanden, die zweiten Lokalbücher
#        if ausgabe != "LÜB":
#            print(df_["Spezial"].unique())
#            df_lok2 = df_[df_["Spezial"]==lok_zwei]
#            print(lok_zwei)
#            print(df_lok2.shape)
#            deckblatt_abschluss(prs, df_lok2, title_text = üs_zwei + " " + ausgaben_dict[ausgabe], lokales=True)
#            tabelle_ressortauswertung(prs, df_lok2, research_object="Darstellungsformen", sort="Lesewert", 
#                                      title_text = "Darstellungsformen " + ausgaben_dict[ausgabe], minimum = 5, lokales=True)
#            darstellungsformen(prs, df_lok2, minimum = 5, geschlecht=False, title_text="Darstellungsformen " +  ausgaben_dict[ausgabe])
#            
#            tabelle_ressortauswertung(prs, df_lok2, research_object="Platzierungen", sort="Lesewert", 
#                                      title_text = "Platzierungen " +  ausgaben_dict[ausgabe], minimum = 5, lokales=True)
#            platzierungen(prs, df_lok2, minimum = 5, geschlecht=False, title_text="Platzierungen " +  ausgaben_dict[ausgabe])
#
#        df_lok3 = df_[df_["Spezial"]=="folgeseiten_1"]
#        deckblatt_abschluss(prs, df_lok3, title_text = "Lokales nach Lok1 " + ausgaben_dict[ausgabe], lokales=True)
#        tabelle_ressortauswertung(prs, df_lok3, research_object="Darstellungsformen", sort="Lesewert", 
#                                  title_text = "Darstellungsformen " + ausgaben_dict[ausgabe], minimum = 5, lokales=True)
#        darstellungsformen(prs, df_lok3, minimum = 5, geschlecht=False, title_text="Darstellungsformen " +  ausgaben_dict[ausgabe])
#        
#        tabelle_ressortauswertung(prs, df_lok3, research_object="Platzierungen", sort="Lesewert", 
#                                  title_text = "Platzierungen " +  ausgaben_dict[ausgabe], minimum = 5, lokales=True)
#        platzierungen(prs, df_lok3, minimum = 5, geschlecht=False, title_text="Platzierungen " +  ausgaben_dict[ausgabe])
#        
#        if ausgabe != "LÜB":
#            df_lok4 = df_[df_["Spezial"]=="folgeseiten_2"]
#            deckblatt_abschluss(prs, df_lok4, title_text = "Sublokale Seiten " + ausgaben_dict[ausgabe], lokales=True)
#            tabelle_ressortauswertung(prs, df_lok4, research_object="Darstellungsformen", sort="Lesewert", 
#                                      title_text = "Darstellungsformen " + ausgaben_dict[ausgabe], minimum = 5, lokales=True)
#            darstellungsformen(prs, df_lok4, minimum = 5, geschlecht=False, title_text="Darstellungsformen " +  ausgaben_dict[ausgabe])
#            
#            tabelle_ressortauswertung(prs, df_lok4, research_object="Platzierungen", sort="Lesewert", 
#                                      title_text = "Platzierungen " +  ausgaben_dict[ausgabe], minimum = 5, lokales=True)
#            platzierungen(prs, df_lok4, minimum = 5, geschlecht=False, title_text="Platzierungen " +  ausgaben_dict[ausgabe])
#    #Sportarten
    
    # LOKALSPORT
#    liste_loksport = ["BI"]         
#    for ausgabe in liste_loksport:
#        df_ = df[df["Ressortbeschreibung"]=="Lokalsport"]
        #df_ = df[(df["ZTG"]==ausgabe) & (df["Ressortbeschreibung"]=="Lokalsport")]
       
#        df_geschlecht_ = df_geschlecht[(df_geschlecht["ZTG"]==ausgabe) & (df_geschlecht["Ressortbeschreibung"]=="Lokalsport")]
#        
#        
##        
##    # 16 Deckblatt
#        deckblatt_abschluss(prs, df_, title_text = "Lokalsport " + ausgaben_dict[ausgabe], lokales=True)
##        
##    # 17 Entwicklung im Messverlauf
#        print(df_.shape)
#        df_1 = df_[(df_["Platzierungen"]=="AA") | (df_["Platzierungen"]=="SK")]
#        print(df_1.shape)
#        for el in werte:
#            grafik_entwicklung(prs, df_, target=el, mean_line=0, legend="large",
#                               grid=True, title_text = el + " Lokalsport " + ausgaben_dict[ausgabe], ma=True)
#            grafik_entwicklung(prs, df_1, target=el, mean_line=0, legend="large",
#                               grid=True, title_text = el + " AA und SK " + ausgaben_dict[ausgabe], ma=True)
#        
#    # 18 Lesewert nach Seitentiteln
##        df_st = df_[df_["Seitentitel"].isin(seitentitel_list)] 
##        
##        grafik_lesewert(prs, df_, target="seitentitel", minimize=5, label_position="xlarge",
##                    ressort_liste=seitentitel_list, special=False, title_text="Lesewerte nach Seitentiteln " + ausgabe, order="ok", legend="normal", sort="Seitennummer", 
##                    article="total")
#        
##    # 19 Lesewert nach Tagen und im Vergleich zum Angebot
#        lesewert_erscheinung(df_, df_, prs, title_text="Lesewert nach Wochentagen " + ausgabe)
##    
##    # 20 TODO Muss-kann-soll
##     
##    # 21 Meldung vs Ankündiger
#        lesewert_ankündigungen(prs, df_)
###    
##    # 22 Top 10
#        top_10(prs, df_, df_berechnet = False, screenshots=True, number_screenshots = 5, 
#           mode="ressort", headline="Top 10 - " + ausgaben_dict[ausgabe], zeitung=True)
##        # nach Geschlechtern
#        top_10(prs, df_geschlecht_, df_berechnet = False, screenshots=False, number_screenshots = 5, 
#           mode="ressort", headline="Top 10 - " + ausgaben_dict[ausgabe] + " - Frauen", zeitung=True, geschlecht="weiblich")
#        top_10(prs, df_geschlecht_, df_berechnet = False, screenshots=False, number_screenshots = 5, 
#           mode="ressort", headline="Top 10 - " + ausgaben_dict[ausgabe] + " - Frauen", zeitung=True, geschlecht="männlich")
#        # nach Altersgruppen
#        top_10(prs, df_geschlecht_, df_berechnet = False, screenshots=False, number_screenshots = 5, 
#           mode="ressort", headline="default", zeitung=True, alter="40 bis 59 Jahre")
#        top_10(prs, df_geschlecht_, df_berechnet = False, screenshots=False, number_screenshots = 5, 
#           mode="ressort", headline="default", zeitung=True, alter="60 bis 79 Jahre") 
#         
##    # 23 Darstellungsformen und Platzierungen
#        tabelle_ressortauswertung(prs, df_, research_object="Darstellungsformen", sort="Lesewert", 
#                                  title_text = "Darstellungsformen " + ausgaben_dict[ausgabe], minimum = 5, lokales=True)
#        darstellungsformen(prs, df_, minimum = 5, geschlecht=False, title_text="Darstellungsformen " +  ausgaben_dict[ausgabe])
##        darstellungsformen(prs, df_geschlecht_, minimum = 5, geschlecht=True, title_text="Darstellungsformen " +  ausgaben_dict[ausgabe])
#        tabelle_ressortauswertung(prs, df_, research_object="Platzierungen", sort="Lesewert", 
#                                  title_text = "Platzierungen " +  ausgaben_dict[ausgabe], minimum = 5, lokales=True)
#        platzierungen(prs, df_, minimum = 5, geschlecht=False, title_text="Platzierungen " +  ausgaben_dict[ausgabe])
#        platzierungen(prs, df_geschlecht_, minimum = 5, geschlecht=True, title_text="Platzierungen " +  ausgaben_dict[ausgabe])
##        
#    # 24 Übersicht Kolumnen + Rubriken
#        print(kolumnen_list)
#        grafik_lesewert(prs, df_, target="rubriken", minimize=5, label_position="xlarge",
#                    ressort_liste=kolumnen_list, special=False, title_text="Rubriken/Kolumnen Lokalsport " + ausgabe, order="ok", legend="normal", sort="Lesewert", 
#                    article="mean", mean_line=0, grid=False)
##    # Übersicht Sportarten im Lokalsport
#        print ("Sportarten starten")
#    
#        df_sportarten = df_.copy()
#        df_sportarten = df_sportarten[df_sportarten["Ressortbeschreibung"]=="Lokalsport"]
#        df_sportarten = df_sportarten[~df_sportarten["Themen"].isnull()]
#    
#        liste_sportarten_nw =["Fußball", "Handball", "Tennis", "Leichtathletik", 
#                          "Radsport", "Reitsport", "Frauenfußball", "Tischtennis", 
#                          "Football", "Handball", "Basketball"]
#    
#        def check_sportarten(elem): 
#        
#            for v in liste_sportarten_nw: 
#                if v in elem:
#                    return v
#        
#    
#    
#        df_sportarten["sport"] = df_sportarten["Themen"].apply(check_sportarten)
#    
#        print(df_sportarten["sport"].value_counts())
#        grafik_lesewert(prs, df_sportarten, target="sport", minimize=5, label_position="large",
#                            ressort_liste=[], special=False, title_text="Beliebteste Sportarten", order="ok", legend="normal", sort="Lesewert", 
#                            article="total", mean_line = 0, grid=False, anzahl_lokales=1, mean_line_title = "", axis=2) 
#        #    
#




 #print("WM startet")
#    # WM-Special 
#    df_wm = df.copy()
#    df_wm_gesch = df_geschlecht.copy()
#    
#    df_wm = df_wm[~df_wm["Themen"].isnull()]
#    df_wm_gesch = df_wm_gesch[~df_wm_gesch["Themen"].isnull()]
#    df_wm = df_wm[~df_wm["Akteure"].isnull()]
#    df_wm_gesch = df_wm_gesch[~df_wm_gesch["Akteure"].isnull()]
#    df_wm = df_wm[~df_wm["Handlungsorte"].isnull()]
#    df_wm_gesch = df_wm_gesch[~df_wm_gesch["Handlungsorte"].isnull()]
#    
#    
#    df_wm = df_wm[df_wm["Ressortbeschreibung"]=="Sport"]
#    df_wm_gesch = df_wm_gesch[df_wm_gesch["Ressortbeschreibung"]=="Sport"]
#    
#    df_wm = df_wm.reset_index(drop=True)
#    df_wm_gesch = df_wm_gesch.reset_index(drop=True)
#    # Idee 1 mit apply
##    def wm_test(elem):
###        if ("Fußball-WM" in elem) | ("FIFA" in elem) |("WM-Qualifikation" in elem) \
###        | ("Länderspiel" in elem)
##       if "Fußball-WM" in elem: 
##            print (elem)
##            return "WM"
##       else: 
##            return "No WM"
##    
##    
##    df_wm["WM"] = df_wm["Themen"].apply(wm_test)
###    
#    
#   # Idee 2 mit get und set_value
#   
##    
#    #
#    df_wm["WM"]= "nein"
#    df_wm_gesch["WM"]= "nein"
#    
##    for i in df_wm.index:  
##    
##        themen = df_wm.get_value(i, "Themen")
##        akteure = df_wm.get_value(i, "Akteure")
##        handlungsorte = df_wm.get_value(i, "Handlungsorte")
##        if "Fußball-WM" in themen:
##            df_wm.set_value(i, "WM", "WM")
##        if "FIFA" in themen:
##            df_wm.set_value(i, "WM", "WM") 
##        if "WM-Qualifikation" in themen:
##            df_wm.set_value(i, "WM", "WM")
##        if "Länderspiel" in themen:
##            df_wm.set_value(i, "WM", "WM")
##        if "Russland" in handlungsorte: 
##            df_wm.set_value(i, "WM", "WM")
##        if "DFB" in akteure:
##            df_wm.set_value(i, "WM", "WMD")
##        if "Kroos" in akteure:
##            df_wm.set_value(i, "WM", "WMD")
##        if "Löw" in akteure:  
##            df_wm.set_value(i, "WM", "WMD")
##        if "Neuer" in akteure:
##            df_wm.set_value(i, "WM", "WMD")
###  
##    
#    for i in df_wm_gesch.index:  
#    
#        themen = df_wm_gesch.get_value(i, "Themen")
#        akteure = df_wm_gesch.get_value(i, "Akteure")
#        handlungsorte = df_wm_gesch.get_value(i, "Handlungsorte")
#        if "Fußball-WM" in themen:
#            df_wm_gesch.set_value(i, "WM", "WM")
#        if "FIFA" in themen:
#            df_wm_gesch.set_value(i, "WM", "WM") 
#        if "WM-Qualifikation" in themen:
#            df_wm_gesch.set_value(i, "WM", "WM")
#        if "Länderspiel" in themen:
#            df_wm_gesch.set_value(i, "WM", "WM")
#        if "Russland" in handlungsorte: 
#            df_wm_gesch.set_value(i, "WM", "WM")
#        if "DFB" in akteure:
#            df_wm_gesch.set_value(i, "WM", "WMD")
#        if "Kroos" in akteure:
#            df_wm_gesch.set_value(i, "WM", "WMD")
#        if "Löw" in akteure:  
#            df_wm_gesch.set_value(i, "WM", "WMD")
#        if "Neuer" in akteure:
#            df_wm_gesch.set_value(i, "WM", "WMD")
#    
#    #df_turnier = df_wm[(df_wm["WM"]=="WM") | (df_wm["WM"]=="WMD")]
#    df_turnier_geschlecht = df_wm_gesch[(df_wm_gesch["WM"]=="WM") | (df_wm_gesch["WM"]=="WMD")]
#    
#    deckblatt_abschluss(prs, df_turnier)
#    wm = df_turnier["Artikel-Lesewert (Erscheinung) in %"].mean()
#    wmd = df_turnier[df_turnier["WM"]=="WMD"]["Artikel-Lesewert (Erscheinung) in %"].mean()
#    sport = df[df["Ressortbeschreibung"]=="Sport"]["Artikel-Lesewert (Erscheinung) in %"].mean()
#    
#    wmb = df_turnier["Artikel-Blickwert (Erscheinung) in %"].mean()
#    wmdb = df_turnier[df_turnier["WM"]=="WMD"]["Artikel-Blickwert (Erscheinung) in %"].mean()
#    sportb = df[df["Ressortbeschreibung"]=="Sport"]["Artikel-Blickwert (Erscheinung) in %"].mean()
#    
#    wmdd = df_turnier["Artikel-Durchlesewerte (Erscheinung) in %"].mean()
#    wmddd = df_turnier[df_turnier["WM"]=="WMD"]["Artikel-Durchlesewerte (Erscheinung) in %"].mean()
#    sportdd = df[df["Ressortbeschreibung"]=="Sport"]["Artikel-Durchlesewerte (Erscheinung) in %"].mean()
#    
#    print("Lesewert alle WM-Artikel: {0:.2f}".format(wm))
#    print("Lesewert nur deutsche WM-Artikel: {0:.2f}".format(wmd))
#    print("Lesewert alle Sport-Artikel: {0:.2f}".format(sport))
#    
#    
#    print("Blickwert alle WM-Artikel: {0:.2f}".format(wmb))
#    print("Blickwert nur deutsche WM-Artikel: {0:.2f}".format(wmdb))
#    print("Blickwert alle Sport-Artikel: {0:.2f}".format(sportb))
#    
#    print("Durchlesewert alle WM-Artikel: {0:.2f}".format(wmdd))
#    print("Durchlesewert nur deutsche WM-Artikel: {0:.2f}".format(wmddd))
#    print("Durchlesewert alle Sport-Artikel: {0:.2f}".format(sportdd))
#    
#    tabelle_ressortauswertung(prs, df_turnier, research_object="Darstellungsformen", sort="Lesewert", 
#                                  title_text = "Darstellungsformen WM", minimum = 5)
#    darstellungsformen(prs, df_turnier_geschlecht, minimum = 5, geschlecht=False, title_text="WM nach Geschlecht")
##    tabelle_ressortauswertung(prs, df_turnier, research_object="Platzierungen", sort="Lesewert", 
##                                  title_text = "Platzierungen ", minimum = 5)
##    platzierungen(prs, df_turnier, minimum = 5, geschlecht=False, title_text="Platzierungen ")
#    platzierungen(prs, df_turnier_geschlecht, minimum = 5, geschlecht=True, title_text="Platzierungen ")
#    
##    top_10(prs, df_turnier, df_berechnet = False, screenshots=True, number_screenshots = 5, 
##           mode="ressort", headline="Top 10 Fußball-WM", zeitung=False, geschlecht="", alter="", 
##           lokales_ressort="Lokales")
#    top_10(prs, df_turnier_geschlecht, df_berechnet = False, screenshots=True, number_screenshots = 5, 
#           mode="ressort", headline="Top 10 Fußball-WM", zeitung=False, geschlecht="weiblich", alter="", 
#           lokales_ressort="Lokales")
#    top_10(prs, df_turnier_geschlecht, df_berechnet = False, screenshots=True, number_screenshots = 5, 
#           mode="ressort", headline="Top 10 Fußball-WM", zeitung=False, geschlecht="männlich", alter="", 
#           lokales_ressort="Lokales")
    
#    print("Fußball startet")
#    
#    
#    vereine_nw = ["Arminia Bielefeld", "Bayern", "Schalke", "SC Paderborn", "BVB", "S04"]
#    
#    df_vereine = df_nk.copy()
#    df_vereine = df_vereine[df_vereine["Ressortbeschreibung"]=="Sport"]
#    
#    df_vereine = df_vereine[~df_vereine["Akteure"].isnull()]
#    def find_verein(elem): 
#        for v in vereine_nw: 
#            if v in elem:
#                return v
#            
#        
#    
#    
#    df_vereine["Verein"] = df_vereine["Akteure"].apply(find_verein)
#    
#    df_vereine = df_vereine[~df_vereine["Verein"].isnull()]
#    print(df_vereine.shape)
#    print(df_vereine["Verein"].value_counts())
#    print(ausgaben_liste)
#    for aus in ausgaben_liste:  
#        df_vereine_aus = df_vereine[df_vereine["ZTG"]==aus]
#        print(aus)
#        print("Hier geht es um Vereine")
#        grafik_lesewert(prs, df_vereine_aus, target="Verein", minimize=5, label_position="xlarge",
#                        ressort_liste=[], special=False, title_text="Vereine " + ausgaben_dict[aus], order="ok", legend="normal", sort="Lesewert", 
#                        article="total", mean_line = 0, grid=False, anzahl_lokales=1, mean_line_title = "", axis=2) 
#    
#    
#    
#    print ("Sportarten starten")
#    
#    df_sportarten = df.copy()
#    df_sportarten = df_sportarten[df_sportarten["Ressortbeschreibung"]=="Sport"]
#    df_sportarten = df_sportarten[~df_sportarten["Themen"].isnull()]
#    
#    liste_sportarten_nw =["Fußball", "Handball", "Tennis", "Leichtathletik", 
#                          "Radsport", "Reitsport", "Frauenfußball", "Tischtennis", 
#                          "Football", "Handball", "Basketball"]
#    
#    def check_sportarten(elem): 
#        
#        for v in liste_sportarten_nw: 
#            if v in elem:
#                return v
#        
#    
#    
#    df_sportarten["sport"] = df_sportarten["Themen"].apply(check_sportarten)
#    
#    print(df_sportarten["sport"].value_counts())
#    grafik_lesewert(prs, df_sportarten, target="sport", minimize=5, label_position="large",
#                        ressort_liste=[], special=False, title_text="Beliebteste Sportarten", order="ok", legend="normal", sort="Lesewert", 
#                        article="total", mean_line = 0, grid=False, anzahl_lokales=1, mean_line_title = "", axis=2) 
##    
##   
#    print("Start Kolumnene")
#    kol_sport = ["Einwurf", "Kurz notiert", "Sport Kompakt", "Wort zum Sport", 
#                 "Fußball kompakt", "WM Kompakt", "Persönlich", 
#                 "Ballgeflüster", "Das besondere WM-Datum", "Beim DSC am Ball", 
#                 "Beim SCP am Ball"]
#                 
#    df_kol = df.copy()            
#    df_kol = df_kol[df_kol["Ressortbeschreibung"]=="Sport"]
#    df_kol = df_kol[~df_kol["Beschreibung"].isnull()]
#    df_kol["Kolumn"] = "nein"
#    
#    def check_kol(elem): 
#        for v in kol_sport:
#            
#            if v in elem: 
#                return elem
#    df_kol["Kolumne"] = df_kol["Beschreibung"].apply(check_kol)
#    
#    grafik_lesewert(prs, df_kol, target="Kolumne", minimize=3, label_position="xxlarge",
#                    ressort_liste=[], special=False, title_text="Kolumnen/Rubriken Sport", order="ok", legend="normal", sort="Lesewert", 
#                    article="total", mean_line = 0, grid=False, anzahl_lokales=1, mean_line_title = "", axis=2)
#    
#    
#    
#    
    
#    
#    df_geschlecht_ = df_geschlecht[df_geschlecht["Ressortbeschreibung"]==elem]
#    df_aufmacher = df_[df_["Darstellungsformen"]=="AA"] 
#    df_großetexte = df_[df_["Platzierungen"].isin(["AA", "SK"])]
#     # 12 - Deckblatt Ressort - CHECK    
#        
#    deckblatt_abschluss(prs, df_)
##        
# # 13 - Ressortübersicht
#    tabelle_ressortauswertung(prs, df_, research_object="Darstellungsformen", sort="Lesewert", 
#                              title_text = "Darstellungsformen " + ressort_dict[elem], minimum = 5)
#    darstellungsformen(prs, df_, minimum = 5, geschlecht=False, title_text="Darstellungsformen " + ressort_dict[elem])
#    darstellungsformen(prs, df_geschlecht_, minimum = 5, geschlecht=True, title_text="Darstellungsformen " + ressort_dict[elem])
#    tabelle_ressortauswertung(prs, df_, research_object="Platzierungen", sort="Lesewert", 
#                              title_text = "Platzierungen " + ressort_dict[elem], minimum = 5)
#    platzierungen(prs, df_, minimum = 5, geschlecht=False, title_text="Platzierungen " + ressort_dict[elem])
#    platzierungen(prs, df_geschlecht_, minimum = 5, geschlecht=True, title_text="Platzierungen " + ressort_dict[elem])
#   
#    
    
    
    
        
    print("-----------------------------------")
    print("Spiele neues PPTX-Dokument aus")
    prs.save("nw_abschluss.pptx")
    print("-----------------------------------")
    print("Analyse abgeschlossen")
    print("-----------------------------------")
    t_1 = time.time()
    endtime = t_1 - t_0 
    print("Dauer des Durchlaufs: {0:.2f}.".format(endtime))
    
    
    
    
#%% Stuttgarter Medienhaus STZ & STN

def stuttgart(df, df_double_split = 0, df_nk=0, df_scans=0, df_umfrage=0,
                       df_geschlecht = 0, liste = ausgaben_liste, df_ereignis=0, 
                       geschlecht_double_split=0):
    
    stutt_dict={"STZ":"Stuttgarter Zeitung", "STN":"Stuttgarter Nachrichten"}
   
    # ZEITMESSUNG Start
    t = time.process_time()
    
    print("Analyse SMH gestartet...")
    
    #prs-Objekt erstellen
    prs = Presentation("LW_neu_2.pptx")
    
    # Stuttgarter Ausgaben identifizieren
    stutt_ausgaben = ["STZ", "STN"]
    
    # Stuttgarter Einzelausgaben über ZTG identifizieren
    stutt_einzelausgaben = ["FI_STZ", "FI_STN", "L_STZ", "S_STZ", "S_STN"]
    werte = ["Lesewert", "Blickwert", "Durchlesewert"]
    # TODO: Diese beiden Artikel für lokale Analysen rausnehmen, sind 
    # Artikel über LW... 
    df = df[df["SplitId"]!="18796"]
    df = df[df["SplitId"]!="19113"]
#    df_double_split = df_double_split[df_double_split["SplitId"]!="18796"]
#    df_double_split = df_double_split[df_double_split["SplitId"]!="19113"]
#    df_nk = df_nk[df_nk["SplitId"]!="18796"]
#    df_nk = df_nk[df_nk["SplitId"]!="19113"]
#    df_geschlecht = df_geschlecht[df_geschlecht["SplitId"]!="18796"]
#    df_geschlecht = df_geschlecht[df_geschlecht["SplitId"]!="19113"]
    
    # ALLGEMEINES
    
#    # 1 Entwicklung Lesewert + Durchlesewert im Messverlauf - nach Titeln (STZ und STN) CHECK
#    # ACHTUNG, Besonderheit: In Stuttgart müssen wir zunächst den Datensatz mit den doppelten IDs nehmen, 
    # dann nach Zeitungen trennen und DANN die doppelten Split_ids eliminieren. Grund: Wir haben es hier 
    # mit zwei Zeitungen zu tun, die sich ansonsten nicht sauber trennen lassen... 
#    
    print ("Analyse Lesewert/Durchlesewert im Messverlauf...")
#    for zt in stutt_ausgaben: 
#        df_ = df_double_split[df_double_split["Zeitung"]==zt]
#        df_ = df_.drop_duplicates(subset="SplitId")
#        for elem in werte:
#            grafik_entwicklung(prs, df_, target=elem, mean_line=0, legend="large", grid=True, 
#                               title_text = "Entwicklung " + elem + " / " +zt, steps=1, ma=True)
             #Hier die WErte nur für Aufmacher und SK - erstmal rausnehmen
#            df__ = df_[df_["Platzierungen"].isin(["AA", "SK"])]
#            grafik_entwicklung(prs, df__, target=elem, mean_line=0, legend="large", grid=True, 
#                               title_text = elem + " AA + SK / " +zt, steps=2)
##    
##    
#    # 2 Fragebogen-Auswertung für alle fünf Panels CHECK
#   
    print ("Initialisiere Auswertung der Fragebögen")
    # Aufruf jeder einzelnen Ausgabe
    for elem in ausgaben_liste:
        
        df_ = df_umfrage[df_umfrage["ZTG"]==elem]
        df_ = df_[df_["Teilnehmernummer"]!="2011442932"]
        #Aufruf gestapelte Bars "Ihre Zeitung ist..."
        zeitungsattribute_berechnung(prs, df_, title = elem) # elem wird später über ausgaben_dict umgewandelt
        
        #Aufruf Bars "Welche Themen
        zeitung_themen(prs, df_, title=elem)
        
        # Aufruf kleine Pie-Charts
        umfrage_pie(prs, df_, title=elem)
        
        # Aufruf kleine Barcharts, TV- und Internetnutzung
        
        mini_bars(prs, df_, title_text = "Fernsehnutzung", liste=umfrage_mini_bars_TV, title=elem)
        mini_bars(prs, df_, title_text = "Internetnutzung", liste=umfrage_mini_bars_netz, title=elem)
        mini_bars(prs, df_, title_text = "Persönliches", liste=umfrage_demografie_bars, title=elem)
    
#    
# 
#    
###    
###    
###    # 4 Kennzahlen der Messung / Marken - Titel, Messwellen, Lokalteile CHECK
####    
##    for i in stutt_ausgaben:
##        df_ = df[df["Zeitung"]==i]
##        df_nichtkum_ = df_nk[df_nk["Zeitung"]==i]
##        df_double_split_ = df_double_split[df_double_split["Zeitung"]==i]
##        print("Werte für " + i)
##        marken_analyse(df_, df_doublesplitid=df_double_split_, df_nichtkum=df_nichtkum_, kunde="MHS")
##    
##    for i in stutt_einzelausgaben: 
##        df_ = df[df["ZTG"]==i]
##        df_nichtkum_ = df_nk[df_nk["ZTG"]==i]
##        df_double_split_ = df_double_split[df_double_split["ZTG"]==i]
##        print()
##        print("WERTE FÜR "+ausgaben_dict[i])
##        marken_analyse(df_, df_doublesplitid=df_double_split, df_nichtkum=df_nichtkum_, kunde="MHS")
##    
##    for i in lokale_liste: 
##        df_ = df[df["Ressortbeschreibung"]==i]
##        df_nichtkum_ = df_nk[df_nk["Ressortbeschreibung"]==i]
##        df_double_split_ = df_double_split[df_double_split["Ressortbeschreibung"]==i]
##        print()
##        print("Werte für "+ i)
##        marken_analyse(df_, df_doublesplitid=df_double_split, df_nichtkum=df_nichtkum_, kunde="MHS")
###        
###    
####    #5 Kennzahlen CHECK
###    kennzahlen(df, df_scans)
####    
####    # 6 Top 5 nach Titeln CHECK
#    for zt in stutt_ausgaben:
#        #df_ = df[df["Zeitung"]==zt]
#        df_ = df_double_split[df_double_split["Zeitung"]==zt]
#        df_ = df_.drop_duplicates(subset="SplitId")
#        top_10(prs, df_, df_berechnet = False, screenshots=True, number_screenshots = 5, 
#           mode="Gesamt", headline="Top 10 " + stutt_dict[zt], zeitung=False)
##    
####    
###    # b 
###    # TODO: WellenteilnahmeId muss noch genutzt werden, um Nutzer zu identifizieren. 
#    for zt in stutt_ausgaben: 
#        df_ = df_scans[df_scans["Zeitung"]==zt]
#        lesetage(prs, df_, title_text = "Lesetage " + stutt_dict[zt])
#        lesezeit(prs, df_, title_text = "Lesezeiten " + stutt_dict[zt])
##       
### 
###    
###    # 8 Übersicht Ressorts nach Lesewert und Umfang - CHECK
#    for zt in stutt_ausgaben: 
#        #df_ = df[df["Zeitung"]==zt]
#        df_ = df_double_split[df_double_split["Zeitung"]==zt]
#        df_ = df_.drop_duplicates(subset="SplitId")
#        mean_line = df_["Artikel-Lesewert (Erscheinung) in %"].mean()
#        print(df_.shape)
#        if zt == "STZ": 
#            sortierung = "Kategorie"
#            liste = ressortliste_stz
#            liste_spezial = ressortliste_spezial_res_stz
#        if zt == "STN": 
#            sortierung = "Kategorie"
#            liste = ressortliste_stn
#            liste_spezial = ressortliste_spezial_res_stn
#        print("erste Grafik")
#        grafik_lesewert(prs, df_, target="ressort", minimize=5, label_position="xxlarge",
#                    ressort_liste=liste, special=False, title_text="Ressorts " + stutt_dict[zt],
#                    legend="normal", sort=sortierung, article="mean", mean_line=mean_line, axis=2)
#        # Übersicht mit verschiedenen Asuwahlen (Ressorts und Seitentitel, z.B. die Lok1)
#        print("zweite Grafik")
#        grafik_lesewert(prs, df_, target="special_res", minimize=5, label_position="xxlarge", ressort_liste = liste_spezial, 
#                    special=False, title_text="Ressorts plus Lokale Eins, " + zt, order="ok",
#                    legend="normal", sort=sortierung, article="mean", mean_line=mean_line, axis=2)
#    # 9 Übersicht Ressorts nach Geschlechtern - CHECK
#    for zt in stutt_ausgaben:
#        print("Grafiken Geschlecht")
#        df_ = geschlecht_double_split[geschlecht_double_split["Zeitung"]==zt]
#        df_ = df_.drop_duplicates(subset="SplitId")
#        multiple_bars_geschlecht(prs, df_, target="Ressorts", grid = True, legend="xlarge", 
#                             ressort_liste = liste, title_text= "Ressorts nach Geschlecht - "+zt, 
#                             sort=sortierung)
    
##    # 10  Lesewert nach Erscheinungstagen - CHECK
#    for zt in stutt_ausgaben:
#        df_ = df_geschlecht[df_geschlecht["Zeitung"]==zt]
#        df_double_split_ = df_double_split[df_double_split["Zeitung"]==zt]
#        lesewert_erscheinung(df_,df_double_split_, prs, title_text="Lesewert nach Wochentagen / "+ zt)
##        grafik_lesewert(prs, df_, target="ressort", minimize=5, label_position="xlarge",
##                    ressort_liste=[], special=False, title_text="Ressorts nach Lesewert " + zt, order="ok",
##                    legend="normal", sort="Seitennummer", article="mean")
##    
##
##    # 11 Leserwert nach Darstellungsformen - CHECK
#    for zt in stutt_ausgaben:   
#        df_ = df[df["Zeitung"]==zt]
#        df_geschlecht_ = df_geschlecht[df_geschlecht["Zeitung"]==zt]
#        darstellungsformen(prs, df_, minimum = 5, geschlecht=False, title_text="LW nach Darstellungsformen - "+zt, special="")
#        darstellungsformen(prs, df_geschlecht_, minimum = 5, geschlecht=True, title_text="Darstellungsformen nach Geschlecht", special="")
#                
#    
##    # 12 Topkolumnen pro Asugabe CHECK
#    for zt in stutt_ausgaben: 
#        #df_ = df[df["Zeitung"]==zt]
#        df_ = df_double_split[df_double_split["Zeitung"]==zt]
#        df_ = df_.drop_duplicates(subset="SplitId")
#        grafik_lesewert(prs, df_, target="rubriken", minimize=5, label_position="xlarge",
#                    ressort_liste=kolumnen_liste, special=False, title_text="Kolumnen " + stutt_dict[zt], order="ok", legend="xlarge", sort="Seitennummer", 
#                    article="total", mean_line = 0, grid=False, anzahl_lokales=1, mean_line_title = "", axis=1)
#    
##   # ÜBERSICHT LOKALTEILE CHECK
### 
    
#     
   # for elem in lokale_liste:
       
#       df_ = df[df["Ressortbeschreibung"]==elem]
#       df_nk_ = df_nk[df_nk["Ressortbeschreibung"]==elem]
#       df_geschlecht_ = df_geschlecht[df_geschlecht["Ressortbeschreibung"]==elem]
#       tabelle_ressortauswertung(prs, df_, research_object="Darstellungsformen", sort="Lesewert", 
#                                  title_text = "Darstellungsformen " + elem, minimum = 5, kunde="MHS")
#       tabelle_ressortauswertung(prs, df_, research_object="Platzierungen", sort="Lesewert", 
#                                  title_text = "Platzierungen " + elem, minimum = 5, kunde="MHS")
#       deckblatt_abschluss(prs, df_, kunde="MHS" )
#       platzierungen(prs, df_, minimum = 5, geschlecht=False, title_text="Platzierungen " + elem, special="")
#       #platzierungen(prs, df_geschlecht_, minimum = 5, geschlecht=True, title_text="Platzierungen nach Geschlecht", special="")
#       darstellungsformen(prs, df_, minimum = 5, geschlecht=False, title_text="Lesewert Darstellungsformen", special="")
#       #darstellungsformen(prs, df_geschlecht_, minimum = 5, geschlecht=True, title_text="Darstellungsformen nach Geschlecht", special="")
##       for el in werte:
##            grafik_entwicklung(prs, df_, target=el, mean_line=0, legend="large",
##                               grid=True, title_text = el + " " + elem, steps=2, ma = True)
##            df__ = df_[(df_["Platzierungen"]=="AA")|(df_["Platzierungen"]=="SK")]
##            grafik_entwicklung(prs, df__, target=el, mean_line=0, legend="large",
##                               grid=True, title_text = el + " Aufmacher und Seitenkeller", steps=2)
#       grafik_lesewert(prs, df_, target="seitentitel", minimize=5, label_position="large",
#                    ressort_liste=[], special=False, order="ok", legend="normal", sort="Seitennummer", 
#                    article="total", mean_line = 0, grid=False, anzahl_lokales=1, axis=2, 
#                    title_text="Seitentitel " + ressort_dict[elem])
#     
#       grafik_lesewert(prs, df_, target="special", minimize=5, label_position="large",
#                    ressort_liste=[], special=False, order="ok", legend="normal", sort="Seitennummer", 
#                    article="total", mean_line = 0, grid=False, anzahl_lokales=1, axis=2, 
#                    title_text="Seitentitel " + ressort_dict[elem])
#       
##       lesewert_erscheinung(df_,df_nk_, prs, title_text="LW nach Wochentagen")
##       lesewert_ankündigungen(prs, df_nk_, title_text="Ankündigungen - " + ressort_dict[elem])
##       top_10(prs, df_, df_berechnet = False, screenshots=True, number_screenshots = 5, 
##           mode="ressort", headline="Top 10 - "+elem, zeitung=True, geschlecht="", alter="", 
##           lokales_ressort="Lokales")
#        # Top 10 Geschlechter
#       top_10(prs, df_geschlecht_, df_berechnet = False, screenshots=False, number_screenshots = 5, 
#           mode="ressort", headline="default", zeitung=True, geschlecht="weiblich")
#       top_10(prs, df_geschlecht_, df_berechnet = False, screenshots=False, number_screenshots = 5, 
#           mode="ressort", headline="default", zeitung=True, geschlecht="männlich")
#        
#        # Top 10 Altersgruppen
#       top_10(prs, df_geschlecht_, df_berechnet = False, screenshots=False, number_screenshots = 5, 
#           mode="ressort", headline="default", zeitung=True, alter="40 bis 59 Jahre")
#       top_10(prs, df_geschlecht_, df_berechnet = False, screenshots=False, number_screenshots = 5, 
#           mode="ressort", headline="default", zeitung=True, alter="60 bis 79 Jahre")
#   

####    
###    
###    
##    ### HÄNDISCHER TEIL LOKALES ######
##    #### In der Betellung unter: Weitere Charts
##    
#     #STZ- Top 5 Kolumnen nach Lesewert, Diagramm 
#    df_lokal = df_double_split[df_double_split["Ressortbeschreibung"].isin(lokale_liste)]
#    for zt in stutt_ausgaben: 
#         df_kol = df_lokal[df_lokal["Zeitung"]==zt]
#         #df_kol = df_kol[df_kol["Beschreibung"].isin(kolumnen_list)]
#         kolumnen_finder(prs, df_kol, title_text = "Lokale Kolumnen/Rubriken " + zt)#
#    
#    
    # Ankündigungen vs anderer Texte - die beiden Lokalteile STN und STZ im Vergleich
#    df_ = df_double_split[df_double_split["Zeitung"]=="STZ"]
#    df_ = df_.drop_duplicates(subset="SplitId")
#    df_ = df_[df_["Ressortbeschreibung"].isin(lokale_liste)]
#    lesewert_ankündigungen(prs, df_, title_text ="ankündigungen stz")                
#     
#    df_ = df_double_split[df_double_split["Zeitung"]=="STN"]
#    df_ = df_.drop_duplicates(subset="SplitId")
#    df_ = df_[df_["Ressortbeschreibung"].isin(lokale_liste)]
#    lesewert_ankündigungen(prs, df_, title_text ="Ankündigungen stn")                    
            
# Region/BaWü STZ
#    df_ = df[df["Zeitung"]=="STZ"]
#    df_= df_[df_["Ressortbeschreibung"]=="Region/Baden-Württ. STZ"]
#    deckblatt_abschluss(prs, df_, kunde="MHS" ) 
#    df_bawü = df_[df["Seitentitel"]=="BADEN-WÜRTTEMBERG"]       
#    df_reg = df_[df["Seitentitel"]=="REGION STUTTGART"]        
#    darstellungsformen(prs, df_reg, minimum = 5, geschlecht=False, title_text="Darstellungsformen Region", special="")
#    darstellungsformen(prs, df_bawü, minimum = 5, geschlecht=False, title_text="Darstellungsformen BaWü", special="")
#    for el in werte:
#        grafik_entwicklung(prs, df_reg, target=el, mean_line=0, legend="large", grid=True, title_text = el + "Region Stutt.", steps=2)
#    for el in werte:
#        grafik_entwicklung(prs, df_bawü, target=el, mean_line=0, legend="large", grid=True, title_text = el + "BaWü", steps=2)
#    
#    top_10(prs, df_, df_berechnet = False, screenshots=True, number_screenshots = 5, 
#           mode="ressort", headline="Top 10 - Region/BaWü", zeitung=True, geschlecht="", alter="", 
#           lokales_ressort="Lokales")
    
    
    
    #Lokalsport
#    df_ = df[df["Ressortbeschreibung"]=="Lokalsport"]
#    deckblatt_abschluss(prs, df_, kunde="MHS" )
#    sport_stz = df_[df_["Seitentitel"]=="Lokalsport"]
#    sport_stn = df_[df_["Seitentitel"]=="Sport vor Ort"]
#    
#    grafik_lesewert(prs, df_, target="seitentitel", minimize=5, label_position="large",
#                    ressort_liste=[], special=False, order="ok", legend="normal", sort="Seitennummer", 
#                    article="total", mean_line = 0, grid=False, anzahl_lokales=1, axis=2, 
#                    title_text="Seitentitel Lokalsport")
#    darstellungsformen(prs, df_, minimum = 5, geschlecht=False, title_text="Darstellungsformen Loksport gesamt", special="")
#    darstellungsformen(prs, sport_stz, minimum = 5, geschlecht=False, title_text="Darstellungsformen STZ", special="")
#    darstellungsformen(prs, sport_stn, minimum = 5, geschlecht=False, title_text="Darstellungsformen STN", special="")
#    top_10(prs, sport_stz, df_berechnet = False, screenshots=True, number_screenshots = 5, 
#           mode="ressort", headline="Top 10 - Loksport STZ", zeitung=True, geschlecht="", alter="", 
#           lokales_ressort="Lokales")
#    top_10(prs, sport_stn, df_berechnet = False, screenshots=True, number_screenshots = 5, 
#           mode="ressort", headline="Top 10 - Loksport STN", zeitung=True, geschlecht="", alter="", 
#           lokales_ressort="Lokales")
    
##    ######## MANTEL #######
##    
##    # Titelseite STZ
#    def mantel_auswertung(prs, df, df_geschlecht):
#        titletext=df.iloc[0]["Ressortbeschreibung"]+" " +df.iloc[0]["Zeitung"]
#        deckblatt_abschluss(prs, df, 
#                            title_text=titletext,
#                            lokales=False, kunde="MHS")
        
#        
#     # 13 - Ressortübersicht
#        tabelle_ressortauswertung(prs, df, research_object="Darstellungsformen", sort="Lesewert", 
#                                  title_text = "Darstellungsformen " + titletext, minimum = 5, kunde="MHS")
#        darstellungsformen(prs, df, minimum = 5, geschlecht=False, title_text="Darstellungsformen " + titletext)
#        darstellungsformen(prs, df_geschlecht, minimum = 5, geschlecht=True, title_text="Darstellungsformen " + titletext)
#        tabelle_ressortauswertung(prs, df, research_object="Platzierungen", sort="Lesewert", 
#                                  title_text = "Platzierungen " + titletext, minimum = 5, kunde="MHS")
#        platzierungen(prs, df, minimum = 5, geschlecht=False, title_text="Platzierungen " + titletext)
#        platzierungen(prs, df_geschlecht, minimum = 5, geschlecht=True, title_text="Platzierungen " + titletext)
#       
        # 14 - Entwicklung LW, BW, DW  CHECK
#        print(werte)
#        for el in werte:
#            grafik_entwicklung(prs, df, target=el, mean_line=0, legend="large", grid=True, title_text = el + " " + titletext, steps_x_label=2)
#        # Entwicklung AA und SK - also nur die großen Stücke
#        df_großetexte = df[(df["Platzierungen"]=="AA") | (df["Platzierungen"]=="SK")]
#        for el in werte: 
#            grafik_entwicklung(prs, df_großetexte, target=el, mean_line=0, legend="large", grid=True, title_text = el + " AA und SK", steps_x_label=2)
        
#        top_10(prs, df, df_berechnet = False, screenshots=True, number_screenshots = 5, 
#           mode="ressort", headline="Top 10 - "+titletext, zeitung=False)
#        top_10(prs, df_geschlecht, df_berechnet = False, screenshots=False, number_screenshots = 5, 
#           mode="ressort", headline="Top 10 - "+titletext, zeitung=False, geschlecht="weiblich")
#        top_10(prs, df_geschlecht, df_berechnet = False, screenshots=False, number_screenshots = 5, 
#           mode="ressort", headline="Top 10 - "+titletext, zeitung=False, geschlecht="männlich")
#        top_10(prs, df_geschlecht, df_berechnet = False, screenshots=False, number_screenshots = 5, 
#           mode="ressort", headline="Top 10 - "+titletext, zeitung=False, alter="40 bis 59 Jahre")        
#        top_10(prs, df_geschlecht, df_berechnet = False, screenshots=False, number_screenshots = 5, 
#           mode="ressort", headline="Top 10 - "+titletext, zeitung=False, alter="60 bis 79 Jahre")
#     
#    
    
#    
#    df_stz = df_double_split[df_double_split["Zeitung"]=="STZ"]
#    df_stz = df_stz.drop_duplicates(subset="SplitId")
#    df_stn = df_double_split[df_double_split["Zeitung"]=="STN"]
#    df_stn = df_stn.drop_duplicates(subset="SplitId")
    
    
    
#    
#    # Titelseite STZ
#    #df_ = df_stz[(df["Ressortbeschreibung"]=="Titelseite") & (df["Zeitung"]=="STZ")]
#    df_ = df_stz[df_stz["Ressortbeschreibung"]=="Titelseite"]
#    df_geschlecht_ = df_geschlecht[(df_geschlecht["Ressortbeschreibung"]=="Titelseite") & (df_geschlecht["Zeitung"]=="STZ")]
#    mantel_auswertung(prs, df_, df_geschlecht_)
#    
#    # Titelseite STN
    #df_ = df[(df["Ressortbeschreibung"]=="Titelseite") & (df["Zeitung"]=="STN")]
#    df_ = df_stn[df_stn["Ressortbeschreibung"]=="Titelseite"]
#    df_geschlecht_ = df_geschlecht[(df_geschlecht["Ressortbeschreibung"]=="Titelseite") & (df_geschlecht["Zeitung"]=="STN")]
#    mantel_auswertung(prs, df_, df_geschlecht_)
#    
    # Tagesthema STZ
    #df_ = df[(df["Ressortbeschreibung"]=="Tagesthema") & (df["Zeitung"]=="STZ")]
#    df_ = df_stz[df_stz["Ressortbeschreibung"]=="Tagesthema"]
#    df_geschlecht_ = df_geschlecht[(df_geschlecht["Ressortbeschreibung"]=="Tagesthema") & (df_geschlecht["Zeitung"]=="STZ")]
#    mantel_auswertung(prs, df_, df_geschlecht_)
    
    # Die Dritte Seite STZ
    #df_ = df[(df["Ressortbeschreibung"]=="Die Dritte Seite") & (df["Zeitung"]=="STZ")]
#    df_ = df_stz[df_stz["Ressortbeschreibung"]=="Die Dritte Seite"]
#    df_geschlecht_ = df_geschlecht[(df_geschlecht["Ressortbeschreibung"]=="Die Dritte Seite") & (df_geschlecht["Zeitung"]=="STZ")]
#    mantel_auswertung(prs, df_, df_geschlecht_)
#    
    #Die Seite Drei STN
#    
#    #df_ = df[(df["Ressortbeschreibung"]=="Die Seite Drei") & (df["Zeitung"]=="STN")]
#    df_ = df_stn[df_stn["Ressortbeschreibung"]=="Die Seite Drei"]
#    df_geschlecht_ = df_geschlecht[(df_geschlecht["Ressortbeschreibung"]=="Die Seite Drei") & (df_geschlecht["Zeitung"]=="STN")]
#    mantel_auswertung(prs, df_, df_geschlecht_)
    
    # Politik STZ
#    #df_ = df[(df["Ressortbeschreibung"]=="Politik") & (df["Zeitung"]=="STZ")]
#    df_ = df_stz[df_stz["Ressortbeschreibung"]=="Politik"]
#    df_geschlecht_ = df_geschlecht[(df_geschlecht["Ressortbeschreibung"]=="Politik") & (df_geschlecht["Zeitung"]=="STZ")]
#    mantel_auswertung(prs, df_, df_geschlecht_)
#    
#    pol_liste = ["INNENPOLITIK", "LANDESPOLITIK", "AUSSENPOLITIK"]
#    grafik_lesewert(prs, df_, target="seitentitel", minimize=5, label_position="large",
#                    ressort_liste= pol_liste, special=False, order="ok", legend="normal", sort="Seitennummer", 
#                    article="total", mean_line = 0, grid=False, anzahl_lokales=1, axis=2, 
#                    title_text="Politik STZ nach Seitentiteln")
#    
    
    # Politik STN
   # df_ = df[(df["Ressortbeschreibung"]=="Politik") & (df["Zeitung"]=="STN")]
#    df_ = df_stn[df_stn["Ressortbeschreibung"]=="Politik"]
#    df_geschlecht_ = df_geschlecht[(df_geschlecht["Ressortbeschreibung"]=="Politik") & (df_geschlecht["Zeitung"]=="STN")]
#    mantel_auswertung(prs, df_, df_geschlecht_)
#    
#    pol_liste = ["Meinung/Nachrichten", "Zeitgeschehen"]
#    grafik_lesewert(prs, df_, target="seitentitel", minimize=5, label_position="large",
#                    ressort_liste= pol_liste, special=False, order="ok", legend="normal", sort="Seitennummer", 
#                    article="total", mean_line = 0, grid=False, anzahl_lokales=1, axis=2, 
#                    title_text="Politik STN nach Seitentiteln")
     
    # Landesnachrichten STN
    #df_ = df[(df["Ressortbeschreibung"]=="Landesnachrichten") & (df["Zeitung"]=="STN")]
#    df_ = df_stn[df_stn["Ressortbeschreibung"]=="Landesnachrichten"]
#    df_geschlecht_ = df_geschlecht[(df_geschlecht["Ressortbeschreibung"]=="Landesnachrichten") & (df_geschlecht["Zeitung"]=="STN")]
#    mantel_auswertung(prs, df_, df_geschlecht_)
#    landes_liste = ["Landesnachrichten", "LANDESPOLITIK", "BADEN-WÜRTTEMBERG"]
#    grafik_lesewert(prs, df, target="seitentitel", minimize=5, label_position="large",
#                    ressort_liste=landes_liste, special=False, order="ok", legend="normal", sort="Seitennummer", 
#                    article="mean", mean_line = 0, grid=False, anzahl_lokales=1, axis=2, 
#                    title_text="Vergleich Berichterstattung Land")
#    
    
    # Aus aller Welt STZ 
    #df_ = df[(df["Ressortbeschreibung"]=="Aus aller Welt") & (df["Zeitung"]=="STZ")]
#    df_ = df_stz[df_stz["Ressortbeschreibung"]=="Aus aller Welt"]
#    df_geschlecht_ = df_geschlecht[(df_geschlecht["Ressortbeschreibung"]=="Aus aller Welt") & (df_geschlecht["Zeitung"]=="STZ")]
#    mantel_auswertung(prs, df_, df_geschlecht_)
#    
    # Panorama STN
    #df_ = df[(df["Ressortbeschreibung"]=="Panorama") & (df["Zeitung"]=="STN")]
#    df_ = df_stn[df_stn["Ressortbeschreibung"]=="Panorama"]
#    df_geschlecht_ = df_geschlecht[(df_geschlecht["Ressortbeschreibung"]=="Panorama") & (df_geschlecht["Zeitung"]=="STN")]
#    mantel_auswertung(prs, df_, df_geschlecht_)
#    
   # Wirtschaft STZ
   # df_ = df[(df["Ressortbeschreibung"]=="Wirtschaft") & (df["Zeitung"]=="STZ")]
#    df_ = df_stz[df_stz["Ressortbeschreibung"]=="Wirtschaft"]
#    df_unternehmen = df_double_split[df_double_split["Ressortbeschreibung"]=="Wirtschaft"]
#    df_geschlecht_ = df_geschlecht[(df_geschlecht["Ressortbeschreibung"]=="Wirtschaft") & (df_geschlecht["Zeitung"]=="STZ")]
#    mantel_auswertung(prs, df_, df_geschlecht_)
#    schlagworte_finden(prs, df_unternehmen, value="unternehmen", sort= "LW")
#    
   # Wirtschaft STN
   #  df_ = df[(df["Ressortbeschreibung"]=="Wirtschaft") & (df["Zeitung"]=="STN")]
#    df_ = df_stn[df_stn["Ressortbeschreibung"]=="Wirtschaft"]
#    df_geschlecht_ = df_geschlecht[(df_geschlecht["Ressortbeschreibung"]=="Wirtschaft") & (df_geschlecht["Zeitung"]=="STN")]
#    mantel_auswertung(prs, df_, df_geschlecht_)
#    
#    
#    # Entdecken STZ
#   # df_ = df[(df["Ressortbeschreibung"]=="Entdecken") & (df["Zeitung"]=="STZ")]
#    df_ = df_stz[df_stz["Ressortbeschreibung"]=="Entdecken"]
#    df_geschlecht_ = df_geschlecht[(df_geschlecht["Ressortbeschreibung"]=="Entdecken") & (df_geschlecht["Zeitung"]=="STZ")]
#    mantel_auswertung(prs, df_, df_geschlecht_)
#    
#    
#    # Wissenswert STN
#   # df_ = df[(df["Ressortbeschreibung"]=="Wissenswert") & (df["Zeitung"]=="STN")]
#    df_ = df_stn[df_stn["Ressortbeschreibung"]=="Wissenswert"]
#    df_geschlecht_ = df_geschlecht[(df_geschlecht["Ressortbeschreibung"]=="Wissenswert") & (df_geschlecht["Zeitung"]=="STN")]
#    mantel_auswertung(prs, df_, df_geschlecht_)
#    
     #Kultur STZ 
    #df_ = df[(df["Ressortbeschreibung"]=="Kultur") & (df["Zeitung"]=="STZ")]
#    df_ = df_stz[df_stz["Ressortbeschreibung"]=="Kultur"]
#    df_geschlecht_ = df_geschlecht[(df_geschlecht["Ressortbeschreibung"]=="Kultur") & (df_geschlecht["Zeitung"]=="STZ")]
#    mantel_auswertung(prs, df_, df_geschlecht_)
#    grafik_lesewert(prs, df_, target="seitentitel", minimize=5, label_position="large",
#                    ressort_liste=[], special=False, order="ok", legend="normal", sort="Seitennummer", 
#                    article="mean", mean_line = 0, grid=False, anzahl_lokales=1, axis=2, 
#                    title_text="Kultur nach Seitentiteln STZ") 
    # Kultur STN 
    #df_ = df[(df["Ressortbeschreibung"]=="Kultur") & (df["Zeitung"]=="STN")]
#    df_ = df_stn[df_stn["Ressortbeschreibung"]=="Kultur"]
#    df_geschlecht_ = df_geschlecht[(df_geschlecht["Ressortbeschreibung"]=="Kultur") & (df_geschlecht["Zeitung"]=="STN")]
#    mantel_auswertung(prs, df_, df_geschlecht_)
#    grafik_lesewert(prs, df_, target="seitentitel", minimize=5, label_position="large",
#                    ressort_liste=[], special=False, order="ok", legend="normal", sort="Seitennummer", 
#                    article="mean", mean_line = 0, grid=False, anzahl_lokales=1, axis=2, 
#                    title_text="Kultur nach Seitentiteln STN") 
    # Sport STZ 
    #df_ = df[(df["Ressortbeschreibung"]=="Sport") & (df["Zeitung"]=="STZ")]
#    df_ = df_stz[df_stz["Ressortbeschreibung"]=="Sport"]
##    df_geschlecht_ = df_geschlecht[(df_geschlecht["Ressortbeschreibung"]=="Sport") & (df_geschlecht["Zeitung"]=="STZ")]
##    mantel_auswertung(prs, df_, df_geschlecht_)
#    schlagworte_finden(prs, df_, value="vereine", sort= "LW")
#    schlagworte_finden(prs, df_, value="sportart", sort= "LW")
#    # Sport STN 
#    #df_ = df[(df["Ressortbeschreibung"]=="Sport") & (df["Zeitung"]=="STN")]
#    df_ = df_stn[df_stn["Ressortbeschreibung"]=="Sport"]
##    df_geschlecht_ = df_geschlecht[(df_geschlecht["Ressortbeschreibung"]=="Sport") & (df_geschlecht["Zeitung"]=="STN")]
##    mantel_auswertung(prs, df_, df_geschlecht_)
#    schlagworte_finden(prs, df_, value="vereine", sort= "LW")
#    schlagworte_finden(prs, df_, value="sportart", sort= "LW")
##    grafik_lesewert(prs, df_, target="seitentitel", minimize=5, label_position="large",
##                    ressort_liste=[], special=False, order="ok", legend="normal", sort="Seitennummer", 
##                    article="mean", mean_line = 0, grid=False, anzahl_lokales=1, axis=2, 
##                    title_text="Sport nach Seitentiteln STN")  
##    #Multimed. Reportage STZ
#    df_ = df[(df["Ressortbeschreibung"]=="Multimed. Reportage")]
#    df_geschlecht_ = df_geschlecht[(df_geschlecht["Ressortbeschreibung"]=="Multimed. Reportage")]
#    mantel_auswertung(prs, df_, df_geschlecht_)
#    
#     #Die Brücke STZ
#    df_ = df[(df["Ressortbeschreibung"]=="Die Brücke")]
#    df_geschlecht_ = df_geschlecht[(df_geschlecht["Ressortbeschreibung"]=="Die Brücke")]
#    mantel_auswertung(prs, df_, df_geschlecht_)
#    
#    # Leserforum STZ
#    #df_ = df[(df["Ressortbeschreibung"]=="Leserforum")]
#    df_ = df_stz[df_stz["Ressortbeschreibung"]=="Leserforum"]
#    df_geschlecht_ = df_geschlecht[(df_geschlecht["Ressortbeschreibung"]=="Leserforum")]
#    mantel_auswertung(prs, df_, df_geschlecht_)
#    
#    # Unsere Leser und wir STN
#    #df_ = df[(df["Ressortbeschreibung"]=="Unsere Leser und wir")]
#    df_ = df_stn[df_stn["Ressortbeschreibung"]=="Unsere Leser und wir"]
#    df_geschlecht_ = df_geschlecht[(df_geschlecht["Ressortbeschreibung"]=="Unsere Leser und wir")]
#    mantel_auswertung(prs, df_, df_geschlecht_)
#    
#    
#    # Wochenende
#    df_ = df[(df["Ressortbeschreibung"]=="Wochenende")]
#    df_geschlecht_ = df_geschlecht[(df_geschlecht["Ressortbeschreibung"]=="Wochenende")]
#    mantel_auswertung(prs, df_, df_geschlecht_)
#    
    
    
    # 13 Übersicht Ressorts inklusive 1. Lokseite S_STZ, S_STN, 5. Bücher
        # und Durchschnittslinie
        
    # 14 Ressorts nach Titeln und Geschlecht
    
    # 15 LW nach Tagen im Verhältnis zum Angebot nach Titeln 
        #lesewert_erscheinung(df_, prs, title_text="Lesewert nach Wochentagen " + zt)
        
    
    # 16 LW nach Darstellungsformen im Verhältnis zum Angebot
       
        #darstellungsformen(prs, df_, minimum = 5)
    # 17 LW nach Darstellungsform und Geschlecht nach Titeln
        #darstellungsformen(prs, df_geschlecht, minimum=5, geschlecht=True, 
             #              title_text = "Darstellungsf. nach Geschlecht " + zt)
    # 18 Übersicht der wichtigsten 5 bis 10 Kolumnen/ Kommentarformen nach Titel
    # wichtige Kolumnen sind: 
    #kommentare_tabelle(prs, df, version="STZ")  
    
    
    # AUSWERTUNG LOKALES
#    
#    for elem in lokale_ressorts:
#        df_ = df[df["Ressortbeschreibung"]==elem]
#        grafik_lesewert(prs, df_, target="ressort", minimize=5, label_position="xlarge",
#                    ressort_liste=ressort_list, special=False, title_text="Ressorts " + stutt_dict[elem], order="ok",
#                    legend="normal", sort="Seitennummer", article="mean", mean_line=mean_line)
#        # Deckblatt Lokalteil
        #TODO: 
        #deckblatt_abschluss(prs, df_)    
        
        # Übersicht LW 
        #for el in werte: 
            #grafik_entwicklung(prs, df_, target=el, mean_line=0, legend="large", grid=True, 
             #                  title_text = el + " " + elem)
            
        #darstellungsformen(prs, df_, minimum = 5, geschlecht=False, title_text="Daratellungsformen " + elem)        
       # platzierungen(prs, df_, minimum = 5, geschlecht=False, title_text="Platzierungen " + elem) 
        #grafik_lesewert(prs, df, target="seitentitel", minimize=5, label_position="normal",
            #        ressort_liste=[], special=False, title_text="", order="ok", legend="normal", sort="Seitennummer", 
            #        article="total")
       # lesewert_erscheinung(df_, prs, title_text="LW/Tage " + elem)
        
        # Muss-kann-soll für einzelne Ressorts
        #mks(df_ereignis, df_)
        
        
        # Ankündigungen vs andere Meldungen
        #lesewert_ankündigungen(prs, df_)
        
        # Top 10 für jeden Lokalteil
#        top_10(prs, df, df_berechnet = False, screenshots=True, number_screenshots = 5, 
#           mode="ressort", headline="default", zeitung=True)
#        
    # weitere Charts STZ
    # Top 5 Kolumnen nach Lesewert
#    df_stz = df[df["Zeitung"]=="STZ"]
#    top_kolumnen(prs, df_stz)
    
    # Vergleich Lokales STZ/STN, Tabellen nebeneinander
        

        
        
    print("Erstelle Powerpoint-Präsentation")
    prs.save("mhs_abschluss.pptx")
    
        
    #ZEITMESSUNG Ende
    elapsed_time = time.process_time() - t
    print("Dauer der Analyse: "+str(elapsed_time)+" sec.")
    
    
    
#%% Schwäbische Zeitung
    
def schwaebische(df, df_double_split = 0, df_nk=0, df_scans=0, df_umfrage=0,
                       df_geschlecht = 0, liste = ausgaben_liste, df_ereignisse=0, 
                       df_geschlecht_nk=0):
    t_0=time.time()
    print("Schwaebische gestartet")
#ALLGEMEINES
    #prs-Objekt erstellen
    prs = Presentation("LW_neu_2.pptx")
    
#    # 1 - Entwicklungsfolie Lesewert/DW/BW im Messverlauf - CHECK
#    print("Starte Auswertung Entwicklung über Messzeitraum...")
##   
#    for elem in werte:
#        grafik_entwicklung(prs, df, target=elem, mean_line=0, legend="large", 
#                           grid=True, title_text = False, ma=False, 
#                           steps_x_label = 2, limit_y=True)
##   
##   
##    
##    
##    # 2 - Fragebogen - CHECK
#    print ("Initialisiere Auswertung der Fragebögen")
#    # Aufruf jeder einzelnen Ausgabe
#    for elem in ausgaben_liste:
#        
#        df_ = df_umfrage[df_umfrage["ZTG"]==elem]
#       
#        #Aufruf gestapelte Bars "Ihre Zeitung ist..."
#        zeitungsattribute_berechnung(prs, df_, title = elem) # elem wird später über ausgaben_dict umgewandelt
#        
#        #Aufruf Bars "Welche Themen
#        zeitung_themen(prs, df_, title=elem)
#        
#        # Aufruf kleine Pie-Charts
#        umfrage_pie(prs, df_, title=elem)
#        
#        # Aufruf kleine Barcharts, TV- und Internetnutzung
#        
#        mini_bars(prs, df_, title_text = "Fernsehnutzung", liste=umfrage_mini_bars_TV, title=elem)
#        mini_bars(prs, df_, title_text = "Internetnutzung", liste=umfrage_mini_bars_netz, title=elem)
#        mini_bars(prs, df_, title_text = "Persönliches", liste=umfrage_demografie_bars, title=elem)
#        
##    # 3 - Kennzahlen der Messung - CHECK
#    
#    kennzahlen(df, df_scans)
#    
#        
#        
#    # 4 - TEXTMARKEN - Analyse LW, BW, DW Gesamt, Ausgaben, Lokales, Mantel - CHECK
#    #print("-------------------------------------------")
#    print("WERTE FÜR LW-MARKEN:")
#    marken_analyse(df, df_doublesplitid=df_double_split, df_nichtkum=df_nk)
#    
#
##     5 -Screenshots Gesamt + Mantel Top 5 - CHECK
#    print("Top 10 Gesamt werden ausgespielt...")
#    ausgaben = ["BIB", "FHA", "RV"]
#    for elem in ausgaben: 
#        df_ = df_nk[df_nk["ZTG"]==elem]
#        top_10(prs, df_, df_berechnet = False, screenshots=True, number_screenshots = 5, 
#           mode="Gesamt", headline="Top 10 Gesamt "+ ausgaben_dict[elem], zeitung=False, lokales_ressort="Lokales")
##    # Top 10 Mantelressorts CHECK
#    print("Top 10 Mantel werden ausgespielt...")
#    df_mantel = df[df["Ressortbeschreibung"].isin(mantel_ressorts)]
#    top_10(prs, df_mantel, df_berechnet = False, screenshots=True, number_screenshots = 5, 
#           mode="Mantel", headline="Top 10 Mantelressorts", zeitung=False)
#
#   
# # 6 - Leseraktivität nach Tageszeit und Wochentag (aus Scandaten) - CHECK
##    TODO Werte für alle vier Wellen ausspielen, dafür muss ich Scan-Daten mit Asugabedaten verknüpfen... 
##     ... muss also wissen, woher der einzelne Scan herkommt. 
#    print("Analyse Lesetage und Lesezeiten gestartet...." )
#    lesetage(prs, df_scans, title_text="Lesetage - Gesamt", multi_line = True)
#    lesezeit(prs, df_scans, title_text="Lesetage - Gesamt")
#    for elem in ausgaben_liste:
#        df_ = df_scans[df_scans["ZTG"]==elem]
#        lesetage(prs, df_, title_text="Lesetage - " + ausgaben_dict[elem]) 
#        lesezeit(prs, df_, title_text="Lesezeiten - " + ausgaben_dict[elem])
## 7 - Ressorts nach Lesewert und Häufigkeit CHECK
#    for elem in ausgaben: 
#        df_ = df_nk[df_nk["ZTG"]==elem]
#        print(elem)
#        mean_line = df_["Artikel-Lesewert (Erscheinung) in %"].mean()
#       
#        grafik_lesewert(prs, df_, target="ressort", minimize=5, label_position="xlarge",
#                    ressort_liste=ressort_list, special=False, title_text="Ressorts " + ausgaben_dict[elem],
#                    legend="normal", sort="Seitennummer", article="mean", mean_line=mean_line, axis=2)
#    
#   # 8 - Ressortnutzung nach Geschlecht CHECK
#    for elem in ausgaben: 
#       df_ = df_geschlecht_nk[df_geschlecht_nk["ZTG"]==elem]
#       print(df_.shape)
#       #df_ = df_.drop_duplicates(subset="SplitId")
#       multiple_bars_geschlecht(prs, df_, target="Ressorts", grid = True, legend="xlarge", 
#                             ressort_liste = ressort_list, title_text= "Ressorts nach Geschlecht - "+elem, 
#                             sort="Seitennummer")
#    # 9 Lesewert nach Tagen  CHECK
#    lesewert_erscheinung(df, df_double_split, prs, title_text="Lesewert nach Wochentagen Gesamt")
#     
#      Darstellungsformen LW, Menge und GEschlecht CHECK
#    darstellungsformen(prs, df, minimum = 5, geschlecht=False, title_text="Darstellungsformen Gesamt")
#    darstellungsformen(prs, df_geschlecht, minimum = 5, geschlecht=True, title_text="Darstellungsformen nach Geschlecht")
    
    # 10 Lesewerte nach Sonderseiten CHECK
#    liste_sonderseiten = ["Garten", "Familie", "Geld & Service", "Wissen",
#                          "Ernährung", "Bauen & Wohnen", "Sternenhimmel", 
#                          "Auto & Verkehr", "Mode", "Reise & Erholung", 
#                          "Medien", "Multimedia", "Literatur", "Kino", "Tiere",
#                          "Szene", "Gesundheit"]
#    grafik_lesewert(prs, df, target="seitentitel", minimize=5, label_position="xlarge",
#                    ressort_liste=liste_sonderseiten, special=False, title_text="Sonderseiten", order="ok",
#                    legend="normal", sort="Lesewert", article="mean", axis=2)
#    
    # 11 Analyse Mantelteil
#    analyse = ['Titelseite', 'Wir im Süden', 'Seite Drei', 
#       'Nachrichten & Hintergrund', 'Meinung & Dialog', 'Wirtschaft',
#       'Panorama/Journal', 'Kultur', 'Sport','Wochenende']
  #  analyse=["Sport"]
    
        
    
#    
#    for elem in analyse:
#        print("-----------------------")
#        elem_groß = elem.upper()
#        print("---------------- " + elem_groß+" ------------------------")
#        df_ = df[df["Ressortbeschreibung"]==elem]
#        df_geschlecht_ = df_geschlecht[df_geschlecht["Ressortbeschreibung"]==elem]
        
        #mantel_auswertung(prs, df_, df_geschlecht_, kunde="SWZ", lokales =False)
#        
#        if elem == "Wir im Süden":  #CHECK
#            df_ereignisse_ = df_ereignisse[df_ereignisse["Ressortbeschreibung"]==elem] 
#            mks(df_ereignisse_, df_, liste=[])
#            
#        if elem=="Kultur": #CHECK
#            seiten = ["Literatur", "Kino", "Szene", "Fernsehen & Freizeit"]
#            grafik_lesewert(prs, df, target="seitentitel", minimize=5, label_position="xlarge",
#                    ressort_liste=seiten, special=False, title_text="Sonderseiten mit Kulturbezug", order="ok",
#                    legend="normal", sort="Lesewert", article="mean", axis=2)
#          
#        if elem == "Sport": #CHECK
#            schlagworte_finden(prs, df_, value="sportart", sort= "LW", min_artikel = 5)
            #darstellungsformen(prs, df, minimum = 5, geschlecht=False, title_text="Darstellungsformen Gesamt")
    
#        if elem == "Wochenende": 
#            seiten = ['Wochenende', 'Menschen', 'Lebensart', 'Szene am Wochenende',
#                      'Meine Seite']
#            grafik_lesewert(prs, df_, target="seitentitel", minimize=5, label_position="xlarge",
#                    ressort_liste=seiten, special=False, title_text="Seitentitel Wochenende", order="ok",
#                    legend="normal", sort="Seitennummer", article="mean", axis=2)
#            for el in werte:
#                grafik_entwicklung(prs, df_, target=el, mean_line=0, legend="large", grid=True, title_text = el + " Wochenende")
    
    
    # 12 LOKALTEIL AUSWERTUNG
    
#    
    for elem in ausgaben:
        titletext = ausgaben_dict[elem]
        kunde = "SWZ"
        lokales = "Lokales"
        df_ = df[(df["Ressortbeschreibung"]=="Lokales") & (df["ZTG"]==elem)]
        df_geschlecht_ = df_geschlecht[(df_geschlecht["Ressortbeschreibung"]=="Lokales") & (df_geschlecht["ZTG"]==elem)]
        df_ereignisse_ = df_ereignisse[(df_ereignisse["Ressortbeschreibung"]=="Lokales") & (df_ereignisse["ZTG"]==elem)]
        deckblatt_abschluss(prs, df_, 
                            title_text="Lokales " +ausgaben_dict[elem],
                            lokales=False, kunde="SWZ")
#       
        for el in werte:
            grafik_entwicklung(prs, df_, target=el, mean_line=0, legend="large", 
                               grid=True, title_text = el + " " + "Lokales " + elem, 
                               steps_x_label=2, limit_y = True)
#       
        if elem == "BIB": 
            seiten_liste = ['Biberach', 'Biberach / Service', 'Rund um Biberach',
                            'Von der Schussen zur Rot', 'Von der Rottum zur Iller',
                            'Wir in Kreis und Region', 'Veranstaltungen', 'Kirchen', 'Lokale Eins BIB']
            lok_1 = "Lokale Eins BIB"
            buch_4 = "Wir in Kreis und Region"
        
        elif elem == "FHA":
            seiten_liste = ['Friedrichshafen', 'Kultur Lokal / Service',
                            'Langenargen / Eriskirch / Kressbronn',
                            'Immenstaad / Oberteuringen / Meckenbeuren', 'Wir am See',
                            'Markdorf und der westliche Bodenseekreis','Veranstaltungen',
                            'Friedrichshafen / Service', 'Region', "Lokale Eins FHA"]
            lok_1 = "Lokale Eins FHA"
            buch_4 = "Wir am See"
            
        elif elem == "RV": 
            seiten_liste = ['Ravensburg / Weingarten', 'Ravensburg', 'Weingarten',
                            'Oberschwaben & Allgäu', 'Gemeinden', 'Veranstaltungen',
                            'Meine Heimat. Mein Verein.', "Lokale Eins RV"]
            lok_1 = "Lokale Eins RV"
            buch_4 = "Gemeinden"
#        
        grafik_lesewert(prs, df_, target="special", minimize=5, label_position="xlarge",
                    ressort_liste=seiten_liste, special=False, title_text="Seitentitel " + ausgaben_dict[elem], order="ok",
                    legend="normal", sort="Seitennummer", article="mean", axis=2)    
#            
#        # Lesewert nach Wochentagen plus Menge der Artikel  
#        lesewert_erscheinung(df_,df_, prs, title_text="Lesewert nach Wochentagen, " + elem)
#        
#        
#        # meldung vs ankündiger
#        lesewert_ankündigungen(prs, df_)
#        
#        # Darstellungsformen und Platzierungen
#        tabelle_ressortauswertung(prs, df_, research_object="Darstellungsformen", sort="Lesewert", 
#                                  title_text = "Darstellungsformen " + titletext, minimum = 5, kunde=kunde, lokales=lokales)
#        darstellungsformen(prs, df_, minimum = 5, geschlecht=False, title_text="Darstellungsformen " + titletext, grid=False)
#        darstellungsformen(prs, df_geschlecht_, minimum = 5, geschlecht=True, title_text="Darstellungsformen " + titletext)
#        tabelle_ressortauswertung(prs, df_, research_object="Platzierungen", sort="Lesewert", 
#                                  title_text = "Platzierungen " + titletext, minimum = 5, kunde=kunde, lokales=lokales)
#        platzierungen(prs, df_, minimum = 5, geschlecht=False, title_text="Platzierungen " + titletext)
#        platzierungen(prs, df_geschlecht_, minimum = 5, geschlecht=True, title_text="Platzierungen " + titletext)
#        
#        
#        
#        # Top 10 
#        top_10(prs, df_, df_berechnet = False, screenshots=True, number_screenshots = 5, 
#           mode="ressort", headline="Top 10 - Lokales "+ausgaben_dict[elem], zeitung=False)
#        top_10(prs, df_geschlecht_, df_berechnet = False, screenshots=False, number_screenshots = 5, 
#           mode="ressort", headline="Top 10 - "+ausgaben_dict[elem] + " Frauen", zeitung=False, geschlecht="weiblich")
#        top_10(prs, df_geschlecht_, df_berechnet = False, screenshots=False, number_screenshots = 5, 
#           mode="ressort", headline="Top 10 - "+ausgaben_dict[elem] + " Männer", zeitung=False, geschlecht="männlich")
#        top_10(prs, df_geschlecht_, df_berechnet = False, screenshots=False, number_screenshots = 5, 
#           mode="ressort", headline="Top 10 - "+ausgaben_dict[elem] + " 40 bis  59 Jahre", zeitung=False, alter="40 bis 59 Jahre")        
#        top_10(prs, df_geschlecht_, df_berechnet = False, screenshots=False, number_screenshots = 5, 
#           mode="ressort", headline="Top 10 - "+ausgaben_dict[elem] + " 60 bis  79 Jahre", zeitung=False, alter="60 bis 79 Jahre")
#        
##     # Übersicht Lok-1
#        lokeins = df_[df_["Spezial"]==lok_1]
#        lokeinsgesch = df_geschlecht_[df_geschlecht_["Spezial"]==lok_1]
#        deckblatt_abschluss(prs, lokeins, 
#                            title_text="Lokale Eins " +ausgaben_dict[elem],
#                            lokales=False, kunde="SWZ")
#        tabelle_ressortauswertung(prs, lokeins, research_object="Darstellungsformen", sort="Lesewert", 
#                                  title_text = "Darstellungsformen Lokale Eins "+elem, minimum = 5, kunde=kunde, lokales=lokales)
#        darstellungsformen(prs, lokeins, minimum = 5, geschlecht=False, title_text="Darstellungsformen Lokale Eins "+elem, grid=False)
#        darstellungsformen(prs, lokeinsgesch, minimum = 5, geschlecht=True, title_text="Darstellungsformen " + titletext)
#        tabelle_ressortauswertung(prs, lokeins, research_object="Platzierungen", sort="Lesewert", 
#                                  title_text = "Platzierungen Lokale Eins " + elem, minimum = 5, kunde=kunde, lokales=lokales)
#        platzierungen(prs, lokeins, minimum = 5, geschlecht=False, title_text="Platzierungen Lokale Eins " + elem)
#        platzierungen(prs, lokeinsgesch, minimum = 5, geschlecht=True, title_text="Platzierungen Lokale Eins " + elem)
#        
#        
#        # Übersicht Buch 4. 
#        buchvier = df_[df_["Spezial"]==buch_4]
#        buchviergesch = df_geschlecht_[df_geschlecht_["Spezial"]==buch_4]
#        titletext = "Aufschlag 4. Buch " + elem
#        deckblatt_abschluss(prs, buchvier, 
#                            title_text="Viertes Buch " +ausgaben_dict[elem],
#                            lokales=False, kunde="SWZ")
#        tabelle_ressortauswertung(prs, buchvier, research_object="Darstellungsformen", sort="Lesewert", 
#                                  title_text = "Darstellungsformen " + titletext, minimum = 5, kunde=kunde, lokales=lokales)
#        darstellungsformen(prs, buchvier, minimum = 5, geschlecht=False, title_text="Darstellungsformen " + titletext, grid=False)
#        darstellungsformen(prs, buchviergesch, minimum = 5, geschlecht=True, title_text="Darstellungsformen " + titletext)
#        tabelle_ressortauswertung(prs, buchvier, research_object="Platzierungen", sort="Lesewert", 
#                                  title_text = "Platzierungen " + titletext, minimum = 5, kunde=kunde, lokales=lokales)
#        platzierungen(prs, buchvier, minimum = 5, geschlecht=False, title_text="Platzierungen " + titletext)
#        platzierungen(prs, buchviergesch, minimum = 5, geschlecht=True, title_text="Platzierungen " + titletext)
##       
#        
#        mks(df_ereignisse_, df_, liste=[])
#        
#        # Übersicht Lokalsport
#        df_lokalsport = df[(df["ZTG"]==elem) & (df["Ressortbeschreibung"]=="Lokalsport")]
#        df_geschlecht_loksport = df_geschlecht[(df_geschlecht["Ressortbeschreibung"]=="Lokalsport") & (df_geschlecht["ZTG"]==elem)]
#        df_ereignisse_loksport = df_ereignisse[(df_ereignisse["Ressortbeschreibung"]=="Lokalsport") & (df_ereignisse["ZTG"]==elem)]
#        
#        # Deckblatt Loksport
#        deckblatt_abschluss(prs, df_lokalsport, 
#                            title_text="Lokalsport " +ausgaben_dict[elem],
#                            lokales=False, kunde="SWZ")
#        
#        for el in werte:
#            grafik_entwicklung(prs, df_lokalsport, target=el, mean_line=0, legend="large", 
#                               grid=True, title_text = el + " " + "Lokales " + elem, 
#                               steps_x_label=2, limit_y = True)
#        
#        schlagworte_finden(prs, df_lokalsport, value="sportart", sort= "LW", 
#                           min_artikel = 5, title = "Sportarten Lokalsport " + elem)
#        
#        titletext = "Lokalsport " + elem
#        tabelle_ressortauswertung(prs, df_lokalsport, research_object="Darstellungsformen", sort="Lesewert", 
#                                  title_text = "Darstellungsformen " + titletext, minimum = 5, kunde=kunde, lokales=lokales)
#        darstellungsformen(prs, df_lokalsport, minimum = 5, geschlecht=False, title_text="Darstellungsformen " + titletext, grid=False)
#        darstellungsformen(prs, df_geschlecht_loksport, minimum = 5, geschlecht=True, title_text="Darstellungsformen " + titletext)
#        tabelle_ressortauswertung(prs, df_lokalsport, research_object="Platzierungen", sort="Lesewert", 
#                                  title_text = "Platzierungen " + titletext, minimum = 5, kunde=kunde, lokales=lokales)
#        platzierungen(prs, df_lokalsport, minimum = 5, geschlecht=False, title_text="Platzierungen " + titletext)
#        platzierungen(prs, df_geschlecht_loksport, minimum = 5, geschlecht=True, title_text="Platzierungen " + titletext)
#        
#        
#         # Top 10 
#        top_10(prs, df_lokalsport, df_berechnet = False, screenshots=True, number_screenshots = 5, 
#           mode="ressort", headline="Top 10 - Lokalsport "+ausgaben_dict[elem], zeitung=False)
#        top_10(prs, df_geschlecht_loksport, df_berechnet = False, screenshots=False, number_screenshots = 5, 
#           mode="ressort", headline="Top 10 - "+ausgaben_dict[elem] + " Frauen", zeitung=False, geschlecht="weiblich")
#        top_10(prs, df_geschlecht_loksport, df_berechnet = False, screenshots=False, number_screenshots = 5, 
#           mode="ressort", headline="Top 10 - "+ausgaben_dict[elem] + " Männer", zeitung=False, geschlecht="männlich")
#        top_10(prs, df_geschlecht_loksport, df_berechnet = False, screenshots=False, number_screenshots = 5, 
#           mode="ressort", headline="Top 10 - "+ausgaben_dict[elem] + " 40 bis  59 Jahre", zeitung=False, alter="40 bis 59 Jahre")        
#        top_10(prs, df_geschlecht_loksport, df_berechnet = False, screenshots=False, number_screenshots = 5, 
#           mode="ressort", headline="Top 10 - "+ausgaben_dict[elem] + " 60 bis  79 Jahre", zeitung=False, alter="60 bis 79 Jahre")
#        # 14 - Entwicklung LW, BW, DW  CHECK
#        
        
        # Entwicklung AA und SK - also nur die großen Stücke
#        df_großetexte = df[(df["Platzierungen"]=="AA") | (df["Platzierungen"]=="SK")]
#        for el in werte: 
#            grafik_entwicklung(prs, df_großetexte, target=el, mean_line=0, legend="large", grid=True, title_text = el + " AA und SK", steps_x_label=2)
#        
#        top_10(prs, df, df_berechnet = False, screenshots=True, number_screenshots = 5, 
#           mode="ressort", headline="Top 10 - "+titletext, zeitung=False)
#        top_10(prs, df_geschlecht, df_berechnet = False, screenshots=False, number_screenshots = 5, 
#           mode="ressort", headline="Top 10 - "+titletext, zeitung=False, geschlecht="weiblich")
#        top_10(prs, df_geschlecht, df_berechnet = False, screenshots=False, number_screenshots = 5, 
#           mode="ressort", headline="Top 10 - "+titletext, zeitung=False, geschlecht="männlich")
#        top_10(prs, df_geschlecht, df_berechnet = False, screenshots=False, number_screenshots = 5, 
#           mode="ressort", headline="Top 10 - "+titletext, zeitung=False, alter="40 bis 59 Jahre")        
#        top_10(prs, df_geschlecht, df_berechnet = False, screenshots=False, number_screenshots = 5, 
#           mode="ressort", headline="Top 10 - "+titletext, zeitung=False, alter="60 bis 79 Jahre")
##     
    
    
    
    print("-----------------------------------")
    print("Spiele neues PPTX-Dokument aus")
    prs.save("schw_abschluss.pptx")
    print("-----------------------------------")
    print("Analyse abgeschlossen")
    print("-----------------------------------")
    t_1 = time.time()
    endtime = t_1 - t_0 
    print("Dauer des Durchlaufs: {0:.2f}.".format(endtime))


#%% Mitteldeutsche Zeitung

def mitteldeutsche(df, df_komplett, df_double_split = 0, df_nk=0, df_scans=0, df_umfrage=0,
                       df_geschlecht = 0, liste = ausgaben_liste, df_ereignisse=0, 
                       df_geschlecht_nk=0):
    print("Mitteldeutsche läuft")
       
    t_0=time.time()
    prs = Presentation("LW_neu_2.pptx")
    
#    # 1 - Entwicklungsfolie Lesewert/DW/BW im Messverlauf - CHECK
#    print("Starte Auswertung Entwicklung über Messzeitraum...")
##   
#    for elem in werte:
#        grafik_entwicklung(prs, df, target=elem, mean_line=0, legend="large", 
#                           grid=True, title_text = False, ma=False, 
#                           steps_x_label = 2, limit_y=True)
##   
# 
#    for elem in werte:
#        grafik_entwicklung(prs, df, target=elem, mean_line=0, legend="large", 
#                           grid=True, title_text = False, ma=False, 
#                           steps_x_label = 2, limit_y=True)
#   

##   
##    
##    
##  # 2 - Fragebogen - CHECK
#    print ("Initialisiere Auswertung der Fragebögen")
#    # Aufruf jeder einzelnen Ausgabe
#    for elem in ausgaben_liste:
#        
#        df_ = df_umfrage[df_umfrage["ZTG"]==elem]
#       
#        #Aufruf gestapelte Bars "Ihre Zeitung ist..."
#        zeitungsattribute_berechnung(prs, df_, title = elem) # elem wird später über ausgaben_dict umgewandelt
#        
#        #Aufruf Bars "Welche Themen
#        zeitung_themen(prs, df_, title=elem)
#        
#        # Aufruf kleine Pie-Charts
#        umfrage_pie(prs, df_, title=elem)
#        
#        # Aufruf kleine Barcharts, TV- und Internetnutzung
#        
#        mini_bars(prs, df_, title_text = "Fernsehnutzung", liste=umfrage_mini_bars_TV, title=elem)
#        mini_bars(prs, df_, title_text = "Internetnutzung", liste=umfrage_mini_bars_netz, title=elem)
#        mini_bars(prs, df_, title_text = "Persönliches", liste=umfrage_demografie_bars, title=elem)
#   
        
        
        # LESERAKTIVITÄT
    lesetage(prs, df_scans, title_text="Lesetage - Gesamt", multi_line = True)
    lesezeit(prs, df_scans, title_text="Lesezeiten - Gesamt")
    for elem in ausgaben_liste:
       df_analyse_ = df_komplett[df_komplett["ZTG"]==elem]
       grafik_lesewert(prs, df_analyse_, target="Textlänge", minimize=1, label_position="large",
                    ressort_liste=textlänge, special=False, order="ok", legend="normal", sort="Kategorie", 
                    article="total", mean_line = 0, grid=False, anzahl_lokales=1, axis=2, 
                    title_text="Textlängen " + ausgaben_dict[elem])
    analyse_mantel = df_komplett[~(df_komplett["Ressortbeschreibung"]=="Lokales") | ~(df_komplett["Ressortbeschreibung"]=="Regionalsport")]
    grafik_lesewert(prs, analyse_mantel, target="Textlänge", minimize=1, label_position="large",
                    ressort_liste=textlänge, special=False, order="ok", legend="normal", sort="Kategorie", 
                    article="total", mean_line = 0, grid=False, anzahl_lokales=1, axis=2, 
                    title_text="Textlängen Mantelressorts")
    df_ = df_scans[df_scans["ZTG"]==elem]
    lesetage(prs, df_, title_text="Lesetage - " + ausgaben_dict[elem]) 
    lesezeit(prs, df_, title_text="Lesezeiten - " + ausgaben_dict[elem])
    
    
    # TOP 10 GESAMT (einzelne Ausgaben und Mantel)
    for elem in ausgaben_liste: 
        df_ = df_nk[df_nk["ZTG"]==elem]
        top_10(prs, df_, df_berechnet = False, screenshots=True, number_screenshots = 5, 
           mode="Gesamt", headline="Top 10 Gesamt "+ ausgaben_dict[elem], zeitung=False, lokales_ressort="Lokales")
##    # Top 10 Mantelressorts CHECK
#    print("Top 10 Mantel werden ausgespielt...")
    df_mantel = df[df["Ressortbeschreibung"].isin(mantel_ressorts)]
    top_10(prs, df_mantel, df_berechnet = False, screenshots=True, number_screenshots = 5, 
            mode="Mantel", headline="Top 10 Mantelressorts", zeitung=False)
#
    
    #RESSORT NACH LW + UMFANG (plus Durchschnittslinie)
  
    # 7 - Ressorts nach Lesewert und Häufigkeit CHECK
    sortierung_liste = ["Titelseite", "Sachsen-Anhalt", "Mitteldeutschland", 
                        "Politik", "Meinung", "Lokale Eins", "Lok 1 HAL", 
                        "Lok 1 JES", "Lok 1 BEB", "Lok 1 ZEI", "Lokales", 
                        "Sport", "Regionalsport", "Kinder", "Leserbriefe", 
                        "Wirtschaft", "Blick", "Kultur& Leben",  "Medien", 
                        "Film", "Ratgeber",  "Wissenschaft", "Panorama"]
    for elem in ausgaben_liste: 
        df_ = df_nk[df_nk["ZTG"]==elem]
        lok_1_df = df_[df_["Seitentitel"]=="Erste Lokalseite"]
        index_df = ["Ressortbeschreibung", "SplitId", "Seitennummer", 
                 "Artikel-Lesewert (Erscheinung) in %"]
        #name = "Lok 1 " + elem
        name = "Lokale Eins"
        splitId = lok_1_df.shape[0]
        seitennummer = lok_1_df.Seitennummer.median()
        lesewert = lok_1_df["Artikel-Lesewert (Erscheinung) in %"].mean()
        lok_1 = pd.Series([name, splitId, seitennummer, lesewert], index=index_df)
        
        
        print(elem)
        mean_line = df_["Artikel-Lesewert (Erscheinung) in %"].mean()
       
        grafik_lesewert(prs, df_, target="ressort", minimize=5, label_position="xlarge",
                    ressort_liste=sortierung_liste, special=False, title_text="Ressorts " + ausgaben_dict[elem],
                    legend="normal", sort="Kategorie", article="mean", mean_line=mean_line, axis=2, 
                    extension_df=lok_1, font_size=27)
 
    
    
    columns_df = ["Ressortbeschreibung", "SplitId", "Seitennummer", 
                 "Artikel-Lesewert (Erscheinung) in %"]
    lok_eins_df = pd.DataFrame(columns=columns_df)
    for elem in ausgaben_liste:
        df_lok1 = df_nk[(df_nk["Seitentitel"]=="Erste Lokalseite") & (df_nk["ZTG"]==elem)]
        
        name = "Lok 1 " + elem
        splitId = df_lok1.shape[0]
        seitennummer = df_lok1.Seitennummer.median()
        lesewert = df_lok1["Artikel-Lesewert (Erscheinung) in %"].mean()
        lok_1 = pd.Series([name, splitId, seitennummer, lesewert], index = columns_df)
        lok_eins_df = lok_eins_df.append(lok_1, ignore_index=True)  
    
    lok_eins_df["SplitId"] = lok_eins_df["SplitId"].astype(float)
    mean_line = df_komplett["Artikel-Lesewert (Erscheinung) in %"].mean()
    grafik_lesewert(prs, df_komplett, target="ressort", minimize=5, label_position="xlarge",
                    ressort_liste=sortierung_liste, special=False, title_text="Ressorts Gesamt",
                    legend="strange", sort="Kategorie", article="mean", mean_line=mean_line, axis=2, 
                    font_size=23, anzahl_lokales=4, extension_df=lok_eins_df)
#   # 8 - Ressortnutzung nach Geschlecht CHECK
    for elem in ausgaben_liste: 
       df_ = df_geschlecht_nk[df_geschlecht_nk["ZTG"]==elem]
       print(df_.shape)
       #df_ = df_.drop_duplicates(subset="SplitId")
       multiple_bars_geschlecht(prs, df_, target="Ressorts", grid = True, legend="xlarge", 
                             ressort_liste = sortierung_liste, title_text= "Ressorts nach Geschlecht - "+elem, 
                             sort="Kategorie", width = "special")
    multiple_bars_geschlecht(prs, df_geschlecht, target="Ressorts", grid = True, legend="xlarge", 
                             ressort_liste = sortierung_liste, title_text= "Ressorts nach Geschlecht GEsamt",  
                             sort="Kategorie", width = "special")
#    
    # 9 Lesewert nach Tagen  CHECK
    lesewert_erscheinung(df_komplett, df_double_split, prs, title_text="Lesewert nach Wochentagen Gesamt")
     
    #  Darstellungsformen LW, Menge und GEschlecht CHECK
    darstellungsformen(prs, df_komplett, minimum = 5, geschlecht=False, title_text="Darstellungsformen Gesamt")
    darstellungsformen(prs, df_geschlecht, minimum = 5, geschlecht=True, grid= True, title_text="Darstellungsformen nach Geschlecht")
    
    # 10 Lesewerte Kolumne und Kommentaren
    liste_kom = ["Aufgepasst", "Ansichtssache", "Empfehlung", "Mein Tipp", 
                 "Nachschlag"]
    df_ = df_komplett[df_komplett["Beschreibung"].isin(liste_kom)]
    lok_liste = ["Lokales", "Regionalsport"]
    df_kommi_lok = df_komplett[(df_komplett["Ressortbeschreibung"].isin(lok_liste)) & (df_komplett["Darstellungsformen"].str.contains("KM"))]
    df_kommi_mantel = df_komplett[(~df_komplett["Ressortbeschreibung"].isin(lok_liste)) & (df_komplett["Darstellungsformen"].str.contains("KM"))]
   
    besch_lok = "Kommentare Lokal"
    besch_mantel = "Kommentare Mantel"
    split_lok = df_kommi_lok.shape[0]
    split_mantel = df_kommi_mantel.shape[0]
    seite_lok = df_kommi_lok.Seitennummer.median()
    seite_mantel = df_kommi_mantel.Seitennummer.median()
    lw_lokal = df_kommi_lok["Artikel-Lesewert (Erscheinung) in %"].mean()
    lw_mantel = df_kommi_mantel["Artikel-Lesewert (Erscheinung) in %"].mean()
    col = ["Beschreibung", "SplitId", "Seitennummer", "Artikel-Lesewert (Erscheinung) in %"]
    extension_df = pd.DataFrame({"Beschreibung":[besch_lok, besch_mantel], 
                                 "SplitId":[split_lok, split_mantel], 
                                 "Seitennummer":[seite_lok, seite_mantel], 
                                 "Artikel-Lesewert (Erscheinung) in %":[lw_lokal, lw_mantel]})
            
          
    print(extension_df)
    grafik_lesewert(prs, df_, target="rubriken", minimize=1, label_position="xlarge", grid=True, 
                    ressort_liste=liste_kom, special=False, title_text="Sonderseiten", order="ok",
                    legend="normal", sort="Lesewert", article="mean", extension_df=extension_df)
    
    
    
    # ÜBERSICHT LOKALES

            
#    for elem in ausgaben_liste:
##        
#        print(elem)
#        print(df_nk.shape)
#        df_ = df_nk[(df_nk["Ressortbeschreibung"]=="Lokales") & (df_nk["ZTG"]==elem)]
#        df_analyse = df_komplett[(df_komplett["Ressortbeschreibung"]=="Lokales") & (df_komplett["ZTG"]==elem)]
#        print(df_.shape)
#        df_nk_ = df_nk[(df_nk["Ressortbeschreibung"]=="Lokales") & (df_nk["ZTG"]==elem)]
#        df_geschlecht_ = df_geschlecht_nk[(df_geschlecht_nk["Ressortbeschreibung"]=="Lokales") & (df_geschlecht_nk["ZTG"]==elem)]
#        deckblatt_abschluss(prs, df_, kunde="MZ", title_text = "Lokales " + ausgaben_dict[elem])
#        tabelle_ressortauswertung(prs, df_, research_object="Darstellungsformen", sort="Lesewert", 
#                                  title_text = "Darstellungsformen " + ausgaben_dict[elem], minimum = 5, kunde="MZ", lokales=True)
#        tabelle_ressortauswertung(prs, df_, research_object="Platzierungen", sort="Lesewert", 
#                                  title_text = "Platzierungen " + ausgaben_dict[elem], minimum = 5, kunde="MZ", lokales=True)
#        
#        platzierungen(prs, df_, minimum = 5, geschlecht=False, title_text="Platzierungen " + elem, special="")
#        platzierungen(prs, df_geschlecht_, minimum = 5, geschlecht=True, title_text="Platzierungen nach Geschlecht", special="")
#        darstellungsformen(prs, df_, minimum = 5, geschlecht=False, title_text="Lesewert Darstellungsformen", special="")
#        darstellungsformen(prs, df_geschlecht_, minimum = 5, geschlecht=True, title_text="Darstellungsformen nach Geschlecht", special="")
#        for el in werte:
#            grafik_entwicklung(prs, df_, target=el, mean_line=0, legend="large",
#                               grid=True, title_text = el + " " + ausgaben_dict[elem],
#                               steps=1,  limit_y=True, ma = False)
#            df__ = df_[(df_["Platzierungen"]=="AA")|(df_["Platzierungen"]=="SK")]
#            grafik_entwicklung(prs, df__, target=el, mean_line=0, legend="large",
#                               grid=True, title_text = el + " Aufmacher und Seitenkeller", steps=2)
#        
### SONDERWÜNSCHE LOKALES
#    
#        #Auswertung Seitentitel
#        grafik_lesewert(prs, df_, target="special", minimize=5, label_position="large",
#                    ressort_liste=[], special=False, order="ok", legend="normal", sort="Seitennummer", 
#                    article="total", mean_line = 0, grid=True, anzahl_lokales=1, axis=1, \
#                    title_text="Lesewert nach Seitentiteln " + ausgaben_dict[elem])
#        
#       # Auswertung nach Textlänge
#        grafik_lesewert(prs, df_analyse, target="Textlänge", minimize=1, label_position="large",
#                    ressort_liste=textlänge, special=False, order="ok", legend="normal", sort="Kategorie", 
#                    article="total", mean_line = 0, grid=False, anzahl_lokales=1, axis=2, 
#                    title_text="Textlängen " + ausgaben_dict[elem])
##        if elem == "HAL": 
##            res_li = ["]
#        grafik_lesewert(prs, df_, target="seitentitel", minimize=5, label_position="large",
#                    ressort_liste=[], special=False, order="ok", legend="normal", sort="Seitennummer", 
#                    article="total", mean_line = 0, grid=False, anzahl_lokales=1, axis=2, 
#                    title_text="Seitentitel " + elem)
#     
#        grafik_lesewert(prs, df_, target="special", minimize=5, label_position="large",
#                    ressort_liste=[], special=False, order="ok", legend="normal", sort="Seitennummer", 
#                    article="total", mean_line = 0, grid=False, anzahl_lokales=1, axis=2, 
#                    title_text="Seitentitel " + elem)
#        
#        lesewert_erscheinung(df_,df_nk_, prs, title_text="LW nach Wochentagen")
#        lesewert_ankündigungen(prs, df_nk_, title_text="Ankündigungen - " + ausgaben_dict[elem])
#        top_10(prs, df_, df_berechnet = False, screenshots=True, number_screenshots = 5, 
#           mode="ressort", headline="Top 10 - Lokales "+ausgaben_dict[elem], zeitung=True, geschlecht="", alter="", 
#           lokales_ressort="Lokales")
#        # Top 10 Geschlechter
#        top_10(prs, df_geschlecht_, df_berechnet = False, screenshots=False, number_screenshots = 5, 
#           mode="ressort", headline="default", zeitung=True, geschlecht="weiblich")
#        top_10(prs, df_geschlecht_, df_berechnet = False, screenshots=False, number_screenshots = 5, 
#           mode="ressort", headline="default", zeitung=True, geschlecht="männlich")
#        
# #        Top 10 Altersgruppen
#        top_10(prs, df_geschlecht_, df_berechnet = False, screenshots=False, number_screenshots = 5, 
#           mode="ressort", headline="default", zeitung=True, alter="40 bis 59 Jahre")
#        top_10(prs, df_geschlecht_, df_berechnet = False, screenshots=False, number_screenshots = 5, 
#           mode="ressort", headline="default", zeitung=True, alter="60 bis 79 Jahre")
#   
## 
##          Kommentare und Kolumnen
#        liste_kommis = ["Guten Morgen", "Moment Mal!", "Moment mal!", "Gedanken zum Tag", "Gedanken zum Sonntag", "Klatschkasten"]
#        df_liste = df_[df_["Beschreibung"].isin(liste_kommis)]
#         
#        df_kommentare = df_[df_["Darstellungsformen"].str.contains("KM")]
#         
#        besch_lok = "Kommentare gesamt"
#        split_lok = df_kommentare.shape[0]
#        seite_lok = df_kommentare.Seitennummer.median()
#        lw_lokal = df_kommentare["Artikel-Lesewert (Erscheinung) in %"].mean()
#        col = ["Beschreibung", "SplitId", "Seitennummer", "Artikel-Lesewert (Erscheinung) in %"]
#        extension_df = pd.DataFrame({"Beschreibung":[besch_lok], 
#                                 "SplitId":[split_lok], 
#                                 "Seitennummer":[seite_lok], 
#                                 "Artikel-Lesewert (Erscheinung) in %":[lw_lokal]})
#            
##          
#        print(extension_df)
#        grafik_lesewert(prs, df_liste, target="rubriken", minimize=1, label_position="xlarge",
#                    ressort_liste=liste_kommis, special=False, title_text="Kolumnen/Rubriken/Meinung " + elem, order="ok",
#                    legend="normal", sort="Lesewert", article="mean", extension_df=extension_df, grid=True)
    
    
        # Vergleich Montags/-Dienstagsausgaben & Moment mal -Auswertungen händisch erstellt
       
        # Entwicklung Bernburg nach Zeiträumen
#        
#        for el in werte:
#            df_zeitraum1 = df_[df_["Erscheinungsdatum"]<="2018-09-20"]
#            df_zeitraum2 = df_[df_["Erscheinungsdatum"]>"2018-09-20"]
#            grafik_entwicklung(prs, df_zeitraum1, target=el, mean_line=0, legend="large",
#                               grid=True, title_text = el + " " + elem + "Zeit1!",
#                               steps=1,  limit_y=True, ma = False)
#            grafik_entwicklung(prs, df_zeitraum2, target=el, mean_line=0, legend="large",
#                               grid=True, title_text = el + " " + elem + "Zeit2",
#                               steps=1,  limit_y=True, ma = False)
#           
            
        
        
        
        # 
        
 
   # AUSWERTUNG MANTEL
    
#    mantel_ressorts_auswertung = ['Titelseite','Sachsen-Anhalt', 'Mitteldeutschland', 'Politik', \
#                   'Meinung', 'Sport', 'Kinder', 'Wirtschaft', 'Kultur& Leben', \
#                   'Ratgeber', 'Wissenschaft', 'Panorama','Blick', 'Film']
#    
#    mantel_ressorts_auswertung_ = ['Mitteldeutschland']
#    
#    df_mantel = df[df["Ressortbeschreibung"].isin(mantel_ressorts_auswertung)]
#    df_geschlecht_mantel = df_geschlecht[df_geschlecht["Ressortbeschreibung"].isin(mantel_ressorts_auswertung)]
#    
#    for el in werte:
#        grafik_entwicklung(prs, df_mantel, target=el, mean_line=0, legend="large",
#                               grid=True, title_text = el + " Mantelressorts",
#                               steps=1,  limit_y=True, ma = False)
#        
#    deckblatt_abschluss(prs, df_mantel, kunde="MZ", title_text = "Mantel Gesamt")
#    tabelle_ressortauswertung(prs, df_mantel, research_object="Darstellungsformen", sort="Lesewert", 
#                                  title_text = "Darstellungsformen Mantel Gesamt", minimum = 5, kunde="MZ", lokales=False)
#    tabelle_ressortauswertung(prs, df_mantel, research_object="Platzierungen", sort="Lesewert", 
#                                  title_text = "Platzierungen Mantel Gesamt", minimum = 5, kunde="MZ", lokales=False)
#        
#    platzierungen(prs, df_mantel, minimum = 5, geschlecht=False, title_text="Platzierungen Mantelressorts", special="")
#    platzierungen(prs, df_geschlecht_mantel, minimum = 5, geschlecht=True, title_text="Platzierungen nach Geschlecht", special="")
#    darstellungsformen(prs, df_mantel, minimum = 5, geschlecht=False, title_text="Lesewert Darstellungsformen", special="")
#    darstellungsformen(prs, df_geschlecht_mantel, minimum = 5, geschlecht=True, title_text="Darstellungsformen nach Geschlecht", special="")
#    
#    top_10(prs, df_mantel, df_berechnet = False, screenshots=True, number_screenshots = 5, 
#           mode="Gesamt", headline="Top 10 - Mantelressorts", zeitung=True, geschlecht="", alter="", 
#           lokales_ressort="Lokales")
#        # Top 10 Geschlechter
#    top_10(prs, df_geschlecht_mantel, df_berechnet = False, screenshots=False, number_screenshots = 5, 
#           mode="Gesamt", zeitung=True, geschlecht="weiblich", 
#           headline="Top 10 - Mantel weiblich")
#    top_10(prs, df_geschlecht_mantel, df_berechnet = False, screenshots=False, number_screenshots = 5, 
#           mode="Gesamt",  zeitung=True, geschlecht="männlich", headline="Top 10 - Mantel männlich")
#        
#        # Top 10 Altersgruppen
#    top_10(prs, df_geschlecht_mantel, df_berechnet = False, screenshots=False, number_screenshots = 5, 
#           mode="Gesamt", headline="default", zeitung=True, alter="40 bis 59 Jahre")
#    top_10(prs, df_geschlecht_mantel, df_berechnet = False, screenshots=False, number_screenshots = 5, 
#           mode="Gesamt", headline="default", zeitung=True, alter="60 bis 79 Jahre")
#    mean_line = df_mantel["Artikel-Lesewert (Erscheinung) in %"].mean()
#    grafik_lesewert(prs, df_mantel, target="ressort", minimize=5, label_position="xlarge",
#                    ressort_liste=mantel_ressorts_auswertung, special=False, title_text="Mantel-Ressorts",
#                    legend="strange", sort="Seitennummer", article="mean", mean_line=mean_line, axis=2, 
#                    anzahl_lokales=4) 
    
#    for elem in mantel_ressorts_auswertung:
#        print(elem)
#        df_ = df[df["Ressortbeschreibung"]==elem]
#        df_geschlecht_ = df_geschlecht[df_geschlecht["Ressortbeschreibung"]==elem]
#        mantel_auswertung(prs, df_, df_geschlecht_, kunde="MZ", lokales=False)
#    
    
#    
    
    # Spezialwünsche Mantel
#    df_titel = df[df["Ressortbeschreibung"]=="Titelseite"]
#    liste_rub = ["Ansichtssache","Tipp des Tages", "Schnell erklärt"]
#    grafik_lesewert(prs, df_titel, target="rubriken", minimize=1, label_position="xlarge",
#                    ressort_liste=liste_rub, special=False, title_text="Rubriken Titelseite", order="ok",
#                    legend="normal", sort="Kategorie", article="mean", grid=True)
    
    # Sport Angebot und Lesewert nach Erscheinungstagen
#    df_sport = df[df["Ressortbeschreibung"]=="Sport"]
#    lesewert_erscheinung(df_sport, df_double_split, prs, title_text="Lesewert nach Wochentagen Sport")
    # Grafik Sportarten
#    schlagworte_finden(prs, df_sport, value="sportart", sort= "LW", min_artikel = 5, title="")
#    schlagworte_finden(prs, df_sport, value="vereine", sort= "LW", min_artikel = 5, title="")
#    grafik_lesewert(prs, df_sport, target="sport", minimize=5, label_position="large",
#                        ressort_liste=[], special=False, title_text="Beliebteste Sportarten", order="ok", legend="normal", sort="Lesewert", 
#                        article="total", mean_line = 0, grid=False, anzahl_lokales=1, mean_line_title = "", axis=2) 
#    
#    REgionalsport
#    df_rs = df_nk[df_nk["Ressortbeschreibung"]=="Regionalsport"]
#    for elem in ausgaben_liste:
#        df_rs_ = df_rs[df_rs["ZTG"]==elem]
#        
#         
#        
#        df_geschlecht_ = df_geschlecht_nk[(df_geschlecht_nk["Ressortbeschreibung"]=="Regionalsport") & (df_geschlecht_nk["ZTG"]==elem)]
#        deckblatt_abschluss(prs, df_rs_, kunde="MZ", title_text = "Regionalsport " + ausgaben_dict[elem] )
#        tabelle_ressortauswertung(prs, df_rs_, research_object="Darstellungsformen", sort="Lesewert", 
#                                  title_text = "Darstellungsformen " + ausgaben_dict[elem], minimum = 5, kunde="MZ", lokales=True)
#        tabelle_ressortauswertung(prs, df_rs_, research_object="Platzierungen", sort="Lesewert", 
#                                  title_text = "Platzierungen " + ausgaben_dict[elem], minimum = 5, kunde="MZ", lokales=True)
#        
#        platzierungen(prs, df_rs_, minimum = 5, geschlecht=False, title_text="Platzierungen " + elem, special="")
#        platzierungen(prs, df_geschlecht_, minimum = 5, geschlecht=True, title_text="Platzierungen nach Geschlecht", special="")
#        darstellungsformen(prs, df_rs_, minimum = 5, geschlecht=False, title_text="Lesewert Darstellungsformen", special="")
#        darstellungsformen(prs, df_geschlecht_, minimum = 5, geschlecht=True, title_text="Darstellungsformen nach Geschlecht", special="")
#        for el in werte:
#            grafik_entwicklung(prs, df_rs_, target=el, mean_line=0, legend="large",
#                               grid=True, title_text = el + " Regionalsport " + ausgaben_dict[elem],
#                               steps=1,  limit_y=True, ma = False)
#        lesewert_erscheinung(df_rs_,df_rs_, prs, title_text="LW nach Wochentagen")
#        
#        top_10(prs, df_rs_, df_berechnet = False, screenshots=True, number_screenshots = 5, 
#           mode="ressort", headline="Top 10 - Regionalsport "+ausgaben_dict[elem], zeitung=True, geschlecht="", alter="", 
#           lokales_ressort="Lokales")
#        # Top 10 Geschlechter
#        top_10(prs, df_geschlecht_, df_berechnet = False, screenshots=False, number_screenshots = 5, 
#           mode="ressort", headline="default", zeitung=True, geschlecht="weiblich")
#        top_10(prs, df_geschlecht_, df_berechnet = False, screenshots=False, number_screenshots = 5, 
#           mode="ressort", headline="default", zeitung=True, geschlecht="männlich")
#        
#        # Top 10 Altersgruppen
#        top_10(prs, df_geschlecht_, df_berechnet = False, screenshots=False, number_screenshots = 5, 
#           mode="ressort", headline="default", zeitung=True, alter="40 bis 59 Jahre")
#        top_10(prs, df_geschlecht_, df_berechnet = False, screenshots=False, number_screenshots = 5, 
#           mode="ressort", headline="default", zeitung=True, alter="60 bis 79 Jahre")
#     
#  Kinder
#    df_kinder = df[df["Ressortbeschreibung"]=="Kinder"]
#    kinder_rub = ["Schnell Schlau","Galaktikus erklärt", "Mannomann", "Bild des Tages",
#                  "Der Buch-Tipp", "Der Spiel-Tipp", "Was macht eigentlich...", "Island-Serie"]
#    grafik_lesewert(prs, df_kinder, target="rubriken", minimize=1, label_position="xlarge",
#                    ressort_liste=kinder_rub, special=False, title_text="Rubriken Kinder", order="ok",
#                    legend="normal", sort="Lesewert", article="mean", grid=True)
#    # blick
#    df_blick = df[df["Ressortbeschreibung"]=="Blick"]
#    grafik_lesewert(prs, df_blick, target="seitentitel", minimize=1, label_position="xlarge",
#                    special=False, title_text="Seitentitel Blick", order="ok",
#                    legend="normal", sort="Lesewert", article="mean", axis=2, grid=True)
#        
    
    prs.save("mzw_abschluss.pptx")
    t_1 = time.time()
    endtime = t_1 - t_0 
    print("Dauer des Durchlaufs: {0:.2f}.".format(endtime))
    
#%% 
    
def vrm(df, df_double_split = 0, df_nk=0, df_scans=0, df_umfrage=0,
                       df_geschlecht = 0, liste = ausgaben_liste, df_ereignisse=0, 
                       df_geschlecht_nk=0):
    print("VRM läuft")  
    t_0=time.time()
    prs = Presentation("LW_neu_2.pptx")
     # 2 - Fragebogen - CHECK
    
    #Aufruf jeder einzelnen Ausgabe
    for elem in ausgaben_liste:
        print(df_umfrage.shape[0])
        df_ = df_umfrage[df_umfrage["Ausgabe"]==elem]
        print(df_.shape[0])
        #Aufruf gestapelte Bars "Ihre Zeitung ist..."
        zeitungsattribute_berechnung(prs, df_, title = elem) # elem wird später über ausgaben_dict umgewandelt
        #zeitungsattribute_berechnung(prs, df_, title = elem, frage_dict = zeitung_attribute_dict2)
        #zeitungsattribute_berechnung(prs, df_, title = elem, frage_dict = zeitung_attribute_dict3)
        
        
        #Aufruf Bars "Welche Themen..." - insgesamt 58 Themen im Fragebogen
        zeitung_themen(prs, df_, title=elem, fragetext="Welche der folgenden Themen interessieren Sie besonders in Ihrer Zeitung?")
        
        # Aufruf kleine Pie-Charts
        umfrage_pie(prs, df_, title=elem)
        
        # Aufruf kleine Barcharts, TV- und Internetnutzung
        
        mini_bars(prs, df_, title_text = "Fernsehnutzung", liste=umfrage_mini_bars_TV, title=elem)
        mini_bars(prs, df_, title_text = "Internetnutzung", liste=umfrage_mini_bars_netz, title=elem)
        mini_bars(prs, df_, title_text = "Persönliches", liste=umfrage_demografie_bars, title=elem)
    
    prs.save("vrm_abschluss.pptx")
    t_1 = time.time()
    endtime = t_1 - t_0 
    print("Dauer des Durchlaufs: {0:.2f} seca.".format(endtime))
    
    
#%%


    
#%% NOZ

def noz(df, df_double_split = 0, df_nk=0, df_scans=0, df_umfrage=0,
                       df_geschlecht = 0, liste = ausgaben_liste, df_ereignisse=0, 
                       df_geschlecht_nk=0):
    print("NOZ läuft")
       
    t_0=time.time()
    prs = Presentation("LW_neu_2.pptx")
    
#    # 1 - Entwicklungsfolie Lesewert/DW/BW im Messverlauf - CHECK
    print("Starte Auswertung Entwicklung über Messzeitraum...")
    blacklist = ["Lokales", "Lokalsport", "Titelseite"]
    
    
    
    df1 = df[df["Spezial_Res"]=="Unterhaltung"]
    gesch1 = df_geschlecht
    mantel_auswertung(prs, df1, gesch1, kunde="NOZ", lokales=False, ztg = "DKB" )

#    for elem in ausgaben_liste:
#        df_ = df[df["ZTG"]==elem]
#        df_geschlecht_ = df_geschlecht[df_geschlecht["ZTG"]==elem]
#        for wert in werte:
#            grafik_entwicklung(prs, df_, target=wert, mean_line=0, legend="large", 
#                               grid=True, title_text = wert + " " + ausgaben_dict[elem], ma=False, 
#                               steps_x_label = 2, limit_y=True)
#    #   
     
      
       

##   
##    
##    
#  # 2 - Fragebogen - CHECK
    
    # Aufruf jeder einzelnen Ausgabe
    #for elem in ausgaben_liste:
#        print(df_umfrage.shape[0])
#        df_ = df_umfrage[df_umfrage["ZTG"]==elem]
#        print(df_.shape[0])
#        #Aufruf gestapelte Bars "Ihre Zeitung ist..."
#        zeitungsattribute_berechnung(prs, df_, title = elem) # elem wird später über ausgaben_dict umgewandelt
#        zeitungsattribute_berechnung(prs, df_, title = elem, frage_dict = zeitung_attribute_dict2)
#        zeitungsattribute_berechnung(prs, df_, title = elem, frage_dict = zeitung_attribute_dict3)
        
#        
#        #Aufruf Bars "Welche Themen..." - insgesamt 58 Themen im Fragebogen
#        zeitung_themen(prs, df_, title=elem, fragetext="Welche der folgenden Themen interessieren Sie besonders in Ihrer Zeitung?")
#        
#        # Aufruf kleine Pie-Charts
#        umfrage_pie(prs, df_, title=elem)
#        
#        # Aufruf kleine Barcharts, TV- und Internetnutzung
#        
#        mini_bars(prs, df_, title_text = "Fernsehnutzung", liste=umfrage_mini_bars_TV, title=elem)
#        mini_bars(prs, df_, title_text = "Internetnutzung", liste=umfrage_mini_bars_netz, title=elem)
#        mini_bars(prs, df_, title_text = "Persönliches", liste=umfrage_demografie_bars, title=elem)
#   
            # LESERAKTIVITÄT
        
        # Entwicklung Werte
#        df_ = df[df["ZTG"]==elem]
#        print(ausgaben_dict[elem] + " läuft. Gemessene Artikel: {}.".format(df_.shape[0]))
#       
#        for e in werte:
#            grafik_entwicklung(prs, df_, target=e, mean_line=0, legend="large", 
#                           grid=True, title_text = e + " " + ausgaben_dict[elem], ma=False, 
#                           steps_x_label = 2, limit_y=True)
#  
#        
#        
#        
#        df_scans_ = df_scans[df_scans["ZTG"]==elem]
#        lesetage(prs, df_scans_, title_text="Lesetage - " + ausgaben_dict[elem], multi_line = True)
#        lesezeit(prs, df_scans_, title_text="Lesezeiten - " + ausgaben_dict[elem])
#        df_analyse_ = df[df["ZTG"]==elem]
#        grafik_lesewert(prs, df_analyse_, target="Textlänge", minimize=1, 
#                        ressort_liste=textlänge, special=False, order="ok", legend="normal", sort="Kategorie", 
#                        article="total", mean_line = 0, grid=False, anzahl_lokales=1, axis=2, 
#                        title_text="Textlängen "+ ausgaben_dict[elem])
#        analyse_mantel = df_analyse_[df_analyse_["Ressortbeschreibung"].isin(mantel_ressorts)]
#        analyse_mantel_ot = df_analyse_[~df_analyse_["Ressortbeschreibung"].isin(blacklist)]
#        grafik_lesewert(prs, analyse_mantel, target="Textlänge", minimize=1, 
#                        ressort_liste=textlänge, special=False, order="ok", legend="normal", sort="Kategorie", 
#                        article="total", mean_line = 0, grid=False, anzahl_lokales=1, axis=2, 
#                        title_text="Textlängen Mantel " + ausgaben_dict[elem])
#        grafik_lesewert(prs, analyse_mantel, target="Textlänge", minimize=1, 
#                        ressort_liste=textlänge, special=False, order="ok", legend="normal", sort="Kategorie", 
#                        article="total", mean_line = 0, grid=False, anzahl_lokales=1, axis=2, 
#                        title_text="Textlängen Mantel ohne Titel " + ausgaben_dict[elem])
#        
        
        # Top 10
    
#        top_10(prs, df_, df_berechnet = False, screenshots=True, number_screenshots = 5, 
#           mode="Gesamt", headline="Top 10 Gesamt "+ ausgaben_dict[aus], zeitung=False, lokales_ressort="Lokales")
###    # Top 10 Mantelressorts CHECK
##    print("Top 10 Mantel werden ausgespielt...")
#        df_mantel = df_[df_["Ressortbeschreibung"].isin(mantel_ressorts)]
#        top_10(prs, df_mantel, df_berechnet = False, screenshots=True, number_screenshots = 5, 
#            mode="Mantel", headline="Top 10 Mantelressorts "+ausgaben_dict[aus], zeitung=False)
        
#        df_mantelohnetitel = df_[~df_["Ressortbeschreibung"].isin(blacklist)]
#        top_10(prs, df_mantelohnetitel, df_berechnet = False, screenshots=True, number_screenshots = 5, 
#            mode="Mantel", headline="Top 10 Mantel ohne Titel "+ausgaben_dict[aus], zeitung=False)
#    
    
    
    
        # Übersicht Ressorts, LW und Artikelzahl
        
#        # erstmal nur die echten Ressorts
#        df_ressorts = df_[df_["Ressortbeschreibung"].isin(ressort_list)]
#        mean_line = df_["Artikel-Lesewert (Erscheinung) in %"].mean()
#        grafik_lesewert(prs, df_ressorts, target="ressort", minimize=5, label_position="xlarge",
#                    ressort_liste=ressort_list, special=False, title_text="Ressorts " + ausgaben_dict[elem],
#                    legend="strange", sort="Seitennummer", article="mean", mean_line=mean_line, axis=2, 
#                    font_size=23, anzahl_lokales=1)
#        columns_df = ["Ressortbeschreibung", "SplitId", "Seitennummer", 
#                 "Artikel-Lesewert (Erscheinung) in %"]
#        lok_eins_df = pd.DataFrame(columns=columns_df)
#        lok_total_df = pd.DataFrame(columns=columns_df)
#        if elem == "DKB": 
#            df_lok = df_[(df_["Ressortbeschreibung"]=="Titelseite") | (df_["Ressortbeschreibung"]=="Lokales")]
#            df_lokrest = df_[df_["Ressortbeschreibung"]=="Lokales"]
#            df_loktotal = df_[(df_["Ressortbeschreibung"]=="Titelseite") | (df_["Ressortbeschreibung"]=="Lokales")]
#            name_lok = "Lokales + Titel"
#        if elem == "NST":
#            df_lok = df_[df_["Spezial_Res"]=="Lokale Eins"]
#            df_lokrest = df_[df_["Spezial_Res"]=="Lokales"]
#            df_loktotal = df_[df_["Ressortbeschreibung"]=="Lokales"]
#            name_lok = "Lokale Eins"
#        
#        
#        splitId = df_lok.shape[0]
#        seitennummer = df_lok.Seitennummer.median()
#        lesewert = df_lok["Artikel-Lesewert (Erscheinung) in %"].mean()
#        lok_ = pd.Series([name_lok, splitId, seitennummer, lesewert], index = columns_df)
#        lok_eins_df = lok_eins_df.append(lok_, ignore_index=True)  
#        lok_total_df = lok_total_df.append(lok_, ignore_index=True)  
#        lok_eins_df["SplitId"] = lok_eins_df["SplitId"].astype(float)
#        lok_total_df["SplitId"] = lok_total_df["SplitId"].astype(float)
#        if elem == "NST": 
#            extension = lok_eins_df
#        elif elem == "DKB": 
#            extension = lok_total_df
#        
#        
#        grafik_lesewert(prs, df_ressorts, target="ressort", minimize=5, label_position="xlarge",
#                    ressort_liste=ressort_list, special=False, title_text="Ressorts Gesamt",
#                    legend="strange", sort="Seitennummer", article="mean", mean_line=mean_line, axis=2, 
#                    font_size=23, anzahl_lokales=1, extension_df=lok_eins_df)
#        
#        grafik_lesewert(prs, df_ressorts, target="special_res", minimize=5, label_position="xlarge",
#                    special=False, title_text="Ressorts Gesamt Spezial "+elem,
#                    legend="strange", sort="Seitennummer", article="mean", mean_line=mean_line, axis=2, 
#                    font_size=23, anzahl_lokales=1)
        
        
        
        
        
        
        # Übersicht nutzung nach Geschlecht / LW nach Darstellungsformen / Platzierungen / Geschle
        
        
#        df_geschlecht_ = df_geschlecht[df_geschlecht["ZTG"]==elem]
#        
#           
#        multiple_bars_geschlecht(prs, df_geschlecht_, target="Ressorts", grid = True, legend="xlarge", 
#                             ressort_liste = ressort_list, title_text= "Ressorts nach Geschlecht - "+elem, 
#                             sort="Seitennummer", width = "special")
#        
#        darstellungsformen(prs, df_, minimum = 5, geschlecht=False, title_text="Darstellungsformen " + ausgaben_dict[elem])
#        darstellungsformen(prs, df_geschlecht_, minimum = 5, geschlecht=True, title_text="Darstellungsformen nach Geschlecht " + elem)
   
        
    
    ######################## Übersicht einzelne Ressorts Mantel
#        if elem == "NST":
#            #ressort_ausw = ["Titelseite"]
##            ressort_ausw = ['Titelseite', 'Politik', 'Einblicke', 'Hintergrund', 'Nordwest', 'Wirtschaft',
##        'Dialog', 'Sport', 'Sport Regional', 'Lokales', 'Tipps und Termine', 'Kultur Regional', 'Weltspiegel',
##       'Gut zu wissen', 'Medien / Fernsehen', 'Kultur', 'Lokalsport',
##       'Film', 'Reportage', 'Leserbriefe', 'Sondervorhaben', 'Wochenendjournal' ]
#           ressort_ausw = [ "Sport Regional"]
#            #ressort_ausw = ["Wochenendjournal"]
#        
#        if elem == "DKB":
#            ressort_ausw = ["Lokalsport"]
#            #ressort_ausw = ["Wochenendjournal"]
##            ressort_ausw = ['Titelseite', 'Lokales', 'Verlagssonderthema', 'Nordwest',
##       'Reportage', 'Wochenendjournal', 'Kultur', 'Politik', 'Einblicke',
##       'Hintergrund', 'Gut zu wissen', 'Wirtschaft', 'Weltspiegel',
##       'Sport', 'Lokalsport','Automarkt',
##       'Immobilienmarkt', 'Stellenmarkt', 'Film']
##            
#        
#        for res in ressort_ausw:
#            df_res = df_[df_["Ressortbeschreibung"]==res]
#            sh = df_res.shape[0]
##            print ("Analyse {} gestartet, {} Datenzeilen gefunden, {} wird ausgewertet".format(ausgaben_dict[elem], df_.shape[0], res))
#            df_gesch = df_geschlecht_[(df_geschlecht_["Ressortbeschreibung"]==res)]
##            if sh >= 10:
##                mantel_auswertung(prs, df_res, df_gesch, kunde="NOZ", lokales=False, ztg = elem)
##            else: 
##                print ("Kein Treffer für Ressort {} in {}".format(res, ausgaben_dict[elem]))
#            if res == "Sport": 
#                mantel_auswertung(prs, df_res, df_gesch, kunde="NOZ", lokales=False, ztg = elem)
#                schlagworte_finden(prs, df_res, value="sportart", sort= "LW", min_artikel = 5, title="")
#                schlagworte_finden(prs, df_res, value="vereine", sort= "LW", min_artikel = 5, title="")
##                grafik_lesewert(prs, df_res, target="sport", minimize=5, label_position="large",
#                        ressort_liste=[], special=False, title_text="Beliebteste Sportarten", order="ok", legend="normal", sort="Lesewert", 
#                        article="total", mean_line = 0, grid=False, anzahl_lokales=1, mean_line_title = "", axis=2) 
                #kolumnen_finder(prs, df_res, anzahl=5, ausnahmem=[], axis = 2, label_position = "large", title_text = "Kolumnen", grid=True, limit=10)
#                
#                df_sport1 = df_res[df_res["Spezial"]=="Sport 1"]
#                mantel_auswertung(prs, df_sport1, df_gesch, kunde="NOZ", lokales=False, ztg = elem)
#                schlagworte_finden(prs, df_sport1, value="sportart", sort= "LW", min_artikel = 5, title="")
#                #kolumnen_finder(prs, df_sport1, anzahl=5, ausnahmem=[], axis = 2, label_position = "large", title_text = "Kolumnen", grid=True, limit=10)
##                
#            
#                
#            if res == "Sport Regional": 
#                if len(df_res) > 0: 
#                    schlagworte_finden(prs, df_res, value="sportart", sort= "LW", min_artikel = 5, title="")
#                    schlagworte_finden(prs, df_res, value="vereine", sort= "LW", min_artikel = 5, title="")
#                   # kolumnen_finder(prs, df_res, anzahl=5, ausnahmem=[], axis = 2, label_position = "large", title_text = "Kolumnen", grid=True, limit=10)
#            
#            if res == "Lokalsport": 
#                if len(df_res) > 0: 
#                    schlagworte_finden(prs, df_res, value="sportart", sort= "LW", min_artikel = 5, title="")
#                    schlagworte_finden(prs, df_res, value="vereine", sort= "LW", min_artikel = 5, title="")
#                    #kolumnen_finder(prs, df_res, anzahl=5, ausnahmem=[], axis = 2, label_position = "large", title_text = "Kolumnen", grid=True, limit=10)
##                
#            if res == "Politik":
#                
#                df_pol1 = df_res[df_res["Spezial"]=="Politik 1"]
#                
#                
#                df_pol2 = df_res[df_res["Spezial"]=="Politik 2"]
#                df_pol3 = df_res[df_res["Spezial"]=="Politik 3"]
#                mantel_auswertung(prs, df_pol1, df_gesch, kunde="NOZ", lokales=False, ztg = elem)
#                if len(df_pol2)>0: 
#                    mantel_auswertung(prs, df_pol2, df_gesch, kunde="NOZ", lokales=False, ztg = elem)
#                if len(df_pol3)>0: 
#                    mantel_auswertung(prs, df_pol3, df_gesch, kunde="NOZ", lokales=False, ztg = elem)
##                    
##            if res == "Dialog": 
##                df_unterhaltung = df_[df_["Seitentitel"]=="unterhaltung"]
##                mantel_auswertung(prs, df_unterhaltung, df_gesch, kunde="NOZ", lokales=False, ztg = elem)
##            
#            
#            if res == "Lokales":
#                liste_kolumnen = ['Till', 'Kommentar', 'Kompakt',
#       'Familienname', 'Kurz notiert', 'Leserbriefe',
#       'Mit Stallgeruch', 'Laga 2018', 'Aus der Politik',
#       'Delmenhorster Chronik', 'Buchtipp', 'Ganderkeseer Chronik',
#       'Wonneproppen', 'Leserbrief', 'Rezept der Woche',
#       'Selbsthilferuppen', 'Lesermeinung', 'Sorry', 'Nostalgie',
#       'Klönschnack', 'Quergedacht', 'Delmenhorst Heute',
#       'Wirtschaft vor Ort', 'Vor 25 Jahren', 'Wir gehen aus',
#       'Worte der Woche', 'De plattdütsche Eck',
#       'Na klar! Erzählnachricht für Kinder', 'Zum Sonntag', 'Stimmen',
#       'Immo. + Wohnungsnot', 'Im Interview', 'Im Gespräch',
#       'Aus kirchlicher Sicht', 'Das tut sich in Osnabrück',
#       'Mein lieber Herr Gesangsverein', 'Pflegefall Altenpflege',
#       'Zeitreise', 'Plattdüütsche Week', 'Delmenhorster Klönschknack',
#       'Persönlich', 'Ganderkeseer Klönschnack', 'NOZ-Telefonaktion',
#       'Jeden Tag ein Türchen', 'Adventskalender Soroptimist Club']
#                lokale_auswertung(prs, df_res, df_geschlecht_, liste_kolumnen=liste_kolumnen, kunde="NOZ", ztg = elem)
#                if elem == "NST": 
#                    df_lokeins = df_res[df_res["Spezial"]=="Lokale Eins"]
#                    df_lokrest = df_res[~(df_res["Spezial_Res"]=="Lokale Eins")]
#                    df_lokrest = df_lokrest[~(df_lokrest["Seitentitel"]=="region")]
#                    
#                    df_region = df_res[(df_res["Seitentitel"]=="region") |(df_res["Seitentitel"]=="Region")]  
#                    list_df = [df_lokeins, df_lokrest, df_region] 
#                if elem == "DKB":
#                    df_lokeins = df_res[df_res["Spezial"]=="Titelseite DKB"]
#                    df_lokrest = df_res[(df_res["Seitentitel"]=="delmenhorst") | (df_res["Seitentitel"]=='delmenhorst / stuhr')]
#                    df_region = df_res[(df_res["Seitentitel"]=="Region") | (df_res["Seitentitel"]=="region")]  
#                    df_gander_komplett = df_res[(df_res["Seitentitel"]=="ganderkesee") | (df_res["Seitentitel"]=="ganderkesee / landkreis") | (df_["Seitentitel"]=="ganderkeseer zeitung") | (df_["Seitentitel"]=="ganderkeseer eins")]  
#                    df_gander_eins = df_res[df_res["Spezial"]=="ganderkesee eins"]
#                    df_gander_rest = df_gander_komplett[~(df_gander_komplett["Spezial"]=="ganderkesee eins")]
#                    list_df = [df_lokeins, df_lokrest, df_region, df_gander_komplett, df_gander_eins, df_gander_rest]
#                for e in list_df: 
#                    if len(e) > 0:
#                        lokale_auswertung(prs, e, df_gesch, liste_kolumnen=liste_kolumnen, kunde="NOZ", ztg = elem)
#                        
#            if res == "Wochenendjournal":
#                df_we = df_res[df_res["weekday"]==5]
#                df_ge = df_gesch[df_gesch["weekday"]==5]
#                print(df_we.shape[0])
#                mantel_auswertung(prs, df_we, df_ge, kunde="NOZ", lokales=False, ztg = elem)
#                
#                df_u = df_res[(df_res["Seitentitel"]=="unterhaltung") & (~df_res["weekday"]==5)]
#                if len(df_u)>0: 
#                    mantel_auswertung(prs, df_u, df_ge, kunde="NOZ", lokales=False, ztg = elem)
    prs.save("noz_abschluss.pptx")
    t_1 = time.time()
    endtime = t_1 - t_0 
    print("Dauer des Durchlaufs: {0:.2f} seca.".format(endtime))
    
#    
#    
#    
##
#%%drp_umfrage
def drp_umfrage(df_umfrage):
    print("DRP Umfrage läuft")  
    t_0=time.time()
    prs = Presentation("LW_neu_2.pptx")
     # 2 - Fragebogen - CHECK
    
    #Aufruf jeder einzelnen Ausgabe
    for elem in ausgaben_liste:
        print(ausgaben_liste)
        print(elem)
        print(df_umfrage.head())
        print(df_umfrage.shape[0])
        df_ = df_umfrage[df_umfrage["ZTG"]==elem]
        print(df_.shape[0])
        #Aufruf gestapelte Bars "Ihre Zeitung ist..."
        zeitungsattribute_berechnung(prs, df_, title = elem) # elem wird später über ausgaben_dict umgewandelt
        #zeitungsattribute_berechnung(prs, df_, title = elem, frage_dict = zeitung_attribute_dict2)
        #zeitungsattribute_berechnung(prs, df_, title = elem, frage_dict = zeitung_attribute_dict3)
        
        
        #Aufruf Bars "Welche Themen..." - insgesamt 58 Themen im Fragebogen
        zeitung_themen(prs, df_, title=elem, fragetext="Welche der folgenden Themen interessieren Sie besonders in Ihrer Zeitung?")
        
        # Aufruf kleine Pie-Charts
        umfrage_pie(prs, df_, title=elem)
        
        # Aufruf kleine Barcharts, TV- und Internetnutzung
        
        mini_bars(prs, df_, title_text = "Fernsehnutzung", liste=umfrage_mini_bars_TV, title=elem)
        mini_bars(prs, df_, title_text = "Internetnutzung", liste=umfrage_mini_bars_netz, title=elem)
        mini_bars(prs, df_, title_text = "Persönliches", liste=umfrage_demografie_bars, title=elem)
    
    prs.save("drp_umfrage.pptx")
    t_1 = time.time()
    endtime = t_1 - t_0 
    print("Dauer des Durchlaufs: {0:.2f} seca.".format(endtime))
    
#%% Umfrage SWZ

def swz_umfrage(df_umfrage):
    print("SWZ Umfrage läuft")  
    t_0=time.time()
    prs = Presentation("LW_neu_2.pptx")
     # 2 - Fragebogen - CHECK
    
    #Aufruf jeder einzelnen Ausgabe
    for elem in ausgaben_liste:
        print(ausgaben_liste)
        print(elem)
        print(df_umfrage.head())
        print(df_umfrage.shape[0])
        df_ = df_umfrage[df_umfrage["ZTG"]==elem]
        print(df_.shape[0])
        #Aufruf gestapelte Bars "Ihre Zeitung ist..."
        zeitungsattribute_berechnung(prs, df_, title = elem) # elem wird später über ausgaben_dict umgewandelt
        #zeitungsattribute_berechnung(prs, df_, title = elem, frage_dict = zeitung_attribute_dict2)
        #zeitungsattribute_berechnung(prs, df_, title = elem, frage_dict = zeitung_attribute_dict3)
        
        
        #Aufruf Bars "Welche Themen..." - insgesamt 58 Themen im Fragebogen
        zeitung_themen(prs, df_, title=elem, fragetext="Welche der folgenden Themen interessieren Sie besonders in Ihrer Zeitung?")
        
        # Aufruf kleine Pie-Charts
        umfrage_pie(prs, df_, title=elem)
        
        # Aufruf kleine Barcharts, TV- und Internetnutzung
        
        mini_bars(prs, df_, title_text = "Fernsehnutzung", liste=umfrage_mini_bars_TV, title=elem)
        mini_bars(prs, df_, title_text = "Internetnutzung", liste=umfrage_mini_bars_netz, title=elem)
        mini_bars(prs, df_, title_text = "Persönliches", liste=umfrage_demografie_bars, title=elem)
    
    prs.save("swz_umfrage.pptx")
    t_1 = time.time()
    endtime = t_1 - t_0 
    print("Dauer des Durchlaufs: {0:.2f} seca.".format(endtime))    
    
    
#%% Zwischenauswertung DRP
def drp_zwischen(df, df_double_split = 0, df_nk=0, df_scans=0, df_umfrage=0,
                       df_geschlecht = 0, liste = ausgaben_liste, df_ereignisse=0, 
                       df_geschlecht_nk=0):
    
    
    t_0=time.time()
    prs = Presentation("LW_neu_2.pptx")
    
##    # 1 - Entwicklungsfolie Lesewert/DW/BW im Messverlauf - CHECK
#    for elem in werte:
#        grafik_entwicklung(prs, df, target=elem, mean_line=0, legend="large", 
#                           grid=True, title_text = False, ma=False, 
#                           steps_x_label = 2, limit_y=True)
    
#    #Lesetage
    
    df_ohne_sonn = df_scans[df_scans["Ressortbeschreibung"]!="Sonntag"]
    df_os = df[df["Ressortbeschreibung"]!="Sonntag"]
    for el in werte:
        grafik_entwicklung(prs, df_os, target=el, mean_line=0, legend="large",
                               grid=True, title_text = "Entwicklung " + el, 
                               steps=1,  limit_y=True, ma = False)
#    lesetage(prs, df_ohne_sonn, title_text="Lesetage") 
#    lesezeit(prs, df_ohne_sonn, title_text="Lesezeiten")
#    lesetage(prs, df_scans, title_text="Lesetage (mit Sonntagszeitung) ") 
#    lesezeit(prs, df_scans, title_text="Lesezeiten (mit Sonntagszeitung) ")
#    
#    # Ressort-Übersichten
#    
#    
#    mean_line = df_os["Artikel-Lesewert (Erscheinung) in %"].mean()
##       
#    grafik_lesewert(prs, df_os, target="ressort", minimize=5, label_position="xlarge",
#                    ressort_liste=ressort_list, special=False, title_text="Ressortübersicht",
#                    legend="normal", sort="Seitennummer", article="mean", mean_line=mean_line, axis=2)
#    
#    #  Darstellungsformen LW, Menge und GEschlecht CHECK
##    darstellungsformen(prs, df, minimum = 5, geschlecht=False, title_text="Darstellungsformen Gesamt")
##    darstellungsformen(prs, df_geschlecht, minimum = 5, geschlecht=True, grid= True, title_text="Darstellungsformen nach Geschlecht")
#    platzierungen(prs, df, minimum = 5, geschlecht=False, title_text="Platzierungen und Lesewerte", special="")
#    platzierungen(prs, df_geschlecht, minimum = 5, geschlecht=True, title_text="Platzierungen nach Geschlecht", special="")
    # TOP 10 GESAMT (einzelne Ausgaben und Mantel)
#    for elem in ausgaben_liste:
         
#        df_ = df_nk[df_nk["ZTG"]==elem]
#        top_10(prs, df_, df_berechnet = False, screenshots=True, number_screenshots = 5, 
#           mode="Gesamt", headline="Top 10 Gesamt "+ ausgaben_dict[elem], zeitung=False, lokales_ressort="Lokales")
##
#   
###    # Top 10 Mantelressorts CHECK
##    print("Top 10 Mantel werden ausgespielt...")
#    df_mantel = df[df["Ressortbeschreibung"].isin(mantel_ressorts)]
#    top_10(prs, df_mantel, df_berechnet = False, screenshots=True, number_screenshots = 5, 
#            mode="Mantel", headline="Top 10 Mantelressorts", zeitung=False)
#    
#    # Übersicht Ressorts nach LW und Häufigkeit
#    mean_line = df_["Artikel-Lesewert (Erscheinung) in %"].mean()
       
#    grafik_lesewert(prs, df_, target="ressort", minimize=5, label_position="xlarge",
#                    ressort_liste=sortierung_liste, special=False, title_text="Ressorts " + ausgaben_dict[elem],
#                    legend="normal", sort="Kategorie", article="mean", mean_line=mean_line, axis=2, 
#                    extension_df=lok_1, font_size=27)
    
    # 8 - Ressortnutzung nach Geschlecht CHECK
#    for elem in ausgaben_liste:
#        df_ = df_geschlecht_nk[df_geschlecht_nk["ZTG"]==elem]
#        print(df_.shape)
#       #df_ = df_.drop_duplicates(subset="SplitId")
#        multiple_bars_geschlecht(prs, df_, target="Ressorts", grid = True, legend="xlarge", title_text= "Ressorts nach Geschlecht - "+elem, 
#                             sort="Kategorie", width = "special")
#        multiple_bars_geschlecht(prs, df_geschlecht, target="Ressorts", grid = True, legend="xlarge", 
#                             title_text= "Ressorts nach Geschlecht GEsamt",  
#                             sort="Kategorie", width = "special")
##  
    
    
    
    
    
#    
#    # Auswertung Mantelressorts
#    print("Starte Auswertung Entwicklung über Messzeitraum...")
#    mantel_ressorts_plus_Sonntag = mantel_ressorts.copy()
#    mantel_ressorts_plus_Sonntag.append("Sonntag")
#    print(mantel_ressorts)
#    print(mantel_ressorts_plus_Sonntag)
#    
#    
#    for elem in mantel_ressorts_plus_Sonntag:
#        print(elem)
#        
#        
#        
#        df_ = df[df["Ressortbeschreibung"]==elem]
#        
#        df_geschlecht_ = df_geschlecht[df_geschlecht["Ressortbeschreibung"]==elem]
#        mantel_auswertung(prs, df_, df_geschlecht_, kunde="DRP", lokales=False)
    
    
    
    
    # Auswertung Lokalteil
    
#    for elem in ausgaben_liste:
#####        
###        
#        df_ = df[(df["Ressortbeschreibung"]=="Lokales") & (df["ZTG"]==elem)]
#        grafik_lesewert(prs, df_, target="seitentitel", minimize=5, label_position="large",
#                            ressort_liste=[], special=False, title_text="Lokalseiten in " + elem, order="ok", legend="normal",  
#                            article="total", mean_line = 0, grid=False, anzahl_lokales=1, mean_line_title = "", axis=2) 
#    
##        
#        
#    
#        df_geschlecht_ = df_geschlecht[(df_geschlecht["Ressortbeschreibung"]=="Lokales") & (df_geschlecht["ZTG"]==elem)]
#        deckblatt_abschluss(prs, df_, kunde="DRP", title_text = "Lokales " + ausgaben_dict[elem])
#        tabelle_ressortauswertung(prs, df_, research_object="Darstellungsformen", sort="Lesewert", 
#                                  title_text = "Darstellungsformen " + ausgaben_dict[elem], minimum = 5, kunde="DRP", lokales=True)
#        tabelle_ressortauswertung(prs, df_, research_object="Platzierungen", sort="Lesewert", 
#                                  title_text = "Platzierungen " + ausgaben_dict[elem], minimum = 5, kunde="DRP", lokales=True)
##        
#        platzierungen(prs, df_, minimum = 5, geschlecht=False, title_text="Platzierungen " + elem, special="")
#        platzierungen(prs, df_geschlecht_, minimum = 5, geschlecht=True, title_text="Platzierungen nach Geschlecht", special="")
#        darstellungsformen(prs, df_, minimum = 5, geschlecht=False, title_text="Lesewert Darstellungsformen", special="")
#        darstellungsformen(prs, df_geschlecht_, minimum = 5, geschlecht=True, title_text="Darstellungsformen nach Geschlecht", special="")
#        for el in werte:
#            grafik_entwicklung(prs, df_, target=el, mean_line=0, legend="large",
#                               grid=True, title_text = el + " " + ausgaben_dict[elem],
#                               steps=1,  limit_y=True, ma = False)
#     
     #19 Lesewert nach Tagen und im Vergleich zum Angebot
#    lesetage(prs, df_scans, title_text="Lesetage - Gesamt", multi_line = True)
#    lesezeit(prs, df_scans, title_text="Lesezeiten - Gesamt")
### 

    # Ressorts nach Geschlechtern Gesamt
#    df_geschlecht_ = df_geschlecht[df_geschlecht["Ressortbeschreibung"]!="Sonntag"]
#    multiple_bars_geschlecht(prs, df_geschlecht_, target="Ressorts", grid = True, legend="xlarge", 
#                             ressort_liste = ressort_list, title_text= "Ressorts nach Geschlecht GEsamt",  
#                             sort="Kategorie", width = "special")
#              
     
    
    # Special Sport 
    #df_sportarten["sport"] = df_sportarten["Themen"].apply(check_sportarten)
#
#    print(df_sportarten["sport"].value_counts())
#        grafik_lesewert(prs, df_sportarten, target="sport", minimize=5, label_position="large",
#                            ressort_liste=[], special=False, title_text="Beliebteste Sportarten", order="ok", legend="normal", sort="Lesewert", 
#                            article="total", mean_line = 0, grid=False, anzahl_lokales=1, mean_line_title = "", axis=2) 
#    df_sport = df[df["Ressortbeschreibung"]=="Sport"]        
#    schlagworte_finden(prs, df_sport, value="sportart", sort= "LW", min_artikel = 10, title="")
#    schlagworte_finden(prs, df_sport, value="vereine", sort= "LW", min_artikel = 10, title="")
#    schlagworte_finden(prs, df_sport, value="sportart", sort= "Artikelanzahl", min_artikel = 10, title="")
#    schlagworte_finden(prs, df_sport, value="vereine", sort= "Artikelanzahl", min_artikel = 10, title="")
#                grafik_lesewert(prs, df_res, target="sport", minimize=5,
    
    # ÜBERSICHT LOKALES und LOKALSPORT
   
            
    #for elem in ausgaben_liste:
##       ressort = "Lokales" 
     #   ressort = "Lokalsport"
#        #df_lok = df[(df["Ressortbeschreibung"]=="Lokales") & (df["ZTG"]==elem)]
    #    df_lok = df[(df["Ressortbeschreibung"]==ressort) & (df["ZTG"]==elem)]
#        #df_lok = df[(df["Ressortbeschreibung"]=="Lokales") & (df["ZTG"]==elem
#        #df_analyse = df_komplett[(df_komplett["Ressortbeschreibung"]=="Lokales") & (df_komplett["ZTG"]==elem)]
#        
#        df_nk_ = df_nk[(df_nk["Ressortbeschreibung"]==ressort) & (df_nk["ZTG"]==elem)]
#        df_geschlecht_ = df_geschlecht[(df_geschlecht["Ressortbeschreibung"]==ressort) & (df_geschlecht["ZTG"]==elem)]
#        deckblatt_abschluss(prs, df_lok, kunde="DRP", title_text = ressort + " " + ausgaben_dict[elem])
#        tabelle_ressortauswertung(prs, df_lok, research_object="Darstellungsformen", sort="Lesewert", 
#                                  title_text = "Darstellungsformen " + ausgaben_dict[elem], minimum = 5, kunde="DRP", lokales=True)
#        tabelle_ressortauswertung(prs, df_lok, research_object="Platzierungen", sort="Lesewert", 
#                                  title_text = "Platzierungen " + ausgaben_dict[elem], minimum = 5, kunde="DRP", lokales=True)
#        
#        platzierungen(prs, df_lok, minimum = 5, geschlecht=False, title_text="Platzierungen " + elem, special="")
#        platzierungen(prs, df_geschlecht_, minimum = 5, geschlecht=True, title_text="Platzierungen nach Geschlecht", special="")
#        darstellungsformen(prs, df_lok, minimum = 5, geschlecht=False, title_text="Lesewert Darstellungsformen", special="")
#        darstellungsformen(prs, df_geschlecht_, minimum = 5, geschlecht=True, title_text="Darstellungsformen nach Geschlecht", special="")
#        for el in werte:
#            grafik_entwicklung(prs, df_lok, target=el, mean_line=0, legend="large",
#                               grid=True, title_text = el + " " + ausgaben_dict[elem],
#                               steps=1,  limit_y=True, ma = False)
##            df__ = df_lok[(df_lok["Platzierungen"]=="AA")|(df_lok["Platzierungen"]=="SK")]
###            grafik_entwicklung(prs, df__, target=el, mean_line=0, legend="large",
###                               grid=True, title_text = el + " Aufmacher und Seitenkeller", steps=2)  
#        top_10(prs, df_lok, df_berechnet = False, screenshots=True, number_screenshots = 5, 
#           headline="Top 10 Lokalsport "+elem, mode="Gesamt", zeitung=False, lokales_ressort="Lokales")
#        schlagworte_finden(prs, df_lok, value="sportart", sort= "LW", min_artikel = 10, title="Höchste Lesewerte nach Sportart")
#        #schlagworte_finden(prs, df_lok, value="vereine", sort= "LW", min_artikel = 10, title="Höchste Lesewerte nach Verein")
#        schlagworte_finden(prs, df_lok, value="sportart", sort= "Artikelanzahl", min_artikel = 10, title="LW der meistgenannten Sportarten")
#        #schlagworte_finden(prs, df_lok, value="vereine", sort= "Artikelanzahl", min_artikel = 10, title="LW der meistgenannten Vereine")
##    
#     
    df_sonntag = df[df["Ressortbeschreibung"]!="Sonntag"]
    df_sonntag_nk = df_nk[df_nk["Ressortbeschreibung"]!="Sonntag"]
#    df_sonntag.shape[0]
#    mean_line = df_sonntag["Artikel-Lesewert (Erscheinung) in %"].mean()
#    grafik_lesewert(prs, df_sonntag, target="seitentitel", minimize=5, label_position="xlarge",special=False, title_text="Seitentitel Sonntag",
#                    legend="normal", sort="Seitennummer", article="total", mean_line=mean_line, axis=2,font_size=27)
#    df_we = df[df["Ressortbeschreibung"]=="Wochenende"]
#
#    mean_line = df_we["Artikel-Lesewert (Erscheinung) in %"].mean()
#    grafik_lesewert(prs, df_we, target="seitentitel", minimize=2, label_position="xlarge",special=False, title_text="Seitentitel Sonntag",
#                    legend="normal", sort="Seitennummer", article="total", mean_line=mean_line, axis=2,font_size=27)
#    
    blacklist = ["57947", "86426", "72408", "73278", "44818", "69220", "44006", "44206", "44364", "44565", "44734"
                 "58154", "58317", "58490", "58651"]      
    df_blacklist_lw = df_sonntag[~df_sonntag["ArtikelId"].isin(blacklist)] 
    df_blacklist_lw_nk = df_sonntag_nk[~df_sonntag_nk["ArtikelId"].isin(blacklist)] 
    df_mantel_bl = df_blacklist_lw[df_blacklist_lw["Ressortbeschreibung"].isin(mantel_ressorts)]
    df_hinter_bl = df_blacklist_lw[df_blacklist_lw["Ressortbeschreibung"]=="Hintergrund"]
    #df_leser_bl = df_blacklist_lw[df_blacklist_lw["Ressortbeschreibung"]=="Leserbriefe"]
    df_pir_nk = df_blacklist_lw_nk[df_blacklist_lw_nk["ZTG"]=="PIR"]
    top_10(prs, df_mantel_bl, df_berechnet = False, screenshots=True, number_screenshots = 5, 
           headline="Top 10 Mantel ", mode="Gesamt", zeitung=False, lokales_ressort="Lokales")
    top_10(prs, df_hinter_bl, df_berechnet = False, screenshots=True, number_screenshots = 5, 
           headline="Top 10 Hintergrund ", mode="Gesamt", zeitung=False, lokales_ressort="Lokales")
    top_10(prs, df_blacklist_lw, df_berechnet = False, screenshots=True, number_screenshots = 5, 
           headline="Top 10 Gesamt ", mode="Gesamt", zeitung=False, lokales_ressort="Lokales")
    top_10(prs, df_pir_nk, df_berechnet = False, screenshots=True, number_screenshots = 5, 
           headline="Top 10 Gesamt Pirmasens ", mode="Gesamt", zeitung=False, lokales_ressort="Lokales")
 
    
    prs.save("drp_zwischen.pptx")
    t_1 = time.time()
    endtime = t_1 - t_0 
    print("Dauer des Durchlaufs: {0:.2f} seca.".format(endtime))
    
    
    
    
#%%    SWZ_2019
    
    
    
    
def swz_top10(df, kunden_id, geschlecht = "", alter=""):
    
    
    t_0=time.time()
    prs = Presentation("LW_neu_2.pptx")
    
#    res = ['Titelseite', 'Wir im Süden', 'Umland/Region',  'Seite Drei', 'Nachrichten & Hintergrund', 
#           'Meinung & Dialog', 'Wirtschaft', 'Panorama/Journal', 'Kultur', 'Fernsehen & Freizeit', 'Ratgeber/Vermischtes',
#           'Sport', 'Wochenende', 'Schwäbische Märkte', 'Reise & Erholung', 'Bauen & Wohnen',  'Auto & Verkehr', 'Immobilien']
           
     
    ausgabe_liste = ['RV', 'WG', 'GR']
    
#    df_lok = df[df["Ressort"]=="Lokales"]
#    res_lok = ["Lokales, Lokalsport"]
    
#    for elem in res:
#        df_ = df[df["Ressort"]==elem]
#        print(elem)
#        top_10(prs, df_,  mode="Gesamt", zeitung=False, headline="Top 10 " + elem, kunden_id = kunden_id, geschlecht = geschlecht, alter=alter, screenshots=False)
#        
    for aus in ausgabe_liste:
        
        print(aus)
        df_neu = df[df["Ausgabename"]==aus]
        print(df_neu.shape)
        top_10(prs, df_neu,  mode="Gesamt", zeitung=False, headline="Top 10 Titelseite" + aus, kunden_id = kunden_id, geschlecht = geschlecht, alter=alter, screenshots=False)
         
    
    
    
    
    prs.save("swz_top10.pptx")
    t_1 = time.time()
    endtime = t_1 - t_0 
    print("Dauer des Durchlaufs: {0:.2f} seca.".format(endtime))
    
    
    
    
    
#%%

def waiblingen(df):
    prs = Presentation("LW_neu_2.pptx")
    ressorts = df.Ressort.unique()
    cleanedList = [x for x in ressorts if str(x) != 'nan']
    for elem in cleanedList: 
        df_ = df[df["Ressort"]==elem]
        if elem == "Lokales":
            aus_liste = ["Lokales Schorndorf", "Lokales Waiblingen"]
            for el in aus_liste: 
                df__ = df_[df_["Ausgabeteil"]==el]
                top_10(prs, df__, df_berechnet = False, screenshots=True, number_screenshots = 5, 
                       headline="Top 10 " + elem, zeitung=False, lokales_ressort="Lokales", kunden_id = 1016, mode="Gesamt")  
            
        else: 
            top_10(prs, df_, df_berechnet = False, screenshots=True, number_screenshots = 5, 
                   headline="Top 10 " + elem, zeitung=False, lokales_ressort="Lokales", kunden_id = 1016, mode="Gesamt")  
    prs.save("waiblingen_zwischen.pptx")